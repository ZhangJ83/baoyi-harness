import json
import os
import re
import sys
import itertools
import threading
from dataclasses import dataclass, field
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Any, Callable

from . import config
from .tools.registry import dispatch, select_tools
from .state import RunState, RuntimePhase, Status, TaskItem
from .hooks import PostToolHook, PreToolHook, ToolEvent
from .runtime import RuntimeController, bounded_tool_result, canonical_call, OBSERVE_TOOLS, TERMINAL_TOOLS
from .events import EventBus, EventKind
from .session import Session, StopReason, TurnOutcome
from .task_profiles import TaskProfile, classify_task
from .task_scope import infer_ppt_mutation_scope, is_scope_continuation, ppt_scope_is_explicit
from .task_compiler import compile_task, brief_json, TaskSpec
from .execution_contract import compile_execution_contract
from .tools.tool_catalog import specialize_tools
from .code_task_compiler import compile_code_task
from .goal_runtime import GoalStore, goal_from_contract
from .planning import plan_next

_READ_ONLY = {"read_file", "read_many", "list_dir", "glob_files", "search_text", "git_status", "git_diff", "verify_files", "run_checks", "deck_info", "ppt_verify"}
_PPT_PLAN_MUTATORS = {"ppt_edit_text", "ppt_style", "ppt_compose", "ppt_arrange"}


@dataclass
class RunControl:
    """Per-turn control plane, mirroring Codex Interrupt / Claude AbortController."""

    tool_limit: int | None
    tool_start: int
    cancelled: threading.Event = field(default_factory=threading.Event)

    def used_tools(self, state: RunState) -> int:
        return state.tool_calls - self.tool_start

    def admits(self, state: RunState, requested: int) -> bool:
        return self.tool_limit is None or self.used_tools(state) + requested <= self.tool_limit


class Harness:
    """Agentic loop: plan -> act -> observe -> verify.

    The LLM returns tool_calls; we execute them, feed results back,
    and continue until the model stops calling tools. Callable repeatedly:
    each `run()` call keeps the session (messages + deck), like a terminal chat.
    """

    def __init__(
        self,
        model: str | None = None,
        max_steps: int | None = None,
        controller_policy: str = "cegar_h",
        interactive: bool = False,
    ):
        config.load_dotenv()
        from .llm import make_client

        self.llm = make_client(model)
        self.max_steps = max_steps or config.max_steps()
        self.messages: list[dict[str, Any]] = []
        self.deck = None
        self.on_tool: Callable[[str, str, str], None] = lambda name, args, out: None
        self.started = False
        self.state = RunState()
        self._done: str | None = None
        self.pre_tool_hooks: list[PreToolHook] = []
        self.post_tool_hooks: list[PostToolHook] = []
        self.loaded_skills: set[str] = set()
        self.skill_allowed_tools: set[str] = set()
        self.recorder = None
        from .controller_policies import PolicyGuard
        self.controller_policy = controller_policy
        self.policy_guard = PolicyGuard(controller_policy)
        self.interactive = interactive
        self._run_control: RunControl | None = None
        self.runtime = RuntimeController()
        self.events = EventBus()
        self.session = Session()
        self.task_profile: TaskProfile | None = None
        self.goal_store = GoalStore(config.sandbox_root())
        self.active_goal = self.goal_store.load()
        self._last_planning_signature = None
        # Product-facing extension points. These belong to the UI/session
        # layer, not to the loop policy, and are deliberately preserved across
        # Harness.reset().
        self.undo_stack: list[bytes] = []
        self.approval_handler: Callable[[str], str] | None = None
        self.stream_callback: Callable[[str], None] | None = None
        # UI extension point: actual reasoning text returned by the provider.
        # The terminal/GUI may display it, clearly labelled as model output;
        # xiaopu never fabricates a hidden chain-of-thought.
        self.reasoning_callback: Callable[[str], None] | None = None
        # Obligation-based progress monitor (CEGAR-H v2):
        # progress = an unresolved obligation was resolved, or the artifact
        # mutation epoch advanced; fresh-but-unchanged evidence is NOT progress.
        self._obligation_checkpoint: tuple | None = None
        self._rejection_history: list[tuple] = []

    def attach_printer(self, printer: Callable[[str, str, str], None]) -> None:
        self.on_tool = printer

    _TASK_LOCAL_FACTS = {
        "task_root", "task_instruction", "task_capability", "task_difficulty",
        "required_output_pptx", "ppt_input_deck", "ppt_input_candidates",
        "official_evaluator", "official_evaluator_present",
        "verification_contract", "verification_contract_terms_present",
        "ppt_source_ir", "task_spec", "selected_skill", "task_intent",
        "execution_capability", "code_language", "code_test_runner",
        "primary_input", "output_path", "manifest_task_id",
        "manifest_batch_goal", "ppt_full_check_downgraded",
        "bound_task_identity", "ppt_repair_observation_calls",
        "repair_observation_grant", "auto_evaluator_coverage_applied",
    }

    def _clear_task_facts(self) -> None:
        """Drop task-scoped intake/compiler facts before a NEW task turn.

        Interactive sessions reuse one Harness, so a plain chat question after
        a benchmark task must not inherit the previous evaluator path, output
        contract, verification contract or task root. Chat history, deck and
        UI state are intentionally preserved.
        """
        state = self.state
        for key in list(state.facts):
            if (
                key in self._TASK_LOCAL_FACTS
                or key.startswith("source:")
                or key.startswith("preflight:")
            ):
                state.facts.pop(key, None)
        state.verification_contract_terms = {}
        state.content_brief = ""
        state.execution_contract = None
        for kind in (
            "task_evaluator", "ppt_contract", "ppt_structural", "ppt_quality",
            "ppt_visual", "file_verification", "code_check",
        ):
            state.unresolved_checks.discard(kind)
        state.repair_attempts = 0
        state.last_verification_failed = False
        self.task_spec = None
        self._last_planning_signature = None

    def undo(self) -> str:
        """Restore the deck snapshot taken before the most recent PPT mutation.

        This is a UI-level convenience, not a loop repair action: it does not
        consume repair budget and it does not touch verification evidence.
        """
        import io

        if not self.undo_stack:
            return "没有可撤销的 PPT 修改（撤销栈为空）。"
        from pptx import Presentation

        payload = self.undo_stack.pop()
        self.deck = Presentation(io.BytesIO(payload))
        self.state.record_change("deck:undo")
        slide_count = len(self.deck.slides) if self.deck is not None else 0
        return f"已撤销最近一次修改，当前内存稿：{slide_count} 页（剩余 {len(self.undo_stack)} 步可撤销）。"

    def subscribe(self, callback) -> Callable[[], None]:
        return self.events.subscribe(callback)

    def _obligation_snapshot(self) -> tuple:
        return (
            tuple(sorted(self.state.unresolved_checks)),
            self.state.mutation_epoch,
            tuple(sorted(record.kind for record in self.state.fresh_evidence())),
        )

    def _obligation_progress(self) -> bool:
        """CEGAR-H v2 progress definition.

        Progress only when an unresolved obligation disappears, the artifact
        mutation epoch advances, or fresh evidence appears for a previously
        unresolved obligation. Re-running the same save/check/render with no
        obligation change is explicitly NOT progress.
        """
        now = self._obligation_snapshot()
        checkpoint = getattr(self, "_obligation_checkpoint", None)
        if checkpoint is None:
            self._obligation_checkpoint = now
            return False
        old_unresolved, old_epoch, old_evidence = checkpoint
        new_unresolved, new_epoch, new_evidence = now
        progress = (
            len(new_unresolved) < len(old_unresolved)
            or any(kind not in new_unresolved for kind in old_unresolved)
            or new_epoch != old_epoch
            or any(kind in new_evidence for kind in old_unresolved if kind not in old_evidence)
        )
        if progress:
            self._obligation_checkpoint = now
            self._rejection_history = []
        return progress

    @staticmethod
    def _rejection_signature(name: str, exc: Exception, blockers) -> tuple:
        """(action, blocker_type, blockers, blocker_target) — the target digest
        keeps 'text not found: 6/20' and 'text not found: 6/26' distinct."""
        return (name, type(exc).__name__, blockers, str(exc)[:120])

    def _mutation_gated(self, tool_name: str) -> bool:
        """CEGAR-H verify-before-continue gate.

        Once the artifact has unverified mutations, further content mutation is
        blocked until fresh structural evidence exists. This makes the runtime
        recommendation an enforceable loop invariant, not a hint.
        """
        from .controller_policies import MUTATION_TOOLS

        if tool_name not in MUTATION_TOOLS:
            return False
        if not getattr(self.state, "mutation_epoch", 0):
            return False
        fresh = {record.kind for record in self.state.fresh_evidence()}
        if "ppt_structural" in fresh:
            return False
        # Interactive sessions and official-evaluator tasks run the full loop;
        # bare unit-test harnesses keep the old permissive path.
        if not (getattr(self, "interactive", False) or self.state.facts.get("official_evaluator_present") == "true"):
            return False
        return True

    def _set_task_profile(self, task: str) -> TaskProfile:
        """Classify the task and expose the selected route as auditable state."""
        profile = classify_task(task)
        previous = getattr(self, "task_profile", None)
        changed = previous is None or previous.name != profile.name
        self.task_profile = profile
        state = getattr(self, "state", None)
        if state is not None:
            state.task_profile = profile.name
            state.task_capabilities = profile.capabilities
            state.verification_plan = profile.verification
            state.design_policy = profile.design_policy
        if changed:
            payload = {
                "profile": profile.name,
                "label": profile.label,
                "description": profile.description,
                "capabilities": list(profile.capabilities),
                "verification": list(profile.verification),
                "design_policy": profile.design_policy,
            }
            recorder = getattr(self, "recorder", None)
            if recorder is not None and callable(getattr(recorder, "event", None)):
                recorder.event("task_profiled", **payload)
            events = getattr(self, "events", None)
            if events is not None:
                events.publish(EventKind.TASK_PROFILED, **payload)
        return profile

    def _bind_ppt_mutation_scope(self, task: str) -> None:
        """Bind an explicit mutation scope, preserving it only for a pure resume turn."""

        if is_scope_continuation(task):
            return
        scope = infer_ppt_mutation_scope(task)
        self.state.ppt_allowed_slides = set(scope.slides)
        self.state.ppt_allowed_shapes = {
            slide: set(shapes) for slide, shapes in scope.shape_targets
        }
        self.state.ppt_scope_global = scope.global_scope
        self.state.ppt_scope_hard = scope.restrictive
        self.state.ppt_scope_explicit = scope.explicit

    def _auto_open_preflight_deck(self) -> str:
        """Open the one exact PPT input selected by deterministic intake.

        A create/template-build task uses any supplied PPTX as a style or
        layout *reference*, not as the editable working copy. Auto-opening it
        would steer the model into inspecting a reference deck instead of
        building the requested new artifact.
        """

        if getattr(self, "deck", None) is not None:
            return ""
        profile_name = getattr(getattr(self, "task_profile", None), "name", "")
        if profile_name in {"create_deck"}:
            return ""
        path = self.state.facts.get("ppt_input_deck", "")
        if not path:
            return ""
        result = dispatch("ppt_open", json.dumps({"path": path}, ensure_ascii=False), self)
        self.state.record_fact("ppt_input_auto_opened", path)
        if getattr(self, "recorder", None):
            self.recorder.event("ppt_input_auto_opened", path=path, result=result)
        return result

    def _publish_phase_change(self, previous: RuntimePhase | None, reason: str = "") -> None:
        current = self.state.phase
        if previous is None or previous == current:
            return
        payload = {
            "from_phase": previous.value,
            "to_phase": current.value,
            "reason": reason or self.state.last_meta_reason,
            "profile": getattr(self.task_profile, "name", ""),
        }
        recorder = getattr(self, "recorder", None)
        if recorder is not None and callable(getattr(recorder, "event", None)):
            recorder.event("phase_changed", **payload)
        self.events.publish(EventKind.PHASE_CHANGED, **payload)

    def request_cancel(self) -> None:
        """Request a cooperative stop at the next safe execution boundary."""
        control = getattr(self, "_run_control", None)
        if control is not None:
            control.cancelled.set()
            cancel_provider = getattr(self.llm, "cancel_current", None)
            if callable(cancel_provider):
                cancel_provider()

    def cancel_requested(self) -> bool:
        control = getattr(self, "_run_control", None)
        return bool(control and control.cancelled.is_set())

    def add_pre_tool_hook(self, hook: PreToolHook) -> None:
        self.pre_tool_hooks.append(hook)

    def add_post_tool_hook(self, hook: PostToolHook) -> None:
        self.post_tool_hooks.append(hook)

    def _system_prompt(self, task: str) -> str:
        from . import memory
        from .context import workspace_context
        from .skills import catalog, discover, match
        from .controller_policies import policy_instruction

        profile = self._set_task_profile(task)
        rules = memory.load_rules()
        notes = memory.load_notes()
        read_only_task = "do not modify" in task.lower() or "read-only" in task.lower()
        skills = [] if read_only_task else discover()
        selected = [] if read_only_task else match(task, skills)
        self.loaded_skills.update(skill.name for skill in selected)
        restrictions = [set(skill.allowed_tools) for skill in selected if skill.allowed_tools]
        self.skill_allowed_tools = set.intersection(*restrictions) if restrictions else set()
        read_only_hint = ""
        if "do not modify" in task.lower() or "read-only" in task.lower():
            read_only_hint = (
                "\n--- Read-only task constraint ---\n"
                "This is a read-only task. Use at most two targeted inspection calls, then call finish immediately with the evidence. "
                "Do not explore unrelated files, run broad scans, or repeat observations.\n"
            )
        fast_path_hint = ""
        if self._is_local_ppt_edit(task):
            fast_path_hint = (
                "\n--- Local PPT edit fast path ---\n"
                "Use the shortest evidence-complete route: ppt_open; inspect the affected scope once only when needed; use the matching semantic edit; ppt_save; ppt_check; finish. "
                "Preserve unrelated content. Use only task-local materials. At most one targeted repair is allowed. "
                "The harness owns task discovery, provenance, official evaluation, final render/visual audit, and trajectory output.\n"
            )
        elif self._is_ppt_task(task) and profile.name == "create_deck":
            fast_path_hint = (
                "\n--- New-deck creation fast path ---\n"
                "There is no existing presentation to open; do not guess filenames or call ppt_open. "
                "Create it from scratch: call new_deck once for the cover, then one ppt_compose (content/comparison/table/quadrant) per additional slide, then ppt_save and ppt_check once the deck is complete.\n"
            )
        return (
            "Identity (non-negotiable): You are 小朴 (Xiaopu), the product's autonomous coding and PowerPoint agent. "
            "Never identify yourself as Claude, Anthropic, Codex, OpenAI, or any other product. "
            "If asked about the active model, state only the configured model identifier when it is available; "
            "do not guess model family, provider deployment, or version.\n"
            "You run inside the 小朴 harness with a phase-scoped tool set, PowerPoint "
            "creation/modification/verification capabilities, and evidence-based completion checks. "
            "Shell or Python execution exists only when its canonical tool is explicitly advertised for this turn.\n"
            "Language policy: use Chinese for user-facing answers, task plans, tool rationales, and concise verification summaries unless the user explicitly requests another language. "
            "Do not reveal hidden chain-of-thought or private reasoning. Instead, communicate a short auditable work summary: goal, actions taken, evidence, and any blocker.\n"
            "Response consistency: if you say items are listed below, above, or as follows, include the actual itemized entries in that same answer. "
            "Never claim that a list was displayed when you only inspected or counted it.\n"
            "You get a task, and you work it step by step using tools. "
            "The harness gives you: tools, execution, and a verification layer.\n"
            "Rules:\n"
            "0. Call only tool names advertised in the current request. Never invent sys_exec, sys_cat, shell aliases, or filenames.\n"
            "1. Inspect before editing. Maintain a short task list for multi-step work.\n"
            "2. Prefer small verifiable steps. After changes, run the cheapest relevant verifier.\n"
            "2a. For office-document tasks, use the harness-provided task brief when present. Source discovery and final output provenance are recorded automatically; do not spend calls recreating them.\n"
            "3. If a tool errors, read the error, adjust, retry. Never repeat the exact same call.\n"
            "4. For PPT work, inspect the affected scope once, choose layouts intentionally, save, and run structural verification. Final rendering, visual inspection, and ContentIR provenance binding are automatic lifecycle duties; do not recreate them with model tool calls.\n"
            "4b. Prefer one semantic PPT composition/edit tool over many primitive text boxes or an ad-hoc script; the harness owns layout and provenance mechanics.\n"
            "4a. Opening an existing deck creates an editable working copy and preserves the original. After a concrete verification defect, repair only the cited slide/shape, stay within the repair budget, then rerun the failed check.\n"
            "5. Call finish only after citing concrete verification evidence.\n"
            "5a. PPT completion requires a saved PPTX plus fresh structural evidence. Call finish next; finish automatically obtains fresh render/pixel evidence when available or records an explicit renderer diagnostic.\n"
            "6. Everything you write must be placed under the workspace directory.\n"
            "7. A successful shell exit is not enough by itself: inspect test output and ensure the command exercised the changed behavior.\n"
            "8. Do not weaken, delete, or bypass tests merely to obtain a passing result.\n"
            "9. Classify the request before acting. For read-only/explanation tasks, inspect only the minimum relevant artifacts, then call finish; do not create a task list or run broad repository scans.\n"
            "10. For a small single-file task, use at most one inspect action, one mutation, one focused verifier, and finish. Avoid re-reading unchanged files.\n"
            "11. Prefer a concise final answer after evidence is sufficient; do not spend turns narrating plans or repeating successful observations.\n"
            f"12. Active controller intervention: {getattr(self, 'controller_policy', 'cegar_h')}. {policy_instruction(getattr(self, 'controller_policy', 'cegar_h'))}\n"
            + ("\n--- Workspace context ---\n" + workspace_context() if not read_only_task else "")
            + "\n"
            + "\n--- Available skills ---\n"
            + catalog(skills)
            + "\n--- Task profile / capability catalog ---\n"
            + f"Selected profile: {profile.name} ({profile.label}). {profile.description}\n"
            + f"Capabilities: {', '.join(profile.capabilities)}\n"
            + f"Verification plan: {', '.join(profile.verification)}\n"
            + f"Design policy: {profile.design_policy}\n"
            + "Route only the capabilities needed by the current phase. This is an adaptive task route, not a rigid script.\n"
            + "\n"
            + ("\n--- Loaded skill instructions ---\n" + "\n\n".join(skill.body for skill in selected) + "\n" if selected else "")
            + (f"\n--- Project rules ---\n{rules}\n" if rules else "")
            + (f"\n--- Session notes ---\n{notes}\n" if notes else "")
            + read_only_hint
            + fast_path_hint
        )

    def reset(self) -> None:
        self.messages = []
        self.deck = None
        self.started = False
        self.state = RunState()
        self._done = None
        self.loaded_skills = set()
        self.skill_allowed_tools = set()
        self.recorder = None
        self.undo_stack = []
        self._obligation_checkpoint = None
        self._rejection_history = []
        from .controller_policies import PolicyGuard
        self.policy_guard = PolicyGuard(getattr(self, "controller_policy", "cegar_h"))
        self._run_control = None
        self.runtime = RuntimeController()
        self.session = Session()
        self.task_profile = None
        self._last_planning_signature = None

    def _bind_task_context(self, task: str) -> str:
        """Resolve a task package and isolate its execution state.

        Conversation UI history is external to model execution state. A new
        task package gets a clean deck, controller, cancellation token, cache,
        recorder and message context, matching clean-session benchmark rules.
        """
        from .intake import task_root_from_prompt

        task_root = task_root_from_prompt(task)
        if task_root is None:
            return task
        relative = str(task_root.relative_to(config.sandbox_root())).replace("/", "\\")
        identity = str(task_root.resolve()).casefold()
        previous = self.state.facts.get("bound_task_identity", "").casefold()
        if previous and previous != identity:
            self._reset_task_local_state()
        elif not previous and (
            self.state.tool_calls
            or self.state.mutation_epoch
            or self.state.phase not in {RuntimePhase.INTAKE, RuntimePhase.UNDERSTAND}
        ):
            self._reset_task_local_state()
        self.state.record_fact("bound_task_identity", identity)
        self.state.record_fact("task_root", relative)
        # Canonicalize the model-visible request so all downstream compilers
        # consume the same explicit package reference.
        if f"tasks\\{task_root.name}".casefold() not in task.replace("/", "\\").casefold():
            return f"完成 tasks\\{task_root.name}。\n用户原始请求：{task}"
        return task

    def _reset_task_local_state(self) -> None:
        """Reset task-scoped state without deleting durable goals or UI chat."""
        self.messages = []
        self.deck = None
        self.started = False
        self.state = RunState()
        self._done = None
        self.loaded_skills = set()
        self.skill_allowed_tools = set()
        self.recorder = None
        self._run_control = None
        self.runtime = RuntimeController()
        self.task_profile = None
        self._last_planning_signature = None
        from .controller_policies import PolicyGuard
        self.policy_guard = PolicyGuard(getattr(self, "controller_policy", "cegar_h"))

    def run(self, task: str) -> str:
        return self.run_turn(task).text

    def start_goal(self, objective: str) -> str:
        """Start one explicit durable goal from the current task contract."""
        if not objective.strip():
            raise ValueError("goal objective cannot be empty")
        # A new explicit goal is a new completion ledger even when the chat and
        # active artifact remain available for continuity.
        self.state.final_summary = None
        self.state.operational_plan = []
        self.state.unresolved_checks.clear()
        self._last_planning_signature = None
        ppt_task = self._is_ppt_task(objective)
        from .intake import prepare_task_brief
        from .task_compiler import portable_contract_for
        prepare_task_brief(objective, self.state, self.recorder) if ppt_task else ""
        spec = compile_task(objective, self.state.facts, self.state.content_brief)
        portable = portable_contract_for(objective, self.state.facts, self.state.content_brief) if ppt_task else None
        contract = compile_execution_contract(spec, ppt_task, None if ppt_task else compile_code_task(objective, config.sandbox_root()), portable=portable)
        self.state.execution_contract = contract
        self.active_goal = goal_from_contract(objective, contract)
        self.goal_store.save(self.active_goal)
        self.events.publish(EventKind.GOAL_UPDATED, goal=self.active_goal.objective, status=self.active_goal.status,
                            milestones=self.active_goal.milestones, completed=[])
        return f"长期目标已启动：{self.active_goal.objective}（{len(self.active_goal.milestones)} 个里程碑）"

    def goal_summary(self) -> str:
        goal = self.active_goal
        if goal is None:
            return "当前没有活动的长期目标。"
        remaining = [item for item in goal.milestones if item not in goal.completed_milestones]
        return (
            f"长期目标：{goal.objective}\n状态：{goal.status}\n"
            f"已完成：{len(goal.completed_milestones)}/{len(goal.milestones)}\n"
            + ("下一里程碑：" + remaining[0] if remaining else "完成条件已满足")
        )

    def _effective_goal_task(self, task: str) -> tuple[str, bool]:
        """Resolve a bare continuation against the durable goal objective."""
        goal = getattr(self, "active_goal", None)
        continuing = bool(goal is not None and goal.status == "active" and is_scope_continuation(task))
        return (goal.objective, True) if continuing else (task, False)

    def _sync_goal(self) -> None:
        goal = getattr(self, "active_goal", None)
        if goal is None:
            return
        completed = [item["content"] for item in self.state.operational_plan if item.get("status") == "completed"]
        for milestone in completed:
            goal.complete(milestone)
        if self.state.final_summary and not self.state.unresolved_checks:
            for milestone in goal.milestones:
                goal.complete(milestone)
            goal.status = "completed"
        store = getattr(self, "goal_store", None)
        if store is not None:
            store.save(goal)
        self.events.publish(EventKind.GOAL_UPDATED, goal=goal.objective, status=goal.status,
                            milestones=goal.milestones, completed=goal.completed_milestones)

    def _publish_planning(self) -> None:
        """Publish only changed, auditable plans; never expose private reasoning."""
        planning = plan_next(self.state, getattr(self.state, "execution_contract", None))
        signature = (planning.stage, planning.next_action, planning.evidence, planning.gaps, planning.reason)
        if signature == getattr(self, "_last_planning_signature", None) and not planning.revised:
            return
        self._last_planning_signature = signature
        payload = {
            "stage": planning.stage,
            "next_action": planning.next_action,
            "evidence": list(planning.evidence),
            "gaps": list(planning.gaps),
            "reason": planning.reason,
            "revised": planning.revised,
        }
        self.events.publish(EventKind.PLANNING_DECISION, **payload)
        recorder = getattr(self, "recorder", None)
        if recorder is not None and callable(getattr(recorder, "event", None)):
            recorder.event("planning_decision", **payload)

    def run_turn(self, task: str) -> TurnOutcome:
        """Execute one typed turn while preserving the legacy `run() -> str` API."""
        if not hasattr(self, "events"):
            self.events = EventBus()
        if not hasattr(self, "session"):
            self.session = Session()
        task = self._bind_task_context(task)
        task_listing = self._workspace_task_listing(task)
        if task_listing is not None:
            outcome = TurnOutcome(
                text=task_listing,
                stop_reason=StopReason.END_TURN,
                tool_calls=self.state.tool_calls,
                total_tokens=self.state.total_tokens,
                mutation_epoch=self.state.mutation_epoch,
                phase=self.state.phase.value,
                artifact=self.state.active_artifact,
            )
            self.session.append(outcome)
            return outcome
        if self._is_progress_question(task):
            text = self.progress_summary()
            outcome = TurnOutcome(
                text=text,
                stop_reason=StopReason.END_TURN,
                tool_calls=self.state.tool_calls,
                total_tokens=self.state.total_tokens,
                mutation_epoch=self.state.mutation_epoch,
                phase=self.state.phase.value,
                artifact=self.state.active_artifact,
            )
            self.session.append(outcome)
            return outcome
        self.events.publish(
            EventKind.TURN_STARTED,
            task=task,
            agent="Xiaopu",
            model=getattr(self.llm, "model", "configured-model"),
            provider=config.provider(),
        )
        try:
            text = self._run_text(task)
            reason = self._infer_stop_reason(text)
        except Exception:
            if config.strict_run_budget():
                raise
            text = f"运行错误：{type(sys.exception()).__name__}: {sys.exception()}\n\n已保留当前状态，可输入“继续”重试。"
            reason = StopReason.ERROR
        outcome = TurnOutcome(
            text=text,
            stop_reason=reason,
            tool_calls=self.state.tool_calls,
            total_tokens=self.state.total_tokens,
            mutation_epoch=self.state.mutation_epoch,
            phase=self.state.phase.value,
            artifact=self.state.active_artifact,
        )
        self.session.append(outcome)
        self._sync_goal()
        self.events.publish(
            EventKind.TURN_COMPLETED,
            stop_reason=reason.value,
            text=text,
            turn_id=outcome.turn_id,
            phase=outcome.phase,
            tool_calls=outcome.tool_calls,
            mutation_epoch=outcome.mutation_epoch,
            artifact=outcome.artifact,
        )
        return outcome

    @staticmethod
    def _workspace_task_listing(task: str) -> str | None:
        """Answer inventory questions from disk, never from model narration."""
        normalized = task.strip().casefold()
        count_cues = ("多少个任务", "几个任务", "任务数量", "how many tasks")
        list_cues = ("列出任务", "任务列表", "任务名称", "列出来", "list tasks", "list the tasks")
        mentions_scope = "任务" in normalized or "task" in normalized
        if not mentions_scope or not any(cue in normalized for cue in (*count_cues, *list_cues)):
            return None
        root = config.sandbox_root().resolve()
        tasks_dir = root / "tasks"
        if not tasks_dir.is_dir():
            return f"当前工作区没有 `tasks` 目录：{root}"
        names = sorted((path.name for path in tasks_dir.iterdir() if path.is_dir()), key=str.casefold)
        if not names:
            return f"当前工作区的 `tasks` 目录为空：{tasks_dir}"
        lines = [f"当前工作区共有 {len(names)} 个任务："]
        lines.extend(f"{index}. {name}" for index, name in enumerate(names, 1))
        return "\n".join(lines)

    @staticmethod
    def _is_progress_question(task: str) -> bool:
        text = task.strip().casefold().rstrip("？?")
        cues = (
            "进展到哪里", "现在进展", "当前进度", "进度怎么样", "做到哪里",
            "what is the progress", "where are we", "status update",
        )
        return any(cue in text for cue in cues)

    def progress_summary(self) -> str:
        """Return operational state without waking the model or running tools."""

        artifacts = getattr(getattr(self, "recorder", None), "manifest", {}).get("artifacts", [])
        final = next((row.get("path") for row in reversed(artifacts) if row.get("role") == "final-pptx"), None)
        evidence = [record for record in self.state.evidence if record.passed]
        lines = [
            "当前进度：",
            f"- 阶段：{self.state.phase.value}",
            f"- 已执行：{self.state.tool_calls} 次工具调用；已产生 {self.state.mutation_epoch} 轮修改",
            f"- 当前文件：{final or self.state.active_artifact or '尚未保存最终产物'}",
        ]
        if self.state.operational_plan:
            lines.append("- 计划：")
            marks = {"completed": "已完成", "in_progress": "进行中", "pending": "待处理", "blocked": "受阻"}
            lines.extend(
                f"  - {marks.get(item['status'], item['status'])}：{item['content']}"
                for item in self.state.operational_plan
            )
        if evidence:
            lines.append("- 最近验证：" + evidence[-1].summary)
        if self.state.unresolved_checks:
            lines.append("- 当前阻塞：" + "、".join(sorted(self.state.unresolved_checks)))
        elif final:
            lines.append("- 下一步：完成最终交付并进入下一题。")
        else:
            lines.append("- 下一步：保存产物并执行增量验证。")
        return "\n".join(lines)

    def _publish_task_plan(self) -> None:
        if not hasattr(self, "events"):
            self.events = EventBus()
        self.events.publish(
            EventKind.TASK_PLAN,
            items=list(self.state.operational_plan),
        )

    def _ensure_execution_plan(self) -> None:
        if self.state.operational_plan:
            return
        contract = getattr(self.state, "execution_contract", None)
        if contract is not None and contract.stages:
            self.state.operational_plan = [
                {
                    "id": stage.id,
                    "content": stage.label,
                    "status": "in_progress" if index == 0 else "pending",
                }
                for index, stage in enumerate(contract.stages)
            ] + [{"id": "deliver", "content": "交付结果", "status": "pending"}]
            self._publish_task_plan()
            return
        compiled = getattr(self, "task_spec", None)
        if compiled is not None and compiled.plan:
            self.state.operational_plan = [
                {"id": f"step_{index}", "content": content,
                 "status": "in_progress" if index == 1 else "pending"}
                for index, content in enumerate(compiled.plan, 1)
            ]
            self._publish_task_plan()
            return
        self.state.operational_plan = [
            {"id": "discover", "content": "确认任务、输入文件与修改范围", "status": "completed" if self.deck is not None else "in_progress"},
            {"id": "edit", "content": "执行最小范围修改", "status": "in_progress" if self.deck is not None else "pending"},
            {"id": "save", "content": "保存最终 PPTX", "status": "pending"},
            {"id": "verify", "content": "增量验证并交付", "status": "pending"},
        ]
        self._publish_task_plan()

    # Compatibility for saved tests/extensions; all new code uses the domain-
    # neutral execution-plan projection above.
    def _ensure_ppt_plan(self) -> None:
        self._ensure_execution_plan()

    def _advance_execution_plan(self, tool: str) -> None:
        if not self.state.operational_plan:
            return
        by_id = {item["id"]: item for item in self.state.operational_plan}
        changed = False
        if {"understand", "produce", "verify", "deliver"}.issubset(by_id):
            if tool in OBSERVE_TOOLS:
                target = "understand"
            elif tool in {"verify_files", "run_checks", "ppt_check", "run_task_evaluator"}:
                target = "verify"
            elif tool == "finish":
                target = "deliver"
            else:
                target = "produce"
            order = ["understand", "produce", "verify", "deliver"]
            target_index = order.index(target)
            for index, key in enumerate(order):
                status = "completed" if index <= target_index else "pending"
                if index == target_index + 1:
                    status = "in_progress"
                if by_id[key]["status"] != status:
                    by_id[key]["status"] = status
                    changed = True
            if changed:
                self.events.publish(
                    EventKind.PROGRESS_UPDATED,
                    tool=tool,
                    items=list(self.state.operational_plan),
                )
            return
        if "discover" not in by_id:
            stages = self.state.operational_plan
            target_index = None
            if tool in _PPT_PLAN_MUTATORS:
                target_index = min(1, len(stages) - 1)
            elif tool == "ppt_save":
                target_index = min(2, len(stages) - 1)
            elif tool == "ppt_check" and self.state.fresh_evidence():
                target_index = len(stages) - 1
            if target_index is not None:
                for index, item in enumerate(stages):
                    status = "completed" if index <= target_index else "pending"
                    if index == target_index + 1 and index < len(stages):
                        status = "in_progress"
                    if item["status"] != status:
                        item["status"] = status
                        changed = True
            if changed:
                self.events.publish(EventKind.PROGRESS_UPDATED, tool=tool, items=list(stages))
            return
        if tool in _PPT_PLAN_MUTATORS:
            for key, status in (("discover", "completed"), ("edit", "completed"), ("save", "in_progress")):
                if by_id[key]["status"] != status:
                    by_id[key]["status"] = status; changed = True
        elif tool == "ppt_save":
            for key, status in (("save", "completed"), ("verify", "in_progress")):
                if by_id[key]["status"] != status:
                    by_id[key]["status"] = status; changed = True
        elif tool == "ppt_check" and self.state.fresh_evidence():
            if by_id["verify"]["status"] != "completed":
                by_id["verify"]["status"] = "completed"; changed = True
        if changed:
            self.events.publish(
                EventKind.PROGRESS_UPDATED,
                tool=tool,
                items=list(self.state.operational_plan),
            )

    def _advance_ppt_plan(self, tool: str) -> None:
        self._advance_execution_plan(tool)

    def _coalesce_atomic_inspect_batch(self, calls):
        """Keep one high-value inspect from an atomic-skill model batch.

        Providers may request overview and targeted inspections concurrently.
        Per-result gates cannot stop the later calls because all were admitted
        before any result existed.  For trajectory-derived atomic skills the
        targeted shapes view dominates a whole-deck summary, so execute only
        that call and let the next model turn act on its evidence.
        """
        if self.state.facts.get("selected_skill") not in {"ppt.atomic_edit", "ppt.atomic_style"}:
            return calls
        inspections = [call for call in calls if call.function.name == "ppt_inspect"]
        if len(inspections) <= 1:
            return calls

        def score(call) -> tuple[int, int]:
            try:
                arguments = json.loads(call.function.arguments or "{}")
            except (TypeError, json.JSONDecodeError):
                arguments = {}
            return (
                2 if arguments.get("detail") == "shapes" else 0,
                1 if arguments.get("slide_number") is not None else 0,
            )

        chosen = max(inspections, key=score)
        reduced = [call for call in calls if call.function.name != "ppt_inspect" or call is chosen]
        removed = len(calls) - len(reduced)
        if removed:
            if not hasattr(self, "events"):
                self.events = EventBus()
            recorder = getattr(self, "recorder", None)
            if recorder is not None:
                recorder.event(
                    "tool_batch_coalesced",
                    skill=self.state.facts.get("selected_skill"),
                    tool="ppt_inspect",
                    removed=removed,
                )
            self.events.publish(
                EventKind.PROGRESS_UPDATED,
                tool="ppt_inspect",
                items=list(self.state.operational_plan),
                note=f"合并 {removed} 个等价检查，只执行最具体的目标页检查。",
            )
        return reduced

    def _constrain_compiled_tool_schemas(self, tools: list[dict]) -> list[dict]:
        """Compatibility delegate; specialization belongs to Tool Runtime."""
        contract = getattr(self.state, "execution_contract", None)
        if contract is None:
            contract = getattr(self, "task_spec", None)
        return specialize_tools(tools, contract)

    def _run_text(self, task: str) -> str:
        from .intake import bind_manifest_task
        from .task_scope import is_scope_continuation
        original_task = task
        task, manifest_task_id = bind_manifest_task(task)
        effective_task, continuing_goal = self._effective_goal_task(task)
        active_goal = getattr(self, "active_goal", None)
        if not continuing_goal and not is_scope_continuation(task):
            self._clear_task_facts()
        if manifest_task_id:
            self.state.record_fact("manifest_task_id", manifest_task_id)
            self.state.record_fact("manifest_batch_goal", original_task)
        # Budgets apply to one user turn. Lifetime totals remain in RunState for
        # audit/cost reporting, but must not make a later question inherit the
        # previous question's allowance.
        turn_total_start = self.state.total_tokens
        turn_generated_start = self.state.generated_output_tokens
        # Like Codex's per-turn Interrupt operation and Claude Code's per-query
        # AbortController, control state never leaks into the next user turn.
        # Claude Code and Codex do not apply a fixed aggregate tool-call cap to
        # their normal interactive REPL. Hard caps remain available only for
        # explicit benchmark/non-interactive execution.
        strict_budget = config.strict_run_budget() and not getattr(self, "interactive", False)
        control = RunControl(config.max_tool_calls() if strict_budget else None, self.state.tool_calls)
        self._run_control = control
        self.current_task = effective_task
        read_only_task = "do not modify" in effective_task.lower() or "read-only" in effective_task.lower()
        if getattr(self, "recorder", None) is None or self.recorder.completed:
            from .lifecycle import RunRecorder, discover_workspace
            self.recorder = RunRecorder(effective_task, self.llm.model, config.provider())
            discovery = discover_workspace()
            self.recorder.event("workspace_discovery", discovery=discovery)
        profile = self._set_task_profile(effective_task)
        from .intake import prepare_task_brief
        preflight_brief = prepare_task_brief(effective_task, self.state, self.recorder) if self._is_ppt_task(effective_task) else ""
        preflight_open = self._auto_open_preflight_deck() if preflight_brief else ""
        self.task_spec = compile_task(effective_task, self.state.facts, preflight_brief)
        is_ppt = self._is_ppt_task(effective_task)
        code_spec = None if is_ppt else compile_code_task(effective_task, config.sandbox_root())
        from .task_compiler import portable_contract_for
        portable = portable_contract_for(effective_task, self.state.facts, preflight_brief) if is_ppt else None
        self.state.execution_contract = compile_execution_contract(self.task_spec, is_ppt, code_spec, portable=portable)
        self.state.record_fact("task_spec", brief_json(self.task_spec))
        self.state.record_fact("selected_skill", self.task_spec.skill)
        self.state.record_fact("task_intent", self.task_spec.intent)
        self.state.record_fact("execution_capability", self.state.execution_contract.capability)
        if code_spec is not None:
            self.state.record_fact("code_language", code_spec.language)
            self.state.record_fact("code_test_runner", code_spec.runner or "none")
        if self.task_spec.primary_input:
            self.state.record_fact("primary_input", self.task_spec.primary_input)
        if self.task_spec.output_path:
            self.state.record_fact("output_path", self.task_spec.output_path)
        if getattr(self, "recorder", None) and callable(getattr(self.recorder, "event", None)):
            self.recorder.event("task_compiled", **self.task_spec.to_dict())
        instruction = self.state.facts.get("task_instruction", "")
        execution_task = effective_task + (("\n\nBound task instruction:\n" + instruction) if instruction else "")
        if is_ppt and self.task_spec.skill == "ppt.template_build":
            # One-shot template builds repeatedly stopped after the cover slide
            # or entered placeholder cleanup before composing any content.
            # Make the completeness contract an explicit part of the task.
            slide_count_hint = ""
            contract_text = self.state.facts.get("verification_contract", "")
            match = re.search(r"slide_count\s*=\s*(\d+)", contract_text)
            if match:
                slide_count_hint = f" Required final deck: {match.group(1)} finished content slides."
            else:
                match = re.search(r"min_slides\s*=\s*(\d+)[^\d]*max_slides\s*=\s*(\d+)", contract_text)
                if match:
                    slide_count_hint = f" Required final deck: {match.group(2)} finished content slides."
            execution_task += (
                "\n\nTemplate-build completion contract: use the supplied template as a scaffold and give every "
                "top-level section of the ContentIR source outline its own finished slide. Compose ALL content "
                f"slides first (one template slide per section, via ppt_compose content/from_outline) and only then "
                f"save/check/evaluate.{slide_count_hint} Do not stop after the cover slide, do not spend the first "
                "pass cleaning template placeholder text, and do not finish until the official evaluator passes."
            )
        if continuing_goal:
            remaining = [m for m in active_goal.milestones if m not in active_goal.completed_milestones]
            execution_task += "\n\nActive durable goal (resume): " + active_goal.objective
            if remaining:
                execution_task += "\nNext milestone: " + remaining[0]
        profile = self._set_task_profile(execution_task)
        if preflight_brief and self.state.phase in {RuntimePhase.INTAKE, RuntimePhase.UNDERSTAND}:
            previous_phase = self.state.phase
            self.state.transition(RuntimePhase.PRODUCE)
            self._publish_phase_change(previous_phase, "preflight ContentIR closed discovery")
        if not self.started:
            self.messages = [{"role": "system", "content": self._system_prompt(execution_task)}]
            self.started = True
        else:
            from .skills import discover, match
            newly_selected = [skill for skill in match(execution_task, discover()) if skill.name not in self.loaded_skills]
            if newly_selected:
                self.loaded_skills.update(skill.name for skill in newly_selected)
                restrictions = [set(skill.allowed_tools) for skill in newly_selected if skill.allowed_tools]
                if restrictions:
                    newly_allowed = set.intersection(*restrictions)
                    existing = getattr(self, "skill_allowed_tools", set())
                    self.skill_allowed_tools = existing & newly_allowed if existing else newly_allowed
                self.messages.append({"role": "system", "content": "Dynamically loaded skill instructions:\n" + "\n\n".join(skill.body for skill in newly_selected)})
        if preflight_brief:
            brief_message = (
                "Harness preflight completed. The following provenance-preserving ContentIR brief is the authoritative "
                "task/source context for this turn. Do not rediscover or reread these inputs. Produce the artifact now, "
                "then save and verify it."
                + (
                    " The unique input PPTX is already open as the active editable working copy; do not call ppt_open again."
                    if preflight_open else ""
                )
                + "\n" + preflight_brief
            )
            if not any(message.get("content") == brief_message for message in self.messages[-4:]):
                self.messages.append({"role": "system", "content": brief_message})
        # Route context is not an execution plan: the harness compiles only the
        # coarse artifact route (new deck vs edit existing) and the verification
        # contract. How to plan and sequence the actual work stays with the
        # model, inside the phase-scoped tool surface and gates.
        if is_ppt:
            route_message = (
                "Route context (not a fixed plan): "
                f"artifact_mode={self.task_spec.artifact_mode}, "
                f"input={self.task_spec.primary_input or 'none'}, "
                f"output={self.task_spec.output_path or 'contract output'}. "
                "Plan your own execution steps; the harness enforces only phase tools, "
                "verification freshness, and the delivery gates."
            )
            if not any(message.get("content") == route_message for message in self.messages[-4:]):
                self.messages.append({"role": "system", "content": route_message})
            # Whole-PPT-class tool guidance by capability (not a task-specific
            # plan): template projects must transform the already-open template
            # instead of replacing it with a blank deck or from_slides synthesis.
            skill_guidance = {
                "ppt.template_build": (
                    "Template capability note: the bound template deck is already open. "
                    "Rewrite its named boxes with ppt_edit_text set_shape_text (names come from ppt_inspect summary), "
                    "add extra content pages with ppt_compose kind='content' up to the slide-count contract, "
                    "and never call new_deck or from_slides for this project."
                ),
                "ppt.source_grounded_build": (
                    "Source-grounded capability note: prefer ppt_compose kind='quadrant' with slide_number=1 "
                    "(or ppt_edit_text set_shape_text on the open template boxes). Record anchor/metric/chart ids "
                    "per quadrant with ppt_metadata on the quadrant's provenance shape. Keep exactly one board slide."
                ),
                "ppt.source_sync": (
                    "Source-sync capability note: rewrite whole cards/tables with set_shape_text/set_table in "
                    "batch_updates; a replace that misses one stale fragment still fails the contract gate. "
                    "Do not reopen files or re-inspect unchanged slides once the per-slide counterexample list is known."
                ),
            }.get(self.task_spec.skill, "")
            if skill_guidance and not any(message.get("content") == skill_guidance for message in self.messages[-4:]):
                self.messages.append({"role": "system", "content": skill_guidance})
        # Compound replacement contract: instructions like
        # "Liability/Liabilities -> Debt/Debts" must be executed as one batch
        # with every variant, not as a single singular-only replace.
        compound_match = re.search(
            r"([A-Za-z]+)\s*/\s*([A-Za-z]+)\s*(?:->|→|to|改为)\s*([A-Za-z]+)\s*/\s*([A-Za-z]+)",
            execution_task,
            re.IGNORECASE,
        )
        if is_ppt and compound_match and self.task_spec.skill == "ppt.atomic_edit":
            old_a, old_b, new_a, new_b = compound_match.groups()
            compound_message = (
                "Deterministic replacement contract: use operation='replace_case_variants' once "
                f"(old='{old_a}' new='{new_a}' new_plural='{new_b}'), not two case-sensitive replaces. "
                "This replaces singular/plural and lowercase/Capitalized/UPPERCASE forms together."
            )
            if not any(message.get("content") == compound_message for message in self.messages[-4:]):
                self.messages.append({"role": "system", "content": compound_message})
        elif is_ppt and self.task_spec.skill == "ppt.atomic_edit":
            # Single capitalized term -> single replacement: case variants still
            # matter (Liability/liability/LIABILITIES -> Debt/debt/DEBTS).
            single_match = re.search(
                r"['\"]?([A-Z][A-Za-z]+)['\"]?\s*(?:替换为|改为|改成|→|->|to)\s*['\"]?([A-Z][A-Za-z]+)",
                execution_task,
                re.IGNORECASE,
            )
            if single_match:
                old_term, new_term = single_match.groups()
                single_message = (
                    "Deterministic replacement contract: use operation='replace_case_variants' once "
                    f"(old='{old_term}' new='{new_term}'), not one case-sensitive replace. "
                    "This replaces singular/plural and lowercase/Capitalized/UPPERCASE forms together."
                )
                if not any(message.get("content") == single_message for message in self.messages[-4:]):
                    self.messages.append({"role": "system", "content": single_message})
        # Give code tasks their deterministic language/test-runner context so the
        # model runs the actual tests (run_checks) instead of stopping at a
        # content-only verify_files assertion.
        if not is_ppt and code_spec is not None and code_spec.runner:
            code_message = (
                f"Code task context: language={code_spec.language}, test runner={code_spec.runner}. "
                "After changing files, run the tests with run_checks using this runner, then finish."
            )
            if not any(message.get("content") == code_message for message in self.messages[-4:]):
                self.messages.append({"role": "system", "content": code_message})
        self.messages.append({"role": "user", "content": task})
        self._maybe_compact(force=len(self.messages) > 12)
        self.state.goal = effective_task
        self._bind_ppt_mutation_scope(execution_task)
        if self.task_spec.mutation_slides and not self.state.ppt_allowed_slides:
            self.state.ppt_allowed_slides = set(self.task_spec.mutation_slides)
        ppt_task = self._is_ppt_task(execution_task)
        # Restored sessions and lightweight test harnesses may predate the
        # runtime field.  Runtime ownership is recoverable and lazily rebuilt.
        if not hasattr(self, "runtime"):
            self.runtime = RuntimeController()
        # Bounded repair budget is contract-owned, not a task-string heuristic:
        # the compiled ExecutionContract already says how many verifier-feedback
        # cycles the current capability may spend (1 for atomic edits, 3 for
        # source-sync/multi-surface work).
        contract = getattr(self.state, "execution_contract", None)
        self.state.max_repairs = max(1, int(getattr(contract, "max_repairs", 3) or 3))
        self._ensure_execution_plan()
        self._done = None
        # A model can satisfy the tool protocol while making no progress (for
        # example, issuing the same failed search/edit forever).  Keep a small
        # circuit breaker independent of the token budget so these stalls do
        # not consume the whole provider allocation.
        last_signature: str | None = None
        same_signature_streak = 0
        no_tool_streak = 0
        # Explicit turn limits belong to benchmark/non-interactive runs. The
        # interactive REPL continues until completion, user interruption, the
        # tool hard-cap, or a deterministic no-progress circuit breaker.
        policy_steps = int(getattr(getattr(self, "policy_guard", None), "spec", None).max_model_steps if getattr(getattr(self, "policy_guard", None), "spec", None) else 0) or self.max_steps
        step_limit = min(self.max_steps, policy_steps) if strict_budget else self.max_steps
        steps = range(step_limit) if strict_budget else itertools.count()
        for step_index in steps:
            if self.cancel_requested():
                return self._interrupt_stop()
            turn_total = self.state.total_tokens - turn_total_start
            turn_generated = self.state.generated_output_tokens - turn_generated_start
            if strict_budget and turn_total >= config.max_total_tokens():
                return self._budget_stop("本轮总预算已经用尽")
            remaining = config.max_total_tokens() - turn_total if strict_budget else config.max_output_tokens()
            generated_remaining = config.max_generated_output_tokens() - turn_generated if strict_budget else config.max_output_tokens()
            # Allocate output budget per turn from the remaining allowance.
            # This prevents a final response from requesting the full default
            # cap and being rejected after it has already consumed the budget.
            if strict_budget and (remaining < 256 or generated_remaining < 256):
                return self._budget_stop("剩余预算不足以安全生成下一次响应")
            if strict_budget and remaining <= config.max_output_tokens() * 2:
                budget_hint = "Budget is nearly exhausted. Do only the minimum final verification needed, then call finish; do not start new exploration."
                if not any(message.get("content") == budget_hint for message in self.messages[-3:] if message.get("role") == "user"):
                    self.messages.append({"role": "user", "content": budget_hint})
            prior_openai_cap = os.environ.get("OPENAI_MAX_TOKENS")
            prior_anthropic_cap = os.environ.get("ANTHROPIC_MAX_TOKENS")
            prior_thinking = os.environ.get("THINKING_ENABLED")
            awaiting_first_action = (
                ppt_task
                and self.state.phase == RuntimePhase.PRODUCE
                and not self.state.changed_files
                and not self.state.fresh_evidence()
            )
            action_cap = config.first_action_output_tokens() if awaiting_first_action else config.max_output_tokens()
            cap = str(min(config.max_output_tokens(), remaining, generated_remaining, action_cap))
            os.environ["OPENAI_MAX_TOKENS"] = cap
            os.environ["ANTHROPIC_MAX_TOKENS"] = cap
            # If a reasoning model consumes the first action window without a
            # tool call, resample once in execution mode. This separates the
            # user-selected reasoning quality from the controller's action
            # deadline and prevents hidden thinking from consuming the run.
            if awaiting_first_action and no_tool_streak:
                os.environ["THINKING_ENABLED"] = "0"
            previous_phase = self.state.phase
            meta_action = self.runtime.decide(
                self.state, ppt_task, getattr(self.state, "execution_contract", None),
                getattr(self, "controller_policy", "cegar_h"),
            )
            self._publish_phase_change(previous_phase, "controller recommendation")
            self.events.publish(
                EventKind.CONTROLLER_DECISION,
                action=meta_action,
                phase=self.state.phase.value,
                reason=self.state.last_meta_reason,
                profile=profile.name,
                capabilities=list(profile.capabilities),
            )
            self._publish_planning()
            recorder = getattr(self, "recorder", None)
            if recorder is not None and callable(getattr(recorder, "event", None)):
                recorder.event(
                    "controller_decision",
                    action=meta_action,
                    phase=self.state.phase.value,
                    reason=self.state.last_meta_reason,
                    profile=profile.name,
                )
            phase_tools = self.runtime.tool_names_for_phase(self.state, ppt_task, profile, execution_task)
            controller_hint = "CEGAR-H runtime decision: " + self.runtime.recommendation_text(meta_action)
            if not any(message.get("content") == controller_hint for message in self.messages[-2:]):
                self.messages.append({"role": "system", "content": controller_hint})
            advertised_tools = select_tools(
                execution_task,
                self.deck is not None,
                getattr(self, "skill_allowed_tools", set()),
                phase_tools,
            )
            advertised_tools = self._constrain_compiled_tool_schemas(advertised_tools)
            advertised_names = {
                tool["function"]["name"] for tool in advertised_tools
            }
            # Read-only compatibility aliases are accepted by dispatch, but
            # arbitrary hidden registered tools must never bypass phase/skill
            # routing merely because a provider hallucinated their names.
            admitted_names = set(advertised_names)
            if "read_file" in advertised_names:
                admitted_names.add("sys_cat")
            if {"read_file", "glob_files"} & advertised_names:
                admitted_names.add("sys_exec")
            try:
                import inspect
                chat_kwargs: dict[str, Any] = {}
                chat_params = inspect.signature(self.llm.chat).parameters
                if "stream" in chat_params:
                    chat_kwargs["stream"] = getattr(self, "stream_callback", None) is not None
                if "on_token" in chat_params:
                    chat_kwargs["on_token"] = getattr(self, "stream_callback", None)
                if "on_reasoning" in chat_params:
                    chat_kwargs["on_reasoning"] = getattr(self, "reasoning_callback", None)
                reply = self.llm.chat(self.messages, advertised_tools, **chat_kwargs)
            except Exception:
                if self.cancel_requested():
                    return self._interrupt_stop()
                raise
            finally:
                if prior_openai_cap is None:
                    os.environ.pop("OPENAI_MAX_TOKENS", None)
                else:
                    os.environ["OPENAI_MAX_TOKENS"] = prior_openai_cap
                if prior_anthropic_cap is None:
                    os.environ.pop("ANTHROPIC_MAX_TOKENS", None)
                else:
                    os.environ["ANTHROPIC_MAX_TOKENS"] = prior_anthropic_cap
                if prior_thinking is None:
                    os.environ.pop("THINKING_ENABLED", None)
                else:
                    os.environ["THINKING_ENABLED"] = prior_thinking
            # Esc can arrive while the provider request is in flight.  Never
            # execute tool calls from a response that completed after cancel.
            if self.cancel_requested():
                return self._interrupt_stop()
            self.state.total_tokens += reply.total_tokens
            self.state.input_tokens += reply.input_tokens
            self.state.generated_output_tokens += reply.output_tokens
            self.state.last_reasoning_chars = reply.reasoning_chars
            self.state.reasoning_chars += reply.reasoning_chars
            self.state.reasoning_observed = self.state.reasoning_observed or reply.reasoning_chars > 0
            self.state.provider_usage_authoritative = self.state.provider_usage_authoritative and reply.usage_authoritative
            reasoning_text = getattr(reply.choices[0].message, "reasoning_content", None) or ""
            self.state.last_reasoning_text = reasoning_text
            self.events.publish(
                EventKind.MODEL_RESPONSE,
                phase=self.state.phase.value,
                has_tool_calls=bool(getattr(reply.choices[0].message, "tool_calls", None)),
                tool_call_count=len(getattr(reply.choices[0].message, "tool_calls", None) or []),
                reasoning_chars=reply.reasoning_chars,
                reasoning_content=reasoning_text,
                output_chars=len(getattr(reply.choices[0].message, "content", "") or ""),
            )
            if recorder is not None and callable(getattr(recorder, "event", None)):
                recorder.event(
                    "model_response",
                    phase=self.state.phase.value,
                    has_tool_calls=bool(getattr(reply.choices[0].message, "tool_calls", None)),
                    tool_call_count=len(getattr(reply.choices[0].message, "tool_calls", None) or []),
                    reasoning_chars=reply.reasoning_chars,
                    output_chars=len(getattr(reply.choices[0].message, "content", "") or ""),
                )
            # A provider may return more tokens than the remaining allowance in
            # one response.  Do not execute any tool calls from an over-budget
            # response; otherwise the harness would violate its own cost gate.
            if strict_budget and ((self.state.total_tokens - turn_total_start) > config.max_total_tokens()
                    or (self.state.generated_output_tokens - turn_generated_start) > config.max_generated_output_tokens()):
                incomplete = [item for item in self.state.tasks if item.status.value in {"pending", "in_progress"}]
                if self.state.fresh_evidence() and not incomplete:
                    self.state.budget_overrun = True
                    self.state.final_summary = "(completed with budget overrun after fresh verification evidence)"
                    return self.state.final_summary
                self._maybe_compact(force=True)
                return self._budget_stop("模型本次响应超过了本轮预算", overrun=True)
            msg = reply.choices[0].message
            rejected = [
                tc.function.name for tc in (getattr(msg, "tool_calls", None) or [])
                if tc.function.name not in admitted_names
            ]
            if rejected:
                # Observation closure is a deterministic loop transition, not a
                # stall: when ppt_inspect has already delivered its bounded
                # evidence and the model asks for more, give ONE explicit
                # action-pass nudge before applying the usual rejection fuse.
                inspect_closed = "ppt_inspect" in rejected and int(self.state.facts.get("ppt_inspect_count", "0")) >= 2
                closure_nudge = int(self.state.facts.get("ppt_observation_closure_nudge", "0"))
                reopen_rejected = "ppt_open" in rejected and getattr(self.state, "mutation_epoch", 0) > 0
                reopen_nudge = int(self.state.facts.get("ppt_open_rejection_nudge", "0"))
                if reopen_rejected and reopen_nudge == 0:
                    self.state.facts["ppt_open_rejection_nudge"] = "1"
                    self.state.no_progress_streak = 0
                    if recorder is not None:
                        recorder.event(
                            "tool_calls_rejected",
                            rejected=rejected,
                            advertised=sorted(advertised_names),
                            reason="progress-preservation: reopening a file is unnecessary",
                        )
                    self.messages.append({"role": "assistant", "content": msg.content or ""})
                    self.messages.append({
                        "role": "user",
                        "content": (
                            "ppt_open was rejected because the ACTIVE deck already contains all current edits; "
                            "reopening any file would risk discarding progress. Do not request file opens. "
                            "Retry the intended edit/save/check sequence on the active deck, or proceed with "
                            "set_shape_text / set_table / batch_updates directly."
                        ),
                    })
                    continue
                if inspect_closed and closure_nudge == 0:
                    self.state.facts["ppt_observation_closure_nudge"] = "1"
                    self.state.no_progress_streak = 0
                    if recorder is not None:
                        recorder.event(
                            "tool_calls_rejected",
                            rejected=rejected,
                            advertised=sorted(advertised_names),
                            reason="ppt observation closed; bounded action pass",
                        )
                    self.messages.append({"role": "assistant", "content": msg.content or ""})
                    if self.state.phase in {RuntimePhase.INTAKE, RuntimePhase.UNDERSTAND}:
                        previous_phase = self.state.phase
                        self.state.transition(RuntimePhase.PRODUCE)
                        self._publish_phase_change(previous_phase, "ppt observation closed; forced production")
                    self.messages.append({
                        "role": "user",
                        "content": (
                            "PPT observation is closed: you already have the full-deck summary with stable shape "
                            "names and table contents, and at most one targeted shapes view. Do not request "
                            "ppt_inspect again. Proceed with the mutation now: use set_shape_text / set_table with "
                            "the shape_name values from the summary (or batch_updates for pure string replacements), "
                            "then save and verify."
                        ),
                    })
                    continue
                self.state.no_progress_streak += 1
                if recorder is not None:
                    recorder.event(
                        "tool_calls_rejected",
                        rejected=rejected,
                        advertised=sorted(advertised_names),
                    )
                self.messages.append({"role": "assistant", "content": msg.content or ""})
                self.messages.append({
                    "role": "user",
                    "content": (
                        "The requested tool name is not available in this phase. Nothing was executed. "
                        f"Use exactly one of: {', '.join(sorted(advertised_names))}."
                    ),
                })
                if self.state.no_progress_streak >= 2:
                    saved = self._save_draft_before_pause()
                    saved_note = f"\n\n暂停前已保存当前草稿：{saved}。" if saved else ""
                    return f"已安全暂停：模型连续请求本阶段未开放的工具，控制器已阻止执行。{saved_note}"
                continue
            if not getattr(msg, "tool_calls", None):
                answer = msg.content or ""
                incomplete = [item for item in self.state.tasks if item.status.value in {"pending", "in_progress"}]
                # Some OpenAI-compatible reasoning models occasionally return
                # the task verbatim after a long thinking pass instead of
                # emitting the requested tool call. An action task must not be
                # reported as complete without any tool/evidence/artifact.
                action_task = self._requires_action(execution_task)
                lacks_completion = action_task and not self._has_completion_evidence(execution_task)
                if lacks_completion:
                    no_tool_streak += 1
                    no_tool_limit = (
                        5 if self.state.facts.get("official_evaluator_present") == "true" else 2
                    )
                    if no_tool_streak <= no_tool_limit:
                        self.messages.append({"role": "assistant", "content": answer})
                        self.messages.append({
                            "role": "user",
                            "content": (
                                "This action task is not complete: required mutation, saved artifact, or fresh verification evidence is missing. "
                                "The execution scheduler is now using a bounded action pass. Emit the next concrete tool call immediately; "
                                "do not spend this pass on further internal deliberation, repeat, summarize, acknowledge, or end the task. "
                                "For a local PPT edit, follow inspect-once -> mutate -> save -> verify -> finish."
                            ),
                        })
                        continue
                    return (
                        "⚠ 模型连续返回文本但没有调用工具，任务尚未开始。\n\n"
                        "小朴没有把任务回显误判为完成。请重试一次，或将 /effort 切换为 high 后继续。"
                    )
                if self.state.changed_files and not self.state.fresh_evidence():
                    self.messages.append({"role": "assistant", "content": answer})
                    self.messages.append({"role": "user", "content": "You changed files but provided no verification evidence. Run the relevant checks, inspect their output, then finish."})
                    continue
                if incomplete:
                    self.messages.append({"role": "assistant", "content": answer})
                    self.messages.append({"role": "user", "content": "The task list still contains unfinished items. Continue working or mark genuinely blocked items with evidence."})
                    continue
                self.state.final_summary = answer
                return answer

            msg.tool_calls = self._coalesce_atomic_inspect_batch(list(msg.tool_calls))

            used_this_turn = control.used_tools(self.state)
            requested_calls = len(msg.tool_calls)
            if not control.admits(self.state, requested_calls):
                assert control.tool_limit is not None
                if getattr(self, "recorder", None):
                    self.recorder.event(
                        "tool_budget_exhausted",
                        used=used_this_turn,
                        limit=control.tool_limit,
                        rejected_tools=[tc.function.name for tc in msg.tool_calls],
                    )
                return self._tool_budget_stop(used_this_turn, control.tool_limit)

            # Append the assistant tool-use turn only after admission. This
            # keeps continuation history valid when a whole batch is rejected
            # by the hard budget gate (no dangling tool_call ids).
            self.messages.append(
                {
                    "role": "assistant",
                    "content": msg.content or "",
                    "tool_calls": [
                        {
                            "id": tc.id,
                            "type": "function",
                            "function": {
                                "name": tc.function.name,
                                "arguments": tc.function.arguments,
                            },
                        }
                        for tc in msg.tool_calls
                    ],
                    # DeepSeek requires this private field to be returned on
                    # subsequent tool-use calls. It is never sent to the UI.
                    **({"reasoning_content": msg.reasoning_content} if msg.reasoning_content else {}),
                }
            )

            results = self._execute_calls(msg.tool_calls)
            if self.cancel_requested():
                return self._interrupt_stop()
            # Obligation-based progress: fresh evidence with unchanged
            # unresolved obligations is not progress and must not reset the
            # stall counter. A changed obligation *set* (e.g. the official
            # evaluator just discovered a new blocker) opens a new bounded
            # CEGAR-H repair iteration and restarts the stall counter, so the
            # loop repairs the counterexample instead of pausing on discovery.
            #
            # `controller_redirects` is a per-iteration repair budget, not a
            # lifetime cap: every productive transition (epoch advance or new
            # obligation set) re-arms the bounded passes for the new iteration.
            if self._obligation_progress():
                self.state.no_progress_streak = 0
                self.state.controller_redirects = 0
                self.state.facts.pop("ppt_repair_observation_calls", None)
                self.state.facts.pop("repair_observation_grant", None)
            else:
                checkpoint = getattr(self, "_obligation_checkpoint", None)
                now = self._obligation_snapshot()
                if checkpoint is not None and set(now[0]) != set(checkpoint[0]):
                    self._obligation_checkpoint = now
                    self.state.no_progress_streak = 0
                    self.state.controller_redirects = 0
                    # A brand-new counterexample set re-arms the bounded
                    # observation budget for its own repair iteration.
                    self.state.facts.pop("ppt_repair_observation_calls", None)
                    self.state.facts.pop("repair_observation_grant", None)
                else:
                    self.state.no_progress_streak += 1
            no_tool_streak = 0
            signatures = [f"{tc.function.name}:{tc.function.arguments}" for tc in msg.tool_calls]
            signature = "|".join(signatures)
            if signature == last_signature:
                same_signature_streak += 1
            else:
                last_signature = signature
                same_signature_streak = 1
            for tc in msg.tool_calls:
                text = results[tc.id]
                self.on_tool(tc.function.name, tc.function.arguments, text)
                # Verifier counterexamples and finish rejections are repair
                # evidence; they get a larger model-visible budget than ordinary
                # tool output so the complete blocker manifest can be planned
                # against, not just its first fragment.
                if tc.function.name == "run_task_evaluator":
                    visible_limit = 22000
                elif tc.function.name == "finish":
                    visible_limit = 12000
                elif tc.function.name == "ppt_inspect":
                    visible_limit = 14000
                elif tc.function.name == "ppt_check":
                    # A failed contract gate is a full per-slide repair manifest;
                    # truncating it hides the slides the model has not fixed yet.
                    visible_limit = 20000 if "FAILED" in text else 8000
                else:
                    visible_limit = 14000 if tc.function.name == "read_many" else 5000
                self.messages.append({"role": "tool", "tool_call_id": tc.id, "content": bounded_tool_result(text, visible_limit)})
            # Observation closure is a transition, not a dead end. If the model
            # answered the closure error with MORE observation calls, replace
            # the next planning pass with an explicit action directive so it
            # emits a mutation from evidence it already holds.
            if (
                ppt_task
                and msg.tool_calls
                and all("observation closed" in results.get(tc.id, "") for tc in msg.tool_calls)
                and self.state.unresolved_checks
            ):
                blockers = sorted(self.state.unresolved_checks)
                self.messages.append({
                    "role": "user",
                    "content": (
                        "Observation stays closed; those observation calls did not execute and will not "
                        "execute on the next attempt. The immediately preceding ppt_inspect / ppt_check / "
                        "run_task_evaluator outputs already contain the CURRENT shape names, table contents "
                        "and exact old/new strings for every failed check. Do not reason about inspecting "
                        f"again. Emit one concrete ppt_edit_text (set_shape_text / set_table / batch_updates) "
                        f"or ppt_metadata call now to repair: {', '.join(blockers)}. Then ppt_save and rerun "
                        "the failed verifier once."
                    ),
                })
            if self._done:
                return self._done
            # Read-only benchmark tasks need evidence, not an open-ended agent
            # loop. After two targeted inspection calls, close deterministically
            # so prompt exploration cannot consume the entire provider budget.
            if read_only_task and self.state.tool_calls >= 2:
                evidence = "\n".join(results.values())
                self.state.final_summary = "Read-only inspection complete. Evidence:\n" + evidence[:6000]
                return self.state.final_summary
            if ppt_task and self.state.no_progress_streak >= 2:
                self.state.controller_redirects += 1
                # Preflight ContentIR can advance the phase to PRODUCE before the
                # first model turn, so a stuck observer is not always in
                # INTAKE/UNDERSTAND. Force a bounded action pass whenever the
                # run has produced nothing yet and is only re-reading unchanged
                # inputs, regardless of the current phase label.
                #
                # The same bounded-pass rule applies after edits: if the model
                # keeps issuing no-op equivalent edits while a saved, freshly
                # verified draft already exists, the shortest path is to run the
                # official evaluator and finish instead of pausing mid-loop.
                official_eval = self.state.facts.get("official_evaluator_present") == "true"
                redirect_limit = 8 if official_eval else 2
                can_redirect = self.state.controller_redirects <= redirect_limit and (
                    self.state.mutation_epoch == 0
                    or (self.state.fresh_evidence() and self.state.mutation_epoch > 0)
                )
                if can_redirect:
                    previous_phase = self.state.phase
                    if self.state.mutation_epoch == 0:
                        self.state.transition(RuntimePhase.PRODUCE)
                        self._publish_phase_change(previous_phase, "no-progress redirect to production")
                        message = (
                            "CEGAR-H detected repeated observations with no new information. Observation is now closed. "
                            "Do not list directories or reread unchanged files. Build or modify the smallest valid artifact, save it, then verify it."
                        )
                    else:
                        self.state.transition(RuntimePhase.VERIFY)
                        self._publish_phase_change(previous_phase, "no-progress redirect to verification")
                        blockers = sorted(self.state.unresolved_checks)
                        if blockers:
                            message = (
                                "CEGAR-H detected repeated equivalent observations/edits with no new information, "
                                f"but these verifier obligations are still unresolved: {', '.join(blockers)}. "
                                "Do not reopen files or re-inspect unchanged slides. Repair ONLY the cited slides/shapes "
                                "with ppt_edit_text (set_shape_text/set_table/batch_updates) or ppt_metadata, then "
                                "ppt_save and rerun ppt_check. finish stays rejected until the blockers clear."
                            )
                        else:
                            message = (
                                "CEGAR-H detected repeated equivalent edits with no new information. "
                                "A saved draft with fresh structural evidence already exists. "
                                "Do not issue another equivalent edit. Run ppt_check if evidence is stale, then call finish "
                                "(the official task evaluator will run automatically and return concrete counterexamples if anything is still wrong)."
                            )
                    self.messages.append({"role": "user", "content": message})
                    self._maybe_compact(force=True)
                    continue
                saved = self._save_draft_before_pause()
                saved_note = f"\n\n暂停前已保存当前草稿：{saved}。" if saved else ""
                return f"已安全暂停：连续三次执行相同操作或等价观察但没有获得新信息，控制器已阻止继续空转。{saved_note}"
            if same_signature_streak >= 3 and self.state.no_progress_streak >= 3:
                saved = self._save_draft_before_pause()
                saved_note = f"\n\n暂停前已保存当前草稿：{saved}。" if saved else ""
                return f"⚠ 已安全暂停：连续三次执行相同操作且 obligations 无进展。{saved_note}\n\n会话与工作状态已保留；请检查工具参数，或输入“继续”让我换一种路径。"
            self._maybe_compact()
            if not strict_budget and (step_index + 1) % self.max_steps == 0:
                self._maybe_compact(force=True)
                self.messages.append({
                    "role": "user",
                    "content": (
                        "Interactive execution checkpoint reached. Continue automatically from the compacted state. "
                        "Do not restart or re-inspect completed work. Prioritize the shortest remaining path to save, verify, and finish."
                    ),
                })

        return "⚠ 评测运行已达到预设最大执行步数，尚未形成最终答案。"

    def _is_ppt_task(self, task: str) -> bool:
        text = task.lower()
        # A deck being open is not, by itself, enough: an interactive user can
        # follow a PPT task with a plain chat question and must not inherit the
        # stale evaluator/output contract. Bare continuations of a deck session
        # still count as PPT.
        if getattr(self, "deck", None) is not None and is_scope_continuation(task):
            return True
        # One-line benchmark requests intentionally contain only a task path.
        # Domain classification must use deterministic workspace facts rather
        # than require the user to repeat "PPT" in every prompt.
        try:
            from .intake import task_root_from_prompt
            task_root = task_root_from_prompt(task)
        except Exception:
            task_root = None
        if task_root is not None:
            input_root = task_root / "input"
            has_pptx = any(task_root.glob("*.pptx"))
            if input_root.is_dir():
                has_pptx = has_pptx or any(input_root.rglob("*.pptx"))
            if has_pptx:
                return True
            card = task_root / "task_card.md"
            if card.is_file():
                card_text = card.read_text(encoding="utf-8-sig", errors="replace").casefold()
                if "source: ppt-eval" in card_text or "ppt" in card_text or "powerpoint" in card_text:
                    return True
        if any(marker in text for marker in ("ppt", "pptx", "powerpoint", "slide", "deck", "演示", "幻灯片", "排版")):
            return True
        office_edit = any(marker in text for marker in ("页", "标题", "项目符号", "文本框", "流程图", "字体", "字号", "图片"))
        action = any(marker in text for marker in ("修改", "新增", "添加", "替换", "调整", "生成", "创建", "改成", "改为", "换成", "变为", "改"))
        return office_edit and action

    @staticmethod
    def _requires_action(task: str) -> bool:
        text = task.lower()
        markers = (
            "create", "edit", "modify", "add", "resize", "save", "render", "execute", "complete",
            "创建", "修改", "新增", "调整", "保存", "渲染", "执行", "完成", "生成", "排版",
        )
        return any(marker in text for marker in markers)

    @staticmethod
    def _is_local_ppt_edit(task: str) -> bool:
        text = task.lower()
        ppt = any(marker in text for marker in ("ppt", "pptx", "powerpoint", "slide", "deck", "幻灯片", "演示文稿"))
        edit = any(marker in text for marker in ("edit", "modify", "add", "resize", "修改", "新增", "调整", "排版"))
        existing = any(marker in text for marker in (".pptx", "existing", "第", "slide"))
        return ppt and edit and existing

    def _has_completion_evidence(self, task: str) -> bool:
        """Completion is an artifact/evidence state, never merely model prose."""
        if self._is_local_ppt_edit(task):
            artifacts = getattr(getattr(self, "recorder", None), "manifest", {}).get("artifacts", [])
            has_pptx = any(item.get("role") == "final-pptx" for item in artifacts)
            kinds = {record.kind for record in self.state.fresh_evidence()}
            # Structural evidence is universal. Render evidence remains
            # conditional on the task contract because some benchmark tasks
            # intentionally specify structural-only evaluation.
            required = {"ppt_structural"}
            lowered = task.lower()
            if any(marker in lowered for marker in ("render", "渲染", "screenshot", "截图", "视觉")):
                required.add("ppt_render")
            return has_pptx and required.issubset(kinds)
        return bool(self.state.changed_files and self.state.fresh_evidence())

    def _save_draft_before_pause(self) -> str:
        """Deliverability invariant: any safety pause keeps an unsaved draft.

        Returns the saved path (or "" when there is nothing to persist). This
        fixes the quick5 failure mode where the circuit breaker stopped a run
        after ``ppt_compose`` but before ``ppt_save``, leaving no output file.
        """
        deck = getattr(self, "deck", None)
        if deck is None or not getattr(self.state, "mutation_epoch", 0):
            return ""
        try:
            from .tools.registry import dispatch

            args: dict = {}
            required = self.state.facts.get("required_output_pptx")
            if required:
                args["path"] = required
            dispatch("ppt_save", json.dumps(args), self)
            return str(required or self.state.facts.get("required_output_pptx", "output/final.pptx"))
        except Exception:
            return ""

    def _budget_stop(self, reason: str, overrun: bool = False) -> str:
        if overrun:
            self.state.budget_overrun = True
        saved = self._save_draft_before_pause()
        saved_note = f"\n\n暂停前已保存当前草稿：{saved}。" if saved else ""
        return (
            f"⚠ 本轮已安全暂停：{reason}。{saved_note}\n\n"
            "未执行超预算响应中携带的工具调用，已有文件、验证证据和任务状态均已保留。"
            "上下文已经压缩；输入“继续”即可从当前状态续接，也可以用 /new 开始新会话。"
        )

    def _tool_budget_stop(self, used: int, limit: int) -> str:
        """Terminal result for a hard tool limit; never feed it back to the model."""
        self._maybe_compact(force=True)
        saved = self._save_draft_before_pause()
        saved_note = f"\n\n暂停前已保存当前草稿：{saved}。" if saved else ""
        return (
            "⚠ 已安全暂停：本轮工具调用额度已用完。\n\n"
            f"本轮已执行 {used}/{limit} 次工具调用。后续工具请求没有执行，也不会继续重试。"
            f"{saved_note}已有文件、任务状态和执行轨迹均已保留；输入“继续”可从当前进度续接。"
        )

    def _interrupt_stop(self) -> str:
        self._maybe_compact(force=True)
        return (
            "⏹ 已按用户要求中止当前任务。\n\n"
            "中止后返回的模型工具请求没有执行；已有文件、验证证据和执行轨迹均已保留。"
        )

    def _execute_calls(self, calls) -> dict[str, str]:
        self.state.tool_calls += len(calls)

        def execute(tc) -> str:
            from .controller_policies import MUTATION_TOOLS as _MUTATION_TOOLS

            arguments = tc.function.arguments
            signature = canonical_call(tc.function.name, arguments)
            try:
                if self.cancel_requested():
                    return "CANCELLED: user interrupted the current task before this tool ran"
                # Repair-phase observation closure: once verifier counterexamples
                # exist, observations are budgeted per repair iteration. A
                # contract/evaluator counterexample usually names text that no
                # longer matches the live deck, so the model gets up to two
                # targeted observations to re-sync before the closure forces a
                # mutation. A failed text replacement additionally grants one
                # fresh observation (the edit proved the cached text is stale).
                if (
                    tc.function.name in {"ppt_inspect", "ppt_open"}
                    and self.state.unresolved_checks
                    and not self.state.facts.get("repair_observation_grant") == "1"
                ):
                    blockers = set(self.state.unresolved_checks)
                    budget = 2 if blockers & {"ppt_contract", "task_evaluator"} else 1
                    used = int(self.state.facts.get("ppt_repair_observation_calls", "0"))
                    if used >= budget:
                        # If the deck has unverified mutations, the shortest
                        # legal path is save+check, not more observations and
                        # not more edits. Resolve the gate deterministically
                        # here so the model returns to concrete repairs.
                        if self._mutation_gated("ppt_edit_text"):
                            try:
                                from .tools.registry import dispatch as _lifecycle_dispatch
                                _lifecycle_dispatch("ppt_save", json.dumps({}), self)
                                _lifecycle_dispatch("ppt_check", json.dumps({"policy": "auto"}), self)
                            except Exception as lifecycle_exc:
                                return (
                                    "TOOL ERROR (RuntimeError): observation closed AND the automatic "
                                    f"save/check also failed ({type(lifecycle_exc).__name__}: {lifecycle_exc}). "
                                    "Fix the reported structural findings before further edits."
                                )
                            if not self._mutation_gated("ppt_edit_text"):
                                self.state.facts.pop("ppt_repair_observation_calls", None)
                                return (
                                    "verify-before-continue gate resolved by the harness: the current draft was "
                                    "saved and ppt_check passed with fresh structural evidence. Now apply the "
                                    "cited contract repairs with ppt_edit_text / ppt_metadata / ppt_arrange."
                                )
                        return (
                            "TOOL ERROR (RuntimeError): observation closed while repair obligations remain "
                            f"({', '.join(sorted(blockers))}). You already used {used}/{budget} targeted "
                            "observations for this repair iteration. Apply the cited repairs now with "
                            "ppt_edit_text / ppt_metadata / ppt_arrange using the exact strings from the "
                            "most recent ppt_check / run_task_evaluator output above; do not inspect or "
                            "reopen again until after a mutation changes the deck."
                        )
                    self.state.record_fact("ppt_repair_observation_calls", str(used + 1))
                elif tc.function.name in {"ppt_inspect", "ppt_open"} and self.state.facts.get("repair_observation_grant") == "1":
                    # One-shot re-observation grant issued after a concrete
                    # mutation failed because its target text no longer exists.
                    self.state.facts.pop("repair_observation_grant", None)
                if self._mutation_gated(tc.function.name):
                    # The verify-before-continue gate is a loop invariant, but
                    # commit + verification are harness-owned lifecycle steps.
                    # Satisfy the gate deterministically when the model forgot
                    # the save/check pair, so a valid repair is never stranded
                    # behind a prompt-only reminder.
                    try:
                        from .tools.registry import dispatch as _lifecycle_dispatch
                        _lifecycle_dispatch("ppt_save", json.dumps({}), self)
                        _lifecycle_dispatch("ppt_check", json.dumps({"policy": "auto"}), self)
                    except Exception as lifecycle_exc:
                        return (
                            "TOOL ERROR (RuntimeError): CEGAR-H verify-before-continue gate. "
                            f"Automatic save/check also failed ({type(lifecycle_exc).__name__}: {lifecycle_exc}). "
                            "Fix the reported structural findings before further content edits."
                        )
                    if not self._mutation_gated(tc.function.name):
                        return (
                            "verify-before-continue gate resolved by the harness: the current draft was saved and "
                            "ppt_check passed with fresh structural evidence. Retry your intended mutation now."
                        )
                    return (
                        "TOOL ERROR (RuntimeError): CEGAR-H verify-before-continue gate. "
                        "The deck has unverified mutations; run ppt_save then ppt_check "
                        "and obtain fresh structural evidence before further content edits."
                    )
                if not hasattr(self, "policy_guard"):
                    from .controller_policies import PolicyGuard
                    self.controller_policy = getattr(self, "controller_policy", "cegar_h")
                    self.policy_guard = PolicyGuard(self.controller_policy)
                self.policy_guard.before_tool(tc.function.name, self.state)
                for hook in self.pre_tool_hooks:
                    arguments = hook(ToolEvent(tc.function.name, arguments))
                    if not isinstance(arguments, str):
                        raise TypeError("pre-tool hooks must return JSON argument text")
                tc.function.arguments = arguments
                fingerprint = f"{tc.function.name}:{arguments}"
                if getattr(self, "recorder", None):
                    self.recorder.event("tool_started", call_id=tc.id, tool=tc.function.name, arguments=arguments[:4000])
                self.events.publish(EventKind.TOOL_STARTED, call_id=tc.id, tool=tc.function.name, arguments=arguments)
                cached = self.runtime.cache.get(signature) if tc.function.name in OBSERVE_TOOLS else None
                if cached is not None:
                    text = cached
                    if getattr(self, "recorder", None):
                        self.recorder.event("tool_cache_hit", call_id=tc.id, tool=tc.function.name, signature=signature[:1000])
                else:
                    out = dispatch(tc.function.name, arguments, self)
                    text = out if isinstance(out, str) else str(out)
                    if tc.function.name in OBSERVE_TOOLS:
                        self.runtime.cache.put(signature, text)
                for hook in self.post_tool_hooks:
                    text = hook(ToolEvent(tc.function.name, arguments, text))
                    if not isinstance(text, str):
                        raise TypeError("post-tool hooks must return text")
                self.state.failures.pop(fingerprint, None)
                previous_phase = self.state.phase
                novel = self.runtime.note_tool_result(self.state, tc.function.name, arguments, text)
                self._advance_execution_plan(tc.function.name)
                # A real mutation re-arms the observation budget: the deck
                # changed, so the next repair iteration may re-observe scope.
                if tc.function.name in _MUTATION_TOOLS:
                    self.state.facts.pop("ppt_repair_observation_calls", None)
                    self.state.facts.pop("repair_observation_grant", None)
                self._publish_phase_change(previous_phase, f"tool result: {tc.function.name}")
                if getattr(self, "recorder", None):
                    self.recorder.event("tool_completed", call_id=tc.id, tool=tc.function.name, novel=novel, output=text)
                self.events.publish(EventKind.TOOL_COMPLETED, call_id=tc.id, tool=tc.function.name, novel=novel, output=text)
                return text
            except Exception as exc:
                fingerprint = f"{tc.function.name}:{arguments}"
                count = self.state.failures.get(fingerprint, 0) + 1
                self.state.failures[fingerprint] = count
                failure_class = f"class:{tc.function.name}:{type(exc).__name__}"
                class_count = self.state.failures.get(failure_class, 0) + 1
                self.state.failures[failure_class] = class_count
                # Stall accounting is owned by the obligation-based progress
                # monitor after each executed turn; do not double-count here.
                if class_count == 1 and self.state.operational_plan:
                    self.events.publish(
                        EventKind.PROGRESS_UPDATED,
                        tool=tc.function.name,
                        items=list(self.state.operational_plan),
                        note=(
                            f"{tc.function.name} 遇到 {type(exc).__name__}；已有进度保留，"
                            "控制器将根据错误证据调整下一步。"
                        ),
                    )
                hint = " Do not repeat this identical call." if count >= 2 else ""
                # A rejected terminal action must reopen verification instead
                # of stranding the loop in DELIVER with only "finish" admitted;
                # otherwise the counterexample the rejection carries can never
                # be repaired.
                if tc.function.name in TERMINAL_TOOLS and not getattr(self, "_done", None):
                    if self.state.phase in {RuntimePhase.DELIVER, RuntimePhase.STOPPED}:
                        previous_phase = self.state.phase
                        self.state.transition(RuntimePhase.VERIFY)
                        self._publish_phase_change(previous_phase, f"rejected {tc.function.name}: verification reopened")
                    hint += " Verification phase was reopened; use the available repair and verification tools."
                known_deck = self.state.facts.get("ppt_input_deck", "")
                if tc.function.name == "ppt_open" and known_deck:
                    hint += f" The harness-discovered input is '{known_deck}'; do not guess another filename."
                if tc.function.name in _MUTATION_TOOLS and "text not found" in str(exc).casefold():
                    # The concrete edit proved the cached observation is stale.
                    # Re-open exactly one targeted observation instead of
                    # trapping the model in a repair-without-sight loop.
                    self.state.record_fact("repair_observation_grant", "1")
                    hint += (
                        " A one-shot observation grant was issued: call ppt_inspect once for the cited "
                        "slide to read the CURRENT text, then retry the edit with those exact strings."
                    )
                if class_count >= 3:
                    hint += (
                        f" Circuit breaker: {class_count} failures of this tool/error class; "
                        "change strategy instead of trying an equivalent call."
                    )
                    if getattr(self, "deck", None) is not None and getattr(self.state, "mutation_epoch", 0):
                        blockers = sorted(self.state.unresolved_checks)
                        if blockers:
                            hint += (
                                " A draft is already in memory; save it, run ppt_check/run_task_evaluator, "
                                "then repair ONLY the cited counterexamples. finish stays rejected while "
                                f"{', '.join(blockers)} remain unresolved."
                            )
                        else:
                            hint += " A draft is already in memory; the shortest valid path is ppt_save, then ppt_check, then finish."
                # Safety fuse (generalized rejection signature): the same
                # (action, error class, blocker set) three times in a row with
                # no obligation progress means the run is stuck on the same
                # counterexample, regardless of which tool is retrying.
                blockers = frozenset(self.state.unresolved_checks)
                signature = self._rejection_signature(tc.function.name, exc, blockers)
                history = getattr(self, "_rejection_history", None)
                if history is None:
                    history = self._rejection_history = []
                history.append(signature)
                del history[:-3]
                stuck = len(history) == 3 and history[0] == history[1] == history[2]
                if stuck:
                    saved = self._save_draft_before_pause()
                    blocker_text = ", ".join(sorted(blockers)) or str(exc)[:120]
                    self.state.record_fact("blocking_obligations", blocker_text)
                    self.state.record_fact("run_status", "paused_unresolved")
                    self._done = (
                        f"⚠ 已进入 STUCK（安全暂停）：同一 rejection signature 连续出现 3 次\n"
                        f"   action={tc.function.name} · error={type(exc).__name__} · blockers=[{blocker_text}]\n"
                        + (f"\n暂停前已保存当前草稿：{saved}。" if saved else "")
                        + "\n\n未解决 obligations 已记录，模型调用已停止；输入“继续”或定向修复 blocker 后续接。"
                    )
                if getattr(self, "recorder", None):
                    self.recorder.event("tool_failed", call_id=tc.id, tool=tc.function.name, error_type=type(exc).__name__, error=str(exc))
                self.events.publish(EventKind.TOOL_FAILED, call_id=tc.id, tool=tc.function.name, error_type=type(exc).__name__, error=str(exc))
                return f"TOOL ERROR ({type(exc).__name__}): {exc}.{hint}"

        parallel = len(calls) > 1 and all(tc.function.name in _READ_ONLY for tc in calls)
        if not parallel:
            return {tc.id: execute(tc) for tc in calls}
        outputs: dict[str, str] = {}
        with ThreadPoolExecutor(max_workers=min(4, len(calls))) as pool:
            futures = {pool.submit(execute, tc): tc.id for tc in calls}
            for future in as_completed(futures):
                outputs[futures[future]] = future.result()
        return outputs

    def _maybe_compact(self, force: bool = False) -> None:
        total = sum(len(json.dumps(message, ensure_ascii=False, default=str)) for message in self.messages)
        # Compact before the history becomes provider-expensive.  The old
        # 48k-character threshold allowed a long tool transcript to consume
        # most of a task's total-token budget before the next useful action.
        if force or total > 20000:
            head = self.messages[:1]
            last_tool_turn = next(
                (index for index in range(len(self.messages) - 1, 0, -1) if self.messages[index].get("role") == "assistant" and self.messages[index].get("tool_calls")),
                None,
            )
            tail = self.messages[last_tool_turn:] if last_tool_turn is not None else self.messages[-4:]
            self.messages = head + [
                {
                    "role": "system",
                    "content": "Long history compacted. Durable working state follows:\n"
                    + self.state.compact()
                    + "\nUse durable facts and cached observations. Do not reread unchanged source artifacts; acquire only specifically missing information.",
                },
                {"role": "user", "content": "Continue the active task from the durable state and most recent tool results below."},
            ] + tail

    @staticmethod
    def _infer_stop_reason(text: str) -> StopReason:
        lowered = text.lower()
        if "中断当前任务" in text or "interrupted" in lowered:
            return StopReason.INTERRUPTED
        if "预算" in text or "budget" in lowered:
            return StopReason.BUDGET_EXHAUSTED
        if "空转" in text or "没有获得新信息" in text or "相同操作" in text:
            return StopReason.NO_PROGRESS
        if "最大执行步数" in text or "max steps" in lowered:
            return StopReason.MAX_STEPS
        return StopReason.FINISHED if text else StopReason.END_TURN


if __name__ == "__main__":
    task = sys.argv[1] if len(sys.argv) > 1 else "create a 3-slide deck about coffee"
    print(Harness().run(task))
