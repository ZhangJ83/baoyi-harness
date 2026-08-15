"""PPT adapter for :mod:`agent.action_transaction`.

This module is intentionally absent from the model-visible tool registry.  PPT
facades can opt into it internally to obtain a durable checkpoint, cooperative
cancellation, scope enforcement, save/reopen verification, and a machine-readable
immutability certificate without teaching the transaction core about PowerPoint.

Closed-world boundary (documented contract): the adapter closes over the PPT
object graph and the source deck, not the entire filesystem. Concretely it
(a) snapshots the deck/source deck, (b) reopens and verifies the deck after
save, (c) issues ``PptImmutabilityCertificate`` for non-target slides, and
(d) rolls the deck/source deck back on failure. The action callback must
restrict side effects to ``harness.deck`` and the harness working copy.
Final-output publication belongs after this transaction commits; arbitrary
external files cannot be compensated by this adapter.
"""

from __future__ import annotations

from copy import deepcopy
from dataclasses import asdict, dataclass
from hashlib import sha256
from io import BytesIO
import json
from pathlib import Path
from typing import Any, Callable, Generic, Iterable, TypeVar
from uuid import uuid4

from lxml import etree
from pptx import Presentation

from .action_transaction import (
    ActionScope,
    ActionTransaction,
    CancellationToken,
    TransactionEvent,
)
from .transaction_journal import DurableTransactionJournal, TransactionJournal


ResultT = TypeVar("ResultT")


def _sha256_bytes(data: bytes) -> str:
    return sha256(data).hexdigest()


def _sha256_file(path: Path | None) -> str | None:
    if path is None or not path.is_file():
        return None
    digest = sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _serialize(deck: Presentation) -> bytes:
    stream = BytesIO()
    deck.save(stream)
    return stream.getvalue()


def _canonical_xml(element: Any) -> bytes:
    return etree.tostring(element, method="c14n", with_comments=False)


def slide_fingerprints(deck: Presentation) -> dict[int, str]:
    """Return persisted-content hashes for every slide in presentation order.

    A fingerprint covers canonical slide XML plus all direct relationships and
    related-part bytes (images, charts, notes, layout, etc.).  It therefore
    catches invisible OOXML changes as well as text and geometry changes.
    """

    fingerprints: dict[int, str] = {}
    for number, slide in enumerate(deck.slides, start=1):
        digest = sha256()
        digest.update(_canonical_xml(slide._element))
        relationships = sorted(
            slide.part.rels.values(),
            key=lambda rel: (rel.reltype, rel.rId),
        )
        for rel in relationships:
            digest.update(rel.rId.encode("utf-8"))
            digest.update(rel.reltype.encode("utf-8"))
            digest.update(b"external" if rel.is_external else b"internal")
            if rel.is_external:
                digest.update(str(rel.target_ref).encode("utf-8", errors="replace"))
            else:
                digest.update(sha256(rel.target_part.blob).digest())
        fingerprints[number] = digest.hexdigest()
    return fingerprints


@dataclass(frozen=True)
class PptTransactionCheckpoint:
    deck_bytes: bytes
    slide_hashes: dict[int, str]
    state_data: dict[str, Any] | None
    working_path: Path | None
    working_bytes: bytes | None
    checkpoint_path: Path | None
    checkpoint_sha256: str


@dataclass(frozen=True)
class PptImmutabilityCertificate:
    schema: str
    transaction_id: str
    allowed_slides: tuple[int, ...]
    requested_slides: tuple[int, ...]
    slide_count_before: int
    slide_count_after: int
    changed_slides: tuple[int, ...]
    unchanged_slides: tuple[int, ...]
    unexpected_changed_slides: tuple[int, ...]
    before_slide_hashes: dict[int, str]
    after_slide_hashes: dict[int, str]
    checkpoint_sha256: str
    committed_deck_sha256: str
    source_path: str | None
    source_sha256: str | None
    source_preserved: bool | None
    passed: bool


@dataclass(frozen=True)
class PptTransactionOutcome(Generic[ResultT]):
    transaction_id: str
    value: ResultT
    certificate: PptImmutabilityCertificate
    checkpoint_path: Path | None
    committed_deck_path: Path | None
    certificate_path: Path | None


class _HarnessCancellation(CancellationToken):
    def __init__(self, harness: Any):
        super().__init__()
        self._harness = harness

    @property
    def is_cancelled(self) -> bool:
        requested = getattr(self._harness, "cancel_requested", None)
        return super().is_cancelled or bool(callable(requested) and requested())


def _transaction_directory(harness: Any, transaction_id: str) -> Path | None:
    recorder = getattr(harness, "recorder", None)
    work = getattr(recorder, "work", None)
    if work is None:
        return None
    return Path(work) / "transactions" / transaction_id


def _event_sink(harness: Any) -> Callable[[TransactionEvent[Any]], None]:
    def emit(event: TransactionEvent[Any]) -> None:
        recorder = getattr(harness, "recorder", None)
        if recorder is not None and callable(getattr(recorder, "event", None)):
            recorder.event(
                "action_transaction",
                transaction_id=event.transaction_id,
                phase=event.phase,
                requested_scope=list(event.requested_scope),
                detail=event.detail,
            )

    return emit


def run_ppt_transaction(
    harness: Any,
    *,
    allowed_slides: Iterable[int],
    requested_slides: Iterable[int],
    action: Callable[[Presentation], ResultT],
    postcondition: Callable[[Presentation, ResultT], bool | None] | None = None,
    transaction_id: str | None = None,
    journal: TransactionJournal | None = None,
) -> PptTransactionOutcome[ResultT]:
    """Run one local PPT mutation as a scoped, recoverable transaction.

    Non-requested slides are immutable by default.  This intentionally supports
    local edits, not insert/delete/reorder operations whose slide identity needs
    a richer stable-ID policy.
    """

    deck = getattr(harness, "deck", None)
    if deck is None:
        raise ValueError("a presentation must be open before starting a PPT transaction")

    allowed = tuple(sorted(set(allowed_slides)))
    requested = tuple(sorted(set(requested_slides)))
    slide_count = len(deck.slides)
    invalid = [number for number in (*allowed, *requested) if number < 1 or number > slide_count]
    if invalid:
        raise ValueError(f"slide number outside current deck: {sorted(set(invalid))}")

    txid = transaction_id or uuid4().hex
    directory = _transaction_directory(harness, txid)
    holder: dict[str, Any] = {}

    # Production runs have a recorder-bound workspace; use it for a durable
    # cross-run state marker.  Lightweight direct adapter clients remain
    # backward compatible and can opt in by passing a journal explicitly.
    if journal is None:
        recorder_workspace = getattr(getattr(harness, "recorder", None), "workspace", None)
        if recorder_workspace is not None:
            journal = DurableTransactionJournal(recorder_workspace)

    source_value = getattr(harness, "deck_source_path", None)
    source_path = Path(source_value).resolve() if source_value is not None else None
    source_hash_before = _sha256_file(source_path)

    def checkpoint() -> PptTransactionCheckpoint:
        # Normalize through save/reopen so before/after hashes are comparable to
        # actual persisted PPTX behavior, rather than only the live object graph.
        live_deck = harness.deck
        deck_bytes = _serialize(live_deck)
        normalized = Presentation(BytesIO(deck_bytes))
        checkpoint_path = None
        if directory is not None:
            directory.mkdir(parents=True, exist_ok=True)
            checkpoint_path = directory / "checkpoint.pptx"
            checkpoint_path.write_bytes(deck_bytes)
        state = getattr(harness, "state", None)
        state_data = deepcopy(state.__dict__) if state is not None else None
        working_path = (
            Path(harness.deck_working_path)
            if getattr(harness, "deck_working_path", None) is not None
            else None
        )
        snapshot = PptTransactionCheckpoint(
            deck_bytes=deck_bytes,
            slide_hashes=slide_fingerprints(normalized),
            state_data=state_data,
            working_path=working_path,
            # A PPTX save can reorder ZIP members even if its semantics do not
            # change.  Preserve the exact pre-transaction working artifact so
            # rollback is byte-for-byte, not merely save/reopen equivalent.
            working_bytes=(working_path.read_bytes() if working_path is not None and working_path.is_file() else None),
            checkpoint_path=checkpoint_path,
            checkpoint_sha256=_sha256_bytes(deck_bytes),
        )
        holder["live_deck"] = live_deck
        holder["checkpoint"] = snapshot
        return snapshot

    def execute() -> ResultT:
        return action(harness.deck)

    def verify(value: ResultT) -> bool:
        snapshot: PptTransactionCheckpoint = holder["checkpoint"]
        committed_bytes = _serialize(harness.deck)
        reopened = Presentation(BytesIO(committed_bytes))
        after_hashes = slide_fingerprints(reopened)

        before_count = len(snapshot.slide_hashes)
        after_count = len(after_hashes)
        all_numbers = sorted(set(snapshot.slide_hashes) | set(after_hashes))
        changed = tuple(
            number
            for number in all_numbers
            if snapshot.slide_hashes.get(number) != after_hashes.get(number)
        )
        requested_set = set(requested)
        unexpected = tuple(number for number in changed if number not in requested_set)
        unchanged = tuple(number for number in all_numbers if number not in changed)
        source_hash_after = _sha256_file(source_path)
        source_preserved = (
            None
            if source_hash_before is None
            else source_hash_after == source_hash_before
        )
        user_passed = True if postcondition is None else postcondition(reopened, value)
        passed = (
            before_count == after_count
            and not unexpected
            and source_preserved is not False
            and user_passed is not False
        )
        certificate = PptImmutabilityCertificate(
            schema="xiaopu-ppt-scope-certificate-v1",
            transaction_id=txid,
            allowed_slides=allowed,
            requested_slides=requested,
            slide_count_before=before_count,
            slide_count_after=after_count,
            changed_slides=changed,
            unchanged_slides=unchanged,
            unexpected_changed_slides=unexpected,
            before_slide_hashes=snapshot.slide_hashes,
            after_slide_hashes=after_hashes,
            checkpoint_sha256=snapshot.checkpoint_sha256,
            committed_deck_sha256=_sha256_bytes(committed_bytes),
            source_path=str(source_path) if source_path is not None else None,
            source_sha256=source_hash_before,
            source_preserved=source_preserved,
            passed=passed,
        )
        holder.update(
            committed_bytes=committed_bytes,
            reopened=reopened,
            certificate=certificate,
        )
        return passed

    def rollback(snapshot: PptTransactionCheckpoint, _error: BaseException) -> None:
        # Preserve object identity when an inner atomic mutator failed before
        # touching the live deck.  Some existing callers retain shape handles;
        # replacing a demonstrably unchanged object would invalidate them.
        # PPTX ZIP serialization is not byte-stable: a second save can change
        # container metadata even when the live object graph is untouched.
        # Preserve identity only when this is still the exact pre-transaction
        # Presentation *and* every slide fingerprint remains unchanged.
        try:
            live_is_unchanged = (
                harness.deck is holder.get("live_deck")
                and slide_fingerprints(harness.deck) == snapshot.slide_hashes
            )
        except Exception:
            live_is_unchanged = False
        if not live_is_unchanged:
            harness.deck = Presentation(BytesIO(snapshot.deck_bytes))
        if snapshot.working_path is not None:
            snapshot.working_path.parent.mkdir(parents=True, exist_ok=True)
            snapshot.working_path.write_bytes(
                snapshot.working_bytes
                if snapshot.working_bytes is not None
                else snapshot.deck_bytes
            )
        state = getattr(harness, "state", None)
        if state is not None and snapshot.state_data is not None:
            state.__dict__.clear()
            state.__dict__.update(deepcopy(snapshot.state_data))

    def commit(_value: ResultT) -> None:
        snapshot: PptTransactionCheckpoint = holder["checkpoint"]
        certificate: PptImmutabilityCertificate = holder["certificate"]
        committed_bytes: bytes = holder["committed_bytes"]
        # Verification uses the save/reopened copy, but a successful commit
        # retains the live Presentation object.  Existing helpers and callers
        # can safely keep shape handles across a canonical local mutation.

        # Existing semantic mutators already advance the epoch.  Direct adapter
        # clients still get one mutation record without double counting.
        state = getattr(harness, "state", None)
        if state is not None:
            old_epoch = snapshot.state_data.get("mutation_epoch", 0) if snapshot.state_data else 0
            if state.mutation_epoch == old_epoch:
                state.record_changes([f"deck:slide:{number}:scoped-transaction" for number in requested])
            else:
                state.ppt_affected_slides.update(requested)

        committed_path = None
        certificate_path = None
        if directory is not None:
            committed_path = directory / "committed.pptx"
            certificate_path = directory / "immutability_certificate.json"
            committed_path.write_bytes(committed_bytes)
            certificate_path.write_text(
                json.dumps(asdict(certificate), ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
        holder["committed_path"] = committed_path
        holder["certificate_path"] = certificate_path

        recorder = getattr(harness, "recorder", None)
        if recorder is not None:
            if callable(getattr(recorder, "event", None)):
                recorder.event(
                    "ppt_scope_verified",
                    transaction_id=txid,
                    requested_slides=list(requested),
                    changed_slides=list(certificate.changed_slides),
                    unexpected_changed_slides=list(certificate.unexpected_changed_slides),
                    certificate=str(certificate_path) if certificate_path else None,
                )
            if callable(getattr(recorder, "check", None)):
                recorder.check(
                    "ppt_scoped_transaction",
                    certificate.passed,
                    f"changed={list(certificate.changed_slides)}; "
                    f"unexpected={list(certificate.unexpected_changed_slides)}; "
                    f"source_preserved={certificate.source_preserved}",
                )

    scope = ActionScope.from_iterables(
        allowed=(f"slide:{number}" for number in allowed),
        requested=(f"slide:{number}" for number in requested),
    )
    transaction = ActionTransaction(
        scope=scope,
        checkpoint=checkpoint,
        execute=execute,
        postcondition=verify,
        rollback=rollback,
        commit=commit,
        cancellation=_HarnessCancellation(harness),
        event_sink=_event_sink(harness),
        journal=journal,
        transaction_id=txid,
    )
    result = transaction.run()
    snapshot = holder["checkpoint"]
    return PptTransactionOutcome(
        transaction_id=result.transaction_id,
        value=result.value,
        certificate=holder["certificate"],
        checkpoint_path=snapshot.checkpoint_path,
        committed_deck_path=holder.get("committed_path"),
        certificate_path=holder.get("certificate_path"),
    )
