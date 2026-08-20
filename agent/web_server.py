"""Lightweight Web GUI Server for Xiaopu Harness.

Serves the modern Claude Desktop / Cowork-inspired HTML5/CSS3 frontend and provides
REST + SSE (Server-Sent Events) APIs for real-time streaming, tool calls, and workspace session management.
"""
from __future__ import annotations

import json
import difflib
import os
import queue
import re
import shutil
import subprocess
import sys
import threading
import time
import urllib.parse
import webbrowser
from datetime import datetime, timezone
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path
import ctypes
from ctypes import wintypes
from typing import Any


class _BROWSEINFOW(ctypes.Structure):
    _fields_ = [
        ("hwndOwner", wintypes.HWND),
        ("pidlRoot", ctypes.c_void_p),
        ("pszDisplayName", wintypes.LPWSTR),
        ("lpszTitle", wintypes.LPCWSTR),
        ("ulFlags", wintypes.UINT),
        ("lpfn", ctypes.c_void_p),
        ("lParam", wintypes.LPARAM),
        ("iImage", ctypes.c_int),
    ]


class _OPENFILENAMEW(ctypes.Structure):
    _fields_ = [
        ("lStructSize", wintypes.DWORD),
        ("hwndOwner", wintypes.HWND),
        ("hInstance", wintypes.HINSTANCE),
        ("lpstrFilter", wintypes.LPCWSTR),
        ("lpstrCustomFilter", wintypes.LPWSTR),
        ("nMaxCustFilter", wintypes.DWORD),
        ("nFilterIndex", wintypes.DWORD),
        ("lpstrFile", wintypes.LPWSTR),
        ("nMaxFile", wintypes.DWORD),
        ("lpstrFileTitle", wintypes.LPWSTR),
        ("nMaxFileTitle", wintypes.DWORD),
        ("lpstrInitialDir", wintypes.LPCWSTR),
        ("lpstrTitle", wintypes.LPCWSTR),
        ("Flags", wintypes.DWORD),
        ("nFileOffset", wintypes.WORD),
        ("nFileExtension", wintypes.WORD),
        ("lpstrDefExt", wintypes.LPCWSTR),
        ("lCustData", wintypes.LPARAM),
        ("lpfnHook", ctypes.c_void_p),
        ("lpTemplateName", wintypes.LPCWSTR),
        ("pvReserved", ctypes.c_void_p),
        ("dwReserved", wintypes.DWORD),
        ("FlagsEx", wintypes.DWORD),
    ]


def _pick_directory_native(initial_dir: str = "") -> str | None:
    # 1. On Windows, use Win32 SHBrowseForFolderW (zero subprocess, zero PowerShell, zero Tkinter)
    if sys.platform == "win32":
        try:
            ole32 = ctypes.windll.ole32
            ole32.CoInitialize(None)
            shell32 = ctypes.windll.shell32

            display_name = ctypes.create_unicode_buffer(260)
            bi = _BROWSEINFOW()
            bi.hwndOwner = None
            bi.pidlRoot = None
            bi.pszDisplayName = ctypes.cast(display_name, wintypes.LPWSTR)
            bi.lpszTitle = "请选择小朴项目 / 工作区目录"
            # BIF_RETURNONLYFSDIRS (0x1) | BIF_NEWDIALOGSTYLE (0x40) | BIF_USENEWUI (0x50)
            bi.ulFlags = 0x00000040 | 0x00000010 | 0x00000001

            pidl = shell32.SHBrowseForFolderW(ctypes.byref(bi))
            if pidl:
                path_buf = ctypes.create_unicode_buffer(1024)
                if shell32.SHGetPathFromIDListW(pidl, path_buf):
                    ole32.CoTaskMemFree(pidl)
                    ole32.CoUninitialize()
                    res_path = path_buf.value
                    if res_path and Path(res_path).is_dir():
                        return str(Path(res_path).resolve())
                ole32.CoTaskMemFree(pidl)
            ole32.CoUninitialize()
        except Exception:
            pass

    # 2. Cross-platform Tkinter fallback
    try:
        import tkinter as tk
        from tkinter import filedialog
        root = tk.Tk()
        root.withdraw()
        root.attributes("-topmost", True)
        selected = filedialog.askdirectory(title="选择工作区/项目目录", initialdir=initial_dir or str(config.sandbox_root()))
        root.destroy()
        if selected and Path(selected).is_dir():
            return str(Path(selected).resolve())
    except Exception:
        pass

    return None


def _pick_save_file_native(initial_dir: str = "", default_name: str = "presentation.pptx") -> str | None:
    # 1. On Windows, use Win32 GetSaveFileNameW
    if sys.platform == "win32":
        try:
            ofn = _OPENFILENAMEW()
            ofn.lStructSize = ctypes.sizeof(_OPENFILENAMEW)
            file_buf = ctypes.create_unicode_buffer(1024)
            if default_name:
                file_buf.value = default_name
            ofn.lpstrFile = ctypes.cast(file_buf, wintypes.LPWSTR)
            ofn.nMaxFile = 1024
            ofn.lpstrFilter = "PowerPoint 演示文稿 (*.pptx)\0*.pptx\0所有文件 (*.*)\0*.*\0\0"
            ofn.lpstrTitle = "另存为 PPT 文件"
            ofn.lpstrDefExt = "pptx"
            if initial_dir:
                ofn.lpstrInitialDir = initial_dir
            ofn.Flags = 0x00000002 | 0x00000800  # OFN_OVERWRITEPROMPT | OFN_PATHMUSTEXIST
            if ctypes.windll.comdlg32.GetSaveFileNameW(ctypes.byref(ofn)):
                res_path = file_buf.value
                if res_path:
                    return str(Path(res_path).resolve())
        except Exception:
            pass

    # 2. Cross-platform Tkinter fallback
    try:
        import tkinter as tk
        from tkinter import filedialog
        root = tk.Tk()
        root.withdraw()
        root.attributes("-topmost", True)
        selected = filedialog.asksaveasfilename(
            title="另存为 PPT 文件",
            defaultextension=".pptx",
            filetypes=[("PowerPoint", "*.pptx")],
            initialdir=initial_dir or str(config.sandbox_root()),
            initialfile=default_name,
        )
        root.destroy()
        if selected:
            return str(Path(selected).resolve())
    except Exception:
        pass

    return None


def _time_ago(iso_str: str) -> str:
    if not iso_str:
        return "now"
    try:
        dt = datetime.fromisoformat(iso_str)
        now = datetime.now(timezone.utc)
        diff = now - dt
        seconds = max(0, int(diff.total_seconds()))
        if seconds < 60:
            return "now"
        elif seconds < 3600:
            return f"{seconds // 60}m"
        elif seconds < 86400:
            return f"{seconds // 3600}h"
        elif seconds < 2592000:
            return f"{seconds // 86400}d"
        elif seconds < 31536000:
            return f"{seconds // 2592000}mo"
        else:
            return f"{seconds // 31536000}y"
    except Exception:
        return "now"

def _session_json(record) -> dict:
    return {
        "id": record.id,
        "title": record.title,
        "created_at": record.created_at,
        "updated_at": record.updated_at,
        "time_ago": _time_ago(record.updated_at),
        "model": record.model,
        "workspace": record.workspace,
        "turn_count": record.turn_count,
        "summary": record.summary,
        "pinned": bool(record.pinned),
        "status": record.status,
    }


def _workspace_json(record) -> dict:
    return {
        "path": record.path,
        "name": record.name,
        "display_name": record.display_name or record.name,
        "last_used": record.last_used,
        "pinned": bool(record.pinned),
        "archived": bool(record.archived),
        "removed_at": record.removed_at or "",
    }


from . import config
from .events import EventKind, RuntimeEvent
from .harness import Harness
from .session_store import (
    archive_session,
    batch_session_action,
    delete_session,
    export_session,
    list_sessions,
    load_session,
    purge_expired_sessions,
    purge_session,
    rename_session,
    restore_harness,
    restore_session,
    save_session,
    set_session_pinned,
    trash_session,
)
from .tools.registry import dispatch
from .workspace_store import (
    archive_workspace,
    list_workspaces,
    purge_workspace,
    register_workspace,
    remove_workspace,
    rename_workspace,
    restore_workspace,
    set_workspace_pinned,
    touch_workspace,
)

WEB_DIR = Path(__file__).resolve().parent / "web"
ASSETS_DIR = Path(__file__).resolve().parent / "assets"


def _find_active_deck(
    ws_path: Path | None = None,
    session_id: str | None = None,
    specific_file: str | Path | None = None,
    harness: Any = None,
) -> Path | None:
    # 1. Explicit file parameter passed from client
    if specific_file:
        sp = Path(specific_file)
        if not sp.is_absolute() and ws_path:
            sp = ws_path / sp
        if sp.exists() and sp.is_file() and sp.suffix.lower() in (".pptx", ".ppt") and sp.stat().st_size > 0:
            return sp

    # 2. If session_id is provided, check if the session has an associated PPT deck
    if session_id:
        try:
            from .session_store import load_session
            snap = load_session(session_id)
            if snap:
                session_root = Path(snap.get("workspace") or ws_path or config.sandbox_root())
                # Check deck_working_path or deck_source_path or required_output_pptx
                deck_values = [snap.get("deck_working_path"), snap.get("deck_source_path")]
                facts = snap.get("facts") or {}
                deck_values.extend((facts.get("required_output_pptx"), facts.get("output_path")))
                for val in deck_values:
                    if val:
                        p = Path(val)
                        if not p.is_absolute():
                            p = session_root / p
                        if p.exists() and p.is_file() and p.suffix.lower() in (".pptx", ".ppt") and p.stat().st_size > 0:
                            return p
                # Check messages for saved pptx paths
                for msg in reversed(snap.get("messages", [])):
                    content = str(msg.get("content", ""))
                    matches = re.findall(r"(?:^|[\s'\"])([^\r\n'\"]+?\.pptx)(?=$|[\s'\"，。；;])", content, re.IGNORECASE)
                    for match in matches:
                        mp = Path(match.strip("\'\""))
                        if not mp.is_absolute():
                            mp = session_root / mp
                        if mp.exists() and mp.is_file() and mp.stat().st_size > 0:
                            return mp
        except Exception:
            pass

    # 3. If harness has an active deck or changed pptx files in current task
    harness_session_id = str(getattr(getattr(harness, "session", None), "id", "") or "") if harness is not None else ""
    harness_matches_session = not session_id or not harness_session_id or harness_session_id == session_id
    if harness is not None and harness_matches_session:
        for deck_path in (
            getattr(harness, "deck_working_path", None),
            getattr(harness, "deck_path", None),
            getattr(harness, "deck_source_path", None),
        ):
            if not deck_path:
                continue
            dp = Path(deck_path)
            if dp.exists() and dp.is_file() and dp.stat().st_size > 0:
                return dp
        changed = getattr(getattr(harness, "state", None), "changed_files", None)
        if changed:
            pptx_changed = [
                Path(p) for p in changed
                if str(p).lower().endswith((".pptx", ".ppt")) and Path(p).exists() and Path(p).is_file() and Path(p).stat().st_size > 0
            ]
            if pptx_changed:
                pptx_changed.sort(key=lambda p: p.stat().st_mtime, reverse=True)
                return pptx_changed[0]

    # A historical session without an associated deck must not silently display
    # an unrelated, newer deck from the same workspace.
    if session_id:
        return None

    # 4. Search in ws_path (or config.sandbox_root())
    root = ws_path or config.sandbox_root()
    if not root.is_dir():
        return None

    # Fast path: normal GUI decks live at the workspace root. Avoid walking a
    # large repository just to discover the conventional deck.pptx.
    conventional = root / "deck.pptx"
    if conventional.is_file() and conventional.stat().st_size > 0:
        return conventional
    direct_candidates = [
        p for p in root.glob("*.pptx")
        if not p.name.startswith("~$") and p.is_file() and p.stat().st_size > 0
    ]
    if direct_candidates:
        direct_candidates.sort(key=lambda p: p.stat().st_mtime_ns, reverse=True)
        return direct_candidates[0]

    # Bounded fallback for nested task outputs. Dependency trees, caches and
    # prior backups are both expensive and incorrect sources for an active deck.
    candidates: list[Path] = []
    try:
        root_depth = len(root.parts)
        ignored = {"node_modules", ".git", ".venv", "venv", "__pycache__", "preview_cache", "ppt_text_backups"}
        for current_dir, dir_names, file_names in os.walk(root):
            current = Path(current_dir)
            depth = len(current.parts) - root_depth
            dir_names[:] = [
                name for name in dir_names
                if name not in ignored and not name.startswith(".") and depth < 6
            ]
            for name in file_names:
                if name.startswith("~$") or not name.lower().endswith(".pptx"):
                    continue
                p = current / name
                if p.is_file() and p.stat().st_size > 0:
                    candidates.append(p)
    except Exception:
        pass

    if candidates:
        # Sort strictly by mtime descending (most recently modified/created first)
        candidates.sort(key=lambda p: p.stat().st_mtime_ns, reverse=True)
        return candidates[0]

    return None


def _get_deck_content_data(pptx_path: Path) -> dict[str, Any]:
    from pptx import Presentation
    try:
        prs = Presentation(str(pptx_path))
    except Exception as exc:
        return {"success": False, "error": str(exc), "total_slides": 0, "slides": [], "text_content": ""}

    slides_info = []
    text_sections = [
        "编辑说明：只修改每个 [Sx.SHx.Px] 标记后面的文字；请保留方括号内的定位标记。",
        "符号 ↵ 表示文本框内的手动换行。",
        "",
    ]
    for idx, slide in enumerate(prs.slides, 1):
        slide_title = ""
        items = []
        records = list(_iter_editable_text_records(slide, idx))
        for record in records:
            ptxt = record["paragraph"].text.strip()
            if not ptxt:
                continue
            if not slide_title and record["top_inches"] < 1.5 and record["kind"] == "shape":
                slide_title = ptxt
            elif ptxt not in items:
                items.append(ptxt)

        title_display = slide_title or f"第 {idx} 页"
        slides_info.append({
            "slide_number": idx,
            "title": title_display,
            "items": items,
        })
        text_sections.append(f"=== 第 {idx} 页 ===")
        for record in records:
            value = record["paragraph"].text.replace("\v", "↵").replace("\n", "↵")
            if value.strip():
                text_sections.append(f"[{record['ref']}] {value}")
        text_sections.append("")

    html_files = []
    try:
        parent_dir = pptx_path.parent
        for html_candidate in sorted(parent_dir.glob("*.html")):
            if html_candidate.is_file() and not html_candidate.name.startswith("."):
                html_files.append({
                    "filename": html_candidate.name,
                    "content": html_candidate.read_text(encoding="utf-8", errors="ignore")[:15000],
                })
    except Exception:
        pass

    return {
        "success": True,
        "deck_name": pptx_path.name,
        "deck_path": str(pptx_path.resolve()),
        "total_slides": len(prs.slides),
        "slides": slides_info,
        "text_content": "\n".join(text_sections).strip(),
        "html_files": html_files,
    }


_EDITABLE_REF_RE = re.compile(r"^\[(S\d+\.(?:SH\d+(?:\.G\d+)*|T\d+\.R\d+\.C\d+)\.P\d+)\]\s?(.*)$", re.IGNORECASE)


def _iter_editable_text_records(slide: Any, slide_number: int):
    """Yield addressable text paragraphs, including grouped shapes and tables."""
    def walk(shapes: Any, group_suffix: str = "", inherited_top: float = 0.0):
        for shape in shapes:
            try:
                shape_id = int(shape.shape_id)
            except Exception:
                continue
            try:
                top_inches = float(shape.top.inches)
            except Exception:
                top_inches = inherited_top
            shape_ref = f"S{slide_number}.SH{shape_id}{group_suffix}"

            if getattr(shape, "has_text_frame", False) and getattr(shape, "text_frame", None) is not None:
                for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs, 1):
                    yield {
                        "ref": f"{shape_ref}.P{paragraph_index}",
                        "paragraph": paragraph,
                        "kind": "shape",
                        "top_inches": top_inches,
                    }

            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows, 1):
                    for column_index, cell in enumerate(row.cells, 1):
                        if getattr(cell, "is_spanned", False):
                            continue
                        for paragraph_index, paragraph in enumerate(cell.text_frame.paragraphs, 1):
                            yield {
                                "ref": f"S{slide_number}.T{shape_id}.R{row_index}.C{column_index}.P{paragraph_index}",
                                "paragraph": paragraph,
                                "kind": "table",
                                "top_inches": top_inches,
                            }

            child_shapes = getattr(shape, "shapes", None)
            if child_shapes is not None:
                yield from walk(child_shapes, f"{group_suffix}.G{shape_id}", top_inches)

    yield from walk(slide.shapes)


def _parse_editable_text(text_content: str) -> dict[str, str]:
    updates: dict[str, str] = {}
    for line in text_content.splitlines():
        match = _EDITABLE_REF_RE.match(line.strip("\r"))
        if not match:
            continue
        ref = match.group(1).upper()
        if ref in updates:
            raise ValueError(f"重复的文本定位标记: {ref}")
        updates[ref] = match.group(2).replace("↵", "\v")
    return updates


def _replace_paragraph_text(paragraph: Any, new_text: str) -> None:
    """Replace paragraph text while retaining unchanged run-level styling."""
    if paragraph.text == new_text:
        return
    runs = list(paragraph.runs)
    if not runs:
        paragraph.text = new_text
        return
    original = "".join(run.text for run in runs)
    if original != paragraph.text:
        runs[0].text = new_text
        for run in runs[1:]:
            run.text = ""
        return

    character_runs = [index for index, run in enumerate(runs) for _ in run.text]
    pieces: list[list[str]] = [[] for _ in runs]
    matcher = difflib.SequenceMatcher(a=original, b=new_text, autojunk=False)
    for tag, old_start, old_end, new_start, new_end in matcher.get_opcodes():
        if tag == "equal":
            for offset, character in enumerate(new_text[new_start:new_end]):
                pieces[character_runs[old_start + offset]].append(character)
        elif tag in {"insert", "replace"} and new_start < new_end:
            if old_start < len(character_runs):
                target_run = character_runs[old_start]
            elif character_runs:
                target_run = character_runs[-1]
            else:
                target_run = 0
            pieces[target_run].append(new_text[new_start:new_end])

    for run, text_pieces in zip(runs, pieces):
        run.text = "".join(text_pieces)


def _apply_deck_text_content(pptx_path: Path, text_content: str) -> dict[str, Any]:
    """Apply marker-addressed text edits atomically to the selected deck."""
    from pptx import Presentation

    updates = _parse_editable_text(text_content)
    if not updates:
        raise ValueError("未找到可识别的 [Sx.SHx.Px] 文本定位标记；请先刷新右侧 PPT 文本")

    prs = Presentation(str(pptx_path))
    records: dict[str, dict[str, Any]] = {}
    for slide_number, slide in enumerate(prs.slides, 1):
        for record in _iter_editable_text_records(slide, slide_number):
            records[record["ref"].upper()] = record

    unknown = sorted(set(updates) - set(records))
    if unknown:
        sample = ", ".join(unknown[:5])
        raise ValueError(f"PPT 已变化，存在失效的文本定位标记: {sample}；请刷新后重试")

    changed_refs: list[str] = []
    changed_slides: set[int] = set()
    for ref, new_text in updates.items():
        paragraph = records[ref]["paragraph"]
        if paragraph.text != new_text:
            _replace_paragraph_text(paragraph, new_text)
            changed_refs.append(ref)
            changed_slides.add(int(ref.split(".", 1)[0][1:]))

    if not changed_refs:
        return {"changed_count": 0, "changed_refs": [], "changed_slides": []}

    backup_dir = pptx_path.parent / ".baoyi" / "ppt_text_backups"
    backup_dir.mkdir(parents=True, exist_ok=True)
    backup_path = backup_dir / f"{pptx_path.stem}-{time.time_ns()}.pptx"
    shutil.copy2(pptx_path, backup_path)

    temp_path = pptx_path.with_name(f".{pptx_path.stem}.{time.time_ns()}.tmp.pptx")
    try:
        prs.save(str(temp_path))
        Presentation(str(temp_path))
        os.replace(temp_path, pptx_path)
    finally:
        if temp_path.exists():
            temp_path.unlink()

    return {
        "changed_count": len(changed_refs),
        "changed_refs": changed_refs,
        "changed_slides": sorted(changed_slides),
        "backup_path": str(backup_path),
    }


_COM_RENDER_LOCK = threading.Lock()


def _render_deck_slide_preview(pptx_path: Path, slide_number: int = 1) -> bytes | None:
    import io
    import hashlib
    # 1. Try Windows PowerPoint COM for 100% native pixel-perfect quality
    if sys.platform == "win32":
        cache_dir = config.sandbox_root() / ".baoyi" / "preview_cache"
        cache_dir.mkdir(parents=True, exist_ok=True)
        path_hash = hashlib.md5(str(pptx_path.resolve()).encode("utf-8", errors="replace")).hexdigest()[:12]
        stat = pptx_path.stat()
        out_png = cache_dir / f"slide_{path_hash}_{slide_number}_{stat.st_mtime_ns}_{stat.st_size}.png"
        if out_png.exists() and out_png.stat().st_size > 0:
            return out_png.read_bytes()

        with _COM_RENDER_LOCK:
            if out_png.exists() and out_png.stat().st_size > 0:
                return out_png.read_bytes()
            try:
                # PowerPoint automation occasionally blocks forever on startup
                # dialogs or add-ins. Isolate it so the HTTP request has a hard
                # deadline and can fall back to the lightweight renderer.
                render_script = """
import sys
import pythoncom
import win32com.client
pythoncom.CoInitialize()
app = None
deck = None
try:
    app = win32com.client.DispatchEx('PowerPoint.Application')
    deck = app.Presentations.Open(sys.argv[1], ReadOnly=True, Untitled=False, WithWindow=False)
    slide = int(sys.argv[3])
    if 1 <= slide <= deck.Slides.Count:
        deck.Slides(slide).Export(sys.argv[2], 'PNG', 1920, 1080)
finally:
    if deck is not None:
        deck.Close()
    if app is not None:
        app.Quit()
    pythoncom.CoUninitialize()
"""
                subprocess.run(
                    [sys.executable, "-c", render_script, str(pptx_path.resolve()), str(out_png), str(slide_number)],
                    stdin=subprocess.DEVNULL,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                    timeout=12,
                    check=False,
                    creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0),
                )
                if out_png.exists() and out_png.stat().st_size > 0:
                    return out_png.read_bytes()
            except (subprocess.SubprocessError, OSError):
                pass

    # 2. Fallback: PIL rasterizer
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw
        prs = Presentation(str(pptx_path))
        if 1 <= slide_number <= len(prs.slides):
            slide = prs.slides[slide_number - 1]
            img = Image.new("RGB", (1920, 1080), color=(15, 23, 42))
            draw = ImageDraw.Draw(img)
            draw.rectangle([(0, 0), (1920, 120)], fill=(30, 41, 59))
            draw.text((60, 45), f"幻灯片预览 · 第 {slide_number} 页", fill=(56, 189, 248))
            y = 180
            for sh in slide.shapes:
                if getattr(sh, "has_text_frame", False) and sh.text_frame:
                    t = (sh.text_frame.text or "").strip()
                    if t:
                        draw.text((80, y), t[:120], fill=(241, 245, 249))
                        y += 50
                        if y > 980:
                            break
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            return buf.getvalue()
    except Exception:
        pass
    return None


class XiaopuWebHandler(BaseHTTPRequestHandler):
    harness: Harness = None  # Class-level shared harness instance
    active_stream_queue: queue.Queue | None = None

    def log_message(self, format, *args):
        # Suppress routine GET logging for clean console output
        pass

    def _send_json(self, data: Any, status: int = 200) -> None:
        payload = json.dumps(data, ensure_ascii=False).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(payload)))
        self.send_header("Access-Control-Allow-Origin", "*")
        self.end_headers()
        self.wfile.write(payload)

    def _send_file(self, file_path: Path, content_type: str) -> None:
        if not file_path.is_file():
            self.send_error(404, "File Not Found")
            return
        content = file_path.read_bytes()
        self.send_response(200)
        self.send_header("Content-Type", content_type)
        self.send_header("Content-Length", str(len(content)))
        self.end_headers()
        self.wfile.write(content)

    def _send_bytes(self, content: bytes, content_type: str = "image/png") -> None:
        self.send_response(200)
        self.send_header("Content-Type", content_type)
        self.send_header("Content-Length", str(len(content)))
        self.send_header("Cache-Control", "no-cache, must-revalidate")
        self.send_header("Access-Control-Allow-Origin", "*")
        self.end_headers()
        self.wfile.write(content)

    def _read_json_body(self) -> dict:
        try:
            length = int(self.headers.get("Content-Length", 0))
            if length <= 0:
                return {}
            body = self.rfile.read(length).decode("utf-8")
            return json.loads(body)
        except Exception:
            return {}

    def do_OPTIONS(self) -> None:
        self.send_response(200)
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "GET, POST, DELETE, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")
        self.end_headers()

    def do_GET(self) -> None:
        parsed = urllib.parse.urlparse(self.path)
        path = parsed.path
        query = urllib.parse.parse_qs(parsed.query)

        # Static asset routing
        if path in {"/", "/index.html"}:
            self._send_file(WEB_DIR / "index.html", "text/html; charset=utf-8")
            return
        elif path == "/style.css":
            self._send_file(WEB_DIR / "style.css", "text/css; charset=utf-8")
            return
        elif path == "/app.js":
            self._send_file(WEB_DIR / "app.js", "application/javascript; charset=utf-8")
            return
        elif path == "/assets/icon.png":
            self._send_file(ASSETS_DIR / "icon.png", "image/png")
            return
        elif path == "/assets/icon.ico":
            self._send_file(ASSETS_DIR / "icon.ico", "image/x-icon")
            return

        # API routing
        if path == "/api/config":
            self._send_json({
                "known_models": config.known_models(),
                "current_model": getattr(self.harness.llm, "model", config.model()),
                "command_policy": config.command_policy(),
                "provider": config.provider(),
                "reasoning_effort": config.reasoning_effort(),
                "thinking_enabled": config.thinking_enabled(),
            })
            return

        if path == "/api/settings":
            raw_key = config.api_key()
            masked_key = (raw_key[:6] + "..." + raw_key[-4:]) if len(raw_key) > 10 else ("*" * len(raw_key))
            self._send_json({
                "provider": config.provider(),
                "api_base": config.api_base(),
                "api_key": raw_key,
                "masked_api_key": masked_key,
                "model": getattr(self.harness.llm, "model", config.model()),
                "known_models": config.known_models(),
                "models_csv": ",".join(config.known_models()),
                "reasoning_effort": config.reasoning_effort(),
                "command_policy": config.command_policy(),
            })
            return

        if path == "/api/workspaces":
            view = query.get("view", ["active"])[0]
            records = list_workspaces(view=view)
            current = str(config.sandbox_root())
            self._send_json({
                "workspaces": [w.path for w in records],
                "current": current,
                "records": [_workspace_json(w) for w in records],
                "archived": [_workspace_json(w) for w in list_workspaces(view="archived")],
                "removed": [_workspace_json(w) for w in list_workspaces(view="removed")],
            })
            return

        if path == "/api/workspaces/manage":
            self._send_json({
                "active": [_workspace_json(w) for w in list_workspaces(view="active")],
                "archived": [_workspace_json(w) for w in list_workspaces(view="archived")],
                "removed": [_workspace_json(w) for w in list_workspaces(view="removed")],
                "current": str(config.sandbox_root()),
            })
            return

        if path == "/api/tree":
            # Trash self-cleans on read; cheap and avoids a background scheduler.
            purge_expired_sessions(days=30)
            view = query.get("view", ["active"])[0]
            if view not in {"active", "archive", "trash", "all"}:
                view = "active"
            q = (query.get("q", [""])[0] or "").strip().casefold()
            current_ws = str(config.sandbox_root())
            known_workspaces = []

            # Workspaces explicitly archived/removed stay hidden even when they
            # are the current workspace or own historical sessions.
            hidden_workspaces: set[str] = set()
            for w in list_workspaces(view="archived") + list_workspaces(view="removed"):
                try:
                    hidden_workspaces.add(str(Path(w.path).resolve()).casefold())
                except Exception:
                    pass

            # 1. Registered active workspaces
            for w in list_workspaces(view="active"):
                w_path = getattr(w, "path", str(w))
                if w_path and w_path not in known_workspaces:
                    try:
                        if Path(w_path).is_dir():
                            known_workspaces.append(str(Path(w_path).resolve()))
                    except Exception:
                        pass

            # 2. Current workspace (unless the user archived/removed it)
            try:
                curr_res = str(Path(current_ws).resolve())
                if curr_res.casefold() not in hidden_workspaces and curr_res not in known_workspaces:
                    known_workspaces.insert(0, curr_res)
            except Exception:
                if current_ws not in known_workspaces:
                    known_workspaces.insert(0, current_ws)

            # 3. Auto-discover valid workspaces from all historical sessions
            active_sessions = list_sessions(view="active")
            for s in active_sessions:
                if s.workspace:
                    try:
                        p = Path(s.workspace)
                        if p.is_dir():
                            resolved = str(p.resolve())
                            # Ignore temporary pytest directories and
                            # explicitly archived/removed workspaces.
                            if "pytest-" not in resolved and "Temp" not in resolved \
                                    and resolved.casefold() not in hidden_workspaces \
                                    and resolved not in known_workspaces:
                                known_workspaces.append(resolved)
                    except Exception:
                        pass

            # Workspace metadata (display alias, pin) keyed by resolved path.
            ws_meta: dict[str, dict] = {}
            for w in list_workspaces(view="all"):
                try:
                    ws_meta[str(Path(w.path).resolve()).casefold()] = _workspace_json(w)
                except Exception:
                    pass

            def matches_query(session) -> bool:
                if not q:
                    return True
                haystack = f"{session.title} {session.summary}".casefold()
                return q in haystack

            assigned_session_ids = set()
            projects = []
            view_sessions = list_sessions(view=view)
            for ws_path in known_workspaces:
                p_path = Path(ws_path)
                p_name = p_path.name or str(ws_path)
                ws_sessions = [s for s in list_sessions(workspace=ws_path, view=view) if matches_query(s)]
                s_list = []
                for r in ws_sessions:
                    assigned_session_ids.add(r.id)
                    s_list.append(_session_json(r))
                try:
                    meta = ws_meta.get(str(Path(ws_path).resolve()).casefold(), {})
                except Exception:
                    meta = {}
                is_curr = False
                try:
                    is_curr = (Path(ws_path).resolve() == Path(current_ws).resolve())
                except Exception:
                    is_curr = (str(ws_path).casefold() == str(current_ws).casefold())

                projects.append({
                    "name": meta.get("display_name") or p_name,
                    "path": ws_path,
                    "is_current": is_curr,
                    "pinned": meta.get("pinned", False),
                    "sessions": s_list,
                })

            general_sessions = []
            for r in view_sessions:
                if r.id not in assigned_session_ids and matches_query(r):
                    general_sessions.append(_session_json(r))

            self._send_json({
                "projects": projects,
                "conversations": general_sessions,
                "current_workspace": current_ws,
                "view": view,
                "query": q,
                "workspace_groups": {
                    "active": [_workspace_json(w) for w in list_workspaces(view="active")],
                    "archived": [_workspace_json(w) for w in list_workspaces(view="archived")],
                    "removed": [_workspace_json(w) for w in list_workspaces(view="removed")],
                },
            })
            return

        if path == "/api/sessions":
            ws = query.get("workspace", [None])[0]
            view = query.get("view", ["active"])[0]
            records = list_sessions(workspace=ws, view=view)
            self._send_json({"sessions": [_session_json(r) for r in records]})
            return

        if path.startswith("/api/session/"):
            session_id = path[len("/api/session/"):]
            payload = load_session(session_id)
            if payload is None:
                self.send_error(404, "Session not found")
                return
            self._send_json(payload)
            return

        if path == "/api/goal":
            self._send_json({"summary": self.harness.goal_summary()})
            return

        if path == "/api/artifacts":
            ws_root = Path(config.sandbox_root())
            artifacts = []
            if ws_root.is_dir():
                candidate_exts = {".pptx", ".ppt", ".md", ".py", ".html", ".json", ".csv", ".xlsx", ".pdf", ".txt"}
                try:
                    # Recursive discovery: task outputs live under tasks/<id>/output,
                    # so a root-only scan hides the artifacts the agent just made.
                    candidates = []
                    root_depth = len(ws_root.parts)
                    ignored = {"node_modules", ".git", ".venv", "venv", "__pycache__", ".baoyi", ".xiaopu"}
                    for current_dir, dir_names, file_names in os.walk(ws_root):
                        current = Path(current_dir)
                        depth = len(current.parts) - root_depth
                        dir_names[:] = [
                            name for name in dir_names
                            if name not in ignored and not name.startswith(".") and depth < 6
                        ]
                        for name in file_names:
                            f = current / name
                            if f.suffix.lower() in candidate_exts:
                                candidates.append(f)
                    candidates.sort(key=lambda p: p.stat().st_mtime if p.exists() else 0, reverse=True)
                    for f in candidates[:200]:
                        stat = f.stat()
                        size_kb = stat.st_size / 1024
                        size_str = f"{stat.st_size} B" if stat.st_size < 1024 else (f"{size_kb:.1f} KB" if size_kb < 1024 else f"{size_kb/1024:.2f} MB")
                        file_type = f.suffix.lower().lstrip(".")

                        slides_count = None
                        if file_type in ("pptx", "ppt"):
                            try:
                                from pptx import Presentation as _Presentation
                                slides_count = len(_Presentation(str(f)).slides)
                            except Exception:
                                slides_count = None

                        artifacts.append({
                            "name": f.name,
                            "path": str(f.resolve()),
                            "type": file_type,
                            "size": stat.st_size,
                            "size_human": size_str,
                            "updated_at": datetime.fromtimestamp(stat.st_mtime, timezone.utc).isoformat(),
                            "time_ago": _time_ago(datetime.fromtimestamp(stat.st_mtime, timezone.utc).isoformat()),
                            "slides_count": slides_count,
                            "is_pptx": file_type in ("pptx", "ppt"),
                        })
                except Exception:
                    pass

            self._send_json({
                "artifacts": artifacts,
                "workspace": str(ws_root),
                "count": len(artifacts),
            })
            return

        if path == "/api/ppt/content":
            ws_param = query.get("workspace", [None])[0]
            session_param = query.get("session_id", [None])[0] or query.get("session", [None])[0]
            file_param = query.get("file", [None])[0] or query.get("path", [None])[0]
            ws_root = Path(ws_param) if ws_param else config.sandbox_root()
            deck_path = _find_active_deck(ws_root, session_id=session_param, specific_file=file_param, harness=self.harness)
            if not deck_path or not deck_path.exists():
                self._send_json({
                    "success": False,
                    "error": "当前会话暂未生成 PPT 文件",
                    "deck_name": "",
                    "total_slides": 0,
                    "slides": [],
                    "text_content": "",
                })
                return
            data = _get_deck_content_data(deck_path)
            self._send_json(data)
            return

        if path == "/api/ppt/preview":
            ws_param = query.get("workspace", [None])[0]
            session_param = query.get("session_id", [None])[0] or query.get("session", [None])[0]
            file_param = query.get("file", [None])[0] or query.get("path", [None])[0]
            ws_root = Path(ws_param) if ws_param else config.sandbox_root()
            slide_idx = int(query.get("slide", ["1"])[0])
            deck_path = _find_active_deck(ws_root, session_id=session_param, specific_file=file_param, harness=self.harness)
            if not deck_path or not deck_path.exists():
                # Return empty fallback image
                from PIL import Image
                import io
                img = Image.new("RGB", (1920, 1080), color=(15, 23, 42))
                buf = io.BytesIO()
                img.save(buf, format="PNG")
                self._send_bytes(buf.getvalue(), "image/png")
                return
            png_bytes = _render_deck_slide_preview(deck_path, slide_idx)
            if png_bytes:
                self._send_bytes(png_bytes, "image/png")
                return
            self.send_error(500, "Failed to render slide preview")
            return

        self.send_error(404, "Not Found")

    def do_POST(self) -> None:
        path = urllib.parse.urlparse(self.path).path
        body = self._read_json_body()

        if path == "/api/reveal_file":
            file_path = body.get("path")
            if file_path and Path(file_path).exists():
                norm = str(Path(file_path).resolve())
                if sys.platform == "win32":
                    try:
                        subprocess.Popen(["explorer.exe", f"/select,{norm}"])
                    except Exception:
                        pass
                self._send_json({"status": "ok", "path": norm})
                return
            self.send_error(400, "File does not exist")
            return

        if path == "/api/choose_directory":
            selected_path = _pick_directory_native(initial_dir=str(config.sandbox_root()))
            if selected_path:
                self._send_json({"status": "ok", "path": selected_path})
            else:
                self._send_json({"status": "cancelled", "path": ""})
            return

        if path == "/api/choose_save_ppt":
            selected_path = _pick_save_file_native(
                initial_dir=str(config.sandbox_root()),
                default_name="presentation.pptx",
            )
            if selected_path:
                self._send_json({"status": "ok", "path": selected_path})
            else:
                self._send_json({"status": "cancelled", "path": ""})
            return

        if path == "/api/workspace":
            new_ws = body.get("workspace")
            if new_ws:
                os.environ["WORKSPACE"] = str(new_ws)
                try:
                    register_workspace(new_ws)
                except Exception:
                    pass
                self.harness.reset()
                self._send_json({"status": "ok", "workspace": str(config.sandbox_root())})
                return
            self.send_error(400, "Missing workspace argument")
            return

        if path == "/api/workspace/action":
            ws_path = body.get("path")
            action = body.get("action")
            if not ws_path or not action:
                self.send_error(400, "Missing path or action")
                return
            handlers = {
                "rename": lambda: rename_workspace(ws_path, body.get("display_name", "")),
                "pin": lambda: set_workspace_pinned(ws_path, bool(body.get("pinned", False))),
                "archive": lambda: archive_workspace(ws_path),
                "restore": lambda: restore_workspace(ws_path),
                "remove": lambda: remove_workspace(ws_path),
                "purge": lambda: purge_workspace(ws_path),
            }
            handler = handlers.get(action)
            if handler is None:
                self.send_error(400, f"Unsupported workspace action: {action}")
                return
            try:
                ok = handler()
            except Exception as exc:
                self._send_json({"status": "error", "error": str(exc)}, status=400)
                return
            self._send_json({"status": "ok" if ok else "not_found", "action": action, "path": ws_path})
            return

        if path == "/api/session/action":
            session_id = body.get("id")
            action = body.get("action")
            if not session_id or not action:
                self.send_error(400, "Missing id or action")
                return
            if action == "rename":
                ok = rename_session(session_id, body.get("title", ""))
                result = {"status": "ok" if ok else "not_found", "action": action}
            elif action == "pin":
                ok = set_session_pinned(session_id, bool(body.get("pinned", False)))
                result = {"status": "ok" if ok else "not_found", "action": action}
            elif action in {"archive", "trash", "restore", "purge"}:
                handlers = {
                    "archive": archive_session,
                    "trash": trash_session,
                    "restore": restore_session,
                    "purge": purge_session,
                }
                ok = handlers[action](session_id)
                result = {"status": "ok" if ok else "not_found", "action": action}
            elif action == "export":
                try:
                    out_file = config.sandbox_root() / f"xiaopu-session-{session_id[:8]}.md"
                    exported = export_session(session_id, out_file)
                    result = {"status": "ok", "action": action, "path": str(exported)}
                except ValueError:
                    result = {"status": "not_found", "action": action}
            else:
                self.send_error(400, f"Unsupported session action: {action}")
                return
            self._send_json(result)
            return

        if path == "/api/sessions/batch":
            ids = [str(i) for i in body.get("ids", [])]
            action = body.get("action", "")
            if not ids or action not in {"archive", "trash", "restore", "purge"}:
                self.send_error(400, "Missing ids or unsupported action")
                return
            result = batch_session_action(ids, action)
            self._send_json({"status": "ok", **result})
            return

        if path == "/api/cancel":
            self.harness.request_cancel()
            self._send_json({"status": "cancelled"})
            return

        if path == "/api/ppt/verify":
            try:
                res = dispatch("ppt_check", json.dumps({"policy": "auto"}), self.harness)
                self._send_json({"result": res})
            except Exception as e:
                self._send_json({"result": f"校验失败: {e}"}, status=500)
            return

        if path == "/api/ppt/save":
            save_path = body.get("path", "presentation.pptx")
            try:
                res = dispatch("ppt_save", json.dumps({"path": save_path}), self.harness)
                self._send_json({"result": res})
            except Exception as e:
                self._send_json({"result": f"保存失败: {e}"}, status=500)
            return

        if path == "/api/ppt/undo":
            try:
                res = self.harness.undo()
                self._send_json({"result": res})
            except Exception as e:
                self._send_json({"result": f"撤销失败: {e}"}, status=500)
            return

        if path == "/api/ppt/apply_content":
            text_content = (body.get("text_content") or "").strip()
            if not text_content:
                self._send_json({"status": "error", "error": "文本内容为空"}, status=400)
                return
            ws_param = body.get("workspace")
            session_param = body.get("session_id") or body.get("session")
            file_param = body.get("deck_path") or body.get("file") or body.get("path")
            ws_root = Path(ws_param) if ws_param else config.sandbox_root()
            deck_path = _find_active_deck(ws_root, session_id=session_param, specific_file=file_param, harness=self.harness)
            if not deck_path or not deck_path.exists():
                self._send_json({"status": "error", "error": "当前会话没有可更新的 PPT 文件"}, status=404)
                return
            try:
                result = _apply_deck_text_content(deck_path, text_content)
                live_paths = {
                    str(Path(value).resolve())
                    for value in (
                        getattr(self.harness, "deck_working_path", None),
                        getattr(self.harness, "deck_path", None),
                    )
                    if value
                }
                if str(deck_path.resolve()) in live_paths:
                    from pptx import Presentation
                    self.harness.deck = Presentation(str(deck_path))
                self._send_json({"status": "ok", "deck_path": str(deck_path.resolve()), **result})
            except ValueError as exc:
                self._send_json({"status": "error", "error": str(exc)}, status=409)
            except Exception as exc:
                self._send_json({"status": "error", "error": f"更新 PPT 失败: {exc}"}, status=500)
            return

        if path == "/api/session/export":
            session_id = body.get("session_id")
            if not session_id:
                rec = save_session(self.harness)
                session_id = rec.id
            out_file = config.sandbox_root() / f"xiaopu-session-{session_id[:8]}.md"
            exported = export_session(session_id, out_file)
            self._send_json({"path": str(exported)})
            return

        if path == "/api/goal":
            objective = body.get("objective", "")
            if objective:
                res = self.harness.start_goal(objective)
                self._send_json({"result": res})
                return
            self.send_error(400, "Missing objective")
            return

        if path == "/api/settings":
            updates = {}
            if "provider" in body and body["provider"].strip():
                updates["PROVIDER"] = str(body["provider"]).strip().lower()
            if "api_base" in body and body["api_base"].strip():
                updates["OPENAI_BASE_URL"] = str(body["api_base"]).strip()
                updates["ANTHROPIC_BASE_URL"] = str(body["api_base"]).strip()
            if "api_key" in body and body["api_key"].strip():
                updates["OPENAI_API_KEY"] = str(body["api_key"]).strip()
                updates["ANTHROPIC_API_KEY"] = str(body["api_key"]).strip()
            if "model" in body and body["model"].strip():
                updates["OPENAI_MODEL"] = str(body["model"]).strip()
                updates["ANTHROPIC_MODEL"] = str(body["model"]).strip()
            if "models_csv" in body:
                updates["XIAOPU_MODELS"] = str(body["models_csv"]).strip()
            if "reasoning_effort" in body and body["reasoning_effort"].strip():
                updates["REASONING_EFFORT"] = str(body["reasoning_effort"]).strip()
            if "command_policy" in body and body["command_policy"].strip():
                updates["XIAOPU_COMMAND_POLICY"] = str(body["command_policy"]).strip()

            if updates:
                config.update_env_settings(updates)
                try:
                    if hasattr(self.harness, "llm"):
                        from agent.llm import create_llm
                        self.harness.llm = create_llm(model=config.model())
                except Exception:
                    pass

            self._send_json({"status": "ok", "settings": body})
            return

        if path == "/api/chat":
            self._handle_chat_stream(body)
            return

        self.send_error(404, "Not Found")

    def do_DELETE(self) -> None:
        parsed = urllib.parse.urlparse(self.path)
        path = parsed.path
        if path.startswith("/api/session/"):
            session_id = path[len("/api/session/"):]
            mode = urllib.parse.parse_qs(parsed.query).get("mode", ["trash"])[0]
            if mode == "purge":
                ok = delete_session(session_id)
                status = "deleted"
            else:
                ok = trash_session(session_id)
                status = "trashed"
            self._send_json({"status": status if ok else "not_found"})
            return
        self.send_error(404, "Not Found")

    def _handle_chat_stream(self, body: dict) -> None:
        # The browser frontend posts `prompt`; older API clients post `task`.
        # Accept both so the web composer always reaches the harness.
        task = body.get("task") or body.get("prompt") or ""
        session_id = body.get("session_id")
        model = body.get("model")
        permission = body.get("permission") or body.get("command_policy")
        reasoning_effort = body.get("reasoning_effort")

        if reasoning_effort:
            os.environ["REASONING_EFFORT"] = str(reasoning_effort).strip()

        if model:
            if config.provider() == "anthropic":
                os.environ["ANTHROPIC_MODEL"] = model
            else:
                os.environ["OPENAI_MODEL"] = model
            if getattr(self.harness, "llm", None):
                self.harness.llm.model = model

        if permission:
            config.set_command_policy(permission)
            os.environ["COMMAND_POLICY"] = permission

        # If resuming a specific session ID
        if session_id and (not hasattr(self.harness, "session") or self.harness.session.id != session_id):
            payload = load_session(session_id)
            if payload:
                restore_harness(self.harness, payload)

        # Set up SSE Streaming headers
        self.send_response(200)
        self.send_header("Content-Type", "text/event-stream; charset=utf-8")
        self.send_header("Cache-Control", "no-cache")
        self.send_header("Connection", "keep-alive")
        self.send_header("Access-Control-Allow-Origin", "*")
        self.end_headers()

        event_queue: queue.Queue = queue.Queue()

        def stream_token(piece: str):
            event_queue.put({"type": "token", "content": piece})

        def reasoning_token(piece: str):
            event_queue.put({"type": "reasoning", "content": piece})

        def runtime_listener(event: RuntimeEvent):
            p = getattr(event, "payload", {})
            kind = getattr(event, "kind", None)
            if kind == EventKind.TOOL_STARTED:
                event_queue.put({
                    "type": "tool_started",
                    "payload": {"tool": p.get("tool"), "arguments": p.get("arguments", "")[:240]},
                })
            elif kind == EventKind.MODEL_RESPONSE:
                event_queue.put({
                    "type": "model_response",
                    "payload": {
                        "content": str(p.get("content") or ""),
                        "reasoning_content": str(p.get("reasoning_content") or ""),
                        "has_tool_calls": bool(p.get("has_tool_calls")),
                        "phase": p.get("phase"),
                    },
                })
            elif kind == EventKind.TOOL_COMPLETED:
                event_queue.put({
                    "type": "tool_completed",
                    "payload": {"tool": p.get("tool"), "output": str(p.get("output", ""))[:400]},
                })
            elif kind == EventKind.TOOL_FAILED:
                event_queue.put({
                    "type": "tool_failed",
                    "payload": {"tool": p.get("tool"), "error": str(p.get("error", ""))[:240]},
                })
            elif kind == EventKind.PHASE_CHANGED:
                event_queue.put({
                    "type": "phase_changed",
                    "payload": {"from_phase": p.get("from_phase"), "to_phase": p.get("to_phase")},
                })

        self.harness.stream_callback = stream_token
        self.harness.reasoning_callback = reasoning_token
        unsub = self.harness.subscribe(runtime_listener)

        worker_done = threading.Event()
        worker_error: list[str] = []

        def worker():
            try:
                res = self.harness.run(task)
                event_queue.put({"type": "result", "content": res})
                rec = save_session(self.harness)
                event_queue.put({"type": "session_saved", "session_id": rec.id})
            except Exception as exc:
                worker_error.append(str(exc))
                event_queue.put({"type": "error", "content": f"{type(exc).__name__}: {exc}"})
            finally:
                worker_done.set()

        threading.Thread(target=worker, name="xiaopu-web-worker", daemon=True).start()

        try:
            while not worker_done.is_set() or not event_queue.empty():
                try:
                    ev = event_queue.get(timeout=0.1)
                    line = f"data: {json.dumps(ev, ensure_ascii=False)}\n\n".encode("utf-8")
                    self.wfile.write(line)
                    self.wfile.flush()
                except queue.Empty:
                    continue
            self.wfile.write(b"data: [DONE]\n\n")
            self.wfile.flush()
        except (BrokenPipeError, ConnectionResetError):
            pass
        finally:
            unsub()
            self.close_connection = True


def _kill_existing_server_on_port(port: int = 8765) -> None:
    """Ensure any stale/zombie process occupying port is terminated before starting."""
    if sys.platform != "win32":
        return
    import subprocess
    import os
    import time
    current_pid = os.getpid()
    try:
        cmd = f'netstat -ano | findstr :{port}'
        res = subprocess.run(cmd, shell=True, capture_output=True, text=True, errors="ignore")
        pids_to_kill = set()
        for line in (res.stdout or "").strip().splitlines():
            parts = line.strip().split()
            if len(parts) >= 5 and "LISTENING" in parts[3].upper():
                try:
                    pid = int(parts[4])
                    if pid != current_pid and pid > 0:
                        pids_to_kill.add(pid)
                except ValueError:
                    pass
        for pid in pids_to_kill:
            try:
                subprocess.run(f"taskkill /F /PID {pid}", shell=True, capture_output=True)
            except Exception:
                pass
        if pids_to_kill:
            time.sleep(0.6)
    except Exception:
        pass


def run_web_gui(host: str = "127.0.0.1", port: int = 8765, open_browser: bool = True, model: str | None = None) -> None:
    _kill_existing_server_on_port(port)
    XiaopuWebHandler.harness = Harness(model=model, interactive=True, max_steps=50)
    server = ThreadingHTTPServer((host, port), XiaopuWebHandler)
    url = f"http://{host}:{port}"
    print(f"\n=======================================================")
    print(f"  报一 Baoyi Web GUI is running at: {url}")
    print(f"=======================================================\n")

    if open_browser:
        try:
            import shutil
            from pathlib import Path
            browser_bin = None
            if sys.platform == "win32":
                candidates = [
                    shutil.which("msedge"),
                    r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
                    r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
                    shutil.which("chrome"),
                    r"C:\Program Files\Google\Chrome\Application\chrome.exe",
                    r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
                ]
                for cand in candidates:
                    if cand and Path(cand).is_file():
                        browser_bin = str(cand)
                        break

            if browser_bin:
                subprocess.Popen([browser_bin, f"--app={url}"])
            else:
                webbrowser.open(url)
        except Exception:
            try:
                webbrowser.open(url)
            except Exception:
                pass

    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\nStopping Baoyi Web GUI server...")
        server.shutdown()


if __name__ == "__main__":
    run_web_gui()
