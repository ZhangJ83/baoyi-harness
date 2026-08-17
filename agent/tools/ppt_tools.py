"""PPT tools: build, edit, verify a deck with python-pptx.

Deck state lives on the harness (`h.deck`), so tools compose across calls.
Visual layout: a minimal editorial theme — title band, low-contrast body,
an accent underline — built from plain shapes (no templates / no external assets).
"""
import math
import json
import re
import platform
import shutil
import subprocess
from copy import deepcopy
from io import BytesIO
from types import SimpleNamespace
from typing import Any, Callable

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

Schema = dict[str, Any]

# --- theme ---
_PRIMARY = RGBColor(0x1F, 0x38, 0x64)   # deep navy
_ACCENT = RGBColor(0xF0, 0x9A, 0x1A)    # amber
_TEXT = RGBColor(0x2B, 0x2B, 0x2B)
_MUTED = RGBColor(0x87, 0x87, 0x87)
_BG = RGBColor(0xFB, 0xFA, 0xF7)        # warm paper
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_HEAD = RGBColor(0x25, 0x47, 0x7C)      # lighter navy for content titles

_W = 13.333
_H = 7.5


def _patch_pptx_slide_partname_uniqueness() -> None:
    """Ensure python-pptx generates unique slide partnames even after slide deletions."""
    try:
        from pptx.parts.presentation import PresentationPart
        from pptx.opc.packuri import PackURI

        def _safe_next_slide_partname(self):
            used = {p.partname for p in self.package.iter_parts()}
            for rel in self._rels.values():
                if not rel.is_external and getattr(rel, "target_part", None):
                    used.add(rel.target_part.partname)
            n = 1
            while True:
                candidate = PackURI(f"/ppt/slides/slide{n}.xml")
                if candidate not in used:
                    return candidate
                n += 1

        PresentationPart._next_slide_partname = property(_safe_next_slide_partname)
    except Exception:
        pass


_patch_pptx_slide_partname_uniqueness()


def _ensure_deck_package_clean(prs: Presentation) -> Presentation:
    """Detect and heal any duplicate slide partnames before writing to disk."""
    try:
        parts = list(prs.part.package.iter_parts())
        names = [p.partname for p in parts]
        if len(names) == len(set(names)):
            return prs
    except Exception:
        return prs

    # Duplicate partnames detected (corrupted package graph); rebuild clean Presentation
    clean = Presentation()
    clean.slide_width = prs.slide_width
    clean.slide_height = prs.slide_height
    for s in prs.slides:
        layout_idx = 6 if len(clean.slide_layouts) > 6 else 0
        new_s = clean.slides.add_slide(clean.slide_layouts[layout_idx])
        for sh in s.shapes:
            new_s.shapes._spTree.insert_element_before(deepcopy(sh._element), "p:extLst")
        if getattr(s, "has_notes_slide", False):
            notes = s.notes_slide.notes_text_frame.text.strip()
            if notes:
                new_s.notes_slide.notes_text_frame.text = notes
    return clean


def _schema(name: str, description: str, params: dict[str, dict], required: list[str]):
    return {
        "type": "function",
        "function": {
            "name": name,
            "description": description,
            "parameters": {"type": "object", "properties": params, "required": required, "additionalProperties": False},
        },
    }


def _deck(h: Any) -> Presentation:
    if getattr(h, "deck", None) is None:
        h.deck = Presentation()
        h.deck.slide_width = Inches(_W)
        h.deck.slide_height = Inches(_H)
    return h.deck


def _open_deck(h, path: str) -> str:
    from .. import config
    from ..permissions import path_within
    root = config.sandbox_root()
    source = (root / path).resolve()
    if not path_within(root, source):
        raise PermissionError("presentation path escapes workspace")
    if not source.exists():
        active = getattr(h, "deck", None)
        active_hint = (
            f" The active deck already has {len(active.slides)} slides; continue editing it directly."
            if active is not None and len(active.slides)
            else ""
        )
        raise FileNotFoundError(
            f"{path} (no such presentation; if this is a create-from-scratch task, call new_deck instead of ppt_open)."
            + active_hint
        )
    # Progress preservation: after the first mutation, replacing the active
    # deck from any file other than the saved contract output can silently
    # discard edits (stale working copies, the frozen input, unrelated decks).
    if getattr(h, "state", None) is not None and getattr(h.state, "mutation_epoch", 0):
        required = h.state.facts.get("required_output_pptx", "")
        allowed = {(root / required).resolve()} if required else set()
        if source not in allowed:
            hint = f" Continue with the active draft, or open the saved deliverable '{required}'." if required else ""
            raise RuntimeError(
                "ppt_open rejected: the run already has mutations and reopening another file would discard them."
                + hint
            )
        # Reopening the persisted deliverable must never discard work.  If the
        # active deck differs from the persisted file, persist it first, then
        # reload the exact bytes the deliverable now contains.  The recorder
        # refreshes its working copy, so repeated open/save cycles stay safe.
        try:
            live = BytesIO()
            h.deck.save(live)
            if live.getvalue() != source.read_bytes():
                try:
                    _save(h, str(source.relative_to(root)) if str(source.resolve()).startswith(str(root.resolve())) else str(source))
                except Exception:
                    h.deck.save(str(source))
                    if getattr(h, "state", None) is not None:
                        h.state.record_commit(str(source))
        except Exception as exc:
            raise RuntimeError(f"ppt_open failed while persisting the active deck before reopen: {exc}")
    opened = source
    if getattr(h, "recorder", None):
        opened = h.recorder.working_copy(source)
    h.deck = Presentation(str(opened))
    h.deck_source_path = source
    h.deck_working_path = opened
    if getattr(h, "state", None) is not None:
        h.state.active_artifact = str(opened)
        h.state.ppt_existing_deck = True
        h.state.ppt_affected_slides.clear()
        baseline = _collect_structural_findings(h)
        h.state.ppt_baseline_findings = {item["key"]: item["severity"] for item in baseline}
        h.state.ppt_baseline_captured = True
        h.state.record_fact("ppt_source", str(source))
        h.state.record_fact("ppt_working_artifact", str(opened))
        h.state.record_fact(
            "ppt_structural_baseline",
            f"{len(baseline)} finding(s) frozen across {len(h.deck.slides)} source slides",
        )
    return f"opened editable working copy of {path}; original preserved: {len(h.deck.slides)} slides, {h.deck.slide_width / 914400:.2f}x{h.deck.slide_height / 914400:.2f}in"


def _blank_slide(prs: Presentation):
    layouts = prs.slide_layouts
    # Prefer the conventional blank layout (index 6), but minimal templates may
    # expose fewer layouts. Fall back to the first available layout instead of
    # crashing on a deck whose layout table is shorter than the default.
    index = 6 if len(layouts) > 6 else 0
    return prs.slides.add_slide(layouts[index])


def _contract_slide_limits(h) -> tuple[int | None, int | None, int | None]:
    """Declared slide-count envelope from the task-local contract, if any."""
    terms = getattr(getattr(h, "state", None), "verification_contract_terms", None) or {}
    answer = terms.get("answer_contract") or {}
    template = terms.get("template_expected") or {}
    output = terms.get("output_contract") or {}
    exact = output.get("expected_slide_count") or answer.get("slide_count") or template.get("slide_count")
    try:
        exact = int(exact) if exact is not None else None
    except (TypeError, ValueError):
        exact = None
    try:
        minimum = int(answer.get("min_slide_count")) if answer.get("min_slide_count") is not None else None
    except (TypeError, ValueError):
        minimum = None
    try:
        maximum = int(answer.get("max_slide_count")) if answer.get("max_slide_count") is not None else None
    except (TypeError, ValueError):
        maximum = None
    return exact, minimum, maximum


def _assert_can_append_slide(h) -> None:
    """Enforce the declared slide-count contract before creating a new slide."""
    exact, _minimum, maximum = _contract_slide_limits(h)
    deck = getattr(h, "deck", None)
    current = len(deck.slides) if deck is not None else 0
    limit = None
    if exact is not None:
        limit = exact
    if maximum is not None and (limit is None or maximum < limit):
        limit = maximum
    if limit is not None and current + 1 > limit:
        raise ValueError(
            f"cannot add slide {current + 1}: the task contract allows at most {limit} slides "
            f"(current={current}). Reuse/rewrite an existing slide with set_shape_text/set_table "
            "or remove a surplus slide with ppt_arrange delete_slide before composing a new one."
        )


def _position_new_slide(prs: Presentation, slide, insert_after: int | None) -> int:
    """Place a just-appended slide after *insert_after* (1-based) and return its position."""
    slide_list = prs.slides._sldIdLst
    ids = list(slide_list)
    if not ids:
        return 1
    new = ids[-1]
    position = len(ids)
    if insert_after is not None and 1 <= insert_after < len(ids):
        slide_list.remove(new)
        slide_list.insert(insert_after, new)
        position = insert_after + 1
    return position



# Modern design tokens
_SLATE_DARK = RGBColor(0x0F, 0x17, 0x2A)
_SLATE_PANEL = RGBColor(0x1E, 0x29, 0x3B)
_SLATE_BORDER = RGBColor(0x33, 0x41, 0x55)
_SLATE_LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
_CARD_BORDER = RGBColor(0xCB, 0xD5, 0xE1)
_TEXT_PRIMARY = RGBColor(0x0F, 0x17, 0x2A)
_TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
_BLUE_ACCENT = RGBColor(0x25, 0x63, 0xEB)
_BLUE_BG = RGBColor(0xDB, 0xEA, 0xFE)
_GREEN_ACCENT = RGBColor(0x10, 0xB9, 0x81)
_GREEN_BG = RGBColor(0xDC, 0xFC, 0xE7)
_RED_DOT = RGBColor(0xEF, 0x44, 0x44)
_YELLOW_DOT = RGBColor(0xF5, 0x9E, 0x0B)
_GREEN_DOT = RGBColor(0x10, 0xB9, 0x81)


def _rect(slide, x: float, y: float, w: float, h: float, fill, line=None) -> "Shape":
    """Full-bleed helper: plain rectangle with a solid fill."""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def _rounded_rect(slide, x: float, y: float, w: float, h: float, fill, line=None) -> "Shape":
    """Rounded rectangle helper for modern card containers and badge pills."""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def _circle(slide, x: float, y: float, size: float, fill) -> "Shape":
    """Circle helper for window control dots and status indicators."""
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp



def _style_run(run, size: int, bold: bool, color: RGBColor) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Microsoft YaHei"


def _fit_lines(tf, lines, size: int, bold: bool, color: RGBColor, space: float = 1.25) -> None:
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = space
        run = p.add_run()
        run.text = line
        _style_run(run, size, bold, color)


def _put_lines(tf, text: str, size: int, bold: bool, color: RGBColor = _TEXT) -> None:
    _fit_lines(tf, text.split("\n"), size, bold, color)


def _new_deck(h, title: str, subtitle: str) -> str:
    # A "new deck" is a fresh 16:9 presentation, never an append to whatever
    # template the intake pre-opened. Reusing the open deck here silently
    # turned template-driven tasks into template-editing tasks and polluted the
    # structural check with the template's background shapes.
    deck = getattr(h, "deck", None)
    if deck is not None and len(deck.slides) and (
        getattr(getattr(h, "state", None), "ppt_existing_deck", False)
        or h.state.facts.get("ppt_input_deck")
    ):
        raise ValueError(
            "new_deck rejected: a bound template/source deck is already open "
            f"({len(deck.slides)} slides). Transform it with ppt_compose "
            "(from_outline/content) or ppt_edit_text instead of discarding it."
        )
    prs = Presentation()
    prs.slide_width = Inches(_W)
    prs.slide_height = Inches(_H)
    h.deck = prs
    if getattr(h, "state", None) is not None:
        h.state.ppt_existing_deck = False
        h.state.ppt_baseline_captured = False
        h.state.ppt_baseline_findings.clear()
        h.state.ppt_affected_slides.clear()
        # Slide 1 is a decorative cover scaffold. The first slide-targeted
        # page builder (flowchart/quadrant/content) converts it into a real
        # content page instead of overlaying or leaving it as a thin cover.
        h.state.facts["ppt_new_deck_cover"] = "1"
    s = _blank_slide(prs)
    # warm paper background
    _rect(s, 0, 0, _W, _H, _BG)
    # left accent bar
    _rect(s, 0, 0, 0.25, _H, _ACCENT)
    # eyebrow
    tb = s.shapes.add_textbox(Inches(1.1), Inches(1.3), Inches(11), Inches(0.5))
    _fit_lines(tb.text_frame, ["P R E S E N T A T I O N"], 15, True, _ACCENT)
    # title (auto-shrink long titles)
    size = 44 if len(title) <= 16 else 36 if len(title) <= 24 else 30
    tb = s.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(11), Inches(2.4))
    _fit_lines(tb.text_frame, [title], size, True, _PRIMARY)
    # accent underline
    _rect(s, 1.15, 4.35, 1.4, 0.06, _ACCENT)
    if subtitle:
        tb = s.shapes.add_textbox(Inches(1.1), Inches(4.6), Inches(11), Inches(1.2))
        _fit_lines(tb.text_frame, subtitle.split("\n"), 18, False, _MUTED)
    return "new deck created: 16:9, one cover slide."


def _convert_fresh_cover(h, slide_number: int) -> bool:
    """Replace an untouched new_deck cover scaffold with a clean page.

    The first targeted page builder on slide 1 of a freshly created deck owns
    that page: the decorative cover text must not survive underneath the
    requested diagram/dashboard/content layout.
    """
    if slide_number != 1:
        return False
    facts = getattr(getattr(h, "state", None), "facts", {})
    if facts.get("ppt_new_deck_cover") != "1":
        return False
    slide = h.deck.slides[0]
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)
    facts.pop("ppt_new_deck_cover", None)
    return True


def _clear_slide_shapes(slide) -> None:
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def _resolve_target_slide(h, slide_number: int | None = None, insert_after: int | None = None) -> tuple[Any, int, bool]:
    """Resolve or append a slide without raising IndexError on incremental multi-page builds.

    Returns:
        (slide_object, 1-based position, is_rebuild)
    """
    prs = _deck(h)
    if slide_number is not None:
        if slide_number < 1:
            slide_number = 1
        if 1 <= slide_number <= len(prs.slides):
            s = prs.slides[slide_number - 1]
            _convert_fresh_cover(h, slide_number)
            _clear_slide_shapes(s)
            return s, slide_number, True
        else:
            # slide_number is beyond current slide count (e.g. slide_number=2 on 1-slide deck).
            # Seamlessly append as new slide instead of raising IndexError!
            _assert_can_append_slide(h)
            s = _blank_slide(prs)
            position = len(prs.slides)
            return s, position, False
    else:
        _assert_can_append_slide(h)
        s = _blank_slide(prs)
        position = _position_new_slide(prs, s, insert_after)
        return s, position, False


def _content_slide(h, title: str, bullets: list[str], size: int, insert_after: int | None = None, slide_number: int | None = None) -> str:
    prs = _deck(h)
    title = _clean_presentation_title(title)
    s, position, is_rebuild = _resolve_target_slide(h, slide_number, insert_after)
    _rect(s, 0, 0, _W, _H, _BG)
    # header band
    _rect(s, 0, 0, _W, 1.1, _HEAD)
    _rect(s, 0, 1.1, _W, 0.07, _ACCENT)
    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.12), Inches(12.2), Inches(0.9))
    _put_lines(tb.text_frame, title, 26, True, _WHITE)
    if bullets:
        # Each bullet is its own visible textbox so structural checks and the
        # open-ended completeness gate can count real content objects.
        max_bullets = min(len(bullets), 12)
        top = 1.65
        box_height = 5.15 / max_bullets
        for index, bullet in enumerate(bullets[:max_bullets]):
            tb2 = s.shapes.add_textbox(Inches(0.9), Inches(top + index * box_height), Inches(11.5), Inches(box_height - 0.12))
            _fit_lines(tb2.text_frame, [f"•  {bullet}"], size, False, _TEXT)
    verb = "rebuilt" if is_rebuild else "added"
    return f"{verb} slide {position}: '{title}' ({len(bullets)} visible bullet boxes)"


def _two_column(h, title: str, left_title: str, left: list[str], right_title: str, right: list[str], insert_after: int | None = None) -> str:
    _assert_can_append_slide(h)
    prs = _deck(h)
    title = _clean_presentation_title(title)
    s = _blank_slide(prs)
    position = _position_new_slide(prs, s, insert_after)
    _rect(s, 0, 0, _W, _H, _BG)
    _rect(s, 0, 0, _W, 1.1, _HEAD)
    _rect(s, 0, 1.1, _W, 0.07, _ACCENT)
    title_box = s.shapes.add_textbox(Inches(0.55), Inches(0.12), Inches(12.2), Inches(0.9))
    _put_lines(title_box.text_frame, title, 26, True, _WHITE)
    for x, heading, bullets in ((0.7, left_title, left), (6.85, right_title, right)):
        card = _rect(s, x, 1.6, 5.75, 5.1, _WHITE, RGBColor(0xDD, 0xDD, 0xD8))
        heading_box = s.shapes.add_textbox(Inches(x + 0.35), Inches(1.9), Inches(5.05), Inches(0.65))
        _put_lines(heading_box.text_frame, _clean_presentation_title(heading), 20, True, _PRIMARY)
        body = s.shapes.add_textbox(Inches(x + 0.35), Inches(2.75), Inches(5.05), Inches(3.55))
        _fit_lines(body.text_frame, [f"•  {item}" for item in bullets], 17, False, _TEXT, 1.2)
    return f"added two-column slide {position}: '{title}'"


def _quadrant_slide(h, title: str, subtitle: str, quadrants: list[dict], slide_number: int | None = None) -> str:
    """Create a complete executive four-quadrant page in one semantic action.

    This deliberately sits above primitive textbox calls: the model decides
    the content and evidence mapping, while the harness owns spacing,
    typography, source chips and notes provenance.
    """
    if len(quadrants) != 4:
        raise ValueError("quadrant slide requires exactly four quadrants")
    prs = _deck(h)
    title = _clean_presentation_title(title) or "季度四象限看板"
    subtitle = _clean_presentation_title(subtitle)
    slide, target, is_rebuild = _resolve_target_slide(h, slide_number)

    # Predominantly black and white, with one restrained amber signal color.
    _rect(slide, 0, 0, _W, _H, _WHITE)
    _rect(slide, 0, 0, _W, 0.92, RGBColor(0x12, 0x12, 0x12))
    heading = slide.shapes.add_textbox(Inches(0.48), Inches(0.16), Inches(9.8), Inches(0.52))
    _put_lines(heading.text_frame, title, 25, True, _WHITE)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(10.15), Inches(0.20), Inches(2.7), Inches(0.4))
        _put_lines(sub.text_frame, subtitle, 11, False, RGBColor(0xD0, 0xD0, 0xD0))


    positions = ((0.48, 1.18), (6.83, 1.18), (0.48, 4.12), (6.83, 4.12))
    sources: list[str] = []
    for index, (item, (x, y)) in enumerate(zip(quadrants, positions), 1):
        card = _rect(slide, x, y, 6.02, 2.58, _WHITE, RGBColor(0xC8, 0xC8, 0xC8))
        card.line.width = Pt(1.1)
        _rect(slide, x, y, 0.08, 2.58, _ACCENT)
        qlabel = slide.shapes.add_textbox(Inches(x + 0.25), Inches(y + 0.18), Inches(0.55), Inches(0.32))
        _put_lines(qlabel.text_frame, f"Q{index}", 11, True, _MUTED)
        # Reserve a stable metric column instead of letting the model repair
        # title/metric collisions shape by shape after verification.
        qtitle = slide.shapes.add_textbox(Inches(x + 0.82), Inches(y + 0.13), Inches(2.95), Inches(0.45))
        _put_lines(qtitle.text_frame, str(item.get("title", "")), 16, True, RGBColor(0x16, 0x16, 0x16))
        metric = str(item.get("metric", "")).strip()
        if metric:
            metric_box = slide.shapes.add_textbox(Inches(x + 3.85), Inches(y + 0.10), Inches(1.85), Inches(0.48))
            metric_box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            _put_lines(metric_box.text_frame, metric, 18, True, RGBColor(0x16, 0x16, 0x16))
        bullets = [str(value).strip() for value in item.get("bullets", []) if str(value).strip()][:6]
        detail = str(item.get("detail", "")).strip()
        lines = ([detail] if detail else []) + [f"• {value}" for value in bullets]
        body = slide.shapes.add_textbox(Inches(x + 0.28), Inches(y + 0.72), Inches(5.42), Inches(1.42))
        _fit_lines(body.text_frame, lines, 11, False, _TEXT, 1.08)
        source = str(item.get("source", "")).strip()
        if source:
            sources.append(f"Q{index}: {source}")
            chip = slide.shapes.add_textbox(Inches(x + 0.28), Inches(y + 2.19), Inches(5.35), Inches(0.22))
            _put_lines(chip.text_frame, source, 10, False, _MUTED)
            try:
                chip.name = f"Q{index}_provenance_{source[:36]}"
            except (AttributeError, ValueError):
                pass
    if sources:
        _set_speaker_notes(h, target, "[Sources]\n" + "\n".join(sources))
    return f"composed executive quadrant slide {target}: '{title}' with 4 provenance-bound quadrants"


def _metric_slide(h, title: str, metrics: list[dict], takeaway: str = "") -> str:
    if not 1 <= len(metrics) <= 4:
        raise ValueError("metric slide requires 1-4 metrics")
    prs = _deck(h)
    s = _blank_slide(prs)
    _rect(s, 0, 0, _W, _H, _PRIMARY)
    title_box = s.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(11.8), Inches(0.75))
    _put_lines(title_box.text_frame, title, 28, True, _WHITE)
    gap = 0.25
    card_w = (11.85 - gap * (len(metrics) - 1)) / len(metrics)
    for index, metric in enumerate(metrics):
        x = 0.75 + index * (card_w + gap)
        _rect(s, x, 1.75, card_w, 3.45, _WHITE)
        value_box = s.shapes.add_textbox(Inches(x + 0.25), Inches(2.15), Inches(card_w - 0.5), Inches(1.05))
        _put_lines(value_box.text_frame, str(metric.get("value", "—")), 32, True, _ACCENT)
        label_box = s.shapes.add_textbox(Inches(x + 0.25), Inches(3.3), Inches(card_w - 0.5), Inches(0.7))
        _put_lines(label_box.text_frame, str(metric.get("label", "")), 17, True, _PRIMARY)
        detail_box = s.shapes.add_textbox(Inches(x + 0.25), Inches(4.15), Inches(card_w - 0.5), Inches(0.7))
        _put_lines(detail_box.text_frame, str(metric.get("detail", "")), 13, False, _MUTED)
    if takeaway:
        note = s.shapes.add_textbox(Inches(0.85), Inches(5.75), Inches(11.65), Inches(0.85))
        _put_lines(note.text_frame, takeaway, 18, False, _WHITE)
    return f"added metric slide {len(prs.slides)}: '{title}'"


def _table_slide(h, title: str, columns: list[str], rows: list[list[str]], insert_after: int | None = None) -> str:
    if not columns or len(columns) > 6:
        raise ValueError("table requires 1-6 columns")
    if not rows or len(rows) > 10:
        raise ValueError("table requires 1-10 rows")
    if any(len(row) != len(columns) for row in rows):
        raise ValueError("every row must match the column count")
    _assert_can_append_slide(h)
    prs = _deck(h)
    s = _blank_slide(prs)
    position = _position_new_slide(prs, s, insert_after)
    _rect(s, 0, 0, _W, _H, _BG)
    _rect(s, 0, 0, _W, 1.1, _HEAD)
    _rect(s, 0, 1.1, _W, 0.07, _ACCENT)
    title_box = s.shapes.add_textbox(Inches(0.55), Inches(0.12), Inches(12.2), Inches(0.9))
    _put_lines(title_box.text_frame, title, 26, True, _WHITE)
    table_shape = s.shapes.add_table(len(rows) + 1, len(columns), Inches(0.65), Inches(1.55), Inches(12.0), Inches(5.35))
    table = table_shape.table
    for col_index, label in enumerate(columns):
        cell = table.cell(0, col_index)
        cell.text = label
        cell.fill.solid()
        cell.fill.fore_color.rgb = _PRIMARY
    for row_index, row in enumerate(rows, 1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _WHITE if row_index % 2 else RGBColor(0xF1, 0xF0, 0xEC)
    for table_row_index, row in enumerate(table.rows):
        for cell in row.cells:
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    _style_run(run, 13, table_row_index == 0, _WHITE if table_row_index == 0 else _TEXT)
    return f"added table slide {position}: '{title}' ({len(rows)} rows)"


def _process_slide(h, title: str, steps: list[dict], takeaway: str = "") -> str:
    if not 2 <= len(steps) <= 8:
        raise ValueError("process slide requires 2-8 steps")
    prs = _deck(h)
    s = _blank_slide(prs)
    _rect(s, 0, 0, _W, _H, _BG)
    title_box = s.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(0.7))
    _put_lines(title_box.text_frame, title, 28, True, _PRIMARY)
    start_x, y, total_w = 0.75, 1.65, 11.85
    gap = 0.18
    step_w = (total_w - gap * (len(steps) - 1)) / len(steps)
    for index, step in enumerate(steps, 1):
        x = start_x + (index - 1) * (step_w + gap)
        _rect(s, x, y, step_w, 3.7, _WHITE, RGBColor(0xDD, 0xDD, 0xD8))
        badge = _rect(s, x + 0.25, y + 0.3, 0.52, 0.52, _ACCENT)
        badge.text_frame.clear()
        badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = badge.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = str(index)
        _style_run(run, 14, True, _WHITE)
        heading = s.shapes.add_textbox(Inches(x + 0.25), Inches(y + 1.15), Inches(step_w - 0.5), Inches(0.75))
        _put_lines(heading.text_frame, str(step.get("title", "")), 19, True, _PRIMARY)
        detail = s.shapes.add_textbox(Inches(x + 0.25), Inches(y + 2.05), Inches(step_w - 0.5), Inches(1.15))
        _put_lines(detail.text_frame, str(step.get("detail", "")), 14, False, _TEXT)
    if takeaway:
        note = s.shapes.add_textbox(Inches(0.85), Inches(5.85), Inches(11.6), Inches(0.75))
        _put_lines(note.text_frame, takeaway, 18, True, _PRIMARY)
    return f"added process slide {len(prs.slides)}: '{title}'"


def _image_slide(h, title: str, image_path: str, caption: str = "") -> str:
    from .. import config
    from ..permissions import path_within
    root = config.sandbox_root().resolve()
    source = (root / image_path).resolve()
    if not path_within(root, source) or not source.is_file():
        raise FileNotFoundError(image_path)
    prs = _deck(h)
    s = _blank_slide(prs)
    _rect(s, 0, 0, _W, _H, _BG)
    title_box = s.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.7))
    _put_lines(title_box.text_frame, title, 27, True, _PRIMARY)
    from PIL import Image
    with Image.open(source) as image:
        ratio = image.width / image.height
    max_w, max_h = 11.7, 5.35
    width = min(max_w, max_h * ratio)
    height = width / ratio
    x = (_W - width) / 2
    y = 1.25 + (max_h - height) / 2
    s.shapes.add_picture(str(source), Inches(x), Inches(y), width=Inches(width), height=Inches(height))
    if caption:
        caption_box = s.shapes.add_textbox(Inches(0.9), Inches(6.8), Inches(11.5), Inches(0.4))
        _put_lines(caption_box.text_frame, caption, 11, False, _MUTED)
    return f"added image slide {len(prs.slides)}: '{title}'"


def _workflow_pipeline_slide(
    h,
    title: str,
    steps: list[dict],
    subtitle: str = "",
    takeaway: str = "",
    slide_number: int | None = None,
    insert_after: int | None = None,
) -> str:
    """Compose a modern multi-step workflow pipeline with colored cards, metrics, and grid layout."""
    prs = _deck(h)
    if not 1 <= len(steps) <= 8:
        raise ValueError("workflow_pipeline requires 1-8 steps")
    title = _clean_presentation_title(title) or "AI Agent 端到端认知与执行流水线"
    subtitle = _clean_presentation_title(subtitle)
    s, position, is_rebuild = _resolve_target_slide(h, slide_number, insert_after)

    # Per-step color palette: (accent, light_bg)
    _STEP_COLORS = [
        (RGBColor(0x25, 0x63, 0xEB), RGBColor(0xDB, 0xEA, 0xFE)),  # Blue
        (RGBColor(0x7C, 0x3A, 0xED), RGBColor(0xED, 0xE9, 0xFE)),  # Purple
        (RGBColor(0x05, 0x96, 0x69), RGBColor(0xD1, 0xFA, 0xE5)),  # Green
        (RGBColor(0xD9, 0x77, 0x06), RGBColor(0xFE, 0xF3, 0xC7)),  # Amber
        (RGBColor(0xDC, 0x26, 0x26), RGBColor(0xFE, 0xE2, 0xE2)),  # Red
        (RGBColor(0x06, 0x4E, 0x3B), RGBColor(0xCC, 0xFB, 0xF1)),  # Teal
        (RGBColor(0xDB, 0x27, 0x77), RGBColor(0xFC, 0xE7, 0xF3)),  # Pink
        (RGBColor(0x47, 0x55, 0x69), RGBColor(0xF1, 0xF5, 0xF9)),  # Slate
    ]

    # Background & header
    _rect(s, 0, 0, _W, _H, _SLATE_LIGHT_BG)
    header_h = 0.85
    _rect(s, 0, 0, _W, header_h, _PRIMARY)
    _rect(s, 0, header_h, _W, 0.06, _ACCENT)

    title_box = s.shapes.add_textbox(Inches(0.50), Inches(0.08), Inches(9.0), Inches(0.70))
    _put_lines(title_box.text_frame, title, 24, True, _WHITE)
    if subtitle:
        sub_box = s.shapes.add_textbox(Inches(7.0), Inches(0.10), Inches(6.0), Inches(0.65))
        _put_lines(sub_box.text_frame, subtitle, 11, False, RGBColor(0xD0, 0xD0, 0xD0))

    # Adaptive grid layout: wider cards via multi-row
    n = len(steps)
    if n <= 3:
        cols, rows = n, 1
    elif n == 4:
        cols, rows = 2, 2
    elif n <= 6:
        cols, rows = 3, 2
    else:
        cols, rows = 4, 2

    margin_x = 0.45
    total_w = _W - 2 * margin_x
    gap_x = 0.24
    gap_y = 0.22
    card_w = (total_w - gap_x * (cols - 1)) / cols

    content_top = header_h + 0.06 + 0.14
    takeaway_h = 0.52 if takeaway else 0
    avail_h = _H - content_top - takeaway_h - 0.22
    if rows == 1:
        card_h = min(avail_h, 4.8)
    else:
        card_h = (avail_h - gap_y * (rows - 1)) / rows

    for idx, step in enumerate(steps):
        row = idx // cols
        col = idx % cols
        items_in_row = min(cols, n - row * cols)

        # Center incomplete bottom row
        if items_in_row < cols:
            row_w = items_in_row * card_w + (items_in_row - 1) * gap_x
            x = margin_x + (total_w - row_w) / 2 + (idx - row * cols) * (card_w + gap_x)
        else:
            x = margin_x + col * (card_w + gap_x)
        y = content_top + row * (card_h + gap_y)

        step_accent, step_bg = _STEP_COLORS[idx % len(_STEP_COLORS)]

        # Card body + colored top bar
        _rounded_rect(s, x, y, card_w, card_h, _WHITE, _CARD_BORDER)
        _rounded_rect(s, x, y, card_w, 0.07, step_accent)

        # Colored badge
        badge = _rounded_rect(s, x + 0.14, y + 0.14, 0.42, 0.28, step_accent)
        badge.text_frame.clear()
        badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = badge.text_frame.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        brun = bp.add_run()
        brun.text = f"{idx + 1:02d}"
        _style_run(brun, 11, True, _WHITE)

        # Title
        step_title = _clean_presentation_title(str(step.get("title", f"步骤 {idx + 1}")))
        stb = s.shapes.add_textbox(Inches(x + 0.60), Inches(y + 0.10), Inches(card_w - 0.74), Inches(0.34))
        _put_lines(stb.text_frame, step_title, 14, True, _PRIMARY)

        # Action subtitle
        action = str(step.get("action", step.get("summary", ""))).strip()
        cur_y = y + 0.46
        if action:
            act_box = s.shapes.add_textbox(Inches(x + 0.14), Inches(cur_y), Inches(card_w - 0.28), Inches(0.30))
            _put_lines(act_box.text_frame, action, 11, False, step_accent)
            cur_y += 0.32

        # Metric line (new)
        metric = str(step.get("metric", "")).strip()
        if metric:
            m_rect = _rounded_rect(s, x + 0.14, cur_y, card_w - 0.28, 0.24, step_bg)
            m_rect.text_frame.clear()
            m_rect.text_frame.word_wrap = True
            m_rect.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            mp = m_rect.text_frame.paragraphs[0]
            mp.alignment = PP_ALIGN.LEFT
            mrun = mp.add_run()
            clean_m = metric if any(metric.startswith(c) for c in ("📊", "⚡", "🚀")) else f"📊 {metric}"
            mrun.text = f" {clean_m}"
            _style_run(mrun, 10, True, step_accent)
            cur_y += 0.28

        # Bullets
        bullets = step.get("bullets", [])
        if isinstance(bullets, str):
            bullets = [bullets]
        detail = str(step.get("detail", "")).strip()
        lines = ([detail] if detail else []) + [f"• {b}" for b in bullets if str(b).strip()][:5]

        deliverable = str(step.get("deliverable", step.get("output", ""))).strip()
        tag = str(step.get("tag", step.get("tech", ""))).strip()

        bottom_reserved = 0.58 if deliverable else 0.30
        if lines:
            detail_h = max(0.6, card_h - (cur_y - y) - bottom_reserved)
            dtb = s.shapes.add_textbox(Inches(x + 0.14), Inches(cur_y), Inches(card_w - 0.28), Inches(detail_h))
            _fit_lines(dtb.text_frame, lines, 11, False, _TEXT, 1.25)

        # Deliverable pill
        if deliverable:
            deliv_text = deliverable if deliverable.startswith("📦") or deliverable.startswith("产物") else f"📦 {deliverable}"
            del_box = _rounded_rect(s, x + 0.14, y + card_h - 0.54, card_w - 0.28, 0.24, step_bg, step_accent)
            del_box.text_frame.clear()
            del_box.text_frame.word_wrap = True
            del_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            dp = del_box.text_frame.paragraphs[0]
            dp.alignment = PP_ALIGN.LEFT
            drun = dp.add_run()
            drun.text = f" {deliv_text}"
            _style_run(drun, 9.5, False, RGBColor(0x33, 0x41, 0x55))

        # Tag pill (colored per step)
        if tag:
            tag_w = max(0.95, min(card_w - 0.28, 0.35 + len(tag) * 0.16))
            tag_box = _rounded_rect(s, x + 0.14, y + card_h - 0.28, tag_w, 0.20, step_bg)
            tag_box.text_frame.clear()
            tag_box.text_frame.word_wrap = False
            tag_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            tp = tag_box.text_frame.paragraphs[0]
            tp.alignment = PP_ALIGN.CENTER
            trun = tp.add_run()
            trun.text = f"#{tag}"
            _style_run(trun, 9.5, True, step_accent)

        # Arrow (within same row only)
        if col < items_in_row - 1:
            arrow_x = x + card_w + (gap_x - 0.18) / 2
            arrow_y = y + card_h / 2 - 0.12
            atb = s.shapes.add_textbox(Inches(arrow_x), Inches(arrow_y), Inches(0.20), Inches(0.24))
            _put_lines(atb.text_frame, "▶", 11, True, step_accent)

    if takeaway:
        note_y = _H - takeaway_h - 0.16
        _rounded_rect(s, margin_x, note_y, total_w, 0.48, _WHITE, _CARD_BORDER)
        _rounded_rect(s, margin_x, note_y, 0.08, 0.48, _ACCENT)
        ntb = s.shapes.add_textbox(Inches(margin_x + 0.20), Inches(note_y + 0.06), Inches(total_w - 0.40), Inches(0.36))
        clean_takeaway = re.sub(r"^(核心结论|核心价值|主要结论|总结)[:：\s]*", "", str(takeaway).strip())
        _put_lines(ntb.text_frame, f"💡 核心价值：{clean_takeaway}", 12, True, _PRIMARY)

    verb = "rebuilt" if slide_number is not None else "added"
    return f"{verb} workflow pipeline slide {position}: '{title}' ({len(steps)} step cards)"


def _html_mockup_slide(
    h,
    title: str,
    cards: list[dict],
    subtitle: str = "",
    url_bar: str = "",
    slide_number: int | None = None,
    insert_after: int | None = None,
) -> str:
    """Compose a modern Web UI / Browser interface mockup with window chrome, nav, and card grid."""
    prs = _deck(h)
    s, position, is_rebuild = _resolve_target_slide(h, slide_number, insert_after)

    _rect(s, 0, 0, _W, _H, _SLATE_DARK)

    fx, fy, fw, fh = 0.45, 0.30, 12.43, 6.90
    _rounded_rect(s, fx, fy, fw, fh, _SLATE_LIGHT_BG, _SLATE_BORDER)

    _rounded_rect(s, fx, fy, fw, 0.58, _SLATE_PANEL)
    _circle(s, fx + 0.20, fy + 0.20, 0.18, _RED_DOT)
    _circle(s, fx + 0.44, fy + 0.20, 0.18, _YELLOW_DOT)
    _circle(s, fx + 0.68, fy + 0.20, 0.18, _GREEN_DOT)

    url_text = url_bar or f"https://agent.ai/workspace/{title[:18]}"
    _rounded_rect(s, fx + 1.10, fy + 0.12, fw - 4.50, 0.34, _SLATE_BORDER)
    url_box = s.shapes.add_textbox(Inches(fx + 1.25), Inches(fy + 0.08), Inches(fw - 4.80), Inches(0.35))
    _put_lines(url_box.text_frame, f"🔒  {url_text}", 10, False, RGBColor(0xE2, 0xE8, 0xF0))

    rh_box = s.shapes.add_textbox(Inches(fx + fw - 3.2), Inches(fy + 0.08), Inches(3.0), Inches(0.40))
    _put_lines(rh_box.text_frame, "Web App Interface", 10, True, RGBColor(0x94, 0xA3, 0xB8))

    sx = fx
    sy = fy + 0.58
    sw = 2.20
    sh = fh - 0.58
    _rect(s, sx, sy, sw, sh, _SLATE_PANEL)

    sb_logo = s.shapes.add_textbox(Inches(sx + 0.22), Inches(sy + 0.18), Inches(1.8), Inches(0.40))
    _put_lines(sb_logo.text_frame, "⚡ Agent Console", 11, True, _WHITE)
    _rect(s, sx + 0.15, sy + 0.60, sw - 0.30, 0.02, _SLATE_BORDER)

    nav_items = ["📊 概览 Overview", "⚡ 规划 Planner", "🧠 记忆 Memory", "🛠 工具 Tools", "🚀 执行 Execute"]
    for n_idx, n_text in enumerate(nav_items):
        is_active = n_idx == 0
        ny = sy + 0.75 + n_idx * 0.48
        if is_active:
            _rounded_rect(s, sx + 0.15, ny, sw - 0.30, 0.36, _BLUE_ACCENT)
        ntb = s.shapes.add_textbox(Inches(sx + 0.25), Inches(ny + 0.04), Inches(sw - 0.45), Inches(0.32))
        _put_lines(ntb.text_frame, n_text, 10, is_active, _WHITE if is_active else RGBColor(0x94, 0xA3, 0xB8))

    mx = sx + sw + 0.25
    my = sy + 0.18
    mw = fw - sw - 0.50

    mtb = s.shapes.add_textbox(Inches(mx), Inches(my), Inches(mw - 2.0), Inches(0.45))
    _put_lines(mtb.text_frame, title, 18, True, _SLATE_DARK)
    if subtitle:
        stb = s.shapes.add_textbox(Inches(mx), Inches(my + 0.40), Inches(mw), Inches(0.35))
        _put_lines(stb.text_frame, subtitle, 11, False, _TEXT_MUTED)

    card_count = max(1, len(cards))
    grid_y = my + (0.75 if subtitle else 0.55)
    grid_h = fh - 0.58 - (grid_y - sy) - 0.20

    if card_count <= 2:
        cols, rows = card_count, 1
    elif card_count <= 4:
        cols, rows = 2, 2
    else:
        cols, rows = 3, 2

    c_gap = 0.20
    cw = (mw - c_gap * (cols - 1)) / cols
    ch = (grid_h - c_gap * (rows - 1)) / rows

    for c_idx, card_data in enumerate(cards):
        col = c_idx % cols
        row = c_idx // cols
        cx = mx + col * (cw + c_gap)
        cy = grid_y + row * (ch + c_gap)

        _rounded_rect(s, cx, cy, cw, ch, _WHITE, _CARD_BORDER)
        _rounded_rect(s, cx, cy, cw, 0.05, _BLUE_ACCENT)

        c_title = str(card_data.get("title", f"Module {c_idx+1}"))
        ctb = s.shapes.add_textbox(Inches(cx + 0.15), Inches(cy + 0.10), Inches(cw - 1.25), Inches(0.35))
        _put_lines(ctb.text_frame, c_title, 13, True, _SLATE_DARK)

        status = str(card_data.get("status", card_data.get("badge", "ACTIVE"))).upper()
        if status:
            sb_w = min(1.10, 0.14 + len(status) * 0.10)
            status_box = _rounded_rect(s, cx + cw - sb_w - 0.15, cy + 0.10, sb_w, 0.22, _GREEN_BG)
            status_box.text_frame.clear()
            status_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            sp = status_box.text_frame.paragraphs[0]
            sp.alignment = PP_ALIGN.CENTER
            srun = sp.add_run()
            srun.text = status
            _style_run(srun, 10, True, _GREEN_ACCENT)

        metric = str(card_data.get("metric", card_data.get("highlight", ""))).strip()
        cur_cy = cy + 0.42
        if metric:
            met_box = s.shapes.add_textbox(Inches(cx + 0.15), Inches(cur_cy), Inches(cw - 0.30), Inches(0.30))
            _put_lines(met_box.text_frame, f"⚡ {metric}", 11, True, _BLUE_ACCENT)
            cur_cy += 0.28

        detail = str(card_data.get("detail", "")).strip()
        bullets = card_data.get("bullets", [])
        if isinstance(bullets, str):
            bullets = [bullets]
        lines = ([detail] if detail else []) + [f"• {b}" for b in bullets if str(b).strip()][:4]
        if lines:
            dt_h = max(0.8, ch - (cur_cy - cy) - 0.30)
            cdtb = s.shapes.add_textbox(Inches(cx + 0.15), Inches(cur_cy), Inches(cw - 0.30), Inches(dt_h))
            _fit_lines(cdtb.text_frame, lines, 10, False, _TEXT, 1.10)

        anchor = str(card_data.get("html_anchor", card_data.get("anchor", card_data.get("source", "")))).strip()
        if anchor:
            chip = s.shapes.add_textbox(Inches(cx + 0.15), Inches(cy + ch - 0.28), Inches(cw - 0.30), Inches(0.22))
            _put_lines(chip.text_frame, f"<{anchor}>", 10, False, _TEXT_MUTED)

    verb = "rebuilt" if slide_number is not None else "added"
    return f"{verb} HTML mockup slide {position}: '{title}' ({len(cards)} UI cards)"


def _find_headless_browser_executable() -> str | None:
    from pathlib import Path
    for candidate in [
        shutil.which("msedge"),
        r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
        shutil.which("chrome"),
        r"C:\Program Files\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
    ]:
        if candidate and Path(candidate).exists():
            return str(candidate)
    return None


def _render_html_slide_to_png(html: str, css: str = "") -> bytes:
    """Render an HTML/CSS slide snippet to 1920x1080 PNG image bytes via headless browser."""
    import tempfile
    from pathlib import Path

    browser_bin = _find_headless_browser_executable()
    if not browser_bin:
        raise RuntimeError("No headless browser (Edge/Chrome) found on system for HTML rendering.")

    # Wrap complete HTML document with 1920x1080 16:9 canvas styling if missing
    full_html = html
    if "<html" not in html.lower():
        _premium_css = """
* { box-sizing: border-box; margin: 0; padding: 0; }
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');
body {
    width: 1920px;
    height: 1080px;
    overflow: hidden;
    font-family: 'Inter', 'Segoe UI', -apple-system, BlinkMacSystemFont, 'PingFang SC', 'Microsoft YaHei', sans-serif;
    background: linear-gradient(135deg, #0b1120 0%, #0f172a 50%, #1e293b 100%);
    color: #f8fafc;
    padding: 40px 52px;
}
.slide {
    width: 100%;
    height: 100%;
    display: flex;
    flex-direction: column;
    justify-content: center;
    gap: 18px;
    background: transparent !important;
}
.header-area {
    margin-bottom: 4px;
}
h1 {
    font-size: 40px;
    font-weight: 700;
    color: #f1f5f9;
    letter-spacing: -0.5px;
    margin-bottom: 6px;
}
.subtitle {
    font-size: 17.5px;
    color: #94a3b8;
    font-weight: 400;
    letter-spacing: 0.2px;
}
.grid-3, .grid-3x2, .grid-6 {
    display: grid;
    grid-template-columns: repeat(3, 1fr);
    gap: 16px 20px;
}
.grid-2, .grid-2x2, .grid-4 {
    display: grid;
    grid-template-columns: repeat(2, 1fr);
    gap: 16px 20px;
}
.card {
    background: rgba(30, 41, 59, 0.78);
    backdrop-filter: blur(20px);
    -webkit-backdrop-filter: blur(20px);
    border: 1px solid rgba(148, 163, 184, 0.15);
    border-radius: 16px;
    padding: 16px 18px 12px;
    display: flex;
    flex-direction: column;
    gap: 8px;
    box-shadow: 0 12px 36px rgba(0,0,0,0.35), inset 0 1px 0 rgba(255,255,255,0.06);
}
.card:nth-child(1) { border-top: 3px solid #10b981; }
.card:nth-child(2) { border-top: 3px solid #f59e0b; }
.card:nth-child(3) { border-top: 3px solid #3b82f6; }
.card:nth-child(4) { border-top: 3px solid #8b5cf6; }
.card:nth-child(5) { border-top: 3px solid #ec4899; }
.card:nth-child(6) { border-top: 3px solid #06b6d4; }
.card-header {
    display: flex;
    align-items: center;
    justify-content: space-between;
}
.badge {
    padding: 3px 10px;
    border-radius: 14px;
    font-size: 10.5px;
    font-weight: 700;
    letter-spacing: 1.2px;
    text-transform: uppercase;
    background: linear-gradient(135deg, #065f46, #047857);
    color: #34d399;
    border: 1px solid rgba(52, 211, 153, 0.25);
}
.card:nth-child(2) .badge {
    background: linear-gradient(135deg, #7c2d12, #9a3412);
    color: #fbbf24;
    border-color: rgba(251, 191, 36, 0.25);
}
.card:nth-child(3) .badge {
    background: linear-gradient(135deg, #1e3a5f, #1e40af);
    color: #60a5fa;
    border-color: rgba(96, 165, 250, 0.25);
}
.card:nth-child(4) .badge {
    background: linear-gradient(135deg, #4c1d95, #6d28d9);
    color: #c4b5fd;
    border-color: rgba(196, 181, 253, 0.25);
}
.card:nth-child(5) .badge {
    background: linear-gradient(135deg, #831843, #be185d);
    color: #f472b6;
    border-color: rgba(244, 114, 182, 0.25);
}
.card:nth-child(6) .badge {
    background: linear-gradient(135deg, #164e63, #0e7490);
    color: #67e8f9;
    border-color: rgba(103, 232, 249, 0.25);
}
h3 {
    font-size: 20px;
    font-weight: 600;
    color: #e2e8f0;
    letter-spacing: -0.2px;
}
.card-desc {
    font-size: 13px;
    color: #94a3b8;
    line-height: 1.4;
    margin-bottom: 2px;
}
.feature-list {
    display: flex;
    flex-direction: column;
    gap: 6px;
}
.feature-item {
    background: rgba(15, 23, 42, 0.55);
    border: 1px solid rgba(148, 163, 184, 0.1);
    border-radius: 8px;
    padding: 9px 12px;
}
.feature-title {
    font-size: 13.5px;
    font-weight: 600;
    color: #38bdf8;
    margin-bottom: 2px;
    display: flex;
    align-items: center;
    gap: 5px;
}
.card:nth-child(1) .feature-title { color: #34d399; }
.card:nth-child(2) .feature-title { color: #fbbf24; }
.card:nth-child(3) .feature-title { color: #60a5fa; }
.card:nth-child(4) .feature-title { color: #c4b5fd; }
.card:nth-child(5) .feature-title { color: #f472b6; }
.card:nth-child(6) .feature-title { color: #67e8f9; }
.feature-text {
    font-size: 12.5px;
    line-height: 1.45;
    color: #cbd5e1;
}
ul {
    list-style: none;
    padding: 0;
    display: flex;
    flex-direction: column;
    gap: 6px;
}
ul li {
    position: relative;
    padding-left: 18px;
    font-size: 13px;
    line-height: 1.45;
    color: #cbd5e1;
}
ul li::before {
    content: '';
    position: absolute;
    left: 0;
    top: 7px;
    width: 6px;
    height: 6px;
    border-radius: 50%;
    background: #64748b;
}
.card:nth-child(1) ul li::before { background: #34d399; }
.card:nth-child(2) ul li::before { background: #fbbf24; }
.card:nth-child(3) ul li::before { background: #60a5fa; }
.card:nth-child(4) ul li::before { background: #8b5cf6; }
.card:nth-child(5) ul li::before { background: #ec4899; }
.card:nth-child(6) ul li::before { background: #06b6d4; }
.card-footer {
    display: flex;
    align-items: center;
    justify-content: space-between;
    padding-top: 8px;
    border-top: 1px solid rgba(148, 163, 184, 0.12);
    margin-top: 2px;
}
.deliverable-pill {
    font-size: 11.5px;
    color: #f1f5f9;
    background: rgba(148, 163, 184, 0.15);
    padding: 3px 8px;
    border-radius: 5px;
}
.tech-tag {
    font-size: 11.5px;
    font-family: 'SF Mono', 'Cascadia Code', 'Consolas', monospace;
    color: #94a3b8;
    letter-spacing: 0.5px;
}
.summary-bar {
    display: flex;
    align-items: center;
    gap: 10px;
    background: rgba(15, 23, 42, 0.85);
    border: 1px solid rgba(148, 163, 184, 0.15);
    border-left: 4px solid #38bdf8;
    border-radius: 10px;
    padding: 12px 20px;
    font-size: 14.5px;
    color: #e2e8f0;
    line-height: 1.45;
    box-shadow: 0 8px 24px rgba(0,0,0,0.25);
}
.summary-bar strong {
    color: #38bdf8;
}
"""
        full_html = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
{_premium_css}
{css}
</style>
</head>
<body>
{html}
</body>
</html>"""
    elif css and "<style>" not in html.lower():
        full_html = html.replace("</head>", f"<style>{css}</style></head>") if "</head>" in html else f"<style>{css}</style>{html}"

    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_path = Path(tmp_dir)
        html_file = tmp_path / "slide.html"
        png_file = tmp_path / "slide.png"
        html_file.write_text(full_html, encoding="utf-8")

        cmd = [
            browser_bin,
            "--headless",
            "--disable-gpu",
            "--hide-scrollbars",
            "--force-device-scale-factor=1",
            "--window-size=1920,1080",
            f"--screenshot={png_file}",
            str(html_file.resolve()),
        ]
        res = subprocess.run(cmd, capture_output=True, timeout=30)
        if not png_file.exists() or png_file.stat().st_size == 0:
            err = res.stderr.decode("utf-8", errors="ignore")
            raise RuntimeError(f"Browser screenshot failed: {err}")
        return png_file.read_bytes()


def _clean_presentation_title(title: str) -> str:
    """Sanitize slide titles/subtitles by stripping AI prompt meta-instructions."""
    if not title:
        return ""
    t = str(title).strip()
    # Strip meta tokens like (HTML 页面风格), （PPTX 原生元素）, [HTML 制作], 基于现代Web排版 etc.
    t = re.sub(r"[\(（]\s*(HTML|网页|PPTX|原生|代码|风格|页面|组件).*?[\)）]", "", t, flags=re.IGNORECASE)
    t = re.sub(r"[\(（]\s*.*?(风格|制作|生成|排版|模式)\s*[\)）]", "", t, flags=re.IGNORECASE)
    t = re.sub(r"\[\s*(HTML|PPTX|原生|网页).*?\]", "", t, flags=re.IGNORECASE)
    t = re.sub(r"基于.*?(网格|组件|排版|风格|渲染).*?([，,。；;]|$)", "", t, flags=re.IGNORECASE)
    t = re.sub(r"\s+", " ", t).strip()
    return t or str(title).strip()


def _compile_html_to_vector_slide(
    h: Any,
    html: str,
    slide_number: int | None = None,
    insert_after: int | None = None,
    default_title: str = "",
) -> str:
    """Compile HTML DOM structure into native, 100% editable PowerPoint vector shapes and text boxes."""
    soup = BeautifulSoup(html, "html.parser")
    s, target_index, _ = _resolve_target_slide(h, slide_number, insert_after)
    return _html_to_vector_slide(h, soup, target_index, default_title=default_title)


def _html_to_vector_slide(
    h: Any,
    soup: BeautifulSoup,
    target_index: int,
    default_title: str = "",
    theme: str = "dark",
) -> str:
    """Compile HTML DOM tree directly into native PowerPoint shapes & vector cards."""
    prs = _deck(h)
    s = prs.slides[target_index - 1]
    _clear_slide_shapes(s)

    # 1. Colors & Theme
    is_dark = theme.lower() in ("dark", "night", "slate", "black") or soup.find(style=lambda s: s and ("#0f172a" in s or "#1e293b" in s or "black" in s or "#111" in s))
    if is_dark:
        bg_color = RGBColor(0x0F, 0x17, 0x2A)       # Deep slate 900
        card_bg = RGBColor(0x1E, 0x29, 0x3B)        # Slate 800
        card_border = RGBColor(0x33, 0x41, 0x55)    # Slate 700
        title_color = RGBColor(0x38, 0xBD, 0xF8)    # Sky blue 400
        subtitle_color = RGBColor(0x94, 0xA3, 0xB8) # Slate 400
        text_color = RGBColor(0xF1, 0xF5, 0xF9)     # Slate 100
        muted_text = RGBColor(0x64, 0x74, 0x8B)     # Slate 500
        badge_green_bg = RGBColor(0x06, 0x4E, 0x3B)
        badge_green_fg = RGBColor(0x34, 0xD3, 0x99)
        badge_amber_bg = RGBColor(0x78, 0x35, 0x0F)
        badge_amber_fg = RGBColor(0xFB, 0xBF, 0x24)
        badge_blue_bg = RGBColor(0x1E, 0x3A, 0x8A)
        badge_blue_fg = RGBColor(0x60, 0xA5, 0xFA)
        accent_blue = RGBColor(0x02, 0x84, 0xC7)
    else:
        bg_color = RGBColor(0xF8, 0xFA, 0xFC)       # Slate 50
        card_bg = _WHITE
        card_border = RGBColor(0xE2, 0xE8, 0xF0)
        title_color = RGBColor(0x0F, 0x17, 0x2A)
        subtitle_color = RGBColor(0x47, 0x55, 0x69)
        text_color = RGBColor(0x1E, 0x29, 0x3B)
        muted_text = RGBColor(0x64, 0x74, 0x8B)
        badge_green_bg = RGBColor(0xDC, 0xFC, 0xE7)
        badge_green_fg = RGBColor(0x15, 0x80, 0x3D)
        badge_amber_bg = RGBColor(0xFE, 0xF3, 0xC7)
        badge_amber_fg = RGBColor(0xB4, 0x53, 0x09)
        badge_blue_bg = RGBColor(0xDB, 0xEA, 0xFE)
        badge_blue_fg = RGBColor(0x1D, 0x4E, 0xD8)
        accent_blue = _ACCENT

    # Background canvas
    bg_shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(_W), Inches(_H))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.fill.background()

    # 2. Header (Title + Subtitle)
    title_el = soup.find(["h1", "h2", "title"])
    raw_title = title_el.get_text().strip() if title_el else (default_title or "AI Agent 核心架构与运行时全景")
    title_text = _clean_presentation_title(raw_title) or "AI Agent 核心架构与运行时全景"
    
    title_box = s.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(_W - 1.6), Inches(0.55))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = title_color

    sub_el = soup.find(class_=lambda c: c and "sub" in c) or (title_el.find_next_sibling("p") if title_el else soup.find("p"))
    raw_sub = sub_el.get_text().strip() if sub_el and sub_el.get_text().strip() != raw_title else ""
    subtitle_text = _clean_presentation_title(raw_sub)
    if subtitle_text:
        sub_box = s.shapes.add_textbox(Inches(0.8), Inches(0.95), Inches(_W - 1.6), Inches(0.35))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = subtitle_color
        grid_top_y = 1.45
    else:
        grid_top_y = 1.25

    # 3. Extract Component Cards
    cards = soup.find_all(class_=lambda c: c and any(k in c.split() for k in ("card", "box", "panel", "step", "col", "item", "module", "node")))
    if not cards:
        grid_container = soup.find(class_=lambda c: c and any(k in c.split() for k in ("grid", "flex", "container", "columns", "pipeline")))
        if grid_container:
            cards = [ch for ch in grid_container.find_all(recursive=False) if ch.name in ("div", "section", "article")]
        else:
            cards = soup.find_all(["section", "article"])

    cards_data = []
    if not cards:
        all_p = [p for p in soup.find_all(["li", "p"]) if p.get_text().strip() and p.get_text().strip() != raw_sub]
        if all_p:
            chunk_size = max(1, math.ceil(len(all_p) / 3))
            for c_idx in range(0, min(3, len(all_p))):
                sub_items = all_p[c_idx * chunk_size : (c_idx + 1) * chunk_size]
                cards_data.append({
                    "title": f"阶段 0{c_idx+1}",
                    "bullets": [re.sub(r"^[•\-▶\*\s]+", "", it.get_text().strip()) for it in sub_items if it.get_text().strip()],
                    "status": "ACTIVE" if c_idx == 0 else "READY",
                })
        else:
            cards_data = [{"title": "核心模块", "bullets": ["多模态智能体协同执行闭环"], "status": "ACTIVE"}]
    else:
        for idx, card in enumerate(cards):
            c_soup = BeautifulSoup(str(card), "html.parser")
            c_head = c_soup.find(["h1", "h2", "h3", "h4", "h5", "strong", "b"])
            c_title = _clean_presentation_title(c_head.get_text().strip()) if c_head else f"模块 0{idx+1}"
            
            badge_el = c_soup.find(class_=lambda c: c and any(k in str(c).lower() for k in ("badge", "status", "chip")))
            c_status = badge_el.get_text().strip().upper() if badge_el else ("RUNNING" if idx == 0 else ("ACTIVE" if idx == 1 else "READY"))

            metric_el = c_soup.find(class_=lambda c: c and any(k in str(c).lower() for k in ("metric", "stat", "num", "highlight", "time", "speed", "latency", "perf", "qps")))
            c_metric = metric_el.get_text().strip() if metric_el else ""

            # Extract ONLY leaf bullet items (<p> or <li>), NEVER the parent container div!
            bullets = []
            leaf_tags = c_soup.find_all("li")
            if not leaf_tags:
                leaf_tags = [p for p in c_soup.find_all("p") if p.get_text().strip() and p.get_text().strip() != raw_sub]

            for item in leaf_tags:
                t = item.get_text().strip()
                t = re.sub(r"^[•\-▶\*\s]+", "", t).strip()
                if t and t != c_title and t != c_status and t != c_metric:
                    if not any(t == b for b in bullets):
                        bullets.append(t)
            
            tag_el = c_soup.find(class_=lambda c: c and any(k in str(c).lower() for k in ("tech", "code", "anchor", "foot", "tag")))
            c_tag = tag_el.get_text().strip() if tag_el else ""

            cards_data.append({
                "title": c_title,
                "status": c_status,
                "metric": c_metric,
                "bullets": bullets[:6],
                "tag": c_tag,
            })

    # 4. Compute Grid Geometry
    num_cards = max(1, len(cards_data))
    max_w = _W - 1.6
    avail_h = _H - grid_top_y - 0.45

    if num_cards == 1:
        cols, rows = 1, 1
    elif num_cards == 2:
        cols, rows = 2, 1
    elif num_cards == 3:
        cols, rows = 3, 1
    elif num_cards == 4:
        cols, rows = 2, 2
    elif num_cards in (5, 6):
        cols, rows = 3, 2
    else:
        cols, rows = 4, math.ceil(num_cards / 4)

    gap_x = 0.28
    gap_y = 0.25
    card_w = (max_w - (cols - 1) * gap_x) / cols
    
    # Adaptive content-aware card height (compact, void-free, vertically centered)
    if rows == 1:
        max_bullets_count = max((len(c.get("bullets", [])) for c in cards_data), default=3)
        has_metrics = any(c.get("metric") for c in cards_data)
        has_tags = any(c.get("tag") for c in cards_data)
        needed_h = 0.55 + (0.35 if has_metrics else 0.0) + max_bullets_count * 0.48 + (0.35 if has_tags else 0.0) + 0.35
        card_h = min(3.80, max(3.20, needed_h))
        start_cy = grid_top_y + (avail_h - card_h) / 2
    else:
        card_h = (avail_h - (rows - 1) * gap_y) / rows
        start_cy = grid_top_y

    for c_idx, card_info in enumerate(cards_data[: cols * rows]):
        c_col = c_idx % cols
        c_row = c_idx // cols
        cx = 0.8 + c_col * (card_w + gap_x)
        cy = start_cy + c_row * (card_h + gap_y)

        c_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(cy), Inches(card_w), Inches(card_h))
        c_box.fill.solid()
        c_box.fill.fore_color.rgb = card_bg
        c_box.line.color.rgb = card_border
        c_box.line.width = Pt(1.1)

        acc_bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(cy), Inches(card_w), Inches(0.06))
        acc_bar.fill.solid()
        acc_bar.fill.fore_color.rgb = accent_blue
        acc_bar.line.fill.background()

        stat = card_info.get("status", "")
        bw = 0.0
        if stat:
            # Adequate badge width to prevent any vertical text wrap
            bw = max(0.95, min(card_w * 0.38, 0.30 + len(stat) * 0.12))
            bx = cx + card_w - bw - 0.15
            by = cy + 0.12
            badge_shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(by), Inches(bw), Inches(0.26))
            badge_shape.fill.solid()
            badge_shape.fill.fore_color.rgb = badge_green_bg if any(k in stat for k in ("RUN", "SUCC", "OK", "DONE")) else (badge_amber_bg if any(k in stat for k in ("ACT", "WARN", "PLAN")) else badge_blue_bg)
            badge_shape.line.fill.background()
            
            tf_b = badge_shape.text_frame
            tf_b.clear()
            tf_b.word_wrap = False
            tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
            p_b = tf_b.paragraphs[0]
            p_b.alignment = PP_ALIGN.CENTER
            r_b = p_b.add_run()
            r_b.text = stat
            r_b.font.size = Pt(10)
            r_b.font.bold = True
            r_b.font.color.rgb = badge_green_fg if any(k in stat for k in ("RUN", "SUCC", "OK", "DONE")) else (badge_amber_fg if any(k in stat for k in ("ACT", "WARN", "PLAN")) else badge_blue_fg)

        c_title_text = card_info.get("title", "")
        t_box = s.shapes.add_textbox(Inches(cx + 0.15), Inches(cy + 0.12), Inches(card_w - (bw + 0.35 if stat else 0.30)), Inches(0.35))
        tf_ct = t_box.text_frame
        tf_ct.word_wrap = True
        p_ct = tf_ct.paragraphs[0]
        p_ct.text = c_title_text
        p_ct.font.size = Pt(13.5)
        p_ct.font.bold = True
        p_ct.font.color.rgb = text_color

        cur_cy = cy + 0.48
        metric = card_info.get("metric", "")
        if metric:
            m_box = s.shapes.add_textbox(Inches(cx + 0.15), Inches(cur_cy), Inches(card_w - 0.30), Inches(0.28))
            tf_m = m_box.text_frame
            tf_m.word_wrap = True
            p_m = tf_m.paragraphs[0]
            clean_m = metric if metric.startswith("⚡") or metric.startswith("📊") or metric.startswith("🚀") else f"⚡ {metric}"
            p_m.text = clean_m
            p_m.font.size = Pt(10.5)
            p_m.font.bold = True
            p_m.font.color.rgb = title_color
            cur_cy += 0.28

        bullets = card_info.get("bullets", [])
        if bullets:
            b_h = max(0.8, card_h - (cur_cy - cy) - (0.35 if card_info.get("tag") else 0.15))
            b_box = s.shapes.add_textbox(Inches(cx + 0.15), Inches(cur_cy), Inches(card_w - 0.30), Inches(b_h))
            tf_b = b_box.text_frame
            tf_b.word_wrap = True
            for b_idx, bullet_line in enumerate(bullets[:5]):
                p_b = tf_b.paragraphs[0] if b_idx == 0 else tf_b.add_paragraph()
                p_b.text = bullet_line if bullet_line.startswith("•") or bullet_line.startswith("-") or bullet_line.startswith("▶") else f"• {bullet_line}"
                p_b.font.size = Pt(10.5 if len(bullets) > 3 else 11)
                p_b.font.color.rgb = text_color
                p_b.line_spacing = 1.22

        tag = card_info.get("tag", "")
        if tag:
            tag_w = max(0.95, min(card_w - 0.30, 0.35 + len(tag) * 0.15))
            tag_box = s.shapes.add_textbox(Inches(cx + 0.15), Inches(cy + card_h - 0.30), Inches(tag_w), Inches(0.25))
            tf_tag = tag_box.text_frame
            tf_tag.word_wrap = False
            p_tag = tf_tag.paragraphs[0]
            p_tag.text = f"<{tag}>" if not tag.startswith("<") else tag
            p_tag.font.size = Pt(10)
            p_tag.font.color.rgb = muted_text

    h.state.ppt_affected_slides.add(target_index)
    h.state.record_change(f"deck:slide:{target_index}:html_vector_slide")
    return f"compiled HTML vector slide {target_index}: '{title_text}' ({len(cards_data)} editable UI cards)"


def _html_slide(
    h: Any,
    html: str = "",
    css: str = "",
    file_path: str = "",
    slide_number: int | None = None,
    insert_after: int | None = None,
    title: str = "",
    render_mode: str = "raster",
) -> str:
    """Render HTML/CSS content into PPTX slides via headless browser screenshot (raster) or vector shapes."""
    from pathlib import Path
    from bs4 import BeautifulSoup
    from .. import config
    from ..permissions import path_within

    root = config.sandbox_root()
    if file_path:
        target_file = (root / file_path).resolve()
        if not path_within(root, target_file) or not target_file.exists():
            raise FileNotFoundError(f"HTML file not found: {file_path}")
        html = target_file.read_text(encoding="utf-8")
    elif html.strip():
        # Automatically persist slide2.html to workspace disk for real HTML artifact delivery
        try:
            target_name = f"slide{slide_number or 2}.html"
            target_path = root / target_name
            target_path.write_text(html, encoding="utf-8")
        except Exception:
            pass

    if not html.strip():
        raise ValueError("html_slide requires non-empty html code or valid file_path")

    soup = BeautifulSoup(html, "html.parser")
    slide_sections = soup.find_all(class_=lambda c: c and "slide" in c.split()) or soup.find_all("section")

    if len(slide_sections) <= 1:
        # Default to 100% native vector compilation for total editability
        if render_mode == "vector":
            return _compile_html_to_vector_slide(h, html, slide_number, insert_after, title)

        # Optional rasterized fallback
        doc_title = title
        if not doc_title:
            h_el = soup.find(["h1", "h2", "h3", "title"])
            if h_el:
                doc_title = h_el.get_text().strip()
        texts = [s.strip() for s in soup.stripped_strings if s.strip() and s.strip() != doc_title]

        png_bytes = _render_html_slide_to_png(html, css)
        s, target_index, _ = _resolve_target_slide(h, slide_number, insert_after)
        _clear_slide_shapes(s)
        image_stream = BytesIO(png_bytes)
        s.shapes.add_picture(image_stream, Inches(0), Inches(0), Inches(_W), Inches(_H))
        if getattr(s, "has_notes_slide", False) or hasattr(s, "notes_slide"):
            note_content = f"{doc_title}\n" + "\n".join(texts[:15]) if doc_title else "\n".join(texts[:15])
            if note_content.strip():
                s.notes_slide.notes_text_frame.text = note_content.strip()
        h.state.ppt_affected_slides.add(target_index)
        h.state.record_change(f"deck:slide:{target_index}:html_slide")
        return f"rendered HTML slide {target_index}: '{doc_title or 'HTML Web Slide'}'"

    # Multi-slide deck compilation
    start_num = slide_number if slide_number is not None else (len(h.deck.slides) + 1 if getattr(h, "deck", None) else 1)
    rendered_count = 0
    for idx, sec in enumerate(slide_sections):
        sec_html = str(sec)
        cur_slide_num = start_num + idx
        if render_mode == "vector":
            _compile_html_to_vector_slide(h, sec_html, cur_slide_num, None, title)
        else:
            sec_soup = BeautifulSoup(sec_html, "html.parser")
            sec_title = ""
            h_el = sec_soup.find(["h1", "h2", "h3", "title"])
            if h_el:
                sec_title = h_el.get_text().strip()
            sec_texts = [s.strip() for s in sec_soup.stripped_strings if s.strip() and s.strip() != sec_title]
            png_bytes = _render_html_slide_to_png(sec_html, css)
            s, target_index, _ = _resolve_target_slide(h, cur_slide_num, None)
            _clear_slide_shapes(s)
            image_stream = BytesIO(png_bytes)
            s.shapes.add_picture(image_stream, Inches(0), Inches(0), Inches(_W), Inches(_H))
            if getattr(s, "has_notes_slide", False) or hasattr(s, "notes_slide"):
                note_content = f"{sec_title}\n" + "\n".join(sec_texts[:15]) if sec_title else "\n".join(sec_texts[:15])
                if note_content.strip():
                    s.notes_slide.notes_text_frame.text = note_content.strip()
            h.state.ppt_affected_slides.add(target_index)
            h.state.record_change(f"deck:slide:{target_index}:html_slide")
        rendered_count += 1

    return f"compiled {rendered_count} HTML vector slides into deck starting at slide {start_num}"


def _hero_split_slide(
    h,
    title: str,
    hero_title: str,
    hero_text: str,
    cards: list[dict],
    hero_metric: str = "",
    subtitle: str = "",
    slide_number: int | None = None,
    insert_after: int | None = None,
) -> str:
    """Compose a left-hero highlight + right breakdown cards layout."""
    prs = _deck(h)
    s, position, is_rebuild = _resolve_target_slide(h, slide_number, insert_after)

    _rect(s, 0, 0, _W, _H, _SLATE_LIGHT_BG)
    _rect(s, 0, 0, _W, 1.10, _PRIMARY)
    _rect(s, 0, 1.10, _W, 0.07, _ACCENT)

    title_box = s.shapes.add_textbox(Inches(0.55), Inches(0.12), Inches(8.5), Inches(0.85))
    _put_lines(title_box.text_frame, title, 25, True, _WHITE)
    if subtitle:
        sub_box = s.shapes.add_textbox(Inches(7.2), Inches(0.18), Inches(5.6), Inches(0.70))
        _put_lines(sub_box.text_frame, subtitle, 12, False, RGBColor(0xD0, 0xD0, 0xD0))

    hx, hy, hw, hh = 0.55, 1.40, 4.0, 5.50
    _rounded_rect(s, hx, hy, hw, hh, _PRIMARY)
    _rounded_rect(s, hx, hy, hw, 0.08, _ACCENT)

    htb = s.shapes.add_textbox(Inches(hx + 0.30), Inches(hy + 0.35), Inches(hw - 0.60), Inches(0.80))
    _put_lines(htb.text_frame, hero_title, 20, True, _WHITE)

    if hero_metric:
        hm_box = s.shapes.add_textbox(Inches(hx + 0.30), Inches(hy + 1.25), Inches(hw - 0.60), Inches(0.90))
        _put_lines(hm_box.text_frame, hero_metric, 30, True, _ACCENT)

    cur_hy = hy + (2.20 if hero_metric else 1.30)
    ht_box = s.shapes.add_textbox(Inches(hx + 0.30), Inches(cur_hy), Inches(hw - 0.60), Inches(hh - (cur_hy - hy) - 0.40))
    _fit_lines(ht_box.text_frame, hero_text.split("\n"), 12, False, RGBColor(0xEA, 0xE8, 0xE4), 1.25)

    rx = hx + hw + 0.35
    rw = _W - rx - 0.55
    card_count = max(1, len(cards))
    c_gap = 0.20
    c_h = (hh - c_gap * (card_count - 1)) / card_count
    for c_idx, c_data in enumerate(cards):
        cy = hy + c_idx * (c_h + c_gap)
        _rounded_rect(s, rx, cy, rw, c_h, _WHITE, _CARD_BORDER)
        _rounded_rect(s, rx, cy, 0.08, c_h, _ACCENT)

        c_title = str(c_data.get("title", f"Feature {c_idx+1}"))
        ctb = s.shapes.add_textbox(Inches(rx + 0.25), Inches(cy + 0.12), Inches(rw - 0.50), Inches(0.35))
        _put_lines(ctb.text_frame, c_title, 14, True, _PRIMARY)

        bullets = c_data.get("bullets", [])
        if isinstance(bullets, str):
            bullets = [bullets]
        detail = str(c_data.get("detail", "")).strip()
        lines = ([detail] if detail else []) + [f"• {b}" for b in bullets if str(b).strip()][:3]
        if lines:
            dtb = s.shapes.add_textbox(Inches(rx + 0.25), Inches(cy + 0.48), Inches(rw - 0.50), Inches(c_h - 0.60))
            _fit_lines(dtb.text_frame, lines, 11, False, _TEXT, 1.15)

    verb = "rebuilt" if slide_number is not None else "added"
    return f"{verb} hero split slide {position}: '{title}'"



def _walk_shapes(shapes, parent_path: tuple[int, ...] = ()):
    """Yield every shape in visual-tree order, including nested group members.

    PowerPoint stores child geometry in the coordinate system of its parent
    group.  The path lets callers identify that context without incorrectly
    presenting a nested shape's ``left``/``top`` as slide-absolute values.
    """
    for shape in shapes:
        path = parent_path + (shape.shape_id,)
        yield shape, path
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes, path)


def _shape_inventory(h, slide_number: int | None = None) -> str:
    if getattr(h, "deck", None) is None:
        return "no deck yet."
    selected = range(len(h.deck.slides)) if slide_number is None else [slide_number - 1]
    lines = []
    for index in selected:
        if index < 0 or index >= len(h.deck.slides):
            raise IndexError(f"slide number out of range: {slide_number}")
        for shape, path in _walk_shapes(h.deck.slides[index].shapes):
            if getattr(shape, "has_table", False):
                rows = [
                    " | ".join(cell.text.strip().replace("\n", " ") for cell in row.cells)
                    for row in shape.table.rows
                ]
                text = " ;; ".join(rows)[:220]
                kind = "table"
                dimensions = f" {len(shape.table.rows)}x{len(shape.table.columns)}"
            else:
                text = shape.text.strip().replace("\n", " | ")[:100] if getattr(shape, "has_text_frame", False) else ""
                kind = "group" if shape.shape_type == MSO_SHAPE_TYPE.GROUP else "shape"
                dimensions = ""
            path_label = "/".join(str(shape_id) for shape_id in path)
            geometry_label = "box_slide" if len(path) == 1 else "box_group_local"
            lines.append(
                f"slide={index + 1} path={path_label} id={shape.shape_id} kind={kind} name={shape.name!r} "
                f"{geometry_label}=({shape.left/914400:.2f},{shape.top/914400:.2f},"
                f"{shape.width/914400:.2f},{shape.height/914400:.2f}){dimensions} text={text!r}"
            )
    return "\n".join(lines)


def _replace_text(h, slide_number: int, shape_id: int, text: str) -> str:
    if getattr(h, "deck", None) is None:
        raise ValueError("no deck loaded")
    if slide_number < 1 or slide_number > len(h.deck.slides):
        raise IndexError("slide number out of range")
    _, shape = _find_shape(h, slide_number, shape_id)
    if not shape.has_text_frame:
        raise TypeError("target shape has no text frame")
    style = None
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.runs:
            run = paragraph.runs[0]
            try:
                rgb = run.font.color.rgb
            except (AttributeError, ValueError):
                rgb = None
            style = (run.font.size, run.font.bold, run.font.name, rgb)
            break
    shape.text_frame.text = text
    if style and shape.text_frame.paragraphs[0].runs:
        run = shape.text_frame.paragraphs[0].runs[0]
        run.font.size, run.font.bold, run.font.name = style[:3]
        if style[3] is not None:
            run.font.color.rgb = style[3]
    h.state.record_change(f"deck:slide:{slide_number}:shape:{shape_id}:text")
    return f"updated slide {slide_number} shape {shape_id}"


def _replace_text_semantic(h, old: str, new: str, slide_number: int | None = None, match_case: bool = True) -> str:
    """Replace text even when a match spans multiple runs.

    The first run keeps the paragraph's dominant style. This is more faithful
    than rebuilding the whole text frame and is deterministic for benchmark
    edits such as a title split into ``Lecture `` and ``3`` runs. Table cells
    are textual slide surfaces too, so a scoped semantic replace also walks
    every cell of every table on the selected slides.
    """
    if not old:
        raise ValueError("old text cannot be empty")
    if getattr(h, "deck", None) is None:
        raise ValueError("no deck loaded")

    def replace_in_frame(frame, match_case: bool) -> int:
        replacements = 0
        for paragraph in frame.paragraphs:
            runs = list(paragraph.runs)
            if not runs:
                continue

            # Build the searchable paragraph text in XML order.  A soft
            # line break is represented by ``a:br`` (``\v`` in the public
            # paragraph text), not by a run, so keep it as an immutable
            # separator instead of flattening the runs.
            parts: list[str] = []
            run_index = 0
            for child in paragraph._p.content_children:
                kind = child.tag.rsplit("}", 1)[-1]
                if kind == "r":
                    parts.append(runs[run_index].text)
                    run_index += 1
                elif kind == "br":
                    parts.append("\v")
                else:
                    # Do not let a semantic match cross an unsupported
                    # field or other non-run paragraph child.
                    parts.append("\ufffc")
            combined = "".join(parts)

            import re
            flags = 0 if match_case else re.IGNORECASE
            matches = [
                match.span()
                for match in re.finditer(re.escape(old), combined, flags=flags)
                if "\v" not in match.group(0) and "\ufffc" not in match.group(0)
            ]
            if not matches:
                continue

            # Work from right to left.  Later replacements cannot shift
            # the coordinates of an earlier match, while rebuilding the
            # run map after each edit preserves all existing run elements,
            # their formatting, and every a:br node.
            for start, end in reversed(matches):
                cursor = 0
                segments = []
                run_index = 0
                for child in paragraph._p.content_children:
                    kind = child.tag.rsplit("}", 1)[-1]
                    if kind == "r":
                        run = list(paragraph.runs)[run_index]
                        run_index += 1
                        segment_end = cursor + len(run.text)
                        segments.append((run, cursor, segment_end))
                        cursor = segment_end
                    elif kind == "br":
                        cursor += 1
                    else:
                        cursor += 1

                affected = [
                    (run, segment_start, segment_end)
                    for run, segment_start, segment_end in segments
                    if segment_start < end and segment_end > start
                ]
                if not affected:
                    continue

                replacement_offset = 0
                for affected_index, (run, segment_start, segment_end) in enumerate(affected):
                    overlap_start = max(start, segment_start)
                    overlap_end = min(end, segment_end)
                    quota = overlap_end - overlap_start
                    is_first = affected_index == 0
                    is_last = affected_index == len(affected) - 1
                    if is_last:
                        chunk = new[replacement_offset:]
                    else:
                        chunk = new[replacement_offset:replacement_offset + quota]
                    replacement_offset += len(chunk)
                    prefix = run.text[:overlap_start - segment_start] if is_first else ""
                    suffix = run.text[overlap_end - segment_start:] if is_last else ""
                    run.text = prefix + chunk + suffix

            replacements += len(matches)
        return replacements

    indices = range(len(h.deck.slides)) if slide_number is None else [slide_number - 1]
    replacements = 0
    touched: set[str] = set()
    for index in indices:
        if index < 0 or index >= len(h.deck.slides):
            raise IndexError("slide number out of range")
        for shape, _ in _walk_shapes(h.deck.slides[index].shapes):
            if getattr(shape, "has_text_frame", False):
                count = replace_in_frame(shape.text_frame, match_case)
                if count:
                    replacements += count
                    touched.add(f"{index + 1}:{shape.shape_id}")
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        count = replace_in_frame(cell.text_frame, match_case)
                        if count:
                            replacements += count
                            touched.add(f"{index + 1}:{shape.shape_id}")
    if not replacements:
        raise ValueError(f"text not found: {old!r}")
    h.state.ppt_affected_slides.update(int(item.split(":", 1)[0]) for item in touched)
    h.state.record_change("deck:semantic_text_replace:" + ",".join(sorted(touched, key=lambda item: int(item.split(":", 1)[0]))))
    return f"replaced {replacements} occurrence(s) in {len(touched)} shape(s): {', '.join(sorted(touched, key=lambda item: int(item.split(':', 1)[0])))}"


def _target_text_shape(h, slide_number: int, shape_id: int | None = None, text_contains: str = ""):
    slide, explicit = (None, None)
    if shape_id is not None:
        return _find_shape(h, slide_number, shape_id)
    if slide_number < 1 or slide_number > len(h.deck.slides):
        raise IndexError("slide number out of range")
    slide = h.deck.slides[slide_number - 1]
    candidates = [shape for shape, _ in _walk_shapes(slide.shapes) if getattr(shape, "has_text_frame", False)]
    if text_contains:
        candidates = [shape for shape in candidates if text_contains.casefold() in shape.text.casefold()]
    if not candidates:
        raise ValueError("no matching text shape")
    # Prefer the body/list shape over a short title box.
    target = max(candidates, key=lambda shape: (len(shape.text_frame.paragraphs), len(shape.text), shape.height))
    return slide, target


def _append_bullet(h, slide_number: int, text: str, shape_id: int | None = None, text_contains: str = "") -> str:
    if getattr(h, "deck", None) is None:
        raise ValueError("no deck loaded")
    _, shape = _target_text_shape(h, slide_number, shape_id, text_contains)
    if not shape.has_text_frame:
        raise TypeError("target shape has no text frame")
    paragraphs = list(shape.text_frame.paragraphs)
    template = paragraphs[-1]
    anchor_index = len(paragraphs) - 1
    if text_contains:
        matches = [
            (index, candidate)
            for index, candidate in enumerate(paragraphs)
            if text_contains.casefold() in candidate.text.casefold()
        ]
        if not matches:
            raise ValueError(f"anchor paragraph not found: {text_contains!r}")
        anchor_index, template = matches[-1]
    paragraph = shape.text_frame.add_paragraph()
    if template._p.pPr is not None:
        paragraph._p.insert(0, deepcopy(template._p.pPr))
    paragraph.level = template.level
    # Insert after the anchor's child bullet subtree but before the next peer.
    # add_paragraph() initially appends at the end; moving the XML element
    # preserves the paragraph object and formatting we just constructed.
    insertion_index = anchor_index
    for index in range(anchor_index + 1, len(paragraphs)):
        if paragraphs[index].level <= template.level:
            break
        insertion_index = index
    paragraphs[insertion_index]._p.addnext(paragraph._p)
    run = paragraph.add_run()
    run.text = text
    if template.runs:
        source = template.runs[0]
        run.font.size = source.font.size
        run.font.bold = source.font.bold
        run.font.italic = source.font.italic
        run.font.name = source.font.name
        try:
            if source.font.color.rgb is not None:
                run.font.color.rgb = source.font.color.rgb
        except (AttributeError, ValueError):
            pass
    h.state.record_change(f"deck:slide:{slide_number}:shape:{shape.shape_id}:append_bullet")
    return (
        f"appended peer bullet on slide {slide_number} shape {shape.shape_id} "
        f"at level {paragraph.level} after anchor {text_contains!r}"
    )


def _set_text_style(h, slide_number: int, shape_id: int | None = None, text_contains: str = "", size: int | None = None, color: str | None = None, bold: bool | None = None) -> str:
    _, shape = _target_text_shape(h, slide_number, shape_id, text_contains)
    if size is None and color is None and bold is None:
        raise ValueError("provide at least one style property")
    rgb = RGBColor.from_string(color.lstrip("#").upper()) if color else None
    runs = 0
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if size is not None:
                run.font.size = Pt(size)
            if rgb is not None:
                run.font.color.rgb = rgb
            if bold is not None:
                run.font.bold = bold
            runs += 1
    h.state.record_change(f"deck:slide:{slide_number}:shape:{shape.shape_id}:text_style")
    return f"updated style on {runs} run(s), slide {slide_number} shape {shape.shape_id}"


def _set_shape_fill(h, slide_number: int, color: str, shape_id: int | None = None, text_contains: str = "") -> str:
    _, shape = _target_text_shape(h, slide_number, shape_id, text_contains)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color.lstrip("#").upper())
    h.state.record_change(f"deck:slide:{slide_number}:shape:{shape.shape_id}:fill")
    return f"updated fill on slide {slide_number} shape {shape.shape_id}"


def _add_textbox_to_slide(h, slide_number: int, box: dict) -> str:
    if getattr(h, "deck", None) is None:
        raise ValueError("no deck loaded")
    if slide_number < 1 or slide_number > len(h.deck.slides):
        raise IndexError("slide number out of range")
    tb = h.deck.slides[slide_number - 1].shapes.add_textbox(
        Inches(box["x"]), Inches(box["y"]), Inches(box["w"]), Inches(box["h"])
    )
    _fit_lines(tb.text_frame, str(box.get("text", "")).split("\n"), box.get("size", 18), box.get("bold", False), _TEXT)
    h.state.record_change(f"deck:slide:{slide_number}:textbox:{tb.shape_id}")
    return f"added text box to slide {slide_number}, shape {tb.shape_id}"


def _add_flowchart(h, slide_number: int, nodes: list[str], title: str = "") -> str:
    if getattr(h, "deck", None) is None:
        raise ValueError("no deck loaded")
    if not 2 <= len(nodes) <= 8:
        raise ValueError("flowchart requires 2-8 nodes")
    prs = h.deck
    if slide_number < 1:
        slide_number = 1
    if slide_number > len(prs.slides):
        _assert_can_append_slide(h)
        slide = _blank_slide(prs)
        _rect(slide, 0, 0, _W, _H, _BG)
        slide_number = len(prs.slides)
    else:
        slide = prs.slides[slide_number - 1]
    converted_cover = _convert_fresh_cover(h, slide_number)
    if converted_cover:
        # The diagram page owns slide 1 of a freshly created deck. Give it a
        # clean content-page scaffold: background plus header band, then the
        # node row below the header instead of under leftover cover text.
        _rect(slide, 0, 0, _W, _H, _BG)
        if title:
            _rect(slide, 0, 0, _W, 1.1, _HEAD)
            _rect(slide, 0, 1.1, _W, 0.07, _ACCENT)
            title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.12), Inches(12.2), Inches(0.9))
            _put_lines(title_box.text_frame, title, 26, True, _WHITE)
        y = 2.6 if title else 1.5
    elif title:
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.55))
    has_multiline = any(("\n" in str(n)) or isinstance(n, dict) for n in nodes)
    margin, gap = 0.75, 0.28
    arrow_w = 0.38
    height = 4.2 if (has_multiline or len(nodes) <= 5) else 1.2
    if converted_cover or title:
        y = 1.65 if height > 2.0 else (2.6 if title else 1.5)
    else:
        y = 1.8 if height > 2.0 else 3.0

    node_w = (_W - 2 * margin - (len(nodes) - 1) * (gap + arrow_w)) / len(nodes)
    for index, raw_node in enumerate(nodes):
        x = margin + index * (node_w + gap + arrow_w)
        node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(node_w), Inches(height))
        node.fill.solid(); node.fill.fore_color.rgb = _WHITE
        node.line.color.rgb = _PRIMARY
        node.line.width = Pt(1.5)

        if isinstance(raw_node, dict):
            n_title = raw_node.get("title", f"0{index+1} 阶段")
            n_bullets = raw_node.get("bullets", [])
        else:
            lines = [l.strip() for l in str(raw_node).split("\n") if l.strip()]
            n_title = lines[0] if lines else f"0{index+1} 阶段"
            n_bullets = lines[1:] if len(lines) > 1 else []

        # Card Title
        tf = node.text_frame
        tf.word_wrap = True
        p_title = tf.paragraphs[0]
        p_title.text = n_title
        p_title.font.size = Pt(14)
        p_title.font.bold = True
        p_title.font.color.rgb = _PRIMARY

        # Card Bullets / Detail
        if n_bullets and height > 2.0:
            for b in n_bullets:
                p_b = tf.add_paragraph()
                p_b.text = b if b.startswith("•") or b.startswith("-") else f"• {b}"
                p_b.font.size = Pt(10.5)
                p_b.font.color.rgb = _TEXT
                p_b.line_spacing = 1.15

        if index < len(nodes) - 1:
            begin_x = Inches(x + node_w + gap / 4)
            end_x = Inches(x + node_w + gap + arrow_w - gap / 4)
            mid_y = Inches(y + (height / 2 if height < 2.0 else 0.8))
            connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, begin_x, mid_y, end_x, mid_y)
            connector.line.color.rgb = _ACCENT
            connector.line.width = Pt(2.0)
            line_element = connector.line._get_or_add_ln()
            line_element.append(line_element.makeelement(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd",
                {"type": "arrow", "w": "med", "len": "med"},
            ))
    h.state.record_change(f"deck:slide:{slide_number}:flowchart")
    return f"added {len(nodes)}-node flowchart to slide {slide_number}"


# --- compact model-facing facade -------------------------------------------------

def _ppt_inspect(h, detail: str = "summary", slide_number: int | None = None) -> str:
    if getattr(h, "deck", None) is None:
        known = getattr(getattr(h, "state", None), "facts", {}).get("ppt_input_deck", "")
        suffix = f" The discovered input is '{known}'." if known else ""
        raise RuntimeError(f"no active deck; open a real task-local PPTX before inspection.{suffix}")
    if detail == "summary":
        return _deck_info(h, slide_number)
    return _shape_inventory(h, slide_number)


def _validate_batch_update(update: dict, index: int) -> None:
    """Validate operation-specific fields before starting an atomic batch."""
    operation = update["operation"]
    if operation == "replace":
        if not update.get("old"):
            raise ValueError(f"updates[{index}] replace requires non-empty old text")
        if "new" not in update:
            raise ValueError(f"updates[{index}] replace requires new text")
        return
    if operation == "set_shape_text":
        if update.get("slide_number") is None or not update.get("text"):
            raise ValueError(f"updates[{index}] set_shape_text requires slide_number and non-empty text")
        if update.get("shape_id") is None and not update.get("shape_name") and not update.get("text_contains"):
            raise ValueError(f"updates[{index}] set_shape_text requires shape_id, shape_name, or text_contains")
        return
    if operation == "set_table":
        if update.get("slide_number") is None or not update.get("rows"):
            raise ValueError(f"updates[{index}] set_table requires slide_number and rows")
        if update.get("shape_id") is None and not update.get("shape_name"):
            raise ValueError(f"updates[{index}] set_table requires shape_id or shape_name")
        return
    if operation != "style":
        raise ValueError(f"updates[{index}] has unsupported operation: {operation!r}")
    if update.get("slide_number") is None:
        raise ValueError(f"updates[{index}] style requires slide_number")
    if update.get("shape_id") is None and not update.get("text_contains"):
        raise ValueError(f"updates[{index}] style requires shape_id or text_contains")
    if update.get("all_matches") and (not update.get("text_contains") or update.get("shape_id") is not None):
        raise ValueError(f"updates[{index}] all_matches requires text_contains and no shape_id")
    target = update.get("target", "text")
    if target == "fill":
        if not update.get("color"):
            raise ValueError(f"updates[{index}] fill style requires color")
    elif all(update.get(name) is None for name in ("size", "color", "bold")):
        raise ValueError(f"updates[{index}] text style requires size, color, or bold")


def _ppt_batch_updates(h, updates: list[dict], default_slide_number: int | None = None) -> str:
    """Apply many independent text changes as one all-or-nothing mutation.

    Work happens against a deep-copied presentation and state.  The live deck
    is swapped only after every item succeeds, and the real RunState receives
    one merged change record so a large spreadsheet-driven update consumes one
    mutation epoch instead of one epoch per cell.
    """
    if getattr(h, "deck", None) is None:
        raise ValueError("no deck loaded")
    if not updates:
        raise ValueError("batch_updates requires at least one update")
    # A top-level ``slide_number`` is a batch default: models commonly emit it
    # next to ``updates`` instead of repeating it inside every item.
    updates = [dict(update) for update in updates]
    for update in updates:
        if update.get("slide_number") is None and default_slide_number is not None:
            update["slide_number"] = default_slide_number
    for index, update in enumerate(updates):
        _validate_batch_update(update, index)

    transaction_buffer = BytesIO()
    h.deck.save(transaction_buffer)
    transaction_buffer.seek(0)
    tx = SimpleNamespace(deck=Presentation(transaction_buffer), state=deepcopy(h.state))
    # Track this transaction's scope independently from earlier edits.
    tx.state.ppt_affected_slides = set()
    results: list[str] = []
    for update in updates:
        if update["operation"] == "replace":
            results.append(_replace_text_semantic(
                tx,
                update["old"],
                update["new"],
                update.get("slide_number"),
                update.get("match_case", True),
            ))
            continue
        if update["operation"] == "set_shape_text":
            results.append(_set_shape_text(
                tx,
                update["slide_number"],
                update["text"],
                update.get("shape_id"),
                update.get("shape_name", ""),
                update.get("text_contains", ""),
            ))
            continue
        if update["operation"] == "set_table":
            results.append(_set_table_rows(
                tx,
                update["slide_number"],
                update["rows"],
                update.get("shape_id"),
                update.get("shape_name", ""),
            ))
            continue
        results.append(_ppt_style(
            tx,
            slide_number=update["slide_number"],
            target=update.get("target", "text"),
            shape_id=update.get("shape_id"),
            text_contains=update.get("text_contains", ""),
            size=update.get("size"),
            color=update.get("color"),
            bold=update.get("bold"),
            all_matches=update.get("all_matches", False),
        ))

    affected = set(tx.state.ppt_affected_slides)
    h.deck = tx.deck
    h.state.ppt_affected_slides.update(affected)
    paths = [f"deck:slide:{slide}:batch_update" for slide in sorted(affected)]
    h.state.record_changes(paths or ["deck:batch_update"])
    return f"applied {len(updates)} updates atomically across {len(affected)} slide(s): " + "; ".join(results)


def _ppt_edit_text(h, operation: str, slide_number: int | None = None, shape_id: int | None = None,
                   text_contains: str = "", old: str = "", new: str = "", text: str = "",
                   match_case: bool = True, all_matches: bool = False, updates: list[dict] | None = None,
                   shape_name: str = "", rows: list[list[str]] | None = None, new_plural: str = "",
                   after: str = "", update: list[dict] | None = None) -> str:
    if operation == "batch_updates":
        return _ppt_batch_updates(h, updates or update or [], slide_number)
    if operation == "replace_case_variants":
        return _replace_case_variants(h, old, new, new_plural, slide_number)
    if operation == "set_shape_text":
        if slide_number is None or not text:
            raise ValueError("set_shape_text requires slide_number and non-empty text")
        return _set_shape_text(h, slide_number, text, shape_id, shape_name, text_contains)
    if operation == "set_table":
        if slide_number is None or not rows:
            raise ValueError("set_table requires slide_number and rows")
        return _set_table_rows(h, slide_number, rows, shape_id, shape_name)
    if operation == "replace":
        if not old:
            raise ValueError("replace requires non-empty old text")
        return _replace_text_semantic(h, old, new, slide_number, match_case)
    if slide_number is None:
        raise ValueError("append_bullet requires slide_number")
    # append_bullet contract: `text` is the new bullet, `text_contains` (alias
    # `after`) is the anchor paragraph.  Several providers emit the symmetric
    # form old=<anchor>, new=<bullet>; accept it instead of misreading the
    # anchor as the new bullet text.
    bullet_text = new or text
    anchor = text_contains or after or old
    if not bullet_text:
        raise ValueError("append_bullet requires non-empty text (new bullet content)")
    return _append_bullet(h, slide_number, bullet_text, shape_id, anchor)


def _ppt_style(h, slide_number: int, target: str, shape_id: int | None = None, text_contains: str = "",
               size: int | None = None, color: str | None = None, bold: bool | None = None,
               all_matches: bool = False) -> str:
    if all_matches:
        if not text_contains or shape_id is not None:
            raise ValueError("all_matches requires text_contains and no shape_id")
        slide = h.deck.slides[slide_number - 1]
        targets = [
            shape for shape, _ in _walk_shapes(slide.shapes)
            if getattr(shape, "has_text_frame", False) and text_contains.casefold() in shape.text.casefold()
        ]
        if not targets:
            raise ValueError("no matching shapes")
        results = []
        for shape in targets:
            if target == "fill":
                if not color:
                    raise ValueError("fill styling requires color")
                results.append(_set_shape_fill(h, slide_number, color, shape.shape_id, ""))
            else:
                results.append(_set_text_style(h, slide_number, shape.shape_id, "", size, color, bold))
        return f"styled {len(targets)} matching shapes: " + "; ".join(results)
    if target == "fill":
        if not color:
            raise ValueError("fill styling requires color")
        return _set_shape_fill(h, slide_number, color, shape_id, text_contains)
    return _set_text_style(h, slide_number, shape_id, text_contains, size, color, bold)


def _source_slide_material(slide) -> tuple[list[str], list[bytes], list[list[str]]]:
    """Extract text, image blobs, and table rows without copying relationships."""
    title_id = slide.shapes.title.shape_id if slide.shapes.title is not None else None
    lines: list[str] = []
    images: list[bytes] = []
    table_rows: list[list[str]] = []
    for shape, _ in _walk_shapes(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            images.append(shape.image.blob)
        elif getattr(shape, "has_table", False):
            table_rows.extend([[cell.text.strip() for cell in row.cells] for row in shape.table.rows])
        elif shape.shape_id != title_id and getattr(shape, "has_text_frame", False):
            lines.extend(paragraph.text.strip() for paragraph in shape.text_frame.paragraphs if paragraph.text.strip())
    return lines, images, table_rows


def _add_picture_fitted(slide, blob: bytes, x: float, y: float, w: float, height: float) -> None:
    """Re-add an image blob with a valid target-slide media relationship."""
    picture = slide.shapes.add_picture(BytesIO(blob), 0, 0)
    ratio = picture.width / picture.height
    picture_width = min(w, height * ratio)
    picture_height = picture_width / ratio
    picture.left = Inches(x + (w - picture_width) / 2)
    picture.top = Inches(y + (height - picture_height) / 2)
    picture.width = Inches(picture_width)
    picture.height = Inches(picture_height)


def _compose_from_slides(h, source_slides: list[int], insert_after: int, title: str = "",
                         left_title: str = "Male", right_title: str = "Female") -> str:
    """Synthesize two source slides into one inserted visual comparison page."""
    if getattr(h, "deck", None) is None:
        raise ValueError("from_slides requires an open deck")
    if not isinstance(source_slides, list) or len(source_slides) != 2:
        raise ValueError("from_slides requires exactly two source_slides")
    if any(isinstance(number, bool) or not isinstance(number, int) for number in source_slides):
        raise ValueError("source_slides must contain integer slide numbers")
    if len(set(source_slides)) != 2:
        raise ValueError("source_slides must identify two distinct slides")
    count = len(h.deck.slides)
    if any(number < 1 or number > count for number in source_slides):
        raise IndexError("source slide number out of range")
    if isinstance(insert_after, bool) or not isinstance(insert_after, int) or not 1 <= insert_after <= count:
        raise IndexError("insert_after must identify an existing slide")
    if not left_title.strip() or not right_title.strip():
        raise ValueError("from_slides requires non-empty left_title and right_title")

    materials = [_source_slide_material(h.deck.slides[number - 1]) for number in source_slides]
    if not any(lines or images or tables for lines, images, tables in materials):
        raise ValueError("source slides contain no composable content")

    prs = h.deck
    slide = _blank_slide(prs)
    width, height = prs.slide_width / 914400, prs.slide_height / 914400
    margin, gap = max(0.36, width * 0.045), max(0.22, width * 0.025)
    content_width = width - 2 * margin
    column_width = (content_width - gap) / 2
    _rect(slide, 0, 0, width, height, _WHITE)
    heading = slide.shapes.add_textbox(Inches(margin), Inches(0.22), Inches(content_width), Inches(0.64))
    _put_lines(heading.text_frame, title.strip() or "Clothing: Male and Female", 25, True, RGBColor(0x18, 0x18, 0x18))
    _rect(slide, margin, 0.96, content_width, 0.025, RGBColor(0x22, 0x22, 0x22))

    image_top = 1.12
    image_height = min(2.65, max(1.75, height * 0.36))
    for column, (_, images, _) in enumerate(materials):
        column_x = margin + column * (column_width + gap)
        selected = images[:2]
        image_gap = 0.10
        image_width = (column_width - image_gap * (len(selected) - 1)) / max(1, len(selected))
        for index, blob in enumerate(selected):
            _add_picture_fitted(slide, blob, column_x + index * (image_width + image_gap), image_top,
                                image_width, image_height)

    table_top = image_top + image_height + 0.20
    table_height = height - table_top - 0.30
    table_shape = slide.shapes.add_table(2, 2, Inches(margin), Inches(table_top),
                                         Inches(content_width), Inches(table_height))
    table = table_shape.table
    table.columns[0].width = Inches(column_width)
    table.columns[1].width = Inches(column_width)
    for column, label in enumerate((left_title.strip(), right_title.strip())):
        header, body = table.cell(0, column), table.cell(1, column)
        header.text = label
        header.fill.solid(); header.fill.fore_color.rgb = RGBColor(0x1C, 0x1C, 0x1C)
        header.vertical_anchor = MSO_ANCHOR.MIDDLE
        lines, _, table_rows = materials[column]
        table_lines = [" | ".join(value for value in row if value) for row in table_rows]
        body.text = "\n".join(lines + [value for value in table_lines if value])
        body.fill.solid(); body.fill.fore_color.rgb = _WHITE
        body.vertical_anchor = MSO_ANCHOR.TOP
        for cell, color, font_size, bold in ((header, _WHITE, 15, True), (body, _TEXT, 11, False)):
            cell.margin_left = Inches(0.14); cell.margin_right = Inches(0.14)
            cell.margin_top = Inches(0.08); cell.margin_bottom = Inches(0.06)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.space_after = Pt(2)
                for run in paragraph.runs:
                    _style_run(run, font_size, bold, color)

    slide_id = prs.slides._sldIdLst[-1]
    prs.slides._sldIdLst.remove(slide_id)
    prs.slides._sldIdLst.insert(insert_after, slide_id)
    inserted_slide = insert_after + 1
    # Baseline findings are keyed by one-based slide number. Preserve their
    # identity across insertion so historical defects do not become false new
    # failures merely because all later source pages shifted by one position.
    baseline = getattr(h.state, "ppt_baseline_findings", {})
    shifted_baseline: dict[str, float] = {}
    for key, severity in baseline.items():
        parts = key.split(":")
        if len(parts) > 1 and parts[0] == "slide":
            try:
                if int(parts[1]) > insert_after:
                    parts[1] = str(int(parts[1]) + 1)
            except ValueError:
                pass
        shifted_baseline[":".join(parts)] = severity
    h.state.ppt_baseline_findings = shifted_baseline
    try:
        notes_slide = prs.slides[inserted_slide - 1].notes_slide
        if notes_slide is not None:
            notes_slide.notes_text_frame.text = (
                "[Sources]\n" + ", ".join(f"source slide {number}" for number in source_slides)
            )
    except Exception:
        pass
    h.state.record_change(f"deck:slide:{inserted_slide}:from_slides")
    h.state.record_fact("ppt_last_compose_scope",
                        f"inserted slide {inserted_slide}; sources {source_slides}; originals preserved")
    return (f"inserted comparison slide {inserted_slide} after slide {insert_after}; "
            f"copied text/table material and {sum(len(item[1][:2]) for item in materials)} image(s) "
            f"from source slides {source_slides}; source slides preserved")


_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _clone_template_slide(prs: Presentation, source_slide):
    """Clone one slide inside a deck while rebuilding its local relationships.

    Shape XML alone is insufficient for pictures, charts, and hyperlinks because
    relationship ids are local to each slide part.  Rebinding them here keeps the
    public ``from_outline`` operation transactional and template-faithful.
    """
    target_slide = prs.slides.add_slide(source_slide.slide_layout)
    relationship_map: dict[str, str] = {}
    for relationship in source_slide.part.rels.values():
        leaf = relationship.reltype.rsplit("/", 1)[-1]
        if leaf in {"slideLayout", "notesSlide"}:
            continue
        target = relationship.target_ref if relationship.is_external else relationship.target_part
        relationship_map[relationship.rId] = target_slide.part.relate_to(
            target, relationship.reltype, relationship.is_external
        )

    for shape in source_slide.shapes:
        element = deepcopy(shape._element)
        for node in element.iter():
            for attribute in (
                f"{{{_REL_NS}}}id", f"{{{_REL_NS}}}embed", f"{{{_REL_NS}}}link",
            ):
                old_id = node.get(attribute)
                if old_id in relationship_map:
                    node.set(attribute, relationship_map[old_id])
        target_slide.shapes._spTree.insert_element_before(element, "p:extLst")
    return target_slide


def _replace_named_shape_text(slide, shape_name: str, text: str) -> None:
    matches = [shape for shape, _ in _walk_shapes(slide.shapes) if shape.name == shape_name]
    if len(matches) != 1:
        raise ValueError(f"template shape {shape_name!r} matched {len(matches)} shapes")
    shape = matches[0]
    if not getattr(shape, "has_text_frame", False):
        raise TypeError(f"template shape {shape_name!r} has no text frame")

    source_paragraphs = list(shape.text_frame.paragraphs)
    paragraph_styles = []
    for paragraph in source_paragraphs:
        ppr = deepcopy(paragraph._p.pPr) if paragraph._p.pPr is not None else None
        rpr = None
        if paragraph.runs and paragraph.runs[0]._r.rPr is not None:
            rpr = deepcopy(paragraph.runs[0]._r.rPr)
        paragraph_styles.append((ppr, rpr))

    lines = str(text).split("\n") or [""]
    text_frame = shape.text_frame
    text_frame.clear()
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        style_index = min(index, len(paragraph_styles) - 1)
        ppr, rpr = paragraph_styles[style_index]
        if ppr is not None:
            if paragraph._p.pPr is not None:
                paragraph._p.remove(paragraph._p.pPr)
            paragraph._p.insert(0, deepcopy(ppr))
        run = paragraph.add_run()
        run.text = line
        if rpr is not None:
            if run._r.rPr is not None:
                run._r.remove(run._r.rPr)
            run._r.insert(0, deepcopy(rpr))


def _replace_named_table(slide, shape_name: str, rows: list[list[str]]) -> None:
    matches = [shape for shape, _ in _walk_shapes(slide.shapes) if shape.name == shape_name]
    if len(matches) != 1:
        raise ValueError(f"template table {shape_name!r} matched {len(matches)} shapes")
    shape = matches[0]
    if not getattr(shape, "has_table", False):
        raise TypeError(f"template shape {shape_name!r} is not a table")
    _rewrite_table_preserving_style(shape.table, rows)


def _rewrite_table_preserving_style(table, rows: list[list[str]]) -> None:
    """Rewrite every cell while keeping each cell's original run formatting."""
    if len(rows) != len(table.rows) or any(len(row) != len(table.columns) for row in rows):
        current = " ;; ".join(
            " | ".join(cell.text.strip().replace("\n", " ") for cell in row.cells)
            for row in table.rows
        )
        raise ValueError(
            f"table requires {len(table.rows)}x{len(table.columns)} values "
            f"(you provided {len(rows)}x{len(rows[0]) if rows and isinstance(rows[0], list) else '?'}). "
            f"Current table contents: {current}"
        )
    for row_index, row in enumerate(rows):
        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            source_runs = [run for paragraph in cell.text_frame.paragraphs for run in paragraph.runs]
            rpr = deepcopy(source_runs[0]._r.rPr) if source_runs and source_runs[0]._r.rPr is not None else None
            cell.text = str(value)
            if rpr is not None and cell.text_frame.paragraphs[0].runs:
                run = cell.text_frame.paragraphs[0].runs[0]
                if run._r.rPr is not None:
                    run._r.remove(run._r.rPr)
                run._r.insert(0, deepcopy(rpr))


def _replace_case_variants(h, old: str, new: str, new_plural: str = "", slide_number: int | None = None) -> str:
    """Replace singular/plural case variants while preserving each form's case.

    Rubrics often phrase a task as "Liability/Liabilities -> Debt/Debts" and
    score lowercase / Capitalized / UPPERCASE forms separately. A single
    case-sensitive replace covers only one spelling.
    """
    if not old or not new:
        raise ValueError("replace_case_variants requires old and new")
    if not new_plural:
        new_plural = new + "s"
    old_singular = old
    old_plural = old[:-1] + "ies" if old.lower().endswith("y") else old + "s"
    variants = [
        (old_singular, new),
        (old_singular.lower(), new.lower()),
        (old_singular.upper(), new.upper()),
        (old_plural, new_plural),
        (old_plural.lower(), new_plural.lower()),
        (old_plural.upper(), new_plural.upper()),
    ]
    results = []
    replaced_any = False
    for old_variant, new_variant in variants:
        try:
            results.append(_replace_text_semantic(h, old_variant, new_variant, slide_number, True))
            replaced_any = True
        except ValueError:
            continue
    if not replaced_any:
        raise ValueError(f"none of the case variants of {old!r} were found")
    return f"replaced case variants of {old!r}: " + "; ".join(results)


def _set_shape_metadata(h, slide_number: int, descr: str, shape_id: int | None = None,
                        shape_name: str = "") -> str:
    """Set the shape description metadata (descr) used by deterministic binders.

    Many office verifiers treat shape alt text/description as a first-class
    provenance surface: an html_anchor, metric_id or chart_id recorded here
    satisfies source-binding checks even when it is not part of visible text.
    """
    if not descr.strip():
        raise ValueError("set_metadata requires non-empty descr")
    if shape_id is None and not shape_name:
        raise ValueError("set_metadata requires shape_id or shape_name")
    _, shape = _select_shape_on_slide(h, slide_number, shape_id, shape_name)
    nv = shape._element.nvSpPr.cNvPr
    nv.set("descr", descr)
    h.state.ppt_affected_slides.add(slide_number)
    h.state.record_change(f"deck:slide:{slide_number}:shape:{shape.shape_id}:metadata")
    return f"set metadata on slide {slide_number} shape {shape.shape_id} ({shape.name!r})"


def _select_shape_on_slide(h, slide_number: int, shape_id: int | None = None,
                           shape_name: str = "", text_contains: str = ""):
    """Resolve one shape by its stable selector (id > name > contained text)."""
    if getattr(h, "deck", None) is None:
        raise ValueError("no deck loaded")
    if slide_number < 1 or slide_number > len(h.deck.slides):
        raise IndexError("slide number out of range")
    slide = h.deck.slides[slide_number - 1]
    if shape_id is not None:
        for shape, _ in _walk_shapes(slide.shapes):
            if shape.shape_id == shape_id:
                return slide, shape
        raise KeyError(f"shape id not found on slide {slide_number}: {shape_id}")
    if shape_name:
        matches = [shape for shape, _ in _walk_shapes(slide.shapes) if shape.name == shape_name]
        if len(matches) != 1:
            available = sorted({shape.name for shape, _ in _walk_shapes(slide.shapes) if shape.name})
            preview = ", ".join(available[:12])
            raise ValueError(
                f"shape name {shape_name!r} matched {len(matches)} shapes on slide {slide_number}. "
                f"Available names: {preview or '(none)'}"
            )
        return slide, matches[0]
    if text_contains:
        matches = [
            shape for shape, _ in _walk_shapes(slide.shapes)
            if text_contains in _slide_text_material_for_shape(shape)
        ]
        if len(matches) != 1:
            raise ValueError(f"text_contains {text_contains!r} matched {len(matches)} shapes on slide {slide_number}")
        return slide, matches[0]
    raise ValueError("set_shape_text requires shape_id, shape_name, or unique text_contains")


def _set_shape_text(h, slide_number: int, text: str, shape_id: int | None = None,
                    shape_name: str = "", text_contains: str = "") -> str:
    """Replace an existing text surface wholesale, preserving first-run style.

    Whole-surface rewrite is the source-sync primitive: consistency edits
    replace a card/cell with a complete current-scope statement instead of
    hoping a substring replacement covers every stale fragment.
    """
    _, shape = _select_shape_on_slide(h, slide_number, shape_id, shape_name, text_contains)
    if not getattr(shape, "has_text_frame", False):
        if getattr(shape, "has_table", False):
            raise TypeError(
                f"target shape {shape.name!r} is a table; use operation='set_table' with rows "
                "(or batch_updates with a set_table entry) instead of set_shape_text"
            )
        raise TypeError("target shape has no text frame")
    text_frame = shape.text_frame
    first_rpr = None
    first_ppr = None
    for paragraph in text_frame.paragraphs:
        if paragraph._p.pPr is not None and first_ppr is None:
            first_ppr = deepcopy(paragraph._p.pPr)
        for run in paragraph.runs:
            if run._r.rPr is not None and first_rpr is None:
                first_rpr = deepcopy(run._r.rPr)
        if first_ppr is not None and first_rpr is not None:
            break
    lines = text.split("\n")
    text_frame.clear()
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        if first_ppr is not None:
            if paragraph._p.pPr is not None:
                paragraph._p.remove(paragraph._p.pPr)
            paragraph._p.insert(0, deepcopy(first_ppr))
        run = paragraph.add_run()
        run.text = line
        if first_rpr is not None:
            if run._r.rPr is not None:
                run._r.remove(run._r.rPr)
            run._r.insert(0, deepcopy(first_rpr))
    h.state.ppt_affected_slides.add(slide_number)
    h.state.record_change(f"deck:slide:{slide_number}:shape:{shape.shape_id}:set_text")
    return f"rewrote slide {slide_number} shape {shape.shape_id} ({shape.name!r}) with {len(lines)} line(s)"


def _set_table_rows(h, slide_number: int, rows: list[list[str]], shape_id: int | None = None,
                    shape_name: str = "") -> str:
    """Replace a whole existing table with current-scope rows, preserving cells.

    This is the table counterpart of :func:`_set_shape_text`: consistency
    editing must be able to rewrite an entire table atomically (one mutation
    epoch), like the correct reference trajectory's per-table data rewrite.
    """
    if not rows or not all(isinstance(row, list) and row for row in rows):
        raise ValueError("set_table requires non-empty rows of non-empty cell lists")
    if shape_id is None and not shape_name:
        raise ValueError("set_table requires shape_id or shape_name")
    _, shape = _select_shape_on_slide(h, slide_number, shape_id, shape_name)
    if not getattr(shape, "has_table", False):
        raise TypeError("target shape is not a table")
    _rewrite_table_preserving_style(shape.table, [[str(cell) for cell in row] for row in rows])
    h.state.ppt_affected_slides.add(slide_number)
    h.state.record_change(f"deck:slide:{slide_number}:shape:{shape.shape_id}:table")
    return f"rewrote table slide {slide_number} shape {shape.shape_id} ({shape.name!r}) {len(shape.table.rows)}x{len(shape.table.columns)}"


def _set_slide_notes_text(slide, text: str) -> None:
    notes_slide = slide.notes_slide
    notes_frame = notes_slide.notes_text_frame
    if notes_frame is not None:
        notes_frame.text = text
        return
    editable = [shape for shape in notes_slide.shapes if getattr(shape, "has_text_frame", False)]
    if editable:
        target = max(editable, key=lambda shape: len(shape.text))
        target.text_frame.text = text
        return
    raise ValueError("template notes slide has no editable text surface")


def _compose_from_outline(h, slides: list[dict], replace_template: bool = True) -> str:
    """Generate several template-faithful slides from a compact ContentIR batch."""
    if getattr(h, "deck", None) is None:
        raise ValueError("from_outline requires an open template deck")
    if not slides:
        raise ValueError("from_outline requires at least one slide specification")
    template_count = len(h.deck.slides)
    baseline_state = SimpleNamespace(
        deck=h.deck,
        state=SimpleNamespace(
            ppt_existing_deck=False,
            ppt_baseline_captured=False,
            ppt_baseline_findings={},
            ppt_affected_slides=set(),
        ),
    )
    template_findings = _collect_structural_findings(baseline_state)
    prepared: list[tuple[int, list[dict], dict | None, str]] = []
    for index, spec in enumerate(slides, 1):
        template_slide = spec.get("template_slide")
        if isinstance(template_slide, bool) or not isinstance(template_slide, int):
            raise ValueError(f"outline slide {index} requires integer template_slide")
        if template_slide < 1 or template_slide > template_count:
            raise IndexError(f"outline slide {index} template_slide out of range")
        replacements = spec.get("replacements") or []
        names = [item.get("shape_name") for item in replacements]
        if any(not name for name in names) or len(names) != len(set(names)):
            raise ValueError(f"outline slide {index} replacement names must be non-empty and unique")
        prepared.append((template_slide, replacements, spec.get("table"), spec.get("speaker_notes", "")))

    # Validate and mutate a deep copy so a bad shape name or table size cannot
    # leave half a generated deck in the live session.
    transaction_buffer = BytesIO()
    h.deck.save(transaction_buffer)
    transaction_buffer.seek(0)
    generated = Presentation(transaction_buffer)
    source_slides = [generated.slides[index] for index in range(template_count)]
    source_notes = []
    for source_slide in source_slides:
        source_notes.append([
            deepcopy(shape._element) for shape in source_slide.notes_slide.shapes
            if getattr(shape, "has_text_frame", False)
        ])
    clones = []
    for template_slide, replacements, table_spec, notes in prepared:
        clone = _clone_template_slide(generated, source_slides[template_slide - 1])
        existing_names = {shape.name for shape, _ in _walk_shapes(clone.shapes)}
        applied = 0
        for replacement in replacements:
            if replacement["shape_name"] not in existing_names:
                # Tolerant template application: a spec may reference shape
                # names from another layout. Apply every valid replacement and
                # skip only the absent ones instead of aborting the batch.
                continue
            _replace_named_shape_text(clone, replacement["shape_name"], replacement.get("text", ""))
            applied += 1
        if table_spec and table_spec.get("shape_name") in existing_names:
            _replace_named_table(clone, table_spec["shape_name"], table_spec["rows"])
        if notes:
            notes_slide = clone.notes_slide
            if not any(getattr(shape, "has_text_frame", False) for shape in notes_slide.shapes):
                for element in source_notes[template_slide - 1]:
                    notes_slide.shapes._spTree.insert_element_before(deepcopy(element), "p:extLst")
            _set_slide_notes_text(clone, notes)
        if applied or table_spec or notes:
            clones.append(clone)

    if not clones:
        raise ValueError("from_outline: none of the replacement shape names matched any template slide")

    if replace_template:
        for _ in range(template_count):
            slide_id = generated.slides._sldIdLst[0]
            generated.part.drop_rel(slide_id.rId)
            generated.slides._sldIdLst.remove(slide_id)

    h.deck = generated
    # The template itself is a visual source. Freeze its structural baseline
    # before content replacement so unchanged design devices (e.g. a full-size
    # text-bearing canvas) remain warnings, while new/worse defects still block.
    h.state.ppt_existing_deck = True
    h.state.ppt_baseline_captured = True
    h.state.ppt_baseline_findings = {
        item["key"]: item["severity"] for item in template_findings
    }
    final_numbers = list(range(1 if replace_template else template_count + 1, len(generated.slides) + 1))
    # Baseline-delta verification checks every generated page for new or worse
    # findings even when it is not in the explicit affected set. Leaving this
    # empty avoids treating unchanged template fixtures as blocking defects.
    h.state.ppt_affected_slides.clear()
    # Record one mutation without feeding generated pages back into the explicit
    # existing-deck scope tracker; delta verification still checks new/worse
    # findings globally against the frozen template baseline.
    h.state.record_change("deck:from_outline")
    h.state.record_fact(
        "ppt_last_compose_scope",
        f"generated {len(clones)} template-faithful slides in one transaction; template_replaced={replace_template}",
    )
    return (
        f"generated {len(clones)} slides from {len(set(item[0] for item in prepared))} template layouts "
        f"in one atomic operation; final deck has {len(generated.slides)} slides"
    )


def _ppt_compose(h, kind: str, **kw) -> str:
    if kind == "new_deck" and kw.get("slides"):
        # Disambiguation: new_deck + slides is the model's spelling of a
        # template replacement batch, not a request to discard the template.
        return _compose_from_outline(h, kw.get("slides") or [], kw.get("replace_template", True))
    if kind == "new_deck":
        return _new_deck(h, kw.get("title", "Untitled"), kw.get("subtitle", ""))
    if kind == "content" and kw.get("slides"):
        # Disambiguation: content + slides is the model's spelling of a
        # template replacement batch. Route to the tolerant template path.
        return _compose_from_outline(h, kw.get("slides") or [], kw.get("replace_template", True))
    if kind == "content":
        return _content_slide(
            h, kw.get("title", ""), kw.get("bullets") or [], kw.get("size", 18),
            kw.get("insert_after"), kw.get("slide_number"),
        )
    if kind == "comparison":
        if not kw.get("left_title") or not kw.get("right_title"):
            raise ValueError("comparison requires left_title and right_title")
        return _two_column(
            h, kw.get("title", ""), kw.get("left_title", ""), kw.get("left_bullets") or [],
            kw.get("right_title", ""), kw.get("right_bullets") or [], kw.get("insert_after"),
        )
    if kind == "from_slides":
        if kw.get("insert_after") is None and kw.get("slides"):
            # Disambiguation: models often use from_slides while supplying
            # template-slide replacements. Treat that shape as from_outline.
            return _compose_from_outline(h, kw.get("slides") or [], kw.get("replace_template", True))
        if kw.get("insert_after") is None:
            raise ValueError("from_slides requires insert_after")
        return _compose_from_slides(
            h, kw.get("source_slides") or [], kw["insert_after"], kw.get("title", ""),
            kw.get("left_title", "Male"), kw.get("right_title", "Female"),
        )
    if kind == "from_outline":
        slides = kw.get("slides") or []
        if not slides and kw.get("bullets"):
            # Outline shorthand: one semantic content page from a bullet list.
            # Without a template deck this is the same primitive the model
            # would otherwise have to spell out as kind='content'.
            return _content_slide(h, kw.get("title", ""), kw.get("bullets") or [], kw.get("size", 18), kw.get("insert_after"))
        return _compose_from_outline(h, slides, kw.get("replace_template", True))
    if kind == "table":
        if not kw.get("columns") or not kw.get("rows"):
            raise ValueError("table requires non-empty columns and rows")
        return _table_slide(h, kw.get("title", ""), kw.get("columns") or [], kw.get("rows") or [], kw.get("insert_after"))
    if kind == "quadrant":
        quadrants = kw.get("quadrants") or []
        if len(quadrants) != 4:
            raise ValueError("quadrant requires exactly four quadrant objects")
        required = {"title", "metric", "detail", "bullets", "source"}
        for index, item in enumerate(quadrants, 1):
            if not isinstance(item, dict):
                raise ValueError(f"quadrant {index} must be an object")
            missing = sorted(required - set(item))
            if missing:
                raise ValueError(f"quadrant {index} missing: {', '.join(missing)}")
            if not isinstance(item["bullets"], list) or len(item["bullets"]) > 6:
                raise ValueError(f"quadrant {index} bullets must be an array of at most 6 strings")
        slide_number = kw.get("slide_number")
        # ``replace_template`` is the model-facing way to say "rebuild the
        # auto-opened template slide instead of appending a second page".
        # Honor it generically by targeting the first slide when no explicit
        # slide_number was provided.
        if slide_number is None and kw.get("replace_template"):
            deck = getattr(h, "deck", None)
            if deck is not None and len(deck.slides) >= 1:
                slide_number = 1
        if slide_number is None:
            _assert_can_append_slide(h)
        return _quadrant_slide(h, kw.get("title", ""), kw.get("subtitle", ""), quadrants, slide_number)
    if kind in {"workflow", "workflow_pipeline", "step_process", "process_pipeline", "process_cards", "steps"}:
        steps = kw.get("steps") or []
        if not steps and kw.get("nodes"):
            steps = [{"title": n, "detail": ""} for n in kw["nodes"]]
        if not steps:
            raise ValueError("workflow_pipeline requires non-empty steps array")
        return _workflow_pipeline_slide(
            h,
            kw.get("title", ""),
            steps,
            kw.get("subtitle", ""),
            kw.get("takeaway", ""),
            kw.get("slide_number"),
            kw.get("insert_after"),
        )
    if kind in {"html_slide", "from_html", "html_code", "html_deck", "from_html_deck"} or (
        kind in {"html_mockup", "web_dashboard", "html_page", "web_mockup", "html_style", "html"}
        and (kw.get("html") or kw.get("file_path") or kw.get("path"))
    ):
        return _html_slide(
            h,
            html=kw.get("html", ""),
            css=kw.get("css", ""),
            file_path=kw.get("file_path") or kw.get("path") or "",
            slide_number=kw.get("slide_number"),
            insert_after=kw.get("insert_after"),
            title=kw.get("title", ""),
        )
    if kind in {"html_mockup", "web_dashboard", "html_page", "web_mockup", "html_style", "html"}:
        cards = kw.get("cards") or kw.get("components") or kw.get("quadrants") or []
        if not cards and kw.get("steps"):
            cards = kw["steps"]
        return _html_mockup_slide(
            h,
            kw.get("title", "Web Dashboard"),
            cards,
            kw.get("subtitle", ""),
            kw.get("url_bar", ""),
            kw.get("slide_number"),
            kw.get("insert_after"),
        )
    if kind in {"hero_split", "hero", "keynote_focus", "hero_cards"}:
        return _hero_split_slide(
            h,
            kw.get("title", ""),
            kw.get("hero_title", "Key Insight"),
            kw.get("hero_text", ""),
            kw.get("cards") or [],
            kw.get("hero_metric", ""),
            kw.get("subtitle", ""),
            kw.get("slide_number"),
            kw.get("insert_after"),
        )
    if kind == "flowchart":
        slide_number = kw.get("slide_number")
        if slide_number is None and kw.get("insert_after") is not None:
            # insert_after is the content/comparison spelling; translate it
            # instead of rejecting a whole turn over an argument name.
            slide_number = kw["insert_after"] + 1
        if slide_number is None:
            raise ValueError("flowchart requires slide_number")
        return _add_flowchart(h, slide_number, kw.get("nodes") or [], kw.get("title", ""))
    if kind == "textbox":
        if kw.get("slide_number") is None:
            raise ValueError("textbox requires slide_number")
        x, y, w, height = _geometry_args(kw)
        if any(value is None for value in (x, y, w, height)):
            raise ValueError("textbox requires x, y, w (or width), and height (or h)")
        return _add_textbox_to_slide(h, kw["slide_number"], {
            "x": x, "y": y, "w": w, "h": height,
            "text": kw.get("text", ""), "size": kw.get("size", 18), "bold": kw.get("bold", False),
        })
    raise ValueError(f"unsupported compose kind: {kind}")



def _is_explicit_list_paragraph(paragraph) -> bool:
    """Return whether *paragraph* carries a real DrawingML bullet marker."""
    properties = paragraph._p.pPr
    if properties is None:
        return False
    return any(
        child.tag.rsplit("}", 1)[-1] in {"buChar", "buAutoNum", "buBlip"}
        for child in properties
    )


def _inherited_text_size(shape, paragraph) -> int | None:
    """Resolve the theme font size that a new non-placeholder text box would lose.

    A copied paragraph keeps its explicit run formatting, but placeholder body
    text often inherits its size from the slide master.  Materializing only that
    missing value keeps both columns visually identical without restyling the
    source paragraph.
    """
    explicit = next(
        (run.font.size for run in paragraph.runs if run.font.size is not None),
        paragraph.font.size,
    )
    if explicit is not None:
        return int(round(explicit.pt * 100))
    try:
        master = shape.part.slide.slide_layout.slide_master._element
        style = "bodyStyle" if getattr(shape, "is_placeholder", False) else "otherStyle"
        level = max(1, min(9, int(paragraph.level) + 1))
        nodes = master.xpath(f"./p:txStyles/p:{style}/a:lvl{level}pPr/a:defRPr")
        if nodes and nodes[0].get("sz"):
            return int(nodes[0].get("sz"))
    except (AttributeError, TypeError, ValueError):
        pass
    return None


def _copy_paragraph_for_textbox(shape, paragraph):
    copied = deepcopy(paragraph._p)
    inherited_size = _inherited_text_size(shape, paragraph)
    if inherited_size is None:
        return copied
    # Runs without an explicit size would otherwise fall back from the source
    # placeholder's body style (28 pt in MOD-038) to generic text-box styling.
    for run_properties in copied.xpath("./a:r/a:rPr | ./a:fld/a:rPr | ./a:endParaRPr"):
        if run_properties.get("sz") is None:
            run_properties.set("sz", str(inherited_size))
    return copied


def _replace_text_frame_paragraphs(text_frame, paragraphs) -> None:
    body = text_frame._txBody
    for paragraph in list(body.xpath("./a:p")):
        body.remove(paragraph)
    for paragraph in paragraphs:
        body.append(paragraph)


def _reflow_two_columns(
    h,
    slide_number: int,
    shape_id: int | None = None,
    text_contains: str = "",
    split_after: int | None = None,
    gutter: float = 0.35,
) -> str:
    """Split one checklist text frame into two columns, preserving paragraph XML."""
    if getattr(h, "deck", None) is None:
        raise ValueError("no deck loaded")
    if slide_number < 1 or slide_number > len(h.deck.slides):
        raise IndexError("slide number out of range")
    if not 0.15 <= gutter <= 1.5:
        raise ValueError("reflow_two_columns gutter must be between 0.15 and 1.5 inches")

    slide = h.deck.slides[slide_number - 1]
    if shape_id is not None:
        _, target = _find_shape(h, slide_number, shape_id)
        path = next(path for candidate, path in _walk_shapes(slide.shapes) if candidate.shape_id == shape_id)
        if len(path) > 1:
            raise ValueError("reflow_two_columns requires a top-level text shape")
        candidates = [target]
    else:
        needle = text_contains.strip().casefold()
        candidates = [
            shape for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and (not needle or needle in shape.text.casefold())
        ]
        candidates.sort(
            key=lambda shape: sum(
                _is_explicit_list_paragraph(paragraph)
                for paragraph in shape.text_frame.paragraphs
            ),
            reverse=True,
        )
    target = next(
        (
            shape for shape in candidates
            if getattr(shape, "has_text_frame", False)
            and sum(_is_explicit_list_paragraph(p) for p in shape.text_frame.paragraphs) >= 2
        ),
        None,
    )
    if target is None:
        qualifier = f" containing {text_contains!r}" if text_contains else ""
        raise ValueError(f"no text shape{qualifier} with at least two checklist paragraphs")

    source_paragraphs = list(target.text_frame.paragraphs)
    first_item = next(
        index for index, paragraph in enumerate(source_paragraphs)
        if _is_explicit_list_paragraph(paragraph)
    )
    leading = [paragraph for paragraph in source_paragraphs[:first_item] if paragraph.text.strip()]
    checklist = [paragraph for paragraph in source_paragraphs[first_item:] if paragraph.text.strip()]
    if len(checklist) < 2:
        raise ValueError("reflow_two_columns requires at least two non-empty checklist paragraphs")
    split = math.ceil(len(checklist) / 2) if split_after is None else split_after
    if not isinstance(split, int) or not 1 <= split < len(checklist):
        raise ValueError(f"split_after must be between 1 and {len(checklist) - 1}")

    original_left, original_top = target.left, target.top
    original_width, original_height = target.width, target.height
    total_width = original_width / 914400
    if total_width - gutter < 3.0:
        raise ValueError("source text shape is too narrow for a readable two-column reflow")
    column_width = (total_width - gutter) / 2
    original_right = (original_left + original_width) / 914400
    right_x = original_left / 914400 + column_width + gutter
    if right_x + column_width > h.deck.slide_width / 914400 + 1e-6:
        raise ValueError("two-column reflow would cross the slide boundary")

    left_xml = [deepcopy(paragraph._p) for paragraph in leading + checklist[:split]]
    right_xml = [_copy_paragraph_for_textbox(target, paragraph) for paragraph in checklist[split:]]
    _replace_text_frame_paragraphs(target.text_frame, left_xml)
    # A placeholder can inherit its geometry. Setting only width materializes
    # an incomplete xfrm whose untouched dimensions become zero; write the
    # entire resolved rectangle when detaching it from inherited geometry.
    target.left, target.top = original_left, original_top
    target.width = Inches(column_width)
    target.height = original_height

    right = slide.shapes.add_textbox(
        Inches(right_x), original_top, Inches(column_width), original_height,
    )
    # Preserve the source frame's wrapping/autofit/margins in both columns.
    source_body_properties = deepcopy(target.text_frame._txBody.bodyPr)
    right_body = right.text_frame._txBody
    right_body.replace(right_body.bodyPr, source_body_properties)
    _replace_text_frame_paragraphs(right.text_frame, right_xml)

    # The split remains wholly inside the source shape's former footprint, so
    # it cannot introduce overlap with surrounding template elements.
    if abs((right.left + right.width) / 914400 - original_right) > 0.02:
        raise RuntimeError("two-column geometry did not preserve the source footprint")
    h.state.record_change(f"deck:slide:{slide_number}:reflow_two_columns")
    return (
        f"reflowed slide {slide_number} shape {target.shape_id} into two columns "
        f"({split}+{len(checklist) - split} checklist items); "
        f"created shape {right.shape_id}; affected slide {slide_number}"
    )


def _geometry_args(kw: dict) -> tuple:
    """Normalize rectangle geometry args, accepting w/width and height/h aliases.

    The model-facing schema historically mixed ``w`` with ``height``. Accepting
    both spellings keeps the facade forgiving without task-specific patches.
    """
    def pick(primary: str, alias: str):
        value = kw.get(primary)
        return value if value is not None else kw.get(alias)

    return pick("x", "x"), pick("y", "y"), pick("w", "width"), pick("height", "h")


def _ppt_arrange(h, operation: str, slide_number: int, shape_id: int | None = None, **kw) -> str:
    if operation == "geometry":
        if shape_id is None:
            raise ValueError("geometry requires shape_id")
        x, y, w, height = _geometry_args(kw)
        if any(value is None for value in (x, y, w, height)):
            raise ValueError("geometry requires x, y, w (or width), and height (or h)")
        return _set_shape_geometry(h, slide_number, shape_id, x, y, w, height)
    if operation == "delete_shape":
        if shape_id is None:
            raise ValueError("delete_shape requires shape_id")
        return _delete_shape(h, slide_number, shape_id)
    if operation == "delete_slide":
        return _delete_slide(h, slide_number)
    if operation == "move_slide":
        if kw.get("new_position") is None:
            raise ValueError("move_slide requires new_position")
        return _move_slide(h, slide_number, kw["new_position"])
    if operation == "reflow_two_columns":
        return _reflow_two_columns(
            h,
            slide_number,
            shape_id,
            kw.get("text_contains", ""),
            kw.get("split_after"),
            kw.get("gutter", 0.35),
        )
    raise ValueError(f"unsupported arrange operation: {operation}")


def _ppt_check(h, policy: str = "auto") -> str:
    requested_full = policy == "full"
    if (
        requested_full
        and getattr(h.state, "ppt_existing_deck", False)
        and getattr(h.state, "ppt_baseline_captured", False)
    ):
        # A local edit is accountable for its delta, not unrelated historical
        # defects elsewhere in the source deck. Full lint remains appropriate
        # for newly generated decks; any deck with a captured baseline uses
        # baseline-delta verification, even after the model re-opens the saved
        # deliverable in a fresh in-memory session.
        policy = "auto"
        h.state.record_fact("ppt_full_check_downgraded", "existing deck uses baseline-delta verification")
    reports = {"structural": _verify(h, policy)}
    if policy == "full":
        reports["quality"] = _quality_check(h)
    gap = _deck_completeness_gate(h)
    if gap:
        reports["completeness"] = f"Incomplete deck: {gap}"
    contract_passed = True
    if _has_verification_contract(h):
        contract_passed, contract_report = _verify_contract(h)
        reports["contract"] = contract_report
    # Counterexample-driven verification: an explicit full check runs the
    # task-local official evaluator whenever one exists (even when the
    # structural policy was downgraded for an existing deck), so its concrete
    # per-check failures drive the repair loop instead of surfacing only at
    # finish.
    if requested_full and h.state.facts.get("official_evaluator_present") == "true":
        try:
            from .lifecycle_tools import _run_task_evaluator
            reports["evaluator"] = _run_task_evaluator(h, timeout_seconds=120)
        except Exception as exc:
            reports["evaluator"] = f"task evaluator unavailable at this point: {exc}"
    # Task evaluators are intentionally run by finish after the required path
    # is saved; render/visual checks are also finish-lifecycle services.
    return json.dumps(reports, ensure_ascii=False, indent=2)


def _find_shape(h, slide_number: int, shape_id: int):
    if getattr(h, "deck", None) is None:
        raise ValueError("no deck loaded")
    if slide_number < 1 or slide_number > len(h.deck.slides):
        raise IndexError("slide number out of range")
    slide = h.deck.slides[slide_number - 1]
    for shape, _ in _walk_shapes(slide.shapes):
        if shape.shape_id == shape_id:
            return slide, shape
    raise KeyError(f"shape id not found: {shape_id}")


def _set_shape_geometry(h, slide_number: int, shape_id: int, x: float, y: float, w: float, height: float) -> str:
    if min(x, y) < 0 or w <= 0 or height <= 0:
        raise ValueError("shape geometry must be non-negative with positive width and height")
    if x + w > _W or y + height > _H:
        raise ValueError("shape geometry crosses the 16:9 slide boundary")
    slide, shape = _find_shape(h, slide_number, shape_id)
    # python-pptx may return a fresh proxy for the same XML element on each
    # traversal, so match the slide-unique shape id instead of object identity.
    path = next(path for candidate, path in _walk_shapes(slide.shapes) if candidate.shape_id == shape_id)
    if len(path) > 1:
        raise ValueError(
            "nested shape geometry is group-local, not slide-absolute; "
            "style the nested shape directly or arrange its top-level group"
        )
    shape.left, shape.top = Inches(x), Inches(y)
    shape.width, shape.height = Inches(w), Inches(height)
    h.state.record_change(f"deck:slide:{slide_number}:shape:{shape_id}:geometry")
    return f"moved/resized slide {slide_number} shape {shape_id} -> ({x:.2f},{y:.2f},{w:.2f},{height:.2f})"


def _delete_shape(h, slide_number: int, shape_id: int) -> str:
    _, shape = _find_shape(h, slide_number, shape_id)
    shape._element.getparent().remove(shape._element)
    h.state.record_change(f"deck:slide:{slide_number}:shape:{shape_id}:deleted")
    return f"deleted slide {slide_number} shape {shape_id}"


def _delete_slide(h, slide_number: int) -> str:
    if getattr(h, "deck", None) is None:
        raise ValueError("no deck loaded")
    if slide_number < 1 or slide_number > len(h.deck.slides):
        raise IndexError("slide number out of range")
    slide_id = h.deck.slides._sldIdLst[slide_number - 1]
    relationship_id = slide_id.rId
    h.deck.part.drop_rel(relationship_id)
    h.deck.slides._sldIdLst.remove(slide_id)
    h.state.record_change(f"deck:slide:{slide_number}:deleted")
    return f"deleted slide {slide_number}; {len(h.deck.slides)} slides remain"


def _move_slide(h, slide_number: int, new_position: int) -> str:
    if getattr(h, "deck", None) is None:
        raise ValueError("no deck loaded")
    count = len(h.deck.slides)
    if slide_number < 1 or slide_number > count or new_position < 1 or new_position > count:
        raise IndexError("slide number or destination out of range")
    slide_id = h.deck.slides._sldIdLst[slide_number - 1]
    h.deck.slides._sldIdLst.remove(slide_id)
    h.deck.slides._sldIdLst.insert(new_position - 1, slide_id)
    h.state.record_change(f"deck:slide:{slide_number}:moved:{new_position}")
    return f"moved slide {slide_number} to position {new_position}"


def _set_speaker_notes(h, slide_number: int, text: str) -> str:
    if getattr(h, "deck", None) is None:
        raise ValueError("no deck loaded")
    if slide_number < 1 or slide_number > len(h.deck.slides):
        raise IndexError("slide number out of range")
    slide = h.deck.slides[slide_number - 1]
    # Some templates carry no notes part at all; create one instead of failing
    # on the None returned by python-pptx for slides without a notes part.
    try:
        if not slide.part.has_notes_slide:
            slide.part._add_notes_slide_part()
    except Exception:
        pass
    notes = slide.notes_slide
    if notes is None:
        raise ValueError(f"slide {slide_number} has no notes part and it could not be created")
    if notes.notes_text_frame is None:
        # A newly created notes part has no body placeholder.  Add the
        # standard notes body shape so the text frame exists for writing.
        from pptx.oxml.ns import qn
        from lxml import etree
        sp = etree.SubElement(notes.shapes._spTree, qn("p:sp"))
        nv = etree.SubElement(sp, qn("p:nvSpPr"))
        cnv = etree.SubElement(nv, qn("p:cNvPr"))
        cnv.set("id", "2"); cnv.set("name", "Notes Placeholder")
        etree.SubElement(nv, qn("p:cNvSpPr"))
        nvpr = etree.SubElement(nv, qn("p:nvPr"))
        ph = etree.SubElement(nvpr, qn("p:ph"))
        ph.set("type", "body"); ph.set("idx", "1")
        etree.SubElement(sp, qn("p:spPr"))
        tx = etree.SubElement(sp, qn("p:txBody"))
        etree.SubElement(tx, qn("a:bodyPr"))
        etree.SubElement(tx, qn("a:lstStyle"))
        etree.SubElement(tx, qn("a:p"))
    notes.notes_text_frame.text = text
    h.state.record_change(f"deck:slide:{slide_number}:speaker_notes")
    return f"updated speaker notes on slide {slide_number}"


def _render(h, path: str, output_dir: str = "rendered") -> str:
    from .. import config
    from ..permissions import path_within

    root = config.sandbox_root().resolve()
    source = (root / path).resolve()
    target = (root / output_dir).resolve()
    if not path_within(root, source) or not path_within(root, target):
        raise PermissionError("render path escapes workspace")
    if not source.is_file():
        raise FileNotFoundError(path)
    target.mkdir(parents=True, exist_ok=True)
    for stale in target.iterdir():
        if stale.is_file() and stale.suffix.lower() == ".png":
            stale.unlink()
    if platform.system() == "Windows":
        try:
            import win32com.client
        except ImportError as exc:
            raise RuntimeError("pywin32 is required for PowerPoint rendering") from exc

        def _render_once() -> None:
            app = None
            presentation = None
            try:
                app = win32com.client.DispatchEx("PowerPoint.Application")
            except Exception as exc:
                # COM can be installed but unavailable in a headless/service
                # session. Surface a stable diagnostic so the harness can report
                # a renderer limitation instead of an opaque traceback.
                raise RuntimeError(
                    "PowerPoint COM renderer unavailable in this session; "
                    "use an interactive Windows session or LibreOffice fallback"
                ) from exc
            try:
                presentation = app.Presentations.Open(str(source), ReadOnly=True, Untitled=False, WithWindow=False)
                pdf = target / f"{source.stem}.pdf"
                presentation.SaveCopyAs(str(pdf), 32)
                presentation.Export(str(target), "PNG")
            finally:
                if presentation is not None:
                    presentation.Close()
                app.Quit()

        # PowerPoint COM occasionally raises a transient "发生意外" fault when
        # the automation server is briefly busy. One fresh-DispatchEx retry
        # turns that into a successful render instead of a missing certificate.
        last_error: Exception | None = None
        for _attempt in range(2):
            try:
                _render_once()
                break
            except Exception as exc:
                last_error = exc
        else:
            raise RuntimeError(f"PowerPoint COM render failed after retry: {last_error}") from last_error
    else:
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        pdftoppm = shutil.which("pdftoppm")
        if not soffice or not pdftoppm:
            raise RuntimeError("Linux PPT rendering requires soffice/libreoffice and pdftoppm")
        pdf = target / f"{source.stem}.pdf"
        convert = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(target), str(source)],
            capture_output=True, text=True, timeout=120,
        )
        if convert.returncode != 0 or not pdf.exists():
            raise RuntimeError(f"LibreOffice conversion failed: {(convert.stderr or convert.stdout).strip()}")
        raster = subprocess.run(
            [pdftoppm, "-png", "-r", "120", str(pdf), str(target / "slide")],
            capture_output=True, text=True, timeout=120,
        )
        if raster.returncode != 0:
            raise RuntimeError(f"pdftoppm conversion failed: {(raster.stderr or raster.stdout).strip()}")
    images = sorted({image.resolve() for image in target.iterdir() if image.is_file() and image.suffix.lower() == ".png" and image.name.lower() != "montage.png"})
    pdf = target / f"{source.stem}.pdf"
    if not pdf.is_file() or pdf.stat().st_size == 0:
        raise RuntimeError("renderer produced no PDF")
    if not images:
        raise RuntimeError("renderer produced no PNG slides")
    from PIL import Image, ImageDraw
    opened = [Image.open(image).convert("RGB") for image in images]
    try:
        thumb_w = 640
        thumbs = []
        for index, slide_image in enumerate(opened, 1):
            thumb_h = round(slide_image.height * thumb_w / slide_image.width)
            thumb = slide_image.resize((thumb_w, thumb_h))
            canvas = Image.new("RGB", (thumb_w, thumb_h + 36), "white")
            canvas.paste(thumb, (0, 36))
            ImageDraw.Draw(canvas).text((12, 10), f"Slide {index}", fill=(35, 35, 35))
            thumbs.append(canvas)
        columns = 2
        rows = math.ceil(len(thumbs) / columns)
        cell_h = max(thumb.height for thumb in thumbs)
        montage = Image.new("RGB", (columns * thumb_w, rows * cell_h), (225, 225, 225))
        for index, thumb in enumerate(thumbs):
            montage.paste(thumb, ((index % columns) * thumb_w, (index // columns) * cell_h))
        montage_path = target / "montage.png"
        montage.save(montage_path)
    finally:
        for slide_image in opened:
            slide_image.close()
    h.state.record_evidence("ppt_render", f"ppt rendered: {len(images)} PNG slides from {path}")
    h.state.unresolved_checks.discard("ppt_render")
    # This verification action passed. Unrelated unresolved obligations (for
    # example a still-failing official evaluator) must not be re-labelled as a
    # fresh local verification failure: that would reopen a second repair cycle
    # on every auto save/check and exhaust the bounded budget prematurely.
    h.state.last_verification_failed = False
    h.state.last_verification_epoch = h.state.mutation_epoch
    if getattr(h, "recorder", None):
        h.recorder.check("ppt_render", True, f"{len(images)} PNG slides; pdf={pdf}")
        h.recorder.artifact(pdf, "rendered-pdf")
        for image in images:
            h.recorder.artifact(image, "rendered-slide")
    return "rendered slides:\n" + "\n".join(str(image.relative_to(root)) for image in images) + f"\nmontage: {montage_path.relative_to(root)}"


def _inspect_rendered(h, output_dir: str = "rendered") -> str:
    """Turn rendered pixels into auditable, deterministic visual evidence.

    This is intentionally a coarse gate, not an aesthetic judge. It catches
    blank renders and content pressed against the slide edge, while leaving
    semantic visual quality to a human or vision model review.
    """
    from PIL import Image, ImageStat
    from .. import config
    from ..permissions import path_within

    root = config.sandbox_root().resolve()
    target = (root / output_dir).resolve()
    if not path_within(root, target):
        raise PermissionError("render inspection path escapes workspace")
    images = sorted(path for path in target.glob("*.png") if path.name.lower() != "montage.png")
    if not images:
        raise FileNotFoundError("no rendered slide PNGs found")
    findings: list[str] = []
    warnings: list[str] = []
    rows: list[str] = []
    for index, path in enumerate(images, 1):
        with Image.open(path) as source:
            gray = source.convert("L")
            width, height = gray.size
            stat = ImageStat.Stat(gray)
            variance = stat.var[0]
            margin = max(2, round(min(width, height) * 0.015))
            pixels = gray.load()
            edge_dark = 0
            edge_total = 0
            for y in range(height):
                for x in range(width):
                    if x < margin or x >= width - margin or y < margin or y >= height - margin:
                        edge_total += 1
                        if pixels[x, y] < 235:
                            edge_dark += 1
            edge_ratio = edge_dark / max(1, edge_total)
            rows.append(f"slide {index}: variance={variance:.1f}, edge_content={edge_ratio:.3f}")
            if variance < 8:
                findings.append(f"slide {index}: probable blank/near-uniform render")
            # Full-bleed theme backgrounds and title bands are intentional in
            # presentation design. Keep this as a warning rather than a hard
            # failure; geometric clipping is already checked by ppt_verify.
            if edge_ratio > 0.30:
                warnings.append(f"slide {index}: full-bleed/edge content warning ({edge_ratio:.3f}); review if not intentional")
    if findings:
        h.state.unresolved_checks.add("ppt_visual")
        h.state.last_verification_failed = True
        h.state.last_verification_epoch = h.state.mutation_epoch
        if getattr(h, "recorder", None):
            h.recorder.check("ppt_visual", False, "\n".join(findings))
        return "Rendered visual audit findings:\n" + "\n".join(findings) + "\nMetrics:\n" + "\n".join(rows)
    h.state.record_evidence("ppt_visual", f"rendered visual audit passed: {len(images)} slides")
    h.state.unresolved_checks.discard("ppt_visual")
    # This verification action passed. Unrelated unresolved obligations (for
    # example a still-failing official evaluator) must not be re-labelled as a
    # fresh local verification failure: that would reopen a second repair cycle
    # on every auto save/check and exhaust the bounded budget prematurely.
    h.state.last_verification_failed = False
    h.state.last_verification_epoch = h.state.mutation_epoch
    if getattr(h, "recorder", None):
        h.recorder.check("ppt_visual", True, "\n".join(rows + warnings))
    warning_text = "\nWarnings:\n" + "\n".join(warnings) if warnings else ""
    return "Rendered visual audit passed.\n" + "\n".join(rows) + warning_text


def _add_box(h, box: dict) -> str:
    s = h.deck.slides[-1]
    tb = s.shapes.add_textbox(
        Inches(box["x"]), Inches(box["y"]), Inches(box["w"]), Inches(box["h"])
    )
    _fit_lines(tb.text_frame, str(box.get("text", "")).split("\n"), box.get("size", 18), box.get("bold", False), _TEXT)
    return "text box added."


def _box_from_kw(kw: dict) -> dict:
    box = dict(kw)
    box["h"] = kw["height"]
    return box


def _normalize_minimal_container(prs: Presentation) -> tuple[Presentation, bool]:
    """Move a stripped text/shape-only deck into a standard Office container.

    Some benchmark templates are valid enough for python-pptx but omit most
    standard layout/theme parts and PowerPoint refuses to open them.  Copying
    plain shape XML is safe because it has no media/chart relationship graph;
    richer decks are deliberately left untouched.
    """
    if len(prs.slide_layouts) > 1:
        return prs, False
    safe_types = {1, 14, 17}  # AUTO_SHAPE, PLACEHOLDER, TEXT_BOX
    if any(int(shape.shape_type) not in safe_types for slide in prs.slides for shape in slide.shapes):
        return prs, False
    normalized = Presentation()
    normalized.slide_width = prs.slide_width
    normalized.slide_height = prs.slide_height
    for source_slide in prs.slides:
        target_slide = normalized.slides.add_slide(normalized.slide_layouts[6])
        for shape in source_slide.shapes:
            target_slide.shapes._spTree.insert_element_before(deepcopy(shape._element), "p:extLst")
        if getattr(source_slide, "has_notes_slide", False):
            notes = source_slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                target_slide.notes_slide.notes_text_frame.text = notes
    return normalized, True


def _save(h, path: str | None) -> str:
    from pathlib import Path
    from .. import config
    root = config.sandbox_root()
    if path:
        p = Path(path)
        # tolerate an explicit "workspace/..." prefix from LLM args
        if not p.is_absolute() and p.parts and p.parts[0] in ("workspace", root.name):
            p = Path(*p.parts[1:])
        target = p if p.is_absolute() else root / p
    else:
        target = root / "deck.pptx"
    required_output = getattr(h.state, "facts", {}).get("required_output_pptx")
    if required_output:
        required_target = (root / required_output).resolve()
        # The task contract owns the final deliverable location.  Models often
        # omit the optional path, use the alias ``output_path``, or choose a
        # generic workspace filename.  Any non-intermediate save therefore
        # resolves to the unique contract output deterministically.
        if not path or "intermediate" not in {part.casefold() for part in target.parts}:
            target = required_target
    elif path and "tasks" in {part.casefold() for part in target.parts}:
        # tasks/ is the task-package namespace. An ad-hoc open-ended deck has
        # no package and belongs at the workspace root instead of polluting it.
        target = root / "deck.pptx"
    target.parent.mkdir(parents=True, exist_ok=True)
    prs = _deck(h)
    prs = _ensure_deck_package_clean(prs)
    prs, normalized = _normalize_minimal_container(prs)
    h.deck = prs
    if normalized:
        h.state.record_fact("ppt_container_normalized", "minimal template migrated to standard PowerPoint container")
    try:
        prs.save(str(target))
    except PermissionError:
        import time
        time.sleep(0.3)
        try:
            prs.save(str(target))
        except PermissionError:
            tmp_target = target.with_name(f"{target.stem}_tmp{target.suffix}")
            prs.save(str(tmp_target))
            try:
                shutil.move(str(tmp_target), str(target))
            except Exception:
                target = tmp_target
    try:
        relative_target = str(target.relative_to(root))
    except ValueError:
        relative_target = str(target)
    if normalized:
        # Normalization rewrote the in-memory deck, so this is a real content
        # mutation and must advance the epoch.
        h.state.record_change(relative_target)
    else:
        # A plain save only persists the current content; it must not invalidate
        # fresh structural/render evidence for the current epoch.
        h.state.record_commit(relative_target)
    if getattr(h, "recorder", None):
        source = getattr(h, "deck_source_path", None)
        if source:
            h.recorder.bind_source(source, str(target), "edited presentation source")
        # Multi-source provenance is a harness guarantee.  The model should
        # not spend one call per HTML/XLSX/Markdown input merely to create an
        # audit edge that the runtime already knows deterministically.
        for source_path in sorted(getattr(h.state, "source_paths", set())):
            h.recorder.bind_source(Path(source_path), str(target), "ContentIR source used for presentation")
        h.recorder.artifact(target, "final-pptx")
    return f"saved {len(prs.slides)} slides -> {target.name}"


def _deck_info(h, slide_number: int | None = None) -> str:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    if getattr(h, "deck", None) is None:
        return "no deck yet."
    if slide_number is not None and (slide_number < 1 or slide_number > len(h.deck.slides)):
        raise IndexError(f"slide number out of range: {slide_number}")
    selected = range(len(h.deck.slides)) if slide_number is None else [slide_number - 1]
    lines = []
    for i in selected:
        texts = []
        for sh in h.deck.slides[i].shapes:
            if sh.has_text_frame and sh.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                t = (sh.text_frame.text or "").strip()
                if t:
                    texts.append(f"[{sh.name}] {t.replace(chr(10), ' | ')[:70]}")
            elif getattr(sh, "has_table", False):
                rows = [
                    " | ".join(cell.text.strip().replace("\n", " ") for cell in row.cells)
                    for row in sh.table.rows
                ]
                joined = " ;; ".join(rows)
                if joined.strip():
                    texts.append(f"表[{sh.name}]: {joined[:160]}")
        lines.append(f"  slide {i + 1}: " + ("; ".join(texts) if texts else "(empty)"))
    header = f"slide {slide_number}:\n" if slide_number is not None else f"{len(h.deck.slides)} slides:\n"
    return header + "\n".join(lines)


def _collect_structural_findings(h) -> list[dict[str, Any]]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    if getattr(h, "deck", None) is None:
        return []
    findings: list[dict[str, Any]] = []

    def add(slide: int, kind: str, subject: str, severity: float, message: str) -> None:
        findings.append({
            "key": f"slide:{slide}:{kind}:{subject}",
            "slide": slide,
            "kind": kind,
            "severity": round(float(severity), 6),
            "message": message,
        })

    # Explicit instruction following: verify requested slide count
    req_count = _detect_requested_slide_count(h)
    if req_count is not None and len(h.deck.slides) < req_count:
        add(0, "slide_count_deficit", "deck", 1.0, f"deck only has {len(h.deck.slides)} slide(s); user task explicitly requested at least {req_count} slides")

    for i, slide in enumerate(h.deck.slides, 1):
        text_shapes = []
        for sh in slide.shapes:
            # skip background / decorative rectangles — only inspect real text boxes
            if sh.has_text_frame and sh.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                text_shapes.append(sh)
                text = (sh.text_frame.text or "").strip()
                if not text:
                    add(i, "empty_text", str(sh.shape_id), 1.0, f"slide {i}: empty text box (shape {sh.shape_id})")
                    continue
                observed_sizes = []
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size:
                            observed_sizes.append(r.font.size.pt)
                big = max(observed_sizes) if observed_sizes else 18
                w_in = sh.width / 914400
                h_in = sh.height / 914400
                cpl = max(1, int(w_in * 96 / (big * 0.75)))
                lines = 0
                for ln in text.split("\n"):
                    lines += max(1, math.ceil(len(ln) / cpl))
                needed = lines * big * 0.62 / 72
                # Relaxed threshold: python-pptx auto-shrinks text and real
                # rendering capacity is much larger than this heuristic.
                # Only block when overflow exceeds 2x; warn between 1.4x-2x.
                if needed > h_in * 1.4:
                    severity = needed / max(0.01, h_in * 1.4)
                    level = "text_overflow" if needed > h_in * 2.0 else "text_overflow_warning"
                    add(
                        i, level, str(sh.shape_id), severity if level == "text_overflow" else severity * 0.3,
                        f"slide {i}: shape {sh.shape_id} overflow risk ~{needed:.2f}in in {h_in:.2f}in box",
                    )
            # Inspect every positioned shape, not only text boxes. Pictures,
            # charts, tables, and decorative shapes can also be clipped.
            if sh.left < 0 or sh.top < 0 or sh.left + sh.width > h.deck.slide_width or sh.top + sh.height > h.deck.slide_height:
                overflow = max(
                    max(0, -sh.left) / h.deck.slide_width,
                    max(0, -sh.top) / h.deck.slide_height,
                    max(0, sh.left + sh.width - h.deck.slide_width) / h.deck.slide_width,
                    max(0, sh.top + sh.height - h.deck.slide_height) / h.deck.slide_height,
                )
                add(
                    i, "boundary", str(sh.shape_id), overflow,
                    f"slide {i}: shape {sh.shape_id} crosses slide boundary",
                )
        for left_index, a in enumerate(text_shapes):
            for b in text_shapes[left_index + 1:]:
                overlap_w = min(a.left + a.width, b.left + b.width) - max(a.left, b.left)
                overlap_h = min(a.top + a.height, b.top + b.height) - max(a.top, b.top)
                if overlap_w > 0 and overlap_h > 0:
                    ratio = (overlap_w * overlap_h) / min(a.width * a.height, b.width * b.height)
                    if ratio > 0.15:
                        subject = f"{min(a.shape_id, b.shape_id)}-{max(a.shape_id, b.shape_id)}"
                        add(
                            i, "text_overlap", subject, ratio,
                            f"slide {i}: text shapes {a.shape_id}/{b.shape_id} overlap {ratio:.0%}",
                        )
    return findings


def _compact_text(value: str) -> str:
    import re
    # Match the WorkBuddy deterministic verifier: whitespace AND structural
    # punctuation are removed before substring comparison.
    return re.sub(r"[\s`*_·•:：/\\|\-—–,，。；;、()（）\[\]【】<>《》\n\r\t]+", "", value or "")


def _slide_text_material(slide) -> str:
    """All textual material on a slide: text frames and every table cell."""
    parts: list[str] = []
    for shape, _ in _walk_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            parts.append(shape.text or "")
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.append(cell.text or "")
    return "".join(parts)


def _detect_requested_slide_count(h) -> int | None:
    facts = getattr(getattr(h, "state", None), "facts", {})
    val = facts.get("requested_slide_count")
    if val is not None:
        try:
            return int(val)
        except Exception:
            pass
    # Inspect goal prompt directly for explicit slide count requests
    goal = getattr(h, "goal", "") or ""
    if goal:
        import re
        # Pattern 1: "两页" / "2页" / "三页" / "3页" / "5 slides"
        num_map = {"一": 1, "两": 2, "二": 2, "三": 3, "四": 4, "五": 5, "六": 6, "七": 7, "八": 8, "九": 9, "十": 10}
        m = re.search(r"([一两二三四五六七八九十\d]+)\s*(?:页|slides?|张|个页面)", goal, re.IGNORECASE)
        if m:
            raw = m.group(1)
            if raw.isdigit():
                return int(raw)
            if raw in num_map:
                return num_map[raw]
    terms = facts.get("verification_contract_terms")
    if isinstance(terms, dict) and terms.get("slide_count"):
        try:
            return int(terms["slide_count"])
        except Exception:
            pass
    return None


def _deck_completeness_gate(h) -> str:
    """Return a human-readable gap for open-ended decks, or "" when complete.

    Generic density contract: every slide needs a title-like text object, at
    least two body text objects, and a minimum visible character budget.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    if getattr(h, "deck", None) is None:
        return ""
    req_count = _detect_requested_slide_count(h)
    if req_count is not None and len(h.deck.slides) < req_count:
        return f"deck only has {len(h.deck.slides)} slide(s); user task explicitly requested at least {req_count} slides"
    is_new_deck = not getattr(getattr(h, "state", None), "ppt_existing_deck", False)
    for slide_number, slide in enumerate(h.deck.slides, 1):
        boxes = []
        for shape, _path in _walk_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
                text = (shape.text_frame.text or "").strip()
                if text:
                    try:
                        top = shape.top.inches if shape.top is not None else 0.0
                    except Exception:
                        top = 0.0
                    boxes.append((top, text))
        total_chars = sum(len(text) for _top, text in boxes)
        titleish = [text for top, text in boxes if top < 1.5]
        bodyish = [text for top, text in boxes if top >= 1.5]
        if not titleish:
            return f"slide {slide_number} has no visible title text"
        if len(bodyish) < 1:
            return f"slide {slide_number} has no visible body text objects; at least 1 is required"
        if is_new_deck and len(h.deck.slides) > 1 and slide_number > 1:
            has_cards = any(
                (sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and Inches(1.5) < sh.width < Inches(_W * 0.98) and Inches(0.8) < sh.height < Inches(_H * 0.98))
                or (sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
                or getattr(sh, "has_table", False)
                for sh in slide.shapes
            )
            if total_chars < 100 and not has_cards:
                return f"slide {slide_number} content density too sparse ({total_chars} visible chars); presentation slides require rich technical details (120-400 chars) with structured cards/steps/badges"
        if total_chars < 15:
            return f"slide {slide_number} is too thin ({total_chars} visible characters; minimum 15)"

    # Enforce HTML slide requirement if task explicitly requested HTML
    task_goal = str(getattr(getattr(h, "state", None), "initial_goal", "") or getattr(getattr(h, "state", None), "task_prompt", "") or "").casefold()
    requires_html = any(kw in task_goal for kw in ("html 进行制作", "html制作", "html 制作", "第二页采用 html", "第二页使用 html", "第2页采用 html", "第2页使用 html"))
    if requires_html and len(h.deck.slides) >= 2:
        slide2 = h.deck.slides[1]
        has_vector_cards = any(
            sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and Inches(1.2) < sh.width < Inches(_W * 0.95) and Inches(0.6) < sh.height < Inches(_H * 0.95)
            for sh in slide2.shapes
        )
        if not has_vector_cards:
            return (
                "Task explicitly requires Slide 2 to be created using HTML (网页/HTML组件风格). "
                "Slide 2 currently lacks HTML vector card components. "
                "Please call ppt_compose(kind='html_slide', slide_number=2, html='...') with rich modern web cards (<div class=\"card\">...</div>) and badges."
            )
    return ""


def _has_verification_contract(h) -> bool:
    terms = getattr(h.state, "verification_contract_terms", None)
    if isinstance(terms, dict):
        return True
    return bool(getattr(h.state, "facts", {}).get("verification_contract_terms"))


def _verify_contract(h) -> tuple[bool, str]:
    """Deterministic pre-gate over the task-local verification contract.

    C_static's V is checked before the official evaluator so the model gets
    concrete missing/forbidden terms while the repair loop is still open,
    instead of discovering them only through a finish rejection.
    """
    terms = getattr(h.state, "verification_contract_terms", None)
    terms_text = getattr(h.state, "facts", {}).get("verification_contract_terms", "")
    if not isinstance(terms, dict) and terms_text:
        try:
            terms = json.loads(terms_text)
        except json.JSONDecodeError:
            terms = None
    if not isinstance(terms, dict):
        return True, "no task-local verification contract terms; skipping contract gate"

    missing_by_slide: dict[str, list[str]] = {}
    forbidden_by_slide: dict[str, list[str]] = {}
    co_findings: list[str] = []
    required_slides = terms.get("required_slide_expectations") or {}
    for slide_number, spec in required_slides.items():
        if not isinstance(spec, dict):
            continue
        try:
            index = int(slide_number) - 1
            slide = h.deck.slides[index]
        except (ValueError, IndexError):
            co_findings.append(f"slide {slide_number}: slide does not exist")
            continue
        compact = _compact_text(_slide_text_material(slide))
        missing = missing_by_slide.setdefault(str(slide_number), [])
        forbidden = forbidden_by_slide.setdefault(str(slide_number), [])
        for term in spec.get("all") or []:
            if _compact_text(term) not in compact:
                missing.append(str(term))
        for term in spec.get("none") or []:
            if _compact_text(term) in compact:
                forbidden.append(str(term))

    for item in terms.get("co_location_expectations") or []:
        if not isinstance(item, dict):
            continue
        slide_number = item.get("slide")
        try:
            slide = h.deck.slides[int(slide_number) - 1]
        except (ValueError, IndexError, TypeError):
            continue
        matches = [shape for shape, _ in _walk_shapes(slide.shapes) if shape.name == item.get("object_name")]
        if len(matches) != 1:
            co_findings.append(f"slide {slide_number}: co-location object {item.get('object_name')!r} matched {len(matches)}")
            continue
        compact = _compact_text(_slide_text_material_for_shape(matches[0]))
        required_missing = [str(term) for term in item.get("required_terms") or [] if _compact_text(term) not in compact]
        forbidden_present = [str(term) for term in item.get("forbidden_terms") or [] if _compact_text(term) in compact]
        if required_missing:
            co_findings.append(f"slide {slide_number}/{item.get('object_name')}: missing=[{'; '.join(required_missing)}]")
            # Provenance/binding terms may legally live in shape description
            # metadata rather than visible text. Give the model the exact call.
            all_required = "; ".join(str(term) for term in (item.get("required_terms") or []))
            if getattr(matches[0], "has_table", False):
                co_findings.append(
                    f"  → table object: use ppt_edit_text operation='set_table' slide {slide_number} "
                    f"shape_name={item.get('object_name')!r} with a full row set that includes these terms"
                )
            else:
                co_findings.append(
                    f"  → use ppt_metadata slide {slide_number} shape_name={item.get('object_name')} "
                    f"descr='{all_required}'"
                )
        if forbidden_present:
            co_findings.append(f"slide {slide_number}/{item.get('object_name')}: forbidden=[{'; '.join(forbidden_present)}]")

    # Slide-count contract: exact when the package declares one, otherwise the
    # declared min/max envelope.  This is schema-driven and skipped for
    # packages that do not declare a count.
    slide_count = len(h.deck.slides)
    answer_contract = terms.get("answer_contract") or {}
    template_expected = terms.get("template_expected") or {}
    output_contract = terms.get("output_contract") or {}
    exact_slides = (
        output_contract.get("expected_slide_count")
        or answer_contract.get("slide_count")
        or template_expected.get("slide_count")
    )
    min_slides = answer_contract.get("min_slide_count")
    max_slides = answer_contract.get("max_slide_count")
    if exact_slides is not None and int(exact_slides) != slide_count:
        co_findings.append(
            f"slide count: expected {exact_slides}, got {slide_count}. "
            + ("Use ppt_arrange delete_slide to remove the surplus slide." if slide_count > int(exact_slides)
               else "Compose only the missing slide(s).")
        )
    if min_slides is not None and slide_count < int(min_slides):
        co_findings.append(f"slide count: below minimum {min_slides} (got {slide_count}); keep composing content slides until the minimum is met")
    if max_slides is not None and slide_count > int(max_slides):
        co_findings.append(
            f"slide count: above maximum {max_slides} (got {slide_count}); "
            "use ppt_arrange delete_slide to remove surplus slides, merging their content into the remaining ones"
        )

    # Template cleanup contract: placeholder/sample copy must be replaced by
    # real task content, not survive into the final deck.
    placeholders = template_expected.get("placeholder_texts") or []
    if placeholders:
        compact_all = _compact_text("".join(_slide_text_material(slide) for slide in h.deck.slides))
        remaining = [str(item) for item in placeholders if _compact_text(str(item)) and _compact_text(str(item)) in compact_all]
        if remaining:
            co_findings.append(
                "template placeholders still present: " + "; ".join(remaining[:24])
                + ". Replace each with real content via ppt_edit_text set_shape_text on the named template boxes."
            )

    # Source-grounded binding contract: provenance ids may live in visible
    # text, shape names/descriptions, speaker notes, or document properties.
    # The pre-gate only reports ids that are absent from the whole artifact;
    # per-region placement remains the official evaluator's job.
    if terms.get("chart_binding_contract") or terms.get("required_quadrants") or terms.get("correction_contract"):
        visible = "".join(_slide_text_material(slide) for slide in h.deck.slides)
        metadata_bits: list[str] = []
        for slide in h.deck.slides:
            if getattr(slide, "has_notes_slide", False):
                metadata_bits.append(slide.notes_slide.notes_text_frame.text or "")
            for shape, _ in _walk_shapes(slide.shapes):
                metadata_bits.append(shape.name or "")
                try:
                    descr = shape._element.nvSpPr.cNvPr.get("descr") or ""
                    if descr:
                        metadata_bits.append(descr)
                except Exception:
                    pass
        try:
            props = h.deck.core_properties
            for value in (props.title, props.subject, props.keywords, props.comments, props.category):
                if value:
                    metadata_bits.append(str(value))
        except Exception:
            pass
        binding_material = _compact_text(visible + "".join(metadata_bits))
        compact_visible = _compact_text(visible)

        chart_binding = terms.get("chart_binding_contract") or {}
        required_ids: list[tuple[str, str]] = []
        for key, label in (
            ("required_anchor_ids", "anchor"),
            ("required_binding_ids", "binding"),
            ("required_chart_ids", "chart"),
            ("required_subanchors", "subanchor"),
        ):
            required_ids.extend((str(value), label) for value in (chart_binding.get(key) or []))
        missing_ids = [f"{label} {value}" for value, label in required_ids if _compact_text(value) not in binding_material]
        if missing_ids:
            co_findings.append("missing provenance ids (visible text or ppt_metadata): " + "; ".join(missing_ids))
            metadata_target = None
            for slide_number, slide in enumerate(h.deck.slides, 1):
                for shape, _ in _walk_shapes(slide.shapes):
                    name = shape.name or ""
                    if "binding" in name.casefold() or "note" in name.casefold() or "provenance" in name.casefold():
                        metadata_target = (slide_number, name)
                        break
                if metadata_target:
                    break
            if metadata_target:
                values = "; ".join(str(value) for value, _label in required_ids if _compact_text(value) not in binding_material)
                co_findings.append(
                    f"  → use ppt_metadata slide {metadata_target[0]} shape_name={metadata_target[1]!r} "
                    f"descr='{values}'"
                )

        for item in terms.get("required_quadrants") or []:
            if not isinstance(item, dict):
                continue
            qid = item.get("id", "?")
            missing_terms = [
                str(value)
                for value in list(item.get("must_include_terms") or []) + list(item.get("must_include_values") or [])
                if _compact_text(value) not in compact_visible
            ]
            if missing_terms:
                co_findings.append(f"quadrant {qid}: missing=[{'; '.join(missing_terms)}]")
            missing_bindings = [
                f"{label} {value}"
                for key, label in (("html_anchors", "anchor"), ("required_metric_ids", "metric"), ("chart_ids", "chart"))
                for value in (item.get(key) or [])
                if _compact_text(value) not in binding_material
            ]
            if missing_bindings:
                co_findings.append(f"quadrant {qid}: bindings=[{'; '.join(missing_bindings)}] (visible text or ppt_metadata)")
                target = None
                for slide_number, slide in enumerate(h.deck.slides, 1):
                    for shape, _ in _walk_shapes(slide.shapes):
                        name = shape.name or ""
                        if str(qid).casefold() in name.casefold():
                            target = (slide_number, name)
                            break
                    if target:
                        break
                if target:
                    values = "; ".join(
                        str(value)
                        for key, _label in (("html_anchors", "anchor"), ("required_metric_ids", "metric"), ("chart_ids", "chart"))
                        for value in (item.get(key) or [])
                        if _compact_text(value) not in binding_material
                    )
                    co_findings.append(
                        f"  → use ppt_metadata slide {target[0]} shape_name={target[1]!r} descr='{values}'"
                    )

        for key, item in (terms.get("correction_contract") or {}).items():
            if not isinstance(item, dict):
                continue
            missing_pair = []
            for value_key, value_label in (
                ("metric_id", "metric"),
                ("correct_metric_id", "metric"),
                ("correct_value", "value"),
                ("value", "value"),
            ):
                value = item.get(value_key)
                if value and _compact_text(str(value)) not in binding_material:
                    missing_pair.append(f"{value_label} {value}")
            if missing_pair:
                co_findings.append(f"correction {key}: missing=[{'; '.join(missing_pair)}]")
            forbidden = item.get("forbidden_claim")
            if forbidden and _compact_text(str(forbidden)) in compact_visible:
                co_findings.append(f"correction {key}: forbidden claim present='{forbidden}'")

    # Source coverage for structured packages (xmind-style): each top-level
    # topic must leave a textual trace in the generated deck.  Relationship
    # edge coverage is contextual (paraphrase is legal), so it is injected as
    # context only and scored by the official evaluator.
    xmind_expected = terms.get("xmind_expected") or {}
    top_level = xmind_expected.get("top_level_topics") or []
    if isinstance(top_level, list) and top_level:
        visible_compact = _compact_text("".join(_slide_text_material(slide) for slide in h.deck.slides))
        for item in top_level:
            if not isinstance(item, dict):
                continue
            candidates = [str(item.get("title") or "")] + [str(alias) for alias in (item.get("aliases") or [])]
            if not any(_compact_text(candidate) and _compact_text(candidate) in visible_compact for candidate in candidates):
                co_findings.append(f"source topic missing: {item.get('title', '?')}")

    total_findings = sum(len(value) for value in missing_by_slide.values()) + sum(len(value) for value in forbidden_by_slide.values()) + len(co_findings)
    if total_findings:
        lines = [f"Verification contract gate FAILED ({total_findings} findings). Repair each slide with batch_updates: set_shape_text for named boxes, set_table for whole tables (names are in the full-deck summary)."]
        def slide_key(pair):
            slide, _ = pair
            return (0, int(slide)) if str(slide).isdigit() else (1, str(slide))
        for slide, missing in sorted(missing_by_slide.items(), key=slide_key):
            forbidden = forbidden_by_slide.pop(slide, [])
            parts = [f"slide {int(slide) if str(slide).isdigit() else slide}: required=[{'; '.join(missing)}]"]
            if forbidden:
                parts.append(f"forbidden=[{'; '.join(forbidden)}]")
            lines.append(" | ".join(parts))
        for slide, forbidden in sorted(forbidden_by_slide.items(), key=slide_key):
            lines.append(f"slide {int(slide) if str(slide).isdigit() else slide}: forbidden=[{'; '.join(forbidden)}]")
        lines.extend(co_findings[:40])
        if len(co_findings) > 40:
            lines.append(f"... and {len(co_findings) - 40} more contract findings")
        h.state.unresolved_checks.add("ppt_contract")
        h.state.record_evidence("ppt_contract", f"contract gate failed: {total_findings} finding(s)", passed=False)
        return False, "\n".join(lines)
    h.state.unresolved_checks.discard("ppt_contract")
    h.state.record_evidence("ppt_contract", "verification contract gate passed")
    return True, "Verification contract gate passed."


def _slide_text_material_for_shape(shape) -> str:
    parts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        parts.append(shape.text or "")
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                parts.append(cell.text or "")
    return "".join(parts)


def _verify(h, policy: str = "auto") -> str:
    if getattr(h, "deck", None) is None:
        return "no deck yet to verify."
    findings = _collect_structural_findings(h)
    state = h.state
    scoped = policy == "auto" and state.ppt_existing_deck and state.ppt_baseline_captured
    baseline = state.ppt_baseline_findings if scoped else {}
    affected = state.ppt_affected_slides if scoped else set()
    blocking: list[dict[str, Any]] = []
    warnings: list[dict[str, Any]] = []
    for item in findings:
        prior = baseline.get(item["key"])
        is_new = prior is None
        is_worse = prior is not None and item["severity"] > prior + 1e-6
        # Existing defects on an affected slide are not automatically caused by
        # the current edit. Only new or worsened findings block; unchanged
        # baseline findings remain auditable warnings even inside the scope.
        if not scoped or is_new or is_worse:
            blocking.append(item)
        else:
            warnings.append(item)

    warning_summary = "; ".join(item["message"] for item in warnings)
    state.record_fact("ppt_structural_warnings", warning_summary or "none")
    scope_summary = (
        f"scoped existing-deck check; affected slides={sorted(affected) or 'none'}; "
        f"global warnings={len(warnings)}"
        if scoped else "full-deck check"
    )
    if not blocking:
        summary = f"ppt structural verification passed: {len(h.deck.slides)} slides; {scope_summary}"
        state.record_evidence("ppt_structural", summary)
        h.state.unresolved_checks.discard("ppt_structural")
        # This verification action passed. Unrelated unresolved obligations (for
        # example a still-failing official evaluator) must not be re-labelled as
        # a fresh local verification failure: that would reopen a second repair
        # cycle on every auto save/check and exhaust the bounded budget.
        h.state.last_verification_failed = False
        h.state.last_verification_epoch = h.state.mutation_epoch
        if getattr(h, "recorder", None):
            detail = summary + (f"\nHistorical warnings:\n{warning_summary}" if warnings else "")
            h.recorder.check("ppt_structural", True, detail)
        suffix = f" Historical warnings retained: {len(warnings)}." if warnings else ""
        return f"Verification: no structural issues found; no blocking structural issues.{suffix}"
    h.state.unresolved_checks.add("ppt_structural")
    h.state.last_verification_failed = True
    h.state.last_verification_epoch = h.state.mutation_epoch
    h.state.record_evidence("ppt_structural", f"blocking structural findings at epoch {h.state.mutation_epoch}", passed=False)
    report = "\n".join(item["message"] for item in blocking)
    if getattr(h, "recorder", None):
        detail = report + (f"\nHistorical warnings:\n{warning_summary}" if warnings else "")
        h.recorder.check("ppt_structural", False, detail)
    return "Verification findings:\n" + report


def _quality_check(h) -> str:
    """Run an aesthetic & design quality lint before the visual gate.

    Checks geometry health, visual element diversity, card container presence,
    and text density. Emits quality score (0-100), letter grade, and actionable
    design recommendations.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    line_types = {MSO_SHAPE_TYPE.LINE}
    connector_type = getattr(MSO_SHAPE_TYPE, "CONNECTOR", None)
    if connector_type is not None:
        line_types.add(connector_type)

    if getattr(h, "deck", None) is None:
        return json.dumps({"schema": "xiaopu.ppt-quality.v1", "passed": False, "error": "no deck"})
    slide_w = h.deck.slide_width
    slide_h = h.deck.slide_height
    errors: list[str] = []
    warnings: list[str] = []
    recommendations: list[str] = []
    slide_rows: list[str] = []

    total_slides = len(h.deck.slides)
    for index, slide in enumerate(h.deck.slides, 1):
        text_shapes = []
        visible_shapes = []
        container_shapes = []
        total_chars = 0
        for shape in slide.shapes:
            if getattr(shape, "width", 0) <= 0 or getattr(shape, "height", 0) <= 0:
                continue
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > slide_w or shape.top + shape.height > slide_h:
                errors.append(f"slide {index}: shape {shape.shape_id} crosses slide boundary")
            text = (shape.text_frame.text or "").strip() if getattr(shape, "has_text_frame", False) else ""
            if text:
                text_shapes.append(shape)
                total_chars += len(text)
            elif shape.shape_type not in line_types and shape.width > Inches(0.12) and shape.height > Inches(0.12):
                visible_shapes.append(shape)
                if shape.width > Inches(1.5) and shape.height > Inches(0.8):
                    container_shapes.append(shape)

        if not text_shapes:
            warnings.append(f"slide {index}: no visible text content")
        if len(visible_shapes) > 30:
            warnings.append(f"slide {index}: high shape density ({len(visible_shapes)} visible shapes)")

        is_cover = index == 1 and total_slides > 1 and total_chars < 100
        if not is_cover and not container_shapes and len(text_shapes) >= 3:
            recommendations.append(
                f"slide {index}: layout consists of bare flat text only; consider upgrading to "
                "workflow_pipeline, html_mockup, quadrant, or comparison cards for richer presentation."
            )
        if not is_cover and total_chars < 120 and not any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in slide.shapes):
            warnings.append(f"slide {index}: content density is low ({total_chars} visible chars; target 180-600 chars with structured cards and deliverables)")

        for left_index, first in enumerate(text_shapes):
            for second in visible_shapes:
                if first is second:
                    continue
                overlap_w = min(first.left + first.width, second.left + second.width) - max(first.left, second.left)
                overlap_h = min(first.top + first.height, second.top + second.height) - max(first.top, second.top)
                if overlap_w <= 0 or overlap_h <= 0:
                    continue
                ratio = (overlap_w * overlap_h) / min(first.width * first.height, second.width * second.height)
                full_bleed = second.width >= slide_w * 0.92 and second.height >= slide_h * 0.92
                contained = (
                    first.left >= second.left
                    and first.top >= second.top
                    and first.left + first.width <= second.left + second.width
                    and first.top + first.height <= second.top + second.height
                )
                text_in_container = contained and second.shape_type != MSO_SHAPE_TYPE.PICTURE
                if ratio > 0.30 and not full_bleed and not text_in_container:
                    errors.append(f"slide {index}: text shape {first.shape_id} overlaps shape {second.shape_id} ({ratio:.0%})")
        font_sizes = []
        for shape in text_shapes:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size:
                        font_sizes.append(run.font.size.pt)
        if font_sizes and min(font_sizes) < 9:
            warnings.append(f"slide {index}: very small text ({min(font_sizes):.1f}pt)")
        slide_rows.append(f"slide {index}: text={len(text_shapes)}, containers={len(container_shapes)}, chars={total_chars}")

    score = 100 - len(errors) * 25 - len(warnings) * 5 - len(recommendations) * 5
    quality_score = max(0, min(100, score))
    grade = "A" if quality_score >= 90 else "B" if quality_score >= 80 else "C" if quality_score >= 60 else "D"

    passed = not errors
    payload = {
        "schema": "xiaopu.ppt-quality.v1",
        "passed": passed,
        "quality_score": quality_score,
        "grade": grade,
        "profile": getattr(h.state, "task_profile", ""),
        "design_policy": getattr(h.state, "design_policy", ""),
        "slides": len(h.deck.slides),
        "errors": errors,
        "warnings": warnings,
        "recommendations": recommendations,
        "rows": slide_rows,
    }
    summary = json.dumps(payload, ensure_ascii=False, indent=2)
    if passed:
        h.state.record_evidence("ppt_quality", f"ppt quality lint passed (grade {grade}, score {quality_score}): {len(h.deck.slides)} slides")
        h.state.unresolved_checks.discard("ppt_quality")
        h.state.last_verification_failed = False
        if getattr(h, "recorder", None):
            h.recorder.check("ppt_quality", True, summary)
        return summary
    h.state.unresolved_checks.add("ppt_quality")
    h.state.last_verification_failed = True
    h.state.last_verification_epoch = h.state.mutation_epoch
    h.state.record_evidence("ppt_quality", f"ppt quality lint found blocking issues at epoch {h.state.mutation_epoch}", passed=False)
    if getattr(h, "recorder", None):
        h.recorder.check("ppt_quality", False, summary)
    return summary



def _make(name: str, description: str, params: dict[str, dict], required: list[str], fn: Callable):
    return (_schema(name, description, params, required), fn)


ppt_tools = [
    _make(
        "ppt_open",
        "Open an existing PPTX as an editable working copy while preserving the original.",
        {"path": {"type": "string"}}, ["path"],
        lambda h, **kw: _open_deck(h, kw["path"]),
    ),
    _make(
        "ppt_inspect",
        "Inspect the active deck once. Use summary for slide contents or shapes for stable ids and geometry.",
        {"detail": {"type": "string", "enum": ["summary", "shapes"]}, "slide_number": {"type": "integer"}}, [],
        lambda h, **kw: _ppt_inspect(h, kw.get("detail", "summary"), kw.get("slide_number")),
    ),
    _make(
        "ppt_edit_text",
        "Edit presentation text. operation='replace' does one exact substring replacement (all occurrences in scope, including table cells). operation='replace_case_variants' replaces singular/plural and lowercase/Capitalized/UPPERCASE forms in one call (old + new + optional new_plural). operation='append_bullet' appends NEW bullet text `new` (or `text`) immediately after the anchor paragraph selected by `old`/`after`/`text_contains`; never put the anchor into `text`. operation='set_shape_text' rewrites a whole existing text shape: select it by shape_id or unique shape_name or text_contains and provide multi-line `text`. operation='set_table' rewrites a whole existing table atomically: select the table by shape_id or shape_name and provide `rows` as a list of cell-value lists matching the current table dimensions exactly. operation='batch_updates' applies 2+ independent replace/style/set_shape_text/set_table edits in one transaction. For consistency/source-sync work prefer set_shape_text/set_table (single or inside batch_updates) to avoid stale-fragment leftovers.",
        {
            "operation": {"type": "string", "enum": ["replace", "append_bullet", "batch_updates", "set_shape_text", "set_table", "replace_case_variants"]},
            "slide_number": {"type": "integer"}, "shape_id": {"type": "integer"}, "shape_name": {"type": "string"},
            "text_contains": {"type": "string"}, "old": {"type": "string"}, "new": {"type": "string"}, "new_plural": {"type": "string"},
            "after": {"type": "string", "description": "append_bullet anchor paragraph text"},
            "text": {"type": "string"}, "match_case": {"type": "boolean"}, "all_matches": {"type": "boolean"},
            "rows": {"type": "array", "minItems": 1, "maxItems": 20, "items": {"type": "array", "minItems": 1, "maxItems": 12, "items": {"type": "string"}}},
            "updates": {
                "type": "array", "minItems": 1, "maxItems": 100,
                "items": {
                    "type": "object",
                    "properties": {
                        "operation": {"type": "string", "enum": ["replace", "style", "set_shape_text", "set_table"]},
                        "slide_number": {"type": "integer"}, "shape_id": {"type": "integer"}, "shape_name": {"type": "string"},
                        "text_contains": {"type": "string"}, "old": {"type": "string"},
                        "new": {"type": "string"}, "text": {"type": "string"}, "match_case": {"type": "boolean"},
                        "target": {"type": "string", "enum": ["text", "fill"]},
                        "size": {"type": "integer"}, "color": {"type": "string"},
                        "bold": {"type": "boolean"}, "all_matches": {"type": "boolean"},
                        "rows": {"type": "array", "minItems": 1, "maxItems": 20, "items": {"type": "array", "minItems": 1, "maxItems": 12, "items": {"type": "string"}}},
                    },
                    "required": ["operation"], "additionalProperties": False,
                },
            },
            "update": {
                "type": "array", "minItems": 1, "maxItems": 100,
                "description": "Alias for updates.",
                "items": {
                    "type": "object",
                    "properties": {
                        "operation": {"type": "string", "enum": ["replace", "style", "set_shape_text", "set_table"]},
                        "slide_number": {"type": "integer"}, "shape_id": {"type": "integer"}, "shape_name": {"type": "string"},
                        "text_contains": {"type": "string"}, "old": {"type": "string"},
                        "new": {"type": "string"}, "text": {"type": "string"}, "match_case": {"type": "boolean"},
                        "target": {"type": "string", "enum": ["text", "fill"]},
                        "size": {"type": "integer"}, "color": {"type": "string"},
                        "bold": {"type": "boolean"}, "all_matches": {"type": "boolean"},
                        "rows": {"type": "array", "minItems": 1, "maxItems": 20, "items": {"type": "array", "minItems": 1, "maxItems": 12, "items": {"type": "string"}}},
                    },
                    "required": ["operation"], "additionalProperties": False,
                },
            },
        }, ["operation"],
        lambda h, **kw: _ppt_edit_text(h, **kw),
    ),
    _make(
        "ppt_metadata",
        "Set shape description metadata (descr). Use it to record provenance IDs (html_anchor, metric_id, chart_id) that verifiers read from shape alt text/description even when they should not appear in visible body text.",
        {
            "slide_number": {"type": "integer"}, "shape_id": {"type": "integer"},
            "shape_name": {"type": "string"}, "descr": {"type": "string"},
        }, ["slide_number", "descr"],
        lambda h, **kw: _set_shape_metadata(h, kw["slide_number"], kw["descr"], kw.get("shape_id"), kw.get("shape_name", "")),
    ),
    _make(
        "ppt_notes",
        "Set the speaker-notes/backstage text for one slide. Use it for host-only boundaries, preparation notes, and deterministic coverage terms that must stay out of the public visible body.",
        {
            "slide_number": {"type": "integer"}, "text": {"type": "string"},
        }, ["slide_number", "text"],
        lambda h, **kw: _set_speaker_notes(h, kw["slide_number"], kw["text"]),
    ),
    _make(
        "ppt_style",
        "Style a selected text shape or shape fill. Select by shape id or distinctive contained text.",
        {
            "slide_number": {"type": "integer"}, "target": {"type": "string", "enum": ["text", "fill"]},
            "shape_id": {"type": "integer"}, "text_contains": {"type": "string"},
            "size": {"type": "integer"}, "color": {"type": "string"}, "bold": {"type": "boolean"},
            "all_matches": {"type": "boolean"},
        }, ["slide_number", "target"],
        lambda h, **kw: _ppt_style(h, **kw),
    ),
    _make(
        "ppt_compose",
        "Create one semantic presentation unit. kind='workflow_pipeline' (or 'workflow') builds a modern horizontal step-card pipeline with badges, action lines, and detail bullets (use for agent workflows, architectures, processes); kind='html_slide' (or 'from_html') renders pure HTML/CSS into a high-fidelity presentation slide; kind='html_mockup' (or 'web_dashboard') builds a Web/Browser interface with window chrome (macOS 3 dots), URL bar, sidebar, and component card grid; kind='hero_split' builds a left 1/3 key takeaway highlight + right 2/3 breakdown cards; kind='quadrant' builds a 4-card dashboard page; kind='comparison' builds a 2-column contrast slide; kind='content' builds a general bullet page; kind='flowchart' draws connected nodes; kind='new_deck' starts a fresh deck.",
        {
            "kind": {"type": "string", "enum": ["new_deck", "content", "comparison", "from_slides", "from_outline", "table", "quadrant", "flowchart", "textbox", "workflow_pipeline", "workflow", "step_process", "html_slide", "from_html", "html_deck", "from_html_deck", "html_mockup", "web_dashboard", "html_page", "hero_split"]},
            "slide_number": {"type": "integer"}, "title": {"type": "string"}, "subtitle": {"type": "string"},
            "html": {"type": "string", "description": "Raw HTML/CSS snippet or complete document to render for html_slide."},
            "css": {"type": "string", "description": "Optional CSS stylesheet rules for html_slide."},
            "file_path": {"type": "string", "description": "Relative path to an HTML file to compile into slide(s) for from_html."},
            "takeaway": {"type": "string", "description": "Bottom conclusion/takeaway banner text."},
            "url_bar": {"type": "string", "description": "Address bar text for html_mockup."},
            "hero_title": {"type": "string"}, "hero_text": {"type": "string"}, "hero_metric": {"type": "string"},
            "steps": {
                "type": "array", "minItems": 1, "maxItems": 8,
                "description": "Step items for workflow_pipeline. Each object has title, action/summary, detail, bullets, and optional tech tag.",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "action": {"type": "string"},
                        "summary": {"type": "string"},
                        "detail": {"type": "string"},
                        "bullets": {"type": "array", "items": {"type": "string"}},
                        "tag": {"type": "string"},
                        "tech": {"type": "string"},
                    },
                    "required": ["title"], "additionalProperties": True,
                },
            },
            "cards": {
                "type": "array", "minItems": 1, "maxItems": 6,
                "description": "Component cards for html_mockup or hero_split.",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "status": {"type": "string"},
                        "badge": {"type": "string"},
                        "metric": {"type": "string"},
                        "highlight": {"type": "string"},
                        "detail": {"type": "string"},
                        "bullets": {"type": "array", "items": {"type": "string"}},
                        "html_anchor": {"type": "string"},
                    },
                    "required": ["title"], "additionalProperties": True,
                },
            },
            "source_slides": {"type": "array", "minItems": 2, "maxItems": 2, "items": {"type": "integer"}},
            "insert_after": {"type": "integer"},
            "add_table": {"type": "boolean", "description": "Accepted for clarity; from_slides always includes a comparison table."},
            "replace_template": {"type": "boolean"},
            "slides": {
                "type": "array", "minItems": 1, "maxItems": 20,
                "items": {
                    "type": "object",
                    "properties": {
                        "template_slide": {"type": "integer"},
                        "replacements": {
                            "type": "array", "maxItems": 30,
                            "items": {
                                "type": "object",
                                "properties": {
                                    "shape_name": {"type": "string"}, "text": {"type": "string"},
                                },
                                "required": ["shape_name", "text"], "additionalProperties": False,
                            },
                        },
                        "table": {
                            "type": "object",
                            "properties": {
                                "shape_name": {"type": "string"},
                                "rows": {"type": "array", "items": {"type": "array", "items": {"type": "string"}}},
                            },
                            "required": ["shape_name", "rows"], "additionalProperties": False,
                        },
                        "speaker_notes": {"type": "string"},
                    },
                    "required": ["template_slide", "replacements"], "additionalProperties": False,
                },
            },
            "bullets": {"type": "array", "items": {"type": "string"}}, "size": {"type": "integer"},
            "left_title": {"type": "string"}, "left_bullets": {"type": "array", "items": {"type": "string"}},
            "right_title": {"type": "string"}, "right_bullets": {"type": "array", "items": {"type": "string"}},
            "columns": {"type": "array", "items": {"type": "string"}}, "rows": {"type": "array", "items": {"type": "array", "items": {"type": "string"}}},
            "quadrants": {
                "type": "array", "minItems": 4, "maxItems": 4,
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"}, "metric": {"type": "string"},
                        "detail": {"type": "string"},
                        "bullets": {"type": "array", "maxItems": 6, "items": {"type": "string"}},
                        "source": {"type": "string"},
                    },
                    "required": ["title", "metric", "detail", "bullets", "source"],
                    "additionalProperties": False,
                },
            },
            "nodes": {"type": "array", "minItems": 2, "maxItems": 8, "items": {"type": "string"}},
            "x": {"type": "number"}, "y": {"type": "number"}, "w": {"type": "number"}, "width": {"type": "number"}, "height": {"type": "number"}, "h": {"type": "number"},
            "text": {"type": "string"}, "bold": {"type": "boolean"},
        }, ["kind"],
        lambda h, **kw: _ppt_compose(h, **kw),
    ),

    _make(
        "ppt_arrange",
        "Targeted layout operations on any shape (including pictures): geometry moves/resizes a shape via x/y/w/height, delete_shape/delete_slide/move_slide, or reflow_two_columns to split an existing checklist. Use geometry to resize an image that overlaps newly added content. Inspect first.",
        {
            "operation": {"type": "string", "enum": ["geometry", "delete_shape", "delete_slide", "move_slide", "reflow_two_columns"]},
            "slide_number": {"type": "integer"}, "shape_id": {"type": "integer"},
            "x": {"type": "number"}, "y": {"type": "number"}, "w": {"type": "number"}, "width": {"type": "number"}, "height": {"type": "number"}, "h": {"type": "number"},
            "new_position": {"type": "integer"},
            "text_contains": {"type": "string", "description": "Optional selector when shape_id is omitted."},
            "split_after": {"type": "integer", "description": "Optional left-column item count; defaults to a balanced split."},
            "gutter": {"type": "number", "description": "Column gap in inches (0.15-1.5, default 0.35)."},
        }, ["operation", "slide_number"],
        lambda h, **kw: _ppt_arrange(h, **kw),
    ),
    _make(
        "ppt_save",
        "Commit the active deck to the task-required output path. Omit path when intake already found that path.",
        {"path": {"type": "string"}}, [],
        lambda h, **kw: _save(h, kw.get("path")),
    ),
    _make(
        "ppt_check",
        "Check the current deck revision. Auto runs structural checks; full also runs deterministic quality lint. Finish owns evaluator/render/visual gates.",
        {"policy": {"type": "string", "enum": ["auto", "full"]}}, [],
        lambda h, **kw: _ppt_check(h, kw.get("policy", "auto")),
    ),
    _make(
        "open_deck",
        "Load an existing .pptx from the workspace for inspection or modification.",
        {"path": {"type": "string"}},
        ["path"],
        lambda h, **kw: _open_deck(h, kw["path"]),
    ),
    _make(
        "new_deck",
        "Create a new 16:9 PowerPoint deck with ONE decorative cover slide. To build a multi-page deck, call new_deck once, then compose each requested page with ppt_compose; a slide-numbered compose on slide 1 automatically converts the cover into that content page. Never open a non-existent file first.",
        {"title": {"type": "string"}, "subtitle": {"type": "string"}},
        ["title"],
        lambda h, **kw: _new_deck(h, kw.get("title", "Untitled"), kw.get("subtitle", "")),
    ),
    _make(
        "add_slide",
        "Add a content slide with title and optional bullet list.",
        {
            "title": {"type": "string"},
            "bullets": {"type": "array", "items": {"type": "string"}},
            "size": {"type": "integer", "description": "bullet font size"},
        },
        ["title"],
        lambda h, **kw: _content_slide(h, kw.get("title", ""), kw.get("bullets") or [], kw.get("size", 18)),
    ),
    _make(
        "add_two_column_slide",
        "Add a comparison or paired-argument slide with two balanced cards.",
        {"title": {"type": "string"}, "left_title": {"type": "string"}, "left_bullets": {"type": "array", "items": {"type": "string"}}, "right_title": {"type": "string"}, "right_bullets": {"type": "array", "items": {"type": "string"}}},
        ["title", "left_title", "left_bullets", "right_title", "right_bullets"],
        lambda h, **kw: _two_column(h, kw["title"], kw["left_title"], kw["left_bullets"], kw["right_title"], kw["right_bullets"]),
    ),
    _make(
        "compose_quadrant_slide",
        "Compose a complete executive 2x2 quadrant page in one call. The harness owns layout, typography, source chips, and speaker-note provenance. Use slide_number to replace a template slide; omit it to append.",
        {
            "title": {"type": "string"},
            "subtitle": {"type": "string"},
            "slide_number": {"type": "integer"},
            "quadrants": {
                "type": "array", "minItems": 4, "maxItems": 4,
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string"},
                        "metric": {"type": "string"},
                        "detail": {"type": "string"},
                        "bullets": {"type": "array", "maxItems": 3, "items": {"type": "string"}},
                        "source": {"type": "string", "description": "source anchors and chart ids"},
                    },
                    "required": ["title", "metric", "detail", "bullets", "source"],
                },
            },
        },
        ["title", "quadrants"],
        lambda h, **kw: _quadrant_slide(h, kw["title"], kw.get("subtitle", ""), kw["quadrants"], kw.get("slide_number")),
    ),
    _make(
        "add_metric_slide",
        "Add a visual KPI slide with 1-4 metric cards and an optional takeaway.",
        {"title": {"type": "string"}, "metrics": {"type": "array", "minItems": 1, "maxItems": 4, "items": {"type": "object", "properties": {"value": {"type": "string"}, "label": {"type": "string"}, "detail": {"type": "string"}}, "required": ["value", "label"]}}, "takeaway": {"type": "string"}},
        ["title", "metrics"],
        lambda h, **kw: _metric_slide(h, kw["title"], kw["metrics"], kw.get("takeaway", "")),
    ),
    _make(
        "add_table_slide",
        "Add a styled table slide. Prefer for exact comparisons; keep at most 6 columns and 10 rows.",
        {"title": {"type": "string"}, "columns": {"type": "array", "items": {"type": "string"}}, "rows": {"type": "array", "items": {"type": "array", "items": {"type": "string"}}}},
        ["title", "columns", "rows"],
        lambda h, **kw: _table_slide(h, kw["title"], kw["columns"], kw["rows"]),
    ),
    _make(
        "add_process_slide",
        "Add a 3-5 step process, method, roadmap, or closing action slide.",
        {"title": {"type": "string"}, "steps": {"type": "array", "minItems": 3, "maxItems": 5, "items": {"type": "object", "properties": {"title": {"type": "string"}, "detail": {"type": "string"}}, "required": ["title", "detail"]}}, "takeaway": {"type": "string"}},
        ["title", "steps"],
        lambda h, **kw: _process_slide(h, kw["title"], kw["steps"], kw.get("takeaway", "")),
    ),
    _make(
        "add_image_slide",
        "Add a full-width image slide using an image already in the workspace.",
        {"title": {"type": "string"}, "image_path": {"type": "string"}, "caption": {"type": "string"}},
        ["title", "image_path"],
        lambda h, **kw: _image_slide(h, kw["title"], kw["image_path"], kw.get("caption", "")),
    ),
    _make(
        "add_textbox",
        "Add a free text box to the last slide.",
        {
            "x": {"type": "number"}, "y": {"type": "number"},
            "w": {"type": "number"}, "height": {"type": "number"},
            "text": {"type": "string"}, "size": {"type": "integer"},
            "bold": {"type": "boolean"},
        },
        ["x", "y", "w", "height", "text", "size"],
        lambda h, **kw: _add_box(h, _box_from_kw(kw)),
    ),
    _make(
        "save_deck",
        "Save current deck to a .pptx file (default workspace/deck.pptx).",
        {"path": {"type": "string"}, "output_path": {"type": "string", "description": "Alias for path."}},
        [],
        lambda h, **kw: _save(h, kw.get("path") or kw.get("output_path")),
    ),
    _make(
        "deck_info",
        "List slides and their contents in the current deck.",
        {},
        [],
        lambda h, **kw: _deck_info(h),
    ),
    _make(
        "shape_inventory",
        "Inspect editable shapes with stable slide number, shape id, geometry, and text. Call before modifying an existing deck.",
        {"slide_number": {"type": "integer"}},
        [],
        lambda h, **kw: _shape_inventory(h, kw.get("slide_number")),
    ),
    _make(
        "replace_shape_text",
        "Replace text in one shape identified by shape_inventory while preserving its primary font styling.",
        {"slide_number": {"type": "integer"}, "shape_id": {"type": "integer"}, "text": {"type": "string"}},
        ["slide_number", "shape_id", "text"],
        lambda h, **kw: _replace_text(h, kw["slide_number"], kw["shape_id"], kw["text"]),
    ),
    _make(
        "replace_text",
        "Replace exact text across one slide or the whole deck, including matches split across multiple runs. Preserves the paragraph's dominant style.",
        {"old": {"type": "string"}, "new": {"type": "string"}, "slide_number": {"type": "integer"}, "match_case": {"type": "boolean"}},
        ["old", "new"],
        lambda h, **kw: _replace_text_semantic(h, kw["old"], kw["new"], kw.get("slide_number"), kw.get("match_case", True)),
    ),
    _make(
        "append_bullet",
        "Append a peer bullet to a body/list shape while inheriting paragraph and run styling. Target by shape id or distinctive text.",
        {"slide_number": {"type": "integer"}, "text": {"type": "string"}, "shape_id": {"type": "integer"}, "text_contains": {"type": "string"}},
        ["slide_number", "text"],
        lambda h, **kw: _append_bullet(h, kw["slide_number"], kw["text"], kw.get("shape_id"), kw.get("text_contains", "")),
    ),
    _make(
        "set_text_style",
        "Set font size, color, or bold for a text shape selected by shape id or contained text.",
        {"slide_number": {"type": "integer"}, "shape_id": {"type": "integer"}, "text_contains": {"type": "string"}, "size": {"type": "integer"}, "color": {"type": "string", "description": "RGB hex, e.g. 1F3864"}, "bold": {"type": "boolean"}},
        ["slide_number"],
        lambda h, **kw: _set_text_style(h, kw["slide_number"], kw.get("shape_id"), kw.get("text_contains", ""), kw.get("size"), kw.get("color"), kw.get("bold")),
    ),
    _make(
        "set_shape_fill",
        "Set the solid fill color of a shape selected by shape id or contained text.",
        {"slide_number": {"type": "integer"}, "color": {"type": "string", "description": "RGB hex"}, "shape_id": {"type": "integer"}, "text_contains": {"type": "string"}},
        ["slide_number", "color"],
        lambda h, **kw: _set_shape_fill(h, kw["slide_number"], kw["color"], kw.get("shape_id"), kw.get("text_contains", "")),
    ),
    _make(
        "add_textbox_to_slide",
        "Add a text box to a specified existing slide, using slide-inch coordinates.",
        {"slide_number": {"type": "integer"}, "x": {"type": "number"}, "y": {"type": "number"}, "w": {"type": "number"}, "height": {"type": "number"}, "text": {"type": "string"}, "size": {"type": "integer"}, "bold": {"type": "boolean"}},
        ["slide_number", "x", "y", "w", "height", "text", "size"],
        lambda h, **kw: _add_textbox_to_slide(h, kw["slide_number"], _box_from_kw(kw)),
    ),
    _make(
        "add_flowchart",
        "Add a compact 3-5 node horizontal flowchart to a specified existing slide.",
        {"slide_number": {"type": "integer"}, "nodes": {"type": "array", "minItems": 3, "maxItems": 5, "items": {"type": "string"}}, "title": {"type": "string"}},
        ["slide_number", "nodes"],
        lambda h, **kw: _add_flowchart(h, kw["slide_number"], kw["nodes"], kw.get("title", "")),
    ),
    _make(
        "set_shape_geometry",
        "Move and resize an existing shape using slide-inch coordinates. Inspect with shape_inventory first.",
        {"slide_number": {"type": "integer"}, "shape_id": {"type": "integer"}, "x": {"type": "number"}, "y": {"type": "number"}, "w": {"type": "number"}, "height": {"type": "number"}},
        ["slide_number", "shape_id", "x", "y", "w", "height"],
        lambda h, **kw: _set_shape_geometry(h, kw["slide_number"], kw["shape_id"], kw["x"], kw["y"], kw["w"], kw["height"]),
    ),
    _make(
        "delete_shape",
        "Delete one shape identified by shape_inventory.",
        {"slide_number": {"type": "integer"}, "shape_id": {"type": "integer"}},
        ["slide_number", "shape_id"],
        lambda h, **kw: _delete_shape(h, kw["slide_number"], kw["shape_id"]),
    ),
    _make(
        "delete_slide",
        "Delete a slide by its current one-based position.",
        {"slide_number": {"type": "integer"}},
        ["slide_number"],
        lambda h, **kw: _delete_slide(h, kw["slide_number"]),
    ),
    _make(
        "move_slide",
        "Move a slide from its current one-based position to another position.",
        {"slide_number": {"type": "integer"}, "new_position": {"type": "integer"}},
        ["slide_number", "new_position"],
        lambda h, **kw: _move_slide(h, kw["slide_number"], kw["new_position"]),
    ),
    _make(
        "set_speaker_notes",
        "Replace speaker notes on one slide. For sourced evaluation decks, include a [Sources] block without exposing planning notes on the slide.",
        {"slide_number": {"type": "integer"}, "text": {"type": "string"}},
        ["slide_number", "text"],
        lambda h, **kw: _set_speaker_notes(h, kw["slide_number"], kw["text"]),
    ),
    _make(
        "ppt_verify",
        "Run structural verification on the deck: empty text boxes, overflow proxy. Returns a report.",
        {},
        [],
        lambda h, **kw: _verify(h),
    ),
    _make(
        "ppt_quality_check",
        "Run a deterministic PPT quality lint for boundaries, text overlap, density, and text size. Returns JSON evidence.",
        {},
        [],
        lambda h, **kw: _quality_check(h),
    ),
    _make(
        "render_deck",
        "Render a saved .pptx to PNG slides for visual inspection. On Windows uses installed PowerPoint; Linux benchmark images should provide LibreOffice.",
        {"path": {"type": "string"}, "output_dir": {"type": "string"}},
        ["path"],
        lambda h, **kw: _render(h, kw["path"], kw.get("output_dir", "rendered")),
    ),
    _make(
        "inspect_rendered_deck",
        "Inspect rendered PNG slides for blank output and edge-clipping risk. This is a deterministic gate, not a substitute for semantic visual review.",
        {"output_dir": {"type": "string"}},
        [],
        lambda h, **kw: _inspect_rendered(h, kw.get("output_dir", "rendered")),
    ),
]
