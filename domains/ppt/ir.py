"""PPT domain IR: structured view of an artifact that the generic core treats as opaque."""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional, Tuple

EMU_PER_INCH = 914400


@dataclass
class ShapeIR:
    id: int
    name: str
    kind: str
    text: str = ""
    left: Optional[float] = None
    top: Optional[float] = None
    width: Optional[float] = None
    height: Optional[float] = None
    subshapes: List["ShapeIR"] = field(default_factory=list)

    def all_text(self) -> str:
        parts = [self.text] if self.text else []
        for sub in self.subshapes:
            parts.append(sub.all_text())
        return "\n".join(p for p in parts if p)


@dataclass
class TextBoxIR(ShapeIR):
    pass


@dataclass
class ImageIR(ShapeIR):
    pass


@dataclass
class TableIR(ShapeIR):
    rows: int = 0
    columns: int = 0


@dataclass
class ChartIR(ShapeIR):
    pass


@dataclass
class SlideIR:
    index: int
    layout: str = ""
    shapes: List[ShapeIR] = field(default_factory=list)
    notes: str = ""

    def all_text(self) -> str:
        return "\n".join(s.all_text() for s in self.shapes if s.all_text())


@dataclass
class PresentationIR:
    path: Path
    slide_width: float
    slide_height: float
    slides: List[SlideIR] = field(default_factory=list)
    kind: str = "presentation"

    def summary(self) -> str:
        n_shapes = sum(len(s.shapes) for s in self.slides)
        return f"slides={len(self.slides)} shapes={n_shapes}"


def _inches(emu: Optional[int]) -> Optional[float]:
    return None if emu is None else round(emu / EMU_PER_INCH, 3)


def _shape_kind(shape) -> str:
    try:
        type_name = str(shape.shape_type or "").upper()
    except Exception:
        type_name = "UNKNOWN"
    if "PICTURE" in type_name:
        return "image"
    if "TABLE" in type_name:
        return "table"
    if "GRAPHIC_FRAME" in type_name or "CHART" in type_name:
        return "chart"
    if "GROUP" in type_name:
        return "group"
    if "TEXT_BOX" in type_name:
        return "text_box"
    if "PLACEHOLDER" in type_name:
        return "placeholder"
    if "AUTO" in type_name:
        return "autoshape"
    return "other"


def _convert_shape(shape) -> ShapeIR:
    text = ""
    if shape.has_text_frame:
        text = shape.text_frame.text
    rows = columns = 0
    if getattr(shape, "has_table", False) and shape.has_table:
        rows, columns = len(shape.table.rows), len(shape.table.columns)
    base = dict(
        id=shape.shape_id,
        name=shape.name,
        kind=_shape_kind(shape),
        text=text,
        left=_inches(shape.left),
        top=_inches(shape.top),
        width=_inches(shape.width),
        height=_inches(shape.height),
    )
    subshapes: List[ShapeIR] = []
    try:
        if base["kind"] == "group":
            subshapes = [_convert_shape(sub) for sub in shape.shapes]
    except Exception:
        pass
    if base["kind"] == "image":
        return ImageIR(**base, subshapes=subshapes)
    if base["kind"] == "table":
        return TableIR(**base, rows=rows, columns=columns, subshapes=subshapes)
    if base["kind"] == "chart":
        return ChartIR(**base, subshapes=subshapes)
    if base["kind"] == "text_box":
        return TextBoxIR(**base, subshapes=subshapes)
    return ShapeIR(**base, subshapes=subshapes)


def from_pptx(path) -> PresentationIR:
    """Build PresentationIR from a file path. Raises on invalid packages."""
    from pptx import Presentation

    prs = Presentation(str(path))
    slides: List[SlideIR] = []
    for index, slide in enumerate(prs.slides, start=1):
        shapes = [_convert_shape(sh) for sh in slide.shapes]
        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
        layout = ""
        try:
            layout = slide.slide_layout.name
        except Exception:
            pass
        slides.append(SlideIR(index=index, layout=layout, shapes=shapes, notes=notes))
    return PresentationIR(
        path=Path(path),
        slide_width=_inches(prs.slide_width) or 0.0,
        slide_height=_inches(prs.slide_height) or 0.0,
        slides=slides,
    )
