"""PPT transaction policy on top of the generic transactional primitives.

PptDelta is attribute-level: every observed change is a record
(slide, shape_id, property, before, after), so immutability checks can verify
that the actual delta is contained in the allowed mutation.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, List, Optional, Set, Tuple

from core.transaction import AllowedMutation, Certificate, Delta, MutationScope


@dataclass
class PptMutationScope(MutationScope):
    """Closed-world mutation scope: only the listed objects may change."""

    slides: Set[int] = field(default_factory=set)
    shapes: Set[int] = field(default_factory=set)
    properties: Tuple[str, ...] = ()

    def __post_init__(self) -> None:
        self.label = "ppt"
        self.fields = ("slides", "shapes", "properties")


@dataclass(frozen=True)
class PptAttributeChange:
    """One property-level change inside a deck."""

    slide: int
    shape_id: Optional[int]  # None marks a slide-level change
    property: str
    before: Any
    after: Any

    def summarize(self) -> str:
        return (f"slide={self.slide} shape={self.shape_id} "
                f"{self.property}: {self.before!r} -> {self.after!r}")


@dataclass
class PptDelta(Delta):
    kind: str = "ppt"
    added_slides: List[int] = field(default_factory=list)
    removed_slides: List[int] = field(default_factory=list)
    changed_shapes: List[int] = field(default_factory=list)
    changed_properties: Tuple[str, ...] = ()
    attribute_changes: Tuple[PptAttributeChange, ...] = ()

    def summarize(self) -> str:
        parts = [
            f"added_slides={self.added_slides}",
            f"removed_slides={self.removed_slides}",
            f"changed_shapes={self.changed_shapes}",
            f"changed_properties={self.changed_properties}",
            f"attribute_changes={len(self.attribute_changes)}",
        ]
        return " ".join(parts)

    def attribute_summary(self) -> str:
        return "; ".join(c.summarize() for c in self.attribute_changes) or "(no attribute changes)"


@dataclass
class PptImmutabilityCertificate(Certificate):
    kind: str = "ppt_immutability"
    changed_slides: List[int] = field(default_factory=list)
    added_shapes: List[int] = field(default_factory=list)
    removed_shapes: List[int] = field(default_factory=list)
    violations: List[str] = field(default_factory=list)


def _shape_attributes(shape) -> dict:
    attrs: dict = {"left": shape.left, "top": shape.top, "width": shape.width, "height": shape.height}
    if shape.has_text_frame:
        attrs["text"] = shape.text_frame.text
        font_sizes = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is not None:
                    font_sizes.append(round(run.font.size.pt, 1))
        attrs["font_sizes"] = tuple(font_sizes)
    if getattr(shape, "has_table", False):
        attrs["table_text"] = ";;".join(
            "|".join(cell.text for cell in row.cells) for row in shape.table.rows
        )
    nv = shape._element.nvSpPr.cNvPr
    descr = nv.get("descr")
    if descr:
        attrs["descr"] = descr
    try:
        if shape.fill.type is not None and shape.fill.type == 1:
            attrs["fill"] = str(shape.fill.fore_color.rgb)
    except Exception:
        pass
    return attrs


def _shape_map(prs) -> dict:
    shapes_by_slide: dict = {}
    for idx, slide in enumerate(prs.slides, start=1):
        shapes_by_slide[idx] = {shape.shape_id: shape for shape in slide.shapes}
    return shapes_by_slide


def diff_decks(baseline_path: Path, modified_path: Path) -> PptDelta:
    """Compute an attribute-level delta between two decks."""
    from pptx import Presentation

    before = Presentation(str(baseline_path))
    after = Presentation(str(modified_path))
    before_slides = list(before.slides)
    after_slides = list(after.slides)

    added_slides = list(range(len(before_slides) + 1, len(after_slides) + 1))
    removed_slides = list(range(len(after_slides) + 1, len(before_slides) + 1))

    changed_shapes: List[int] = []
    changed_properties: List[str] = []
    attribute_changes: List[PptAttributeChange] = []

    for idx in range(min(len(before_slides), len(after_slides))):
        old_ids = {shape.shape_id for shape in before_slides[idx].shapes}
        new_ids = {shape.shape_id for shape in after_slides[idx].shapes}
        added = sorted(new_ids - old_ids)
        removed = sorted(old_ids - new_ids)
        for shape_id in added:
            changed_shapes.append(shape_id)
            attribute_changes.append(PptAttributeChange(
                slide=idx + 1, shape_id=shape_id, property="shape_added",
                before=None, after="present"))
        for shape_id in removed:
            changed_shapes.append(shape_id)
            attribute_changes.append(PptAttributeChange(
                slide=idx + 1, shape_id=shape_id, property="shape_removed",
                before="present", after=None))

        old_shapes = {shape.shape_id: shape for shape in before_slides[idx].shapes}
        new_shapes = {shape.shape_id: shape for shape in after_slides[idx].shapes}
        for shape_id in sorted(old_ids & new_ids):
            old_attrs = _shape_attributes(old_shapes[shape_id])
            new_attrs = _shape_attributes(new_shapes[shape_id])
            for key in sorted(set(old_attrs) | set(new_attrs)):
                old_value = old_attrs.get(key)
                new_value = new_attrs.get(key)
                if old_value != new_value:
                    if key not in changed_properties:
                        changed_properties.append(key)
                    if shape_id not in changed_shapes:
                        changed_shapes.append(shape_id)
                    attribute_changes.append(PptAttributeChange(
                        slide=idx + 1, shape_id=shape_id, property=key,
                        before=old_value, after=new_value))

    return PptDelta(
        added_slides=added_slides,
        removed_slides=removed_slides,
        changed_shapes=changed_shapes,
        changed_properties=tuple(changed_properties),
        attribute_changes=tuple(attribute_changes),
    )


def delta_within_mutation(delta: PptDelta, mutation: AllowedMutation) -> Tuple[bool, List[str]]:
    """Check delta(PPT) subset of AllowedMutation. Returns (passed, reasons)."""
    reasons: List[str] = []
    scope = mutation.scope
    allowed_slides = set(scope.slides) if isinstance(scope, PptMutationScope) else set()
    allowed_shapes = set(scope.shapes) if isinstance(scope, PptMutationScope) else set()
    allowed_properties = set(scope.properties) if isinstance(scope, PptMutationScope) else set()

    for slide in (*delta.added_slides, *delta.removed_slides):
        if allowed_slides and slide not in allowed_slides:
            reasons.append(f"slide {slide} added/removed outside allowed slides {sorted(allowed_slides)}")

    for change in delta.attribute_changes:
        if allowed_slides and change.slide not in allowed_slides:
            reasons.append(f"change on non-allowed slide {change.slide}: {change.summarize()}")
            continue
        if change.shape_id is not None and allowed_shapes and change.shape_id not in allowed_shapes:
            reasons.append(f"change on non-allowed shape {change.shape_id}: {change.summarize()}")

    if mutation.policy is not None:
        deny = set(getattr(mutation.policy, "deny", ()))
        allow = set(getattr(mutation.policy, "allow", ()))
        touched = set(delta.changed_properties)
        if touched & deny:
            reasons.append(f"denied properties changed: {sorted(touched & deny)}")
        if allow and not touched <= allow:
            reasons.append(f"properties outside allowed set changed: {sorted(touched - allow)}")

    return (not reasons), reasons


def verify_immutability(
    baseline_path: Path,
    modified_path: Path,
    allowed_slide_indexes: Optional[Set[int]] = None,
    epoch: int = 0,
) -> PptImmutabilityCertificate:
    """Compare two decks at attribute granularity and certify non-target slides unchanged."""
    allowed = allowed_slide_indexes or set()
    delta = diff_decks(baseline_path, modified_path)

    violations: List[str] = []
    for slide in (*delta.added_slides, *delta.removed_slides):
        if slide not in allowed:
            violations.append(f"slide {slide} added/removed outside allowed slides")
    for change in delta.attribute_changes:
        if change.slide not in allowed:
            violations.append(change.summarize())

    return PptImmutabilityCertificate(
        artifact_ref=str(modified_path),
        epoch=epoch,
        passed=not violations,
        changed_slides=sorted({c.slide for c in delta.attribute_changes}),
        added_shapes=sorted({c.shape_id for c in delta.attribute_changes if c.property == "shape_added" and c.shape_id is not None}),
        removed_shapes=sorted({c.shape_id for c in delta.attribute_changes if c.property == "shape_removed" and c.shape_id is not None}),
        violations=violations,
        detail={
            "baseline": str(baseline_path),
            "allowed_slides": sorted(allowed),
            "delta": delta.summarize(),
            "attribute_delta": delta.attribute_summary(),
        },
    )
