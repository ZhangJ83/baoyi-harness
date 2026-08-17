"""Unit tests for Xiaopu Web GUI Server & REST/SSE APIs."""
import json
import os
import threading
import time
import urllib.request
import urllib.parse
import pytest

from agent.harness import Harness
from agent.web_server import ThreadingHTTPServer, XiaopuWebHandler


@pytest.fixture(scope="module")
def web_test_server(tmp_path_factory):
    # Isolate every session/workspace side effect from the real user history.
    isolated_home = tmp_path_factory.mktemp("webgui-home")
    previous_home = os.environ.get("XIAOPU_HOME")
    os.environ["XIAOPU_HOME"] = str(isolated_home)
    XiaopuWebHandler.harness = Harness(interactive=True)
    server = ThreadingHTTPServer(("127.0.0.1", 8769), XiaopuWebHandler)
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()
    time.sleep(0.5)
    yield "http://127.0.0.1:8769"
    server.shutdown()
    if previous_home is None:
        os.environ.pop("XIAOPU_HOME", None)
    else:
        os.environ["XIAOPU_HOME"] = previous_home


def test_web_static_index_html(web_test_server):
    req = urllib.request.urlopen(f"{web_test_server}/")
    assert req.status == 200
    assert "text/html" in req.headers.get("Content-Type", "")
    content = req.read().decode("utf-8")
    assert "报一" in content or "Baoyi" in content
    assert "activity-drawer" in content
    assert "prompt-input" in content
    # Sidebar management v2 controls
    assert "sidebar-search" in content
    assert "workspaces-modal" in content
    assert "sidebar-batch-bar" in content


def test_web_static_css_and_js(web_test_server):
    req_css = urllib.request.urlopen(f"{web_test_server}/style.css")
    assert req_css.status == 200
    assert "text/css" in req_css.headers.get("Content-Type", "")

    req_js = urllib.request.urlopen(f"{web_test_server}/app.js")
    assert req_js.status == 200
    assert "application/javascript" in req_js.headers.get("Content-Type", "")
    content_js = req_js.read().decode("utf-8")
    assert "Thought process" in content_js
    # Regression guards for the browser<->server contract discovered by the
    # Playwright flow: the composer posts `prompt`, the server streams token
    # events as event.content, and history needs appendAssistantMessage.
    assert "prompt: prompt," in content_js
    assert "payload.text || directText" in content_js
    assert "function appendAssistantMessage" in content_js


def test_web_api_config(web_test_server):
    req = urllib.request.urlopen(f"{web_test_server}/api/config")
    assert req.status == 200
    data = json.loads(req.read().decode("utf-8"))
    assert "known_models" in data
    assert "current_model" in data


def test_web_api_workspaces(web_test_server):
    req = urllib.request.urlopen(f"{web_test_server}/api/workspaces")
    assert req.status == 200
    data = json.loads(req.read().decode("utf-8"))
    assert "workspaces" in data
    assert "current" in data


def test_web_api_sessions_and_chat_stream(web_test_server):
    # Test getting sessions
    req = urllib.request.urlopen(f"{web_test_server}/api/sessions")
    assert req.status == 200
    data = json.loads(req.read().decode("utf-8"))
    assert "sessions" in data

    # Test sending a mock chat request with SSE streaming
    post_data = json.dumps({"task": "只回答两个字：OK"}).encode("utf-8")
    req_chat = urllib.request.Request(
        f"{web_test_server}/api/chat",
        data=post_data,
        headers={"Content-Type": "application/json"}
    )
    with urllib.request.urlopen(req_chat, timeout=30) as resp:
        assert resp.status == 200
        assert "text/event-stream" in resp.headers.get("Content-Type", "")
        # Read lines until [DONE] or first events
        collected = []
        for _ in range(50):
            line = resp.readline().decode("utf-8")
            if not line:
                break
            collected.append(line)
            if "[DONE]" in line:
                break
        raw_text = "".join(collected)
        assert "data:" in raw_text


def test_web_api_chat_accepts_prompt_field_from_browser(web_test_server, monkeypatch):
    """The web composer posts `prompt` (not `task`); the server must forward it."""
    import agent.web_server as web_server
    from types import SimpleNamespace

    captured = {}

    class FakeHarness:
        def __init__(self):
            self.session = SimpleNamespace(id="fake-session-id")
            self.stream_callback = None
            self.reasoning_callback = None

        def run(self, task):
            captured["task"] = task
            if self.stream_callback:
                self.stream_callback("OK")
            return "OK"

        def subscribe(self, callback):
            return lambda: None

    monkeypatch.setattr(web_server, "save_session",
                        lambda harness: SimpleNamespace(id="fake-session-saved"))
    monkeypatch.setattr(web_server.config, "set_command_policy",
                        lambda value: captured.__setitem__("policy", value))
    monkeypatch.setattr(XiaopuWebHandler, "harness", FakeHarness())

    policy = web_server.config.command_policy()
    post_data = json.dumps({
        "prompt": "只回答两个字：OK",
        "command_policy": policy,
    }).encode("utf-8")
    req = urllib.request.Request(
        f"{web_test_server}/api/chat",
        data=post_data,
        headers={"Content-Type": "application/json"}
    )
    with urllib.request.urlopen(req, timeout=30) as resp:
        assert resp.status == 200
        assert "text/event-stream" in resp.headers.get("Content-Type", "")
        raw_text = ""
        for _ in range(50):
            line = resp.readline().decode("utf-8")
            if not line:
                break
            raw_text += line
            if "[DONE]" in line:
                break

    assert captured["task"] == "只回答两个字：OK"
    assert captured["policy"] == policy
    assert "session_saved" in raw_text


def test_web_api_tree(web_test_server):
    req = urllib.request.urlopen(f"{web_test_server}/api/tree")
    assert req.status == 200
    data = json.loads(req.read().decode("utf-8"))
    assert "projects" in data
    assert "conversations" in data
    assert "current_workspace" in data
    assert len(data["projects"]) >= 1
    p = data["projects"][0]
    assert "name" in p
    assert "path" in p
    assert "sessions" in p
    assert "workspace_groups" in data


def test_web_api_tree_view_and_search_params(web_test_server):
    req = urllib.request.urlopen(f"{web_test_server}/api/tree?view=archive&q=no-such-session")
    assert req.status == 200
    data = json.loads(req.read().decode("utf-8"))
    assert data["view"] == "archive"
    assert data["query"] == "no-such-session"


def test_web_api_workspace_manage(web_test_server):
    req = urllib.request.urlopen(f"{web_test_server}/api/workspaces/manage")
    assert req.status == 200
    data = json.loads(req.read().decode("utf-8"))
    for key in ("active", "archived", "removed", "current"):
        assert key in data


def test_web_api_session_action_routes(web_test_server, monkeypatch):
    import agent.web_server as web_server
    captured = {}

    monkeypatch.setattr(web_server, "rename_session", lambda sid, title: captured.setdefault("rename", (sid, title)) or True)
    monkeypatch.setattr(web_server, "set_session_pinned", lambda sid, pinned: captured.setdefault("pin", (sid, pinned)) or True)
    monkeypatch.setattr(web_server, "archive_session", lambda sid: captured.setdefault("archive", sid) or True)
    monkeypatch.setattr(web_server, "trash_session", lambda sid: captured.setdefault("trash", sid) or True)
    monkeypatch.setattr(web_server, "restore_session", lambda sid: captured.setdefault("restore", sid) or True)
    monkeypatch.setattr(web_server, "purge_session", lambda sid: captured.setdefault("purge", sid) or True)

    def post(payload):
        req = urllib.request.Request(
            f"{web_test_server}/api/session/action",
            data=json.dumps(payload).encode("utf-8"),
            headers={"Content-Type": "application/json"}
        )
        with urllib.request.urlopen(req, timeout=10) as resp:
            return json.loads(resp.read().decode("utf-8"))

    assert post({"id": "s1", "action": "rename", "title": "新标题"})["status"] == "ok"
    assert post({"id": "s1", "action": "pin", "pinned": True})["status"] == "ok"
    assert post({"id": "s1", "action": "archive"})["status"] == "ok"
    assert post({"id": "s1", "action": "restore"})["status"] == "ok"
    assert post({"id": "s1", "action": "trash"})["status"] == "ok"
    assert post({"id": "s1", "action": "purge"})["status"] == "ok"
    assert captured["rename"] == ("s1", "新标题")
    assert captured["pin"] == ("s1", True)
    assert captured["archive"] == "s1"


def test_web_api_session_batch_and_workspace_action(web_test_server, monkeypatch):
    import agent.web_server as web_server

    captured = {}

    def fake_batch(ids, action):
        captured["batch"] = (ids, action)
        return {"ok": ids, "missing": []}

    def fake_rename_workspace(path, name):
        captured["ws"] = (path, name)
        return True

    monkeypatch.setattr(web_server, "batch_session_action", fake_batch)
    monkeypatch.setattr(web_server, "rename_workspace", fake_rename_workspace)

    batch_req = urllib.request.Request(
        f"{web_test_server}/api/sessions/batch",
        data=json.dumps({"ids": ["a", "b"], "action": "archive"}).encode("utf-8"),
        headers={"Content-Type": "application/json"}
    )
    with urllib.request.urlopen(batch_req, timeout=10) as resp:
        data = json.loads(resp.read().decode("utf-8"))
    assert data["status"] == "ok"
    assert captured["batch"] == (["a", "b"], "archive")

    ws_req = urllib.request.Request(
        f"{web_test_server}/api/workspace/action",
        data=json.dumps({"path": "C:/work", "action": "rename", "display_name": "Work"}).encode("utf-8"),
        headers={"Content-Type": "application/json"}
    )
    with urllib.request.urlopen(ws_req, timeout=10) as resp:
        data = json.loads(resp.read().decode("utf-8"))
    assert data["status"] == "ok"
    assert captured["ws"] == ("C:/work", "Work")


def test_web_api_choose_directory(web_test_server, monkeypatch, tmp_path):
    import agent.web_server
    test_dir = str(tmp_path / "mock_workspace")
    monkeypatch.setattr(agent.web_server, "_pick_directory_native", lambda **kwargs: test_dir)

    req = urllib.request.Request(
        f"{web_test_server}/api/choose_directory",
        data=b"{}",
        headers={"Content-Type": "application/json"}
    )
    with urllib.request.urlopen(req) as resp:
        assert resp.status == 200
        data = json.loads(resp.read().decode("utf-8"))
        assert data["status"] == "ok"
        assert "mock_workspace" in data["path"]


def test_web_api_artifacts(web_test_server):
    req = urllib.request.urlopen(f"{web_test_server}/api/artifacts")
    assert req.status == 200
    data = json.loads(req.read().decode("utf-8"))
    assert "artifacts" in data
    assert "workspace" in data
    assert "count" in data
    assert isinstance(data["artifacts"], list)


def test_web_api_reveal_file(web_test_server, tmp_path):
    f = tmp_path / "test_artifact.pptx"
    f.write_text("dummy", encoding="utf-8")

    req = urllib.request.Request(
        f"{web_test_server}/api/reveal_file",
        data=json.dumps({"path": str(f)}).encode("utf-8"),
        headers={"Content-Type": "application/json"}
    )
    with urllib.request.urlopen(req) as resp:
        assert resp.status == 200
        data = json.loads(resp.read().decode("utf-8"))
        assert data["status"] == "ok"


def test_web_api_settings(web_test_server):
    # GET settings
    req = urllib.request.urlopen(f"{web_test_server}/api/settings")
    assert req.status == 200
    data = json.loads(req.read().decode("utf-8"))
    assert "provider" in data
    assert "api_base" in data
    assert "model" in data
    assert "reasoning_effort" in data

    # POST settings
    post_req = urllib.request.Request(
        f"{web_test_server}/api/settings",
        data=json.dumps({
            "provider": "openai",
            "model": "deepseek-v4-flash",
            "reasoning_effort": "max",
            "command_policy": "ask"
        }).encode("utf-8"),
        headers={"Content-Type": "application/json"}
    )
    with urllib.request.urlopen(post_req) as resp:
        assert resp.status == 200
        res = json.loads(resp.read().decode("utf-8"))
        assert res["status"] == "ok"


def test_web_gui_index_ppt_elements(web_test_server):
    req = urllib.request.urlopen(f"{web_test_server}/")
    assert req.status == 200
    content = req.read().decode("utf-8")
    assert "tab-ppt-btn" in content
    assert "tab-ppt-panel" in content
    assert "ppt-preview-img" in content
    assert "ppt-content-txt" in content
    assert "btn-apply-ppt-text" in content


def test_web_api_ppt_content_and_preview(web_test_server):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from agent.config import sandbox_root

    # Create test deck.pptx in active workspace sandbox
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    tb1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(2))
    p1 = tb1.text_frame.paragraphs[0]
    p1.text = "AI Agent 端到端流水线"
    p1.font.size = Pt(24)

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    tb2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(2))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "AI Agent 运行时架构"
    p2.font.size = Pt(24)

    deck_file = sandbox_root() / "deck.pptx"
    prs.save(str(deck_file))

    # Test GET /api/ppt/content
    req_content = urllib.request.urlopen(f"{web_test_server}/api/ppt/content")
    assert req_content.status == 200
    content_data = json.loads(req_content.read().decode("utf-8"))
    assert content_data["success"] is True
    assert content_data["total_slides"] == 2
    assert "AI Agent 端到端流水线" in content_data["text_content"]
    assert "AI Agent 运行时架构" in content_data["text_content"]

    # Test GET /api/ppt/preview
    req_preview = urllib.request.urlopen(f"{web_test_server}/api/ppt/preview?slide=1")
    assert req_preview.status == 200
    assert req_preview.headers.get("Content-Type", "") == "image/png"
    img_bytes = req_preview.read()
    assert len(img_bytes) > 50

    # Test POST /api/ppt/apply_content
    post_apply = urllib.request.Request(
        f"{web_test_server}/api/ppt/apply_content",
        data=json.dumps({"text_content": "=== 第 1 页: 新标题 ===\n• 新要点 1"}).encode("utf-8"),
        headers={"Content-Type": "application/json"}
    )
    with urllib.request.urlopen(post_apply) as resp:
        assert resp.status == 200
        apply_data = json.loads(resp.read().decode("utf-8"))
        assert apply_data["status"] == "ok"
        assert "新标题" in apply_data["instruction"]


