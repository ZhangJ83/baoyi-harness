import json
import os
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches
from PIL import Image, ImageDraw

from agent.state import RunState
from agent.tools.registry import dispatch
from agent.tools.registry import _PPT_MUTATORS
from agent.tools.code_tools import _finish
from agent.tools.ppt_tools import _normalize_minimal_container


class DummyHarness:
    def __init__(self):
        self.state = RunState()
        self.deck = None


class PowerPointTests(unittest.TestCase):
    def test_append_bullet_inherits_anchor_level_after_its_child_subtree(self):
        h = DummyHarness()
        deck = Presentation()
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        shape = slide.shapes.add_textbox(0, 0, 5000000, 3000000)
        frame = shape.text_frame
        frame.paragraphs[0].text = "Thread"
        anchor = frame.add_paragraph(); anchor.text = "Dual Mode"; anchor.level = 0
        child1 = frame.add_paragraph(); child1.text = "Protection"; child1.level = 1
        child2 = frame.add_paragraph(); child2.text = "Kernel"; child2.level = 1
        next_peer = frame.add_paragraph(); next_peer.text = "Address Space"; next_peer.level = 0
        h.deck = deck

        result = dispatch("ppt_edit_text", json.dumps({
            "operation": "append_bullet", "slide_number": 1,
            "shape_id": shape.shape_id, "text_contains": "Dual Mode", "text": "File System",
        }), h)

        paragraphs = list(shape.text_frame.paragraphs)
        texts = [paragraph.text for paragraph in paragraphs]
        inserted = paragraphs[texts.index("File System")]
        self.assertEqual(inserted.level, 0)
        self.assertEqual(texts, ["Thread", "Dual Mode", "Protection", "Kernel", "File System", "Address Space"])
        self.assertIn("level 0", result)

    def test_build_verify_save(self):
        with tempfile.TemporaryDirectory() as tmp, patch.dict(os.environ, {"WORKSPACE": tmp}):
            h = DummyHarness()
            dispatch("new_deck", json.dumps({"title": "Test", "subtitle": "Evidence"}), h)
            dispatch("add_slide", json.dumps({"title": "Result", "bullets": ["One", "Two"]}), h)
            report = dispatch("ppt_verify", "{}", h)
            self.assertIn("no structural issues", report)
            dispatch("save_deck", json.dumps({"path": "deck.pptx"}), h)
            self.assertEqual(len(Presentation(Path(tmp) / "deck.pptx").slides), 2)

    def test_quality_check_returns_structured_evidence(self):
        h = DummyHarness()
        dispatch("new_deck", json.dumps({"title": "Quality", "subtitle": "Evidence"}), h)
        report = json.loads(dispatch("ppt_quality_check", "{}", h))
        self.assertEqual(report["schema"], "xiaopu.ppt-quality.v1")
        self.assertTrue(report["passed"])
        self.assertIn("ppt_quality", {record.kind for record in h.state.fresh_evidence()})

    def test_quality_check_allows_text_inside_container_autoshape(self):
        # A quadrant slide places its text boxes inside card autoshapes; that is
        # normal text-in-container layout, not an overlap defect.
        h = DummyHarness()
        dispatch("new_deck", json.dumps({"title": "Quadrant"}), h)
        quadrants = [
            {"title": "A", "metric": "M1", "detail": "d", "bullets": ["b"], "source": "H-117"},
            {"title": "B", "metric": "M2", "detail": "d", "bullets": ["b"], "source": "H-205"},
            {"title": "C", "metric": "M3", "detail": "d", "bullets": ["b"], "source": "H-304"},
            {"title": "D", "metric": "M4", "detail": "d", "bullets": ["b"], "source": "H-402"},
        ]
        dispatch("ppt_compose", json.dumps({"kind": "quadrant", "title": "T", "quadrants": quadrants}), h)
        report = json.loads(dispatch("ppt_quality_check", "{}", h))
        self.assertTrue(report["passed"], report.get("errors"))
        self.assertNotIn("ppt_quality", h.state.unresolved_checks)

    def test_open_inventory_replace_and_two_column(self):
        with tempfile.TemporaryDirectory() as tmp, patch.dict(os.environ, {"WORKSPACE": tmp}):
            source = Presentation()
            slide = source.slides.add_slide(source.slide_layouts[6])
            box = slide.shapes.add_textbox(0, 0, 2000000, 500000)
            box.text = "Before"
            source.save(Path(tmp) / "source.pptx")
            h = DummyHarness()
            dispatch("open_deck", json.dumps({"path": "source.pptx"}), h)
            inventory = dispatch("shape_inventory", json.dumps({"slide_number": 1}), h)
            self.assertIn("Before", inventory)
            dispatch("replace_shape_text", json.dumps({"slide_number": 1, "shape_id": box.shape_id, "text": "After"}), h)
            dispatch("add_two_column_slide", json.dumps({"title": "Compare", "left_title": "A", "left_bullets": ["one"], "right_title": "B", "right_bullets": ["two"]}), h)
            self.assertIn("After", dispatch("shape_inventory", json.dumps({"slide_number": 1}), h))
            self.assertEqual(len(h.deck.slides), 2)

    def test_semantic_edit_tools_cover_common_trajectory_failures(self):
        h = DummyHarness()
        source = Presentation()
        slide = source.slides.add_slide(source.slide_layouts[6])
        title = slide.shapes.add_textbox(0, 0, 3000000, 500000)
        paragraph = title.text_frame.paragraphs[0]
        paragraph.add_run().text = "Lecture "
        paragraph.add_run().text = "3"
        body = slide.shapes.add_textbox(0, 700000, 5000000, 2000000)
        body.text_frame.paragraphs[0].text = "First bullet"
        h.deck = source

        dispatch("replace_text", json.dumps({"old": "Lecture 3", "new": "Lecture 1"}), h)
        dispatch("append_bullet", json.dumps({"slide_number": 1, "text_contains": "First bullet", "text": "Second bullet"}), h)
        dispatch("set_text_style", json.dumps({"slide_number": 1, "text_contains": "Lecture 1", "size": 28, "color": "1F3864"}), h)
        dispatch("set_shape_fill", json.dumps({"slide_number": 1, "text_contains": "First bullet", "color": "F09A1A"}), h)
        dispatch("add_textbox_to_slide", json.dumps({"slide_number": 1, "x": 8.0, "y": 6.0, "w": 4.0, "height": 0.5, "text": "Caption", "size": 12}), h)
        dispatch("add_flowchart", json.dumps({"slide_number": 1, "nodes": ["Assets", "Liabilities", "Equity"]}), h)

        inventory = dispatch("shape_inventory", json.dumps({"slide_number": 1}), h)
        self.assertIn("Lecture 1", inventory)
        self.assertIn("Second bullet", inventory)
        self.assertIn("Caption", inventory)
        self.assertIn("Equity", inventory)

    def test_semantic_replace_preserves_soft_breaks_run_boundaries_and_styles(self):
        h = DummyHarness()
        source = Presentation()
        slide = source.slides.add_slide(source.slide_layouts[6])
        title = slide.shapes.add_textbox(0, 0, 5000000, 1500000)
        paragraph = title.text_frame.paragraphs[0]
        first = paragraph.add_run()
        first.text = "Lecture "
        first.font.bold = True
        number = paragraph.add_run()
        number.text = "3"
        number.font.italic = True
        paragraph.add_line_break()
        paragraph.add_run().text = "Agenda: "
        paragraph.add_run().text = "Lecture "
        paragraph.add_run().text = "3"
        h.deck = source

        child_kinds_before = [child.tag.rsplit("}", 1)[-1] for child in paragraph._p.content_children]
        break_nodes_before = list(paragraph._p.br_lst)
        dispatch(
            "ppt_edit_text",
            json.dumps({"operation": "replace", "old": "Lecture 3", "new": "Lecture 1"}),
            h,
        )

        self.assertEqual(paragraph.text, "Lecture 1\vAgenda: Lecture 1")
        self.assertEqual([run.text for run in paragraph.runs], ["Lecture ", "1", "Agenda: ", "Lecture ", "1"])
        self.assertEqual([child.tag.rsplit("}", 1)[-1] for child in paragraph._p.content_children], child_kinds_before)
        self.assertEqual(list(paragraph._p.br_lst), break_nodes_before)
        self.assertTrue(paragraph.runs[0].font.bold)
        self.assertTrue(paragraph.runs[1].font.italic)

    def test_compact_facade_routes_semantic_edit_and_compose(self):
        h = DummyHarness()
        source = Presentation()
        slide = source.slides.add_slide(source.slide_layouts[6])
        body = slide.shapes.add_textbox(0, 0, 5000000, 1500000)
        body.text = "Lecture 3"
        h.deck = source
        dispatch("ppt_edit_text", json.dumps({"operation": "replace", "old": "Lecture 3", "new": "Lecture 1"}), h)
        dispatch("ppt_style", json.dumps({"slide_number": 1, "target": "text", "text_contains": "Lecture 1", "size": 28}), h)
        dispatch("ppt_compose", json.dumps({"kind": "flowchart", "slide_number": 1, "nodes": ["Inspect", "Edit", "Verify"]}), h)
        self.assertIn("Lecture 1", dispatch("ppt_inspect", json.dumps({"detail": "shapes", "slide_number": 1}), h))
        self.assertIn("structural", dispatch("ppt_check", "{}", h))

    def test_arrange_reflows_checklist_to_two_columns_atomically(self):
        h = DummyHarness()
        source = Presentation()
        source.slide_width = Inches(13.333)
        source.slide_height = Inches(7.5)
        slide = source.slides.add_slide(source.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(11.5), Inches(4.7))
        frame = box.text_frame
        frame.paragraphs[0].text = "Poetry"
        frame.paragraphs[0].runs[0].font.size = Pt(30)
        items = [f"Checklist item {index}" for index in range(1, 7)]
        for text in items:
            paragraph = frame.add_paragraph()
            paragraph.text = text
            properties = paragraph._p.get_or_add_pPr()
            bullet_font = OxmlElement("a:buFont")
            bullet_font.set("typeface", "Wingdings")
            bullet = OxmlElement("a:buChar")
            bullet.set("char", "q")
            properties.append(bullet_font)
            properties.append(bullet)
            paragraph.runs[0].font.size = Pt(24)
        h.deck = source
        epoch = h.state.mutation_epoch

        report = dispatch(
            "ppt_arrange",
            json.dumps({"operation": "reflow_two_columns", "slide_number": 1, "shape_id": box.shape_id}),
            h,
        )

        self.assertIn("3+3 checklist items", report)
        self.assertEqual(h.state.mutation_epoch, epoch + 1)
        self.assertEqual(h.state.ppt_affected_slides, {1})
        self.assertEqual(len(slide.shapes), 2)
        right = slide.shapes[1]
        self.assertEqual([p.text for p in box.text_frame.paragraphs], ["Poetry", *items[:3]])
        self.assertEqual([p.text for p in right.text_frame.paragraphs], items[3:])
        self.assertGreater(right.left, box.left + box.width)
        self.assertLessEqual(right.left + right.width, source.slide_width)
        for paragraph in [*box.text_frame.paragraphs[1:], *right.text_frame.paragraphs]:
            self.assertTrue(paragraph._p.xpath("./a:pPr/a:buChar"))
        self.assertIn("no structural issues", dispatch("ppt_check", "{}", h))

    def test_arrange_reflow_rejects_non_checklist_or_invalid_split(self):
        h = DummyHarness()
        source = Presentation()
        slide = source.slides.add_slide(source.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
        box.text = "Plain text\nwithout DrawingML bullets"
        h.deck = source
        with self.assertRaisesRegex(ValueError, "at least two checklist paragraphs"):
            dispatch(
                "ppt_arrange",
                json.dumps({"operation": "reflow_two_columns", "slide_number": 1, "shape_id": box.shape_id}),
                h,
            )

    def test_compose_from_slides_inserts_media_safe_comparison_and_preserves_sources(self):
        with tempfile.TemporaryDirectory() as tmp, patch.dict(os.environ, {"WORKSPACE": tmp}):
            source = Presentation()
            source.slide_width = Inches(10)
            source.slide_height = Inches(7.5)
            for number in range(1, 5):
                slide = source.slides.add_slide(source.slide_layouts[6])
                slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.6)).text = f"Source {number}"
            male = source.slides[1]
            male.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4), Inches(3)).text = (
                "Male\nUpper part: collarless jacket\nLower part: bahag\nHeadgear: putong"
            )
            female = source.slides[2]
            female.shapes.add_textbox(Inches(4.5), Inches(1.1), Inches(4), Inches(3)).text = (
                "Female\nUpper part: baro or kamisa\nLower part: saya or patadyong\nTapis"
            )
            image_paths = []
            for name, color in (("male.png", "navy"), ("female.png", "gold")):
                path = Path(tmp) / name
                Image.new("RGB", (180, 260), color).save(path)
                image_paths.append(path)
            male.shapes.add_picture(str(image_paths[0]), Inches(5.4), Inches(1.2), width=Inches(1.5))
            female.shapes.add_picture(str(image_paths[1]), Inches(1.2), Inches(1.2), width=Inches(1.5))
            table = female.shapes.add_table(2, 2, Inches(4.5), Inches(4.3), Inches(4.5), Inches(1.2)).table
            table.cell(0, 0).text, table.cell(0, 1).text = "Item", "Name"
            table.cell(1, 0).text, table.cell(1, 1).text = "Wrap", "Tapis"
            source.save(Path(tmp) / "source.pptx")

            h = DummyHarness()
            dispatch("ppt_open", json.dumps({"path": "source.pptx"}), h)
            original_two = [shape.text for shape in h.deck.slides[1].shapes if getattr(shape, "has_text_frame", False)]
            original_three = [shape.text for shape in h.deck.slides[2].shapes if getattr(shape, "has_text_frame", False)]
            result = dispatch("ppt_compose", json.dumps({
                "kind": "from_slides", "source_slides": [2, 3], "insert_after": 3,
                "title": "Traditional Clothing", "left_title": "Male", "right_title": "Female",
            }), h)

            self.assertIn("inserted comparison slide 4", result)
            self.assertEqual(len(h.deck.slides), 5)
            self.assertEqual(original_two, [shape.text for shape in h.deck.slides[1].shapes if getattr(shape, "has_text_frame", False)])
            self.assertEqual(original_three, [shape.text for shape in h.deck.slides[2].shapes if getattr(shape, "has_text_frame", False)])
            inserted = h.deck.slides[3]
            self.assertEqual(sum(shape.shape_type == 13 for shape in inserted.shapes), 2)
            self.assertEqual(sum(shape.has_table for shape in inserted.shapes), 1)
            self.assertIn("source slide 2", inserted.notes_slide.notes_text_frame.text)
            self.assertEqual(h.state.ppt_affected_slides, {4})
            self.assertIn("no blocking structural issues", dispatch("ppt_check", "{}", h))

            dispatch("ppt_save", json.dumps({"path": "output/final.pptx"}), h)
            saved = Presentation(Path(tmp) / "output" / "final.pptx")
            self.assertEqual(len(saved.slides), 5)
            self.assertEqual(sum(shape.shape_type == 13 for shape in saved.slides[3].shapes), 2)

    def test_compose_from_slides_rejects_ambiguous_scope_before_mutation(self):
        h = DummyHarness()
        h.deck = Presentation()
        h.deck.slides.add_slide(h.deck.slide_layouts[6])
        with self.assertRaisesRegex(ValueError, "at least 2|exactly two"):
            dispatch("ppt_compose", json.dumps({
                "kind": "from_slides", "source_slides": [1], "insert_after": 1,
            }), h)
        self.assertEqual(len(h.deck.slides), 1)

    def test_batch_updates_commit_once_and_merge_affected_slides(self):
        h = DummyHarness()
        source = Presentation()
        for text in ("Status: Draft", "Decision: Pending"):
            slide = source.slides.add_slide(source.slide_layouts[6])
            slide.shapes.add_textbox(0, 0, 5000000, 1000000).text = text
        h.deck = source
        epoch_before = h.state.mutation_epoch

        report = dispatch("ppt_edit_text", json.dumps({
            "operation": "batch_updates",
            "updates": [
                {"operation": "replace", "slide_number": 1, "old": "Draft", "new": "Final"},
                {"operation": "replace", "slide_number": 2, "old": "Pending", "new": "Approved"},
                {"operation": "style", "slide_number": 2, "text_contains": "Approved", "size": 24},
            ],
        }), h)

        self.assertIn("applied 3 updates atomically", report)
        self.assertEqual(h.state.mutation_epoch, epoch_before + 1)
        self.assertEqual(h.state.ppt_affected_slides, {1, 2})
        self.assertEqual(h.deck.slides[0].shapes[0].text, "Status: Final")
        self.assertEqual(h.deck.slides[1].shapes[0].text, "Decision: Approved")
        self.assertEqual(h.deck.slides[1].shapes[0].text_frame.paragraphs[0].runs[0].font.size, Pt(24))
        with tempfile.TemporaryDirectory() as tmp:
            persisted = Path(tmp) / "batch.pptx"
            h.deck.save(persisted)
            reopened = Presentation(persisted)
            self.assertEqual(reopened.slides[0].shapes[0].text, "Status: Final")
            self.assertEqual(reopened.slides[1].shapes[0].text, "Decision: Approved")

    def test_batch_updates_roll_back_entire_deck_when_one_item_fails(self):
        h = DummyHarness()
        source = Presentation()
        slide = source.slides.add_slide(source.slide_layouts[6])
        slide.shapes.add_textbox(0, 0, 5000000, 1000000).text = "Status: Draft"
        h.deck = source
        live_deck = h.deck
        epoch_before = h.state.mutation_epoch

        with self.assertRaisesRegex(ValueError, "text not found"):
            dispatch("ppt_edit_text", json.dumps({
                "operation": "batch_updates",
                "updates": [
                    {"operation": "replace", "old": "Draft", "new": "Final"},
                    {"operation": "replace", "old": "Missing value", "new": "Approved"},
                ],
            }), h)

        self.assertIs(h.deck, live_deck)
        self.assertEqual(h.deck.slides[0].shapes[0].text, "Status: Draft")
        self.assertEqual(h.state.mutation_epoch, epoch_before)
        self.assertEqual(h.state.ppt_affected_slides, set())

    def test_from_outline_clones_template_relationships_and_commits_once(self):
        h = DummyHarness()
        template = Presentation()
        first = template.slides.add_slide(template.slide_layouts[6])
        title = first.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(0.8))
        title.name = "title"
        title.text = "Template title"
        footer = first.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(5), Inches(0.3))
        footer.name = "footer"
        footer.text = "Template footer"
        image_path = Path(tempfile.gettempdir()) / "xiaopu-outline-image.png"
        Image.new("RGB", (40, 30), "#70AD47").save(image_path)
        first.shapes.add_picture(str(image_path), Inches(8), Inches(1), Inches(2), Inches(1.5))
        second = template.slides.add_slide(template.slide_layouts[6])
        heading = second.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(0.8))
        heading.name = "title"
        heading.text = "Second template"
        h.deck = template
        epoch_before = h.state.mutation_epoch

        report = dispatch("ppt_compose", json.dumps({
            "kind": "from_outline",
            "slides": [
                {"template_slide": 1, "replacements": [{"shape_name": "title", "text": "Community workshop"}]},
                {"template_slide": 2, "replacements": [{"shape_name": "title", "text": "Follow-up"}],
                 "speaker_notes": "Private facilitator cue"},
            ],
        }), h)

        self.assertIn("generated 2 slides", report)
        self.assertEqual(len(h.deck.slides), 2)
        self.assertEqual(h.state.mutation_epoch, epoch_before + 1)
        self.assertEqual(h.state.ppt_affected_slides, set())
        self.assertEqual(h.deck.slides[0].shapes[0].text, "Community workshop")
        self.assertIn("Template footer", [shape.text for shape in h.deck.slides[0].shapes if shape.has_text_frame])
        self.assertEqual(sum(shape.shape_type == 13 for shape in h.deck.slides[0].shapes), 1)
        self.assertIn("Private facilitator cue", h.deck.slides[1].notes_slide.notes_text_frame.text)

    def test_from_outline_rolls_back_bad_shape_name(self):
        h = DummyHarness()
        template = Presentation()
        slide = template.slides.add_slide(template.slide_layouts[6])
        box = slide.shapes.add_textbox(0, 0, 2000000, 500000)
        box.name = "title"
        box.text = "Original"
        h.deck = template
        epoch_before = h.state.mutation_epoch

        with self.assertRaisesRegex(ValueError, "none of the replacement"):
            dispatch("ppt_compose", json.dumps({
                "kind": "from_outline", "slides": [{"template_slide": 1, "replacements": [
                    {"shape_name": "missing_shape", "text": "New"},
                ]}],
            }), h)

        self.assertEqual(len(h.deck.slides), 1)
        self.assertEqual(h.deck.slides[0].shapes[0].text, "Original")
        self.assertEqual(h.state.mutation_epoch, epoch_before)

    def test_compact_facade_inspects_and_styles_nested_group_shapes(self):
        h = DummyHarness()
        source = Presentation()
        slide = source.slides.add_slide(source.slide_layouts[6])
        outer = slide.shapes.add_group_shape()
        first = outer.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.2), Inches(0.2), Inches(1.5), Inches(0.6)
        )
        first.text = "vCPU1"
        inner = outer.shapes.add_group_shape()
        second = inner.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.1), Inches(0.1), Inches(1.5), Inches(0.6)
        )
        second.text = "vCPU1"
        third = inner.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.8), Inches(0.1), Inches(1.5), Inches(0.6)
        )
        third.text = "vCPU1"
        h.deck = source

        inventory = dispatch("ppt_inspect", json.dumps({"detail": "shapes", "slide_number": 1}), h)
        self.assertIn(f"path={outer.shape_id}/{first.shape_id}", inventory)
        self.assertIn(f"path={outer.shape_id}/{inner.shape_id}/{second.shape_id}", inventory)
        self.assertIn("box_group_local=", inventory)
        self.assertNotIn(f"path={outer.shape_id}/{first.shape_id} id={first.shape_id} kind=shape name={first.name!r} box_slide=", inventory)

        result = dispatch(
            "ppt_style",
            json.dumps({
                "slide_number": 1, "target": "fill", "text_contains": "vCPU1",
                "color": "70AD47", "all_matches": True,
            }),
            h,
        )
        self.assertIn("styled 3 matching shapes", result)
        for shape in (first, second, third):
            self.assertEqual(shape.fill.fore_color.rgb, RGBColor(0x70, 0xAD, 0x47))
        with self.assertRaisesRegex(ValueError, "group-local"):
            dispatch(
                "ppt_arrange",
                json.dumps({
                    "operation": "geometry", "slide_number": 1, "shape_id": second.shape_id,
                    "x": 1.0, "y": 1.0, "w": 2.0, "height": 1.0,
                }),
                h,
            )

        dispatch(
            "replace_shape_text",
            json.dumps({"slide_number": 1, "shape_id": second.shape_id, "text": "nested"}),
            h,
        )
        self.assertEqual(second.text, "nested")
        dispatch(
            "ppt_edit_text",
            json.dumps({"operation": "replace", "old": "nested", "new": "nested updated", "slide_number": 1}),
            h,
        )
        self.assertEqual(second.text, "nested updated")

    def test_facade_schemas_reject_unknown_or_incomplete_arguments(self):
        h = DummyHarness()
        dispatch("new_deck", json.dumps({"title": "Strict"}), h)
        # Unknown args are now silently dropped (tolerant to LLM typos)
        result = dispatch("ppt_check", json.dumps({"polciy": "auto"}), h)
        self.assertIsInstance(result, str)  # should succeed, not raise
        with self.assertRaisesRegex(ValueError, "at least 4|exactly four"):
            dispatch("ppt_compose", json.dumps({"kind": "quadrant", "quadrants": []}), h)

    def test_metric_and_table_layouts_verify(self):
        h = DummyHarness()
        dispatch("add_metric_slide", json.dumps({"title": "KPIs", "metrics": [{"value": "42%", "label": "Win rate", "detail": "+5pp"}], "takeaway": "Improving"}), h)
        dispatch("add_table_slide", json.dumps({"title": "Options", "columns": ["Name", "Score"], "rows": [["A", "9"], ["B", "7"]]}), h)
        dispatch("add_process_slide", json.dumps({"title": "Method", "steps": [{"title": "A", "detail": "Inspect"}, {"title": "B", "detail": "Edit"}, {"title": "C", "detail": "Verify"}]}), h)
        self.assertIn("no structural issues", dispatch("ppt_verify", "{}", h))

    def test_semantic_quadrant_tool_replaces_template_and_binds_sources(self):
        h = DummyHarness()
        dispatch("new_deck", json.dumps({"title": "Template"}), h)
        quadrants = [
            {"title": f"Q{i}", "metric": f"{i}0%", "detail": "Current fact", "bullets": ["Judgment", "Action"], "source": f"H-{i} / CH-0{i}"}
            for i in range(1, 5)
        ]
        result = dispatch("compose_quadrant_slide", json.dumps({
            "title": "Quarterly board", "subtitle": "2026 Q2",
            "slide_number": 1, "quadrants": quadrants,
        }), h)
        self.assertIn("quadrant slide 1", result)
        self.assertEqual(len(h.deck.slides), 1)
        inventory = dispatch("shape_inventory", json.dumps({"slide_number": 1}), h)
        self.assertIn("Quarterly board", inventory)
        self.assertIn("Q4", inventory)
        self.assertIn("[Sources]", h.deck.slides[0].notes_slide.notes_text_frame.text)
        self.assertIn("no structural issues", dispatch("ppt_verify", "{}", h))
        metric_shapes = [shape for shape in h.deck.slides[0].shapes if getattr(shape, "has_text_frame", False) and "%" in shape.text]
        self.assertEqual(len(metric_shapes), 4)
        for shape in metric_shapes:
            self.assertGreaterEqual(shape.width / 914400, 1.8)

    def test_edit_primitives_invalidate_prior_evidence(self):
        h = DummyHarness()
        dispatch("new_deck", json.dumps({"title": "Cover"}), h)
        dispatch("add_slide", json.dumps({"title": "First", "bullets": ["A"]}), h)
        dispatch("add_slide", json.dumps({"title": "Second", "bullets": ["B"]}), h)
        dispatch("ppt_verify", "{}", h)
        self.assertTrue(h.state.fresh_evidence())

        target = next(shape for shape in h.deck.slides[1].shapes if shape.has_text_frame and shape.text_frame.text == "First")
        dispatch("set_shape_geometry", json.dumps({"slide_number": 2, "shape_id": target.shape_id, "x": 1.0, "y": 0.5, "w": 8.0, "height": 0.7}), h)
        self.assertFalse(h.state.fresh_evidence())
        dispatch("ppt_verify", "{}", h)
        self.assertTrue(h.state.fresh_evidence())

        dispatch("move_slide", json.dumps({"slide_number": 3, "new_position": 2}), h)
        self.assertFalse(h.state.fresh_evidence())
        dispatch("delete_slide", json.dumps({"slide_number": 3}), h)
        self.assertEqual(len(h.deck.slides), 2)

    def test_rendered_pixel_audit_records_fresh_visual_evidence(self):
        with tempfile.TemporaryDirectory() as tmp, patch.dict(os.environ, {"WORKSPACE": tmp}):
            rendered = Path(tmp) / "rendered"
            rendered.mkdir()
            image = Image.new("RGB", (640, 360), "white")
            draw = ImageDraw.Draw(image)
            draw.rectangle((100, 80, 540, 280), fill=(30, 60, 100))
            image.save(rendered / "slide-1.png")
            h = DummyHarness()
            report = dispatch("inspect_rendered_deck", json.dumps({"output_dir": "rendered"}), h)
            self.assertIn("audit passed", report)
            self.assertEqual(h.state.fresh_evidence()[0].kind, "ppt_visual")

    def test_every_layout_mutation_invalidates_evidence(self):
        h = DummyHarness()
        dispatch("new_deck", json.dumps({"title": "One"}), h)
        dispatch("ppt_verify", "{}", h)
        verified_epoch = h.state.mutation_epoch
        self.assertTrue(h.state.fresh_evidence())
        dispatch("add_slide", json.dumps({"title": "Two", "bullets": ["new"]}), h)
        self.assertGreater(h.state.mutation_epoch, verified_epoch)
        self.assertFalse(h.state.fresh_evidence())

    def test_speaker_notes_are_mutating_and_persisted(self):
        with tempfile.TemporaryDirectory() as tmp, patch.dict(os.environ, {"WORKSPACE": tmp}):
            h = DummyHarness()
            dispatch("new_deck", json.dumps({"title": "Sources"}), h)
            dispatch("ppt_verify", "{}", h)
            self.assertTrue(h.state.fresh_evidence())
            dispatch("set_speaker_notes", json.dumps({"slide_number": 1, "text": "[Sources] Internal evaluation task packet"}), h)
            self.assertFalse(h.state.fresh_evidence())
            dispatch("save_deck", json.dumps({"path": "deck.pptx"}), h)
            saved = Presentation(Path(tmp) / "deck.pptx")
            self.assertIn("[Sources]", saved.slides[0].notes_slide.notes_text_frame.text)

    def test_save_binds_matching_filename_to_required_task_output(self):
        with tempfile.TemporaryDirectory() as tmp, patch.dict(os.environ, {"WORKSPACE": tmp}):
            h = DummyHarness()
            h.state.facts["required_output_pptx"] = str(Path("tasks/demo/output/final.pptx"))
            dispatch("new_deck", json.dumps({"title": "Bound"}), h)
            dispatch("save_deck", json.dumps({"path": "output/final.pptx"}), h)
            self.assertTrue((Path(tmp) / "tasks/demo/output/final.pptx").is_file())
            self.assertFalse((Path(tmp) / "output/final.pptx").exists())

    def test_save_without_path_uses_required_task_output(self):
        with tempfile.TemporaryDirectory() as tmp, patch.dict(os.environ, {"WORKSPACE": tmp}):
            h = DummyHarness()
            h.state.facts["required_output_pptx"] = str(Path("tasks/demo/output/final.pptx"))
            dispatch("new_deck", json.dumps({"title": "Bound"}), h)
            dispatch("save_deck", json.dumps({"output_path": "generic.pptx"}), h)
            self.assertTrue((Path(tmp) / "tasks/demo/output/final.pptx").is_file())
            self.assertFalse((Path(tmp) / "generic.pptx").exists())

    def test_finish_automatically_buys_render_and_visual_evidence(self):
        with tempfile.TemporaryDirectory() as tmp, patch.dict(os.environ, {"WORKSPACE": tmp}):
            h = DummyHarness()
            h._done = None
            h.recorder = type("Recorder", (), {
                "manifest": {"artifacts": [{"role": "final-pptx", "path": str(Path(tmp) / "final.pptx")}]},
                "evidence": Path(tmp) / "evidence",
                "check": lambda self, *args: None,
                "finish": lambda self, *args: None,
            })()
            h.recorder.evidence.mkdir()
            h.state.record_change("final.pptx")
            h.state.record_evidence("ppt_structural", "pass")

            def fake_render(harness, path, output_dir):
                harness.state.record_evidence("ppt_render", "pass")

            def fake_inspect(harness, output_dir):
                harness.state.record_evidence("ppt_visual", "pass")

            with patch("agent.tools.ppt_tools._render", side_effect=fake_render), patch("agent.tools.ppt_tools._inspect_rendered", side_effect=fake_inspect):
                _finish(h, "done")
            self.assertEqual({record.kind for record in h.state.fresh_evidence()}, {"ppt_structural", "ppt_render", "ppt_visual"})

    def test_finish_buys_render_before_requiring_render_certificate(self):
        # finish owns render/visual lifecycle. It must buy those certificates
        # *before* the contract certificate gate, otherwise a skill whose
        # finish_certificates include render/visual can never complete.
        with tempfile.TemporaryDirectory() as tmp, patch.dict(os.environ, {"WORKSPACE": tmp}):
            h = DummyHarness()
            h._done = None
            h.recorder = type("Recorder", (), {
                "manifest": {"artifacts": [{"role": "final-pptx", "path": str(Path(tmp) / "final.pptx")}]},
                "evidence": Path(tmp) / "evidence",
                "check": lambda self, *args: None,
                "finish": lambda self, *args: None,
            })()
            h.recorder.evidence.mkdir()
            h.state.record_change("final.pptx")
            h.state.record_evidence("ppt_structural", "pass")
            from agent.task_compiler import TaskSpec
            from agent.execution_contract import compile_execution_contract
            spec = TaskSpec(skill="ppt.source_grounded_build", verification=("ppt_structural", "ppt_render", "ppt_visual"))
            h.state.execution_contract = compile_execution_contract(spec, True, None)

            def fake_render(harness, path, output_dir):
                harness.state.record_evidence("ppt_render", "pass")

            def fake_inspect(harness, output_dir):
                harness.state.record_evidence("ppt_visual", "pass")

            with patch("agent.tools.ppt_tools._render", side_effect=fake_render), patch("agent.tools.ppt_tools._inspect_rendered", side_effect=fake_inspect):
                _finish(h, "done")
            self.assertEqual({record.kind for record in h.state.fresh_evidence()}, {"ppt_structural", "ppt_render", "ppt_visual"})

    def test_save_is_commit_not_content_mutation(self):
        # Saving the contract output must satisfy the finish gate's changed_files
        # check without advancing the mutation epoch (and thus invalidating
        # fresh evidence produced before the save).
        with tempfile.TemporaryDirectory() as tmp, patch.dict(os.environ, {"WORKSPACE": tmp}):
            h = DummyHarness()
            h.state.facts["required_output_pptx"] = "tasks/demo/output/final.pptx"
            dispatch("new_deck", json.dumps({"title": "T"}), h)
            epoch_after_compose = h.state.mutation_epoch
            dispatch("ppt_save", json.dumps({}), h)
            self.assertEqual(h.state.mutation_epoch, epoch_after_compose)
            self.assertTrue(any(str(p).endswith(".pptx") for p in h.state.changed_files))

    def test_geometry_accepts_h_and_width_aliases(self):
        # The facade historically mixed `w` with `height`; models naturally pair
        # `w` with `h`. Both spellings must normalize to x/y/w/height.
        from agent.tools.ppt_tools import _geometry_args
        self.assertEqual(_geometry_args({"x": 1, "y": 2, "w": 3, "h": 4}), (1, 2, 3, 4))
        self.assertEqual(_geometry_args({"x": 1, "y": 2, "width": 3, "height": 4}), (1, 2, 3, 4))
        self.assertEqual(_geometry_args({"x": 1, "y": 2, "w": 3, "height": 4}), (1, 2, 3, 4))

    def test_finish_rejects_unsaved_in_memory_ppt_even_when_structure_passes(self):
        h = DummyHarness()
        h._done = None
        h.recorder = type("Recorder", (), {
            "manifest": {"artifacts": []},
            "finish": lambda self, *args: None,
        })()
        h.state.record_change("deck:compose_quadrant_slide")
        h.state.record_evidence("ppt_structural", "pass")
        with self.assertRaisesRegex(ValueError, "no saved final-pptx"):
            _finish(h, "done")

    def test_minimal_text_only_template_is_normalized_to_standard_container(self):
        source = Presentation()
        # Exercise the one-layout branch with a facade; actual stripped
        # benchmark templates are covered by the real acceptance test.
        slide = source.slides.add_slide(source.slide_layouts[0])
        slide.shapes.add_textbox(0, 0, 1000000, 300000).text = "Portable"
        with patch.object(type(source.slide_layouts), "__len__", return_value=1):
            normalized, changed = _normalize_minimal_container(source)
        self.assertTrue(changed)
        self.assertGreater(len(normalized.slide_layouts), 1)
        self.assertIn("Portable", " ".join(shape.text for shape in normalized.slides[0].shapes if getattr(shape, "has_text_frame", False)))

    def test_verify_checks_non_text_shapes_for_boundary_crossing(self):
        h = DummyHarness()
        dispatch("new_deck", json.dumps({"title": "Boundary"}), h)
        slide = h.deck.slides[0]
        shape = slide.shapes.add_shape(1, 0, 0, 100000, 100000)
        shape.left = h.deck.slide_width - 100000
        shape.width = 200000
        report = dispatch("ppt_verify", "{}", h)
        self.assertIn("crosses slide boundary", report)

    def test_scoped_auto_check_keeps_unrelated_source_defect_as_warning(self):
        with tempfile.TemporaryDirectory() as tmp, patch.dict(os.environ, {"WORKSPACE": tmp}):
            source = Presentation()
            for index in range(19):
                slide = source.slides.add_slide(source.slide_layouts[6])
                if index == 1:
                    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(0.8)).text = "Before"
            historical = source.slides[18].shapes.add_shape(1, 0, 0, 100000, 100000)
            historical.left = source.slide_width - 100000
            historical.width = 200000
            source.save(Path(tmp) / "source.pptx")

            h = DummyHarness()
            dispatch("ppt_open", json.dumps({"path": "source.pptx"}), h)
            dispatch("ppt_edit_text", json.dumps({
                "operation": "replace", "slide_number": 2, "old": "Before", "new": "After",
            }), h)
            report = dispatch("ppt_check", json.dumps({"policy": "auto"}), h)

            self.assertIn("no blocking structural issues", report)
            self.assertIn("Historical warnings retained: 1", report)
            self.assertEqual(h.state.ppt_affected_slides, {2})
            self.assertNotIn("ppt_structural", h.state.unresolved_checks)
            self.assertIn("slide 19", h.state.facts["ppt_structural_warnings"])

            full_report = dispatch("ppt_check", json.dumps({"policy": "full"}), h)
            self.assertIn("no blocking structural issues", full_report)
            self.assertNotIn("ppt_structural", h.state.unresolved_checks)
            self.assertIn("baseline-delta", h.state.facts["ppt_full_check_downgraded"])

    def test_scoped_auto_check_blocks_new_defect_on_affected_slide(self):
        with tempfile.TemporaryDirectory() as tmp, patch.dict(os.environ, {"WORKSPACE": tmp}):
            source = Presentation()
            for index in range(19):
                slide = source.slides.add_slide(source.slide_layouts[6])
                if index == 1:
                    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2.0), Inches(0.25)).text = "Short"
            historical = source.slides[18].shapes.add_shape(1, 0, 0, 100000, 100000)
            historical.left = source.slide_width - 100000
            historical.width = 200000
            source.save(Path(tmp) / "source.pptx")

            h = DummyHarness()
            dispatch("ppt_open", json.dumps({"path": "source.pptx"}), h)
            dispatch("ppt_edit_text", json.dumps({
                "operation": "replace", "slide_number": 2, "old": "Short",
                "new": "This newly inserted sentence is deliberately much too long for the tiny text box " * 4,
            }), h)
            report = dispatch("ppt_check", json.dumps({"policy": "auto"}), h)

            self.assertIn("slide 2", report)
            self.assertIn("overflow risk", report)
            self.assertIn("ppt_structural", h.state.unresolved_checks)
            self.assertFalse(any(record.kind == "ppt_structural" for record in h.state.fresh_evidence()))

    def test_workflow_pipeline_slide_composes_and_verifies(self):
        h = DummyHarness()
        dispatch("new_deck", json.dumps({"title": "AI Workflow"}), h)
        steps = [
            {"title": "Goal Parsing", "action": "Deconstruct Intent", "bullets": ["Natural Language Understanding", "Task DAG Generation"], "tag": "NLU"},
            {"title": "Tool Selection", "action": "Dynamic Routing", "bullets": ["MCP Tool Registry", "Permission Gating"], "tag": "MCP"},
            {"title": "Execution Loop", "action": "ReAct Deliberation", "bullets": ["Atomic Transactions", "Self-Correction"], "tag": "Engine"},
            {"title": "Verification", "action": "Evidence Audit", "bullets": ["Deterministic Checks", "Rollback Boundary"], "tag": "CEGAR-H"},
        ]
        result = dispatch("ppt_compose", json.dumps({
            "kind": "workflow_pipeline",
            "slide_number": 1,
            "title": "Autonomous Agent Pipeline",
            "steps": steps,
            "takeaway": "Evidence-driven loop guarantees deterministic convergence.",
        }), h)
        self.assertIn("workflow pipeline slide 1", result)
        check = dispatch("ppt_check", json.dumps({"policy": "full"}), h)
        self.assertIn("no structural issues", check)
        self.assertEqual(len(h.deck.slides), 1)

    def test_html_mockup_slide_composes_and_verifies(self):
        h = DummyHarness()
        dispatch("new_deck", json.dumps({"title": "App Console"}), h)
        cards = [
            {"title": "Planner", "status": "ACTIVE", "metric": "99.8% Precision", "bullets": ["Multi-turn DAG", "Dependency Resolver"], "html_anchor": "div.planner"},
            {"title": "Memory Buffer", "status": "SYNCED", "metric": "128k Context", "bullets": ["Micro-compaction", "Episodic Recall"], "html_anchor": "div.memory"},
            {"title": "Tool Registry", "status": "READY", "metric": "42 Capabilities", "bullets": ["Schema Validation", "ACID Rollback"], "html_anchor": "div.tools"},
            {"title": "Runtime Loop", "status": "RUNNING", "metric": "<120ms Latency", "bullets": ["Event Bus Stream", "State Journal"], "html_anchor": "div.runtime"},
        ]
        result = dispatch("ppt_compose", json.dumps({
            "kind": "html_mockup",
            "title": "Agent System Workbench",
            "subtitle": "Real-time Component Status & Telemetry",
            "cards": cards,
            "url_bar": "https://console.agent.internal/dashboard",
        }), h)
        self.assertIn("HTML mockup slide 2", result)
        check = dispatch("ppt_check", json.dumps({"policy": "full"}), h)
        self.assertIn("no structural issues", check)
        self.assertEqual(len(h.deck.slides), 2)

    def test_hero_split_slide_composes_and_verifies(self):
        h = DummyHarness()
        dispatch("new_deck", json.dumps({"title": "Architecture Focus"}), h)
        cards = [
            {"title": "Safety Sandbox", "detail": "Containerized isolation with strict path jail and stripped credentials.", "bullets": ["Dockerized Runner", "Subprocess Isolation"]},
            {"title": "Freshness Epoch", "detail": "Every mutating action invalidates prior verification certificates.", "bullets": ["Monotonic Counter", "Anti-Stale Gating"]},
        ]
        result = dispatch("ppt_compose", json.dumps({
            "kind": "hero_split",
            "slide_number": 1,
            "title": "Xiaopu Core Invariants",
            "hero_title": "100% Deterministic Guarantee",
            "hero_metric": "0 False Passes",
            "hero_text": "By binding task completion to fresh, un-mutated certificates, Xiaopu eliminates speculative hallucination.",
            "cards": cards,
        }), h)
        self.assertIn("hero split slide 1", result)
        check = dispatch("ppt_check", json.dumps({"policy": "full"}), h)
        self.assertIn("no structural issues", check)

    def test_requested_slide_count_gate_blocks_undercomplete_deck(self):
        h = DummyHarness()
        h.goal = "请制作一个两页的“AI Agent 工作流程” PPT"
        dispatch("new_deck", json.dumps({"title": "AI Workflow"}), h)
        dispatch("ppt_compose", json.dumps({
            "kind": "workflow_pipeline", "slide_number": 1, "title": "Pipeline",
            "steps": [{"title": "Step 1", "action": "Act", "bullets": ["b1"]}],
        }), h)
        check = dispatch("ppt_check", "{}", h)
        self.assertIn("deck only has 1 slide(s)", check)
        self.assertIn("ppt_structural", h.state.unresolved_checks)

        # Adding second slide satisfies the requirement
        dispatch("ppt_compose", json.dumps({
            "kind": "html_mockup", "slide_number": 2, "title": "Dashboard",
            "cards": [{"title": "Card 1", "metric": "OK", "bullets": ["c1"]}],
        }), h)
        check_pass = dispatch("ppt_check", "{}", h)
        self.assertIn("no structural issues found", check_pass)
        self.assertNotIn("ppt_structural", h.state.unresolved_checks)

    def test_mutation_registry_covers_all_stateful_ppt_operations(self):
        expected = {
            "new_deck", "add_slide", "add_two_column_slide", "compose_quadrant_slide", "add_metric_slide",
            "add_table_slide", "add_process_slide", "add_image_slide", "add_textbox",
            "replace_shape_text", "set_shape_geometry", "delete_shape", "delete_slide",
            "replace_text", "append_bullet", "set_text_style", "set_shape_fill", "add_textbox_to_slide", "add_flowchart",
            "move_slide", "set_speaker_notes", "save_deck",
            "ppt_edit_text", "ppt_style", "ppt_metadata", "ppt_notes", "ppt_compose", "ppt_arrange", "ppt_save",
        }
        self.assertEqual(_PPT_MUTATORS, expected)



if __name__ == "__main__":
    unittest.main()
