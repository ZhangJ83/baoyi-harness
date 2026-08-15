import json
import random
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from benchmarks.lock_pptbench_review_form import lock
from benchmarks.prepare_pptbench_blind_review import prepare


def _fixture(tmp_path: Path, leaked: bool = False):
    protocol = {
        "systems":["sys_a","sys_b"],
        "blind_review":{"dimensions":["content_fidelity","layout"]},
        "tasks":[
            {"id":"t1","audience":"leaders","deck_job":"decide","prompt":"make a decision deck","facts":["A"],"required_text":["A"],"min_slides":1,"max_slides":1},
            {"id":"t2","audience":"staff","deck_job":"learn","prompt":"make a training deck","facts":["B"],"required_text":["B"],"min_slides":1,"max_slides":1},
        ],
    }
    protocol_path = tmp_path / "protocol.json"
    protocol_path.write_text(json.dumps(protocol), encoding="utf-8")
    raw = tmp_path / "raw"
    for system in protocol["systems"]:
        for task in protocol["tasks"]:
            folder = raw / system / task["id"]
            (folder / "slides").mkdir(parents=True)
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            box.text = "Xiaopu" if leaked and system == "sys_a" and task["id"] == "t1" else task["facts"][0]
            prs.save(folder / "deck.pptx")
            image = Image.new("RGB", (100, 60), "white")
            image.save(folder / "montage.png")
            image.save(folder / "slides/slide-1.png")
    return protocol, protocol_path, raw


def test_blind_bundles_hide_systems_and_randomize_orders(tmp_path: Path):
    protocol, protocol_path, raw = _fixture(tmp_path)
    out = tmp_path / "blind"
    result = prepare(protocol_path, raw, out, random.Random(7))
    assert result["valid"] is True and result["n_artifacts"] == 4
    assert result["reviewer_orders_distinct"] is True
    private = json.loads((out / "private_mapping.json").read_text(encoding="utf-8"))
    assert {row["system"] for row in private["artifacts"]} == {"sys_a", "sys_b"}
    public = "\n".join(path.read_text(encoding="utf-8") for reviewer in out.glob("reviewer_*") for path in reviewer.rglob("*.json"))
    assert "sys_a" not in public and "sys_b" not in public


def test_visible_system_identity_refuses_blind_bundle(tmp_path: Path):
    _, protocol_path, raw = _fixture(tmp_path, leaked=True)
    with pytest.raises(ValueError, match="identity leak"):
        prepare(protocol_path, raw, tmp_path / "blind", random.Random(1))


def test_review_form_locker_requires_attestation_and_complete_scores():
    draft = {
        "reviewer_id":"R-A", "independent_from_generation":True,
        "reviewer_non_author":True, "conflicts_declared":True,
        "review_completed_at":"2026-08-11T06:00:00+08:00",
        "locked_before_adjudication":False,
        "scores":[{"anonymous_id":"D-1", "dimensions":{"content_fidelity":4,"layout":5}, "blocking_comment":""}],
    }
    result = lock(draft, ["content_fidelity","layout"], {"D-1"}, "a" * 64)
    assert result["locked_before_adjudication"] is True
    bad = dict(draft, independent_from_generation=False)
    with pytest.raises(ValueError, match="independent_from_generation"):
        lock(bad, ["content_fidelity","layout"], {"D-1"}, "a" * 64)
