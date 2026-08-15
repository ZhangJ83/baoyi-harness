from __future__ import annotations

import json
from pathlib import Path
from tempfile import TemporaryDirectory

import pytest
from pptx import Presentation
from pptx.util import Inches

from agent.action_transaction import ScopeViolation
from agent.lifecycle import RunRecorder
from agent.ppt_transaction_adapter import slide_fingerprints
from agent.state import RunState
from agent.tools.registry import dispatch


class HarnessStub:
    def __init__(self, deck: Presentation, recorder: RunRecorder):
        self.deck = deck
        self.recorder = recorder
        self.state = RunState()

    def cancel_requested(self) -> bool:
        return False


def _deck() -> Presentation:
    deck = Presentation()
    for text in ("Alpha draft", "Beta draft", "Gamma draft"):
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text = text
    return deck


def _harness(root: Path) -> HarnessStub:
    source = root / "source.pptx"
    _deck().save(source)
    recorder = RunRecorder("scoped edit", "fake", "test", root)
    harness = HarnessStub(Presentation(source), recorder)
    harness.deck_source_path = source
    harness.deck_working_path = recorder.working_copy(source)
    return harness


def _only_certificate(harness: HarnessStub) -> tuple[Path, dict]:
    paths = list((harness.recorder.work / "transactions").glob("*/immutability_certificate.json"))
    assert len(paths) == 1
    return paths[0], json.loads(paths[0].read_text(encoding="utf-8"))


def test_real_dispatch_writes_certificate_and_preserves_other_slides() -> None:
    with TemporaryDirectory() as tmp:
        harness = _harness(Path(tmp))
        before = slide_fingerprints(harness.deck)

        result = dispatch(
            "ppt_edit_text",
            json.dumps({
                "operation": "replace",
                "slide_number": 2,
                "old": "draft",
                "new": "final",
            }),
            harness,
        )

        after = slide_fingerprints(harness.deck)
        assert "replaced" in result
        assert after[1] == before[1]
        assert after[3] == before[3]
        assert after[2] != before[2]
        certificate_path, certificate = _only_certificate(harness)
        assert certificate_path.is_file()
        assert certificate["requested_slides"] == [2]
        assert certificate["changed_slides"] == [2]
        assert certificate["unexpected_changed_slides"] == []
        assert certificate["passed"] is True


def test_dispatch_permission_violation_has_zero_side_effects() -> None:
    with TemporaryDirectory() as tmp:
        harness = _harness(Path(tmp))
        harness.state.ppt_allowed_slides = {1}
        harness.state.unresolved_checks.add("ppt_structural")
        before = slide_fingerprints(harness.deck)
        epoch = harness.state.mutation_epoch
        attempts = harness.state.repair_attempts

        with pytest.raises(ScopeViolation):
            dispatch(
                "ppt_style",
                json.dumps({
                    "slide_number": 2,
                    "target": "text",
                    "text_contains": "Beta",
                    "bold": True,
                }),
                harness,
            )

        assert slide_fingerprints(harness.deck) == before
        assert harness.state.mutation_epoch == epoch
        assert harness.state.repair_attempts == attempts
        assert not (harness.recorder.work / "transactions").exists()


def test_failed_batch_dispatch_rolls_back_deck_state_and_working_copy() -> None:
    with TemporaryDirectory() as tmp:
        harness = _harness(Path(tmp))
        before = slide_fingerprints(harness.deck)
        working_before = Path(harness.deck_working_path).read_bytes()
        live_deck = harness.deck

        with pytest.raises(ValueError, match="text not found"):
            dispatch(
                "ppt_edit_text",
                json.dumps({
                    "operation": "batch_updates",
                    "updates": [
                        {"operation": "replace", "slide_number": 1, "old": "draft", "new": "final"},
                        {"operation": "replace", "slide_number": 2, "old": "missing", "new": "final"},
                    ],
                }),
                harness,
            )

        assert harness.deck is live_deck
        assert slide_fingerprints(harness.deck) == before
        assert Path(harness.deck_working_path).read_bytes() == working_before
        assert harness.state.mutation_epoch == 0
        assert harness.state.ppt_affected_slides == set()
        transaction_dirs = list((harness.recorder.work / "transactions").iterdir())
        assert len(transaction_dirs) == 1
        assert (transaction_dirs[0] / "checkpoint.pptx").is_file()
        assert not (transaction_dirs[0] / "immutability_certificate.json").exists()


def test_global_replace_freezes_scope_to_all_current_slides() -> None:
    with TemporaryDirectory() as tmp:
        harness = _harness(Path(tmp))
        dispatch(
            "ppt_edit_text",
            json.dumps({"operation": "replace", "old": "draft", "new": "final"}),
            harness,
        )
        _, certificate = _only_certificate(harness)
        assert certificate["requested_slides"] == [1, 2, 3]
        assert certificate["changed_slides"] == [1, 2, 3]


def test_slide_reordering_arrange_operation_is_not_transaction_wrapped() -> None:
    with TemporaryDirectory() as tmp:
        harness = _harness(Path(tmp))
        dispatch(
            "ppt_arrange",
            json.dumps({"operation": "move_slide", "slide_number": 1, "new_position": 3}),
            harness,
        )
        assert not (harness.recorder.work / "transactions").exists()

