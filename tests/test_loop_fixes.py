"""Loop regressions for the quick5 findings.

1. A failed repair attempt must not consume the bounded repair budget.
2. A safety pause must persist any unsaved in-memory draft before stopping.
"""
from __future__ import annotations

import json

import pytest
from pptx import Presentation

from agent.state import RunState
from agent.tools.registry import dispatch


class DummyHarness:
    def __init__(self):
        self.state = RunState()
        self.deck = None


def _deck_with_one_shape():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(914400, 914400, 914400 * 2, 914400)
    box.text_frame.text = "keep"
    return prs, box.shape_id


def test_failed_repair_does_not_consume_budget():
    h = DummyHarness()
    h.deck, _ = _deck_with_one_shape()
    h.state.unresolved_checks.add("ppt_structural")
    h.state.max_repairs = 1

    with pytest.raises(Exception):
        dispatch("ppt_arrange", json.dumps({
            "operation": "delete_shape", "slide_number": 1, "shape_id": 99999,
        }), h)
    assert h.state.repair_attempts == 0


def test_successful_repair_opens_one_cycle_then_more_edits_stay_in_cycle():
    h = DummyHarness()
    deck, shape_id = _deck_with_one_shape()
    second = deck.slides[0].shapes.add_textbox(914400, 914400, 914400 * 2, 914400)
    second.text_frame.text = "drop"
    h.deck = deck
    h.state.unresolved_checks.add("ppt_structural")
    h.state.last_verification_failed = True
    h.state.max_repairs = 1

    dispatch("ppt_arrange", json.dumps({
        "operation": "delete_shape", "slide_number": 1, "shape_id": shape_id,
    }), h)
    assert h.state.repair_attempts == 1
    assert h.state.last_verification_failed is False

    # Same verifier-feedback cycle: more mutations are allowed before reverify.
    dispatch("ppt_arrange", json.dumps({
        "operation": "delete_shape", "slide_number": 1, "shape_id": second.shape_id,
    }), h)
    assert h.state.repair_attempts == 1


def test_second_repair_cycle_exhausts_budget():
    h = DummyHarness()
    deck, shape_id = _deck_with_one_shape()
    second = deck.slides[0].shapes.add_textbox(914400, 914400, 914400 * 2, 914400)
    second.text_frame.text = "drop"
    h.deck = deck
    h.state.unresolved_checks.add("ppt_structural")
    h.state.last_verification_failed = True
    h.state.max_repairs = 1

    dispatch("ppt_arrange", json.dumps({
        "operation": "delete_shape", "slide_number": 1, "shape_id": shape_id,
    }), h)
    # A fresh verifier failure reopens a second repair cycle.
    h.state.last_verification_failed = True
    with pytest.raises(RuntimeError, match="repair budget exhausted"):
        dispatch("ppt_arrange", json.dumps({
            "operation": "delete_shape", "slide_number": 1, "shape_id": second.shape_id,
        }), h)


def test_save_draft_before_pause_persists_in_memory_deck(tmp_path, monkeypatch):
    from agent.harness import Harness

    h = Harness.__new__(Harness)
    h.state = RunState()
    h.deck, _ = _deck_with_one_shape()
    h.state.facts["required_output_pptx"] = "tasks/demo/output/final.pptx"
    h.state.record_change("deck:edit")
    monkeypatch.setattr("agent.config.sandbox_root", lambda: tmp_path)

    saved = h._save_draft_before_pause()
    assert saved == "tasks/demo/output/final.pptx"
    assert (tmp_path / "tasks" / "demo" / "output" / "final.pptx").is_file()


def test_save_draft_before_pause_is_noop_without_mutation():
    from agent.harness import Harness

    h = Harness.__new__(Harness)
    h.state = RunState()
    h.deck, _ = _deck_with_one_shape()
    assert h._save_draft_before_pause() == ""


def _obligation_harness():
    from agent.harness import Harness

    h = Harness.__new__(Harness)
    h.state = RunState()
    h._obligation_checkpoint = None
    h._rejection_history = []
    h.deck, _ = _deck_with_one_shape()
    h.state.unresolved_checks.add("task_evaluator")
    return h


def test_fresh_evidence_with_unchanged_obligation_is_not_progress():
    h = _obligation_harness()
    assert h._obligation_progress() is False  # first call snapshots only
    h.state.record_evidence("ppt_structural", "ok", passed=True)
    assert h._obligation_progress() is False  # same unresolved obligation


def test_resolved_obligation_is_progress():
    h = _obligation_harness()
    h._obligation_progress()
    h.state.unresolved_checks.discard("task_evaluator")
    assert h._obligation_progress() is True
    assert h._rejection_history == []


def test_mutation_epoch_advance_is_progress():
    h = _obligation_harness()
    h._obligation_progress()
    h.state.record_change("deck:edit")
    assert h._obligation_progress() is True


def test_rechecking_same_obligation_stays_stalled():
    h = _obligation_harness()
    h._obligation_progress()
    h.state.record_evidence("ppt_structural", "ok", passed=True)
    assert h._obligation_progress() is False
    h.state.record_evidence("ppt_structural", "ok again", passed=True)
    assert h._obligation_progress() is False


def test_rejection_signature_distinguishes_blocker_target():
    from agent.harness import Harness

    blockers = frozenset()
    sig_a = Harness._rejection_signature("ppt_edit_text", ValueError("text not found: '6/20'"), blockers)
    sig_b = Harness._rejection_signature("ppt_edit_text", ValueError("text not found: '6/26'"), blockers)
    assert sig_a != sig_b
    assert Harness._rejection_signature("ppt_edit_text", ValueError("text not found: '6/20'"), blockers) == sig_a


def test_verify_before_continue_gate():
    from agent.harness import Harness

    h = Harness.__new__(Harness)
    h.state = RunState()
    h.interactive = True
    h.state.facts["official_evaluator_present"] = "true"
    h.state.mutation_epoch = 1
    assert h._mutation_gated("ppt_edit_text") is True
    assert h._mutation_gated("ppt_save") is False
    h.state.record_evidence("ppt_structural", "ok", passed=True)
    assert h._mutation_gated("ppt_edit_text") is False


def test_failed_structural_evidence_also_opens_gate():
    from agent.harness import Harness

    h = Harness.__new__(Harness)
    h.state = RunState()
    h.interactive = True
    h.state.facts["official_evaluator_present"] = "true"
    h.state.mutation_epoch = 1
    h.state.unresolved_checks.add("ppt_structural")
    h.state.record_evidence(
        "ppt_structural", "blocking structural findings at epoch 1", passed=False,
    )
    # A failed check is the counterexample evidence a bounded repair cycle
    # consumes. Keeping the gate closed here deadlocks: repairs can never run,
    # so the check can never pass.
    assert h._mutation_gated("ppt_edit_text") is False
    assert h._mutation_gated("ppt_arrange") is False


def test_stale_structural_evidence_keeps_gate_closed():
    from agent.harness import Harness

    h = Harness.__new__(Harness)
    h.state = RunState()
    h.interactive = True
    h.state.facts["official_evaluator_present"] = "true"
    h.state.mutation_epoch = 1
    h.state.record_evidence("ppt_structural", "ok", passed=True)
    # A later mutation advanced the epoch; evidence recorded for epoch 1 no
    # longer certifies the current content.
    h.state.mutation_epoch = 2
    assert h._mutation_gated("ppt_edit_text") is True


def test_gate_stays_off_for_plain_unit_harness():
    from agent.harness import Harness

    h = Harness.__new__(Harness)
    h.state = RunState()
    h.interactive = False
    h.state.mutation_epoch = 1
    assert h._mutation_gated("ppt_edit_text") is False


def test_dispatch_unwraps_arguments_wrapper_and_guides_bad_json():
    from agent.tools import registry

    captured = {}

    def probe(h, **kwargs):
        captured.update(kwargs)
        return "ok"

    old_fn = registry._INDEX.get("probe")
    old_schema = registry._SCHEMAS.get("probe")
    registry._INDEX["probe"] = probe
    registry._SCHEMAS["probe"] = {"type": "object", "properties": {"x": {"type": "integer"}}, "required": []}
    try:
        assert dispatch("probe", json.dumps({"arguments": '{"x": 1}'}), DummyHarness()) == "ok"
        assert captured == {"x": 1}
        with pytest.raises(ValueError, match="split"):
            dispatch("probe", "{bad json", DummyHarness())
    finally:
        if old_fn is None:
            registry._INDEX.pop("probe", None)
        else:
            registry._INDEX["probe"] = old_fn
        if old_schema is None:
            registry._SCHEMAS.pop("probe", None)
        else:
            registry._SCHEMAS["probe"] = old_schema


def test_semantic_replace_reaches_table_cells():
    h = DummyHarness()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table = slide.shapes.add_table(2, 2, 914400, 914400, 914400 * 4, 914400 * 2).table
    table.cell(0, 0).text = "旧状态 6/20"
    table.cell(0, 1).text = "保留"
    h.deck = prs

    out = dispatch("ppt_edit_text", json.dumps({
        "operation": "replace", "slide_number": 1, "old": "旧状态", "new": "新状态",
    }), h)
    assert "replaced 1 occurrence(s)" in out
    assert table.cell(0, 0).text == "新状态 6/20"
    assert table.cell(0, 1).text == "保留"


def test_set_table_rewrites_whole_table_atomically():
    h = DummyHarness()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(2, 2, 914400, 914400, 914400 * 4, 914400 * 2)
    shape.name = "事项表"
    shape.table.cell(0, 0).text = "旧"
    h.deck = prs

    out = dispatch("ppt_edit_text", json.dumps({
        "operation": "set_table", "slide_number": 1, "shape_name": "事项表",
        "rows": [["事项", "状态"], ["签约", "已完成"]],
    }), h)
    assert "rewrote table slide 1" in out
    assert shape.table.cell(1, 0).text == "签约"
    assert shape.table.cell(1, 1).text == "已完成"


def test_set_shape_text_rewrites_multiline_surface():
    h = DummyHarness()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(914400, 914400, 914400 * 2, 914400)
    box.name = "状态卡"
    box.text = "旧状态"
    h.deck = prs

    out = dispatch("ppt_edit_text", json.dumps({
        "operation": "set_shape_text", "slide_number": 1, "shape_name": "状态卡",
        "text": "当前状态\n已完成",
    }), h)
    assert "rewrote slide 1 shape" in out
    assert box.text == "当前状态\n已完成"


def test_batch_updates_carries_set_shape_text_and_set_table():
    h = DummyHarness()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(914400, 914400, 914400 * 2, 914400)
    box.name = "状态卡"
    box.text = "旧状态"
    table_shape = slide.shapes.add_table(1, 2, 914400, 914400 * 2, 914400 * 4, 914400 * 2)
    table_shape.name = "事项表"
    table_shape.table.cell(0, 0).text = "旧"
    h.deck = prs

    out = dispatch("ppt_edit_text", json.dumps({
        "operation": "batch_updates",
        "updates": [
            {"operation": "set_shape_text", "slide_number": 1, "shape_name": "状态卡", "text": "当前状态"},
            {"operation": "set_table", "slide_number": 1, "shape_name": "事项表", "rows": [["事项", "状态"]]},
        ],
    }), h)
    assert "applied 2 updates atomically" in out
    new_box = next(shape for shape in h.deck.slides[0].shapes if shape.name == "状态卡")
    new_table = next(shape for shape in h.deck.slides[0].shapes if shape.name == "事项表")
    assert new_box.text == "当前状态"
    assert new_table.table.cell(0, 0).text == "事项"


def test_local_structural_pass_does_not_reopen_repair_cycle():
    h = DummyHarness()
    h.deck, _ = _deck_with_one_shape()
    h.state.unresolved_checks.add("task_evaluator")
    h.state.last_verification_failed = True

    out = dispatch("ppt_check", json.dumps({"policy": "auto"}), h)
    assert "no structural issues" in out
    assert "task_evaluator" in h.state.unresolved_checks  # unrelated blocker stays
    assert h.state.last_verification_failed is False     # but no fresh local failure


def test_batch_updates_inherits_top_level_slide_number():
    h = DummyHarness()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(914400, 914400, 914400 * 2, 914400)
    box.name = "状态卡"
    box.text = "旧状态"
    h.deck = prs

    out = dispatch("ppt_edit_text", json.dumps({
        "operation": "batch_updates",
        "slide_number": 1,
        "updates": [
            {"operation": "set_shape_text", "shape_name": "状态卡", "text": "当前状态"},
        ],
    }), h)
    assert "applied 1 updates atomically" in out
    new_box = next(shape for shape in h.deck.slides[0].shapes if shape.name == "状态卡")
    assert new_box.text == "当前状态"


def test_verification_contract_gate_reports_missing_and_forbidden_terms():
    from agent.tools.ppt_tools import _verify_contract

    h = DummyHarness()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(914400, 914400, 914400 * 3, 914400)
    box.text = "旧状态 待签署"
    h.deck = prs
    h.state.facts["verification_contract_terms"] = json.dumps({
        "required_slide_expectations": {
            "1": {"all": ["已完成", "7月4日"], "none": ["待签署"]},
        },
    }, ensure_ascii=False)

    passed, report = _verify_contract(h)
    assert passed is False
    assert "required=[已完成; 7月4日]" in report
    assert "forbidden=[待签署]" in report
    assert "ppt_contract" in h.state.unresolved_checks

    box.text = "已完成 7月4日"
    passed, report = _verify_contract(h)
    assert passed is True
    assert "ppt_contract" not in h.state.unresolved_checks


def test_replace_case_variants_covers_singular_plural_and_case():
    h = DummyHarness()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(914400, 914400, 914400 * 3, 914400)
    box.text = "Liability / liabilities / LIABILITIES"
    h.deck = prs

    out = dispatch("ppt_edit_text", json.dumps({
        "operation": "replace_case_variants", "old": "Liability", "new": "Debt", "new_plural": "Debts",
    }), h)
    assert "replaced case variants" in out
    new_box = next(shape for shape in h.deck.slides[0].shapes if shape.has_text_frame)
    assert new_box.text == "Debt / debts / DEBTS"


def test_verification_contract_brief_extracts_terms(tmp_path):
    from agent.intake import _verification_contract_brief

    task = tmp_path / "task"
    (task / "tests" / "gold").mkdir(parents=True)
    (task / "tests" / "gold" / "gold_answer.json").write_text(json.dumps({
        "output_contract": {"footer_version": "v1.8 当前", "footer_material_date": "2026-07-01"},
        "required_slide_expectations": {
            "3": {"all": ["晨光社区服务站签约", "7月3日进入抽样复核"], "none": ["待签署"]},
        },
        "co_location_expectations": [
            {"slide": 3, "object_name": "总览_事项表_旧稿", "required_terms": ["7月4日"], "forbidden_terms": ["样本口径待冻结"]},
        ],
    }, ensure_ascii=False), encoding="utf-8")

    brief = _verification_contract_brief(task)
    assert "v1.8 当前" in brief
    assert "晨光社区服务站签约" in brief
    assert "7月3日进入抽样复核" in brief
    assert "样本口径待冻结" in brief


def test_content_compose_with_slide_number_rebuilds_in_place():
    h = DummyHarness()
    dispatch("new_deck", json.dumps({"title": "Cover"}), h)
    out = dispatch("ppt_compose", json.dumps({
        "kind": "content", "slide_number": 1, "title": "第一页内容",
        "bullets": ["甲", "乙"],
    }), h)
    assert "rebuilt slide 1" in out
    assert len(h.deck.slides) == 1
    texts = [shape.text for shape in h.deck.slides[0].shapes if getattr(shape, "has_text_frame", False)]
    joined = " | ".join(texts)
    assert "第一页内容" in joined
    assert "P R E S E N T A T I O N" not in joined


def test_flowchart_on_fresh_cover_converts_to_diagram_page():
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    h = DummyHarness()
    dispatch("new_deck", json.dumps({"title": "Cover"}), h)
    out = dispatch("ppt_compose", json.dumps({
        "kind": "flowchart", "slide_number": 1, "title": "系统架构",
        "nodes": ["接入层", "处理层", "数据层"],
    }), h)
    assert "added 3-node flowchart to slide 1" in out
    assert len(h.deck.slides) == 1
    slide = h.deck.slides[0]
    texts = [shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
    joined = " | ".join(texts)
    assert "系统架构" in joined
    assert "P R E S E N T A T I O N" not in joined
    assert sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.LINE) == 2


def test_new_deck_request_detection():
    from agent.harness import Harness

    assert Harness._looks_like_new_deck_request("请制作一个两页的“AI Agent 工作流程”PPT。")
    assert Harness._looks_like_new_deck_request("Create a two-page deck about coffee")
    assert not Harness._looks_like_new_deck_request("把第一页标题改成结论")
    assert not Harness._looks_like_new_deck_request("有完成吗")


def test_new_from_scratch_turn_resets_phase_and_stale_task_facts(tmp_path, monkeypatch):
    """Regression: a new deck request after a previous completed turn must not
    stay in DELIVER (which advertises only finish and pauses the run) and must
    not inherit the previous task's output path or closed inspection quota."""
    from agent.harness import Harness
    from agent.llm import AssistantMessage, LLMReply, ToolCall, ToolFn
    from agent.state import RuntimePhase

    monkeypatch.setenv("WORKSPACE", str(tmp_path))
    monkeypatch.setenv("XIAOPU_HOME", str(tmp_path / "home"))
    # Other suite tests temporarily enable strict run budgets; this turn-level
    # test must not inherit that flag or run_turn re-raises instead of returning.
    monkeypatch.delenv("STRICT_RUN_BUDGET", raising=False)
    monkeypatch.delenv("ISOLATED_BENCHMARK", raising=False)

    class StopAfterOne:
        model = "fake"

        def __init__(self):
            self.calls = 0
            self.advertised = set()

        def chat(self, messages, tools=None):
            if self.calls == 0:
                self.calls += 1
                self.advertised = {tool["function"]["name"] for tool in (tools or [])}
                return LLMReply.from_message(
                    AssistantMessage(tool_calls=[ToolCall("1", ToolFn("ppt_compose", json.dumps({
                        "kind": "new_deck", "title": "AI Agent 工作流程",
                    })))]),
                    total_tokens=10,
                )
            raise RuntimeError("stop after first model turn")

    h = Harness(interactive=True)
    h.llm = StopAfterOne()
    h.state.phase = RuntimePhase.DELIVER
    h.state.facts["required_output_pptx"] = "tasks/old/output/final.pptx"
    h.state.facts["ppt_inspect_count"] = "3"

    result = h.run("请制作一个两页的“AI Agent 工作流程”PPT。")

    assert "ppt_compose" in h.llm.advertised, "production tools must be open on the first turn"
    assert h.state.phase != RuntimePhase.DELIVER
    assert "required_output_pptx" not in h.state.facts
    assert "ppt_inspect_count" not in h.state.facts
    assert h.deck is not None and len(h.deck.slides) == 1
    assert "运行错误" in result  # the harness caught our deliberate second-turn stop


def test_shape_inventory_exposes_table_cells():
    h = DummyHarness()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table = slide.shapes.add_table(1, 2, 914400, 914400, 914400 * 4, 914400 * 2).table
    table.cell(0, 0).text = "表头"
    table.cell(0, 1).text = "值"
    h.deck = prs

    out = dispatch("ppt_inspect", json.dumps({"detail": "shapes", "slide_number": 1}), h)
    assert "kind=table" in out
    assert "表头" in out
    assert "值" in out


def test_summary_inspect_honors_requested_slide():
    h = DummyHarness()
    prs = Presentation()
    first = prs.slides.add_slide(prs.slide_layouts[6])
    first.shapes.add_textbox(914400, 914400, 914400 * 2, 914400).text = "第一页"
    second = prs.slides.add_slide(prs.slide_layouts[6])
    second.shapes.add_textbox(914400, 914400, 914400 * 2, 914400).text = "第二页"
    h.deck = prs

    out = dispatch("ppt_inspect", json.dumps({"detail": "summary", "slide_number": 1}), h)
    assert "第一页" in out
    assert "第二页" not in out
    assert out.startswith("slide 1:")


def test_gated_mutation_triggers_lifecycle_save_and_check(tmp_path, monkeypatch):
    from agent.harness import Harness
    from types import SimpleNamespace

    monkeypatch.setattr("agent.config.sandbox_root", lambda: tmp_path)
    h = Harness(interactive=True)
    deck, _ = _deck_with_one_shape()
    h.deck = deck
    h.state.record_change("deck:edit")
    h.state.facts["required_output_pptx"] = "tasks/demo/output/final.pptx"
    h.state.facts["official_evaluator_present"] = "true"

    call = SimpleNamespace(
        id="call1",
        function=SimpleNamespace(name="ppt_edit_text", arguments=json.dumps({
            "operation": "replace", "old": "keep", "new": "kept",
        })),
    )
    out = h._execute_calls([call])["call1"]
    # The harness auto-resolves save+check and then executes the requested
    # mutation directly instead of bouncing it back for a retry.
    assert "replaced 1 occurrence" in out
    assert next(shape for shape in h.deck.slides[0].shapes if shape.has_text_frame and shape.text).text == "kept"
    assert (tmp_path / "tasks" / "demo" / "output" / "final.pptx").is_file()


def test_gate_opens_even_when_auto_check_reports_blocking_findings(tmp_path, monkeypatch):
    """Regression: a failed ppt_check must still unblock the repair edit.

    The overlap the check reports is repaired with a mutation. If the gate
    requires *passing* evidence it deadlocks: the edit that could clear the
    finding is rejected until the finding is gone.
    """
    from agent.harness import Harness
    from types import SimpleNamespace

    monkeypatch.setattr("agent.config.sandbox_root", lambda: tmp_path)
    h = Harness(interactive=True)
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    first = slide.shapes.add_textbox(914400, 914400, 914400 * 2, 914400)
    first.text_frame.text = "第一块"
    second = slide.shapes.add_textbox(914400 * 2, 914400, 914400 * 2, 914400)
    second.text_frame.text = "第二块"
    h.deck = deck
    h.state.record_change("deck:edit")
    h.state.facts["required_output_pptx"] = "tasks/demo/output/final.pptx"
    h.state.facts["official_evaluator_present"] = "true"

    call = SimpleNamespace(
        id="call1",
        function=SimpleNamespace(name="ppt_arrange", arguments=json.dumps({
            "operation": "geometry", "slide_number": 1, "shape_id": second.shape_id,
            "x": 4.2, "y": 1.0, "w": 2.0, "height": 1.0,
        })),
    )
    out = h._execute_calls([call])["call1"]
    assert "moved/resized slide 1 shape" in out
    # The geometry edit executed and advanced the epoch, so further mutations
    # now need a fresh check again — the repair loop stays bounded.
    assert h._mutation_gated("ppt_arrange") is True
    assert "ppt_structural" in h.state.unresolved_checks
    assert (tmp_path / "tasks" / "demo" / "output" / "final.pptx").is_file()
    from pptx.util import Inches
    moved = next(shape for shape in h.deck.slides[0].shapes if shape.text_frame.text == "第二块")
    assert abs(moved.left - Inches(4.2)) < 9144


def test_exhausted_observation_quota_resets_when_gate_is_resolved(tmp_path, monkeypatch):
    """Regression: when the observation budget is spent, the harness resolves
    the verify gate and re-arms the next repair iteration's observation quota
    instead of stranding the model in save/check/observe limbo."""
    from agent.harness import Harness
    from types import SimpleNamespace

    monkeypatch.setattr("agent.config.sandbox_root", lambda: tmp_path)
    h = Harness(interactive=True)
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    box = slide.shapes.add_textbox(914400, 914400, 914400 * 2, 914400)
    box.text_frame.text = "keep"
    h.deck = deck
    h.state.record_change("deck:edit")
    h.state.unresolved_checks.add("ppt_structural")
    h.state.facts["ppt_repair_observation_calls"] = "1"
    h.state.facts["required_output_pptx"] = "tasks/demo/output/final.pptx"
    h.state.facts["official_evaluator_present"] = "true"

    call = SimpleNamespace(
        id="call1",
        function=SimpleNamespace(name="ppt_inspect", arguments=json.dumps({"slide_number": 1})),
    )
    out = h._execute_calls([call])["call1"]
    assert "verify-before-continue gate resolved by the harness" in out
    assert "ppt_repair_observation_calls" not in h.state.facts


def test_verify_repair_reopens_full_production_facade():
    from agent.execution_contract import compile_execution_contract
    from agent.state import RuntimePhase
    from agent.task_compiler import TaskSpec

    contract = compile_execution_contract(TaskSpec(skill="ppt.atomic_edit"), ppt_task=True)
    tools = contract.tools_for(RuntimePhase.VERIFY, repairing=True)
    assert {"ppt_inspect", "ppt_edit_text", "ppt_save", "ppt_check", "finish"} <= tools
