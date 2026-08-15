"""Tests for the runnable Trajectory-Guided Domain Specialization pipeline (non-loop)."""
from __future__ import annotations

import json
from pathlib import Path

from research.trajectory_method import (
    build_static_contract,
    claim_gate,
    extract_static_requirements,
    induce,
    load_runs,
)


def _run(root: Path, agent: str, task_id: str, lines: list[dict]) -> None:
    traj = root / "agent_workspaces" / "full13" / agent / "tasks" / task_id / "trajectory"
    traj.mkdir(parents=True)
    (traj / "steps.jsonl").write_text(
        "\n".join(json.dumps(line, ensure_ascii=False) for line in lines), encoding="utf-8"
    )


def test_pipeline_reports_prevalence_not_only_frequency(tmp_path):
    _run(tmp_path, "alpha", "t1", [{"step": 1, "action": "read"}, {"step": 2, "action": "render"}, {"step": 3, "action": "stop"}])
    _run(tmp_path, "beta", "t1", [{"step": 1, "action": "read"}, {"step": 2, "action": "verify"}])
    _run(tmp_path, "beta", "t2", [{"step": 1, "action": "repair"}, {"step": 2, "action": "stop"}])

    runs = load_runs(tmp_path)
    assert len(runs) == 3
    stats = {s.behavior: s for s in induce(runs, all_tasks=["t1", "t2"])}
    assert stats["read"].frequency == 2
    assert stats["read"].run_prevalence == 2 / 3
    assert stats["read"].task_prevalence == 0.5
    assert stats["render"].run_prevalence == 1 / 3
    assert stats["repair"].can_fix_before_loop is False
    assert stats["read"].can_fix_before_loop is True


def test_claim_gate_is_orthogonal_to_design_layer():
    assert claim_gate("render") == "B"
    assert claim_gate("read") == "A"
    assert claim_gate("plan") == "B"


def test_static_requirement_extraction_drops_loop_behaviors(tmp_path):
    _run(tmp_path, "alpha", "t1", [{"step": 1, "action": "read"}, {"step": 2, "action": "repair"}, {"step": 3, "action": "stop"}])
    stats = induce(load_runs(tmp_path), all_tasks=["t1"])
    rows = extract_static_requirements(stats)
    slots = {r["behavior"]: r["slot"] for r in rows}
    assert slots["read"] == ["I"]
    assert slots["plan"] == ["T"]
    assert "repair" not in slots
    assert "detect_defect" not in slots


def test_build_static_contract_is_the_landed_c_static(tmp_path):
    from pptx import Presentation

    task = tmp_path / "tasks" / "demo"
    (task / "input").mkdir(parents=True)
    (task / "instruction.md").write_text("convert the checklist into a two-column layout", encoding="utf-8")
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(task / "input" / "deck.pptx"))

    contract = build_static_contract(task)
    assert contract.task_type == "layout_reflow"
    assert contract.verification.required_kinds() == ("structural", "render", "visual")
    assert contract.mutation is not None
    assert "content_text" in contract.mutation.deny
    assert contract.output.path.name == "final.pptx"
