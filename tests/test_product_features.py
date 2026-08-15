"""Tests for the final product features: streaming, approvals, undo, themes."""
from __future__ import annotations

import io
import json
import os
from pathlib import Path
from unittest.mock import patch

import pytest
from pptx import Presentation

from agent import config
from agent.state import RunState


def test_openai_chat_streams_tokens():
    from agent.llm import LLM

    class ChunkChoice:
        def __init__(self, content=None, reasoning=None):
            self.delta = type("Delta", (), {"content": content, "reasoning_content": reasoning})()

    chunks = [
        type("C", (), {"choices": [ChunkChoice("Hello")], "usage": None})(),
        type("C", (), {"choices": [ChunkChoice(" world")], "usage": None})(),
        type("C", (), {"choices": [], "usage": type("U", (), {"prompt_tokens": 2, "completion_tokens": 2, "total_tokens": 4})()})(),
    ]

    class Completions:
        def create(self, **kwargs):
            assert kwargs["stream"] is True
            return iter(chunks)

    llm = LLM.__new__(LLM)
    llm.model = "deepseek-v4"
    llm._cancelled = type("E", (), {"set": lambda self: None, "clear": lambda self: None, "is_set": lambda self: False})()
    llm._client = type("Client", (), {"chat": type("Chat", (), {"completions": Completions()})()})()
    seen = []
    reply = llm.chat([{"role": "user", "content": "hi"}], stream=True, on_token=seen.append)
    assert "".join(seen) == "Hello world"
    assert reply.total_tokens == 4


def test_approval_handler_allows_and_denies_shell(tmp_path, monkeypatch):
    from agent.tools import fs_tools

    class H:
        def __init__(self, decision):
            self.state = RunState()
            self.approval_handler = lambda command: decision

    monkeypatch.setattr(fs_tools.config, "command_policy", lambda: "ask")
    monkeypatch.setattr(fs_tools, "_root", lambda: tmp_path)
    allowed = fs_tools._run(H("allow"), "echo hello")
    assert "exit_code=0" in allowed
    denied = fs_tools._run(H("deny"), "echo hello")
    assert denied.startswith("PERMISSION DENY")


def test_approval_without_handler_keeps_legacy_ask_message(tmp_path, monkeypatch):
    from agent.tools import fs_tools

    class H:
        def __init__(self):
            self.state = RunState()

    monkeypatch.setattr(fs_tools.config, "command_policy", lambda: "ask")
    monkeypatch.setattr(fs_tools, "_root", lambda: tmp_path)
    result = fs_tools._run(H(), "echo hello")
    assert result.startswith("PERMISSION ASK")


def test_harness_undo_restores_previous_deck(tmp_path):
    from agent.harness import Harness
    from agent.tools.registry import dispatch

    h = Harness.__new__(Harness)
    h.state = RunState()
    h.undo_stack = []
    h.deck = Presentation()
    slide = h.deck.slides.add_slide(h.deck.slide_layouts[6])
    box = slide.shapes.add_textbox(914400, 914400, 914400 * 2, 914400)
    box.text_frame.text = "before"

    dispatch("ppt_edit_text", json.dumps({
        "operation": "replace", "slide_number": 1, "shape_id": box.shape_id,
        "old": "before", "new": "after",
    }), h)
    assert box.text_frame.text == "after"
    assert len(h.undo_stack) == 1

    result = h.undo()
    assert "已撤销" in result
    restored = h.deck.slides[0].shapes[0].text_frame.text
    assert restored == "before"


def test_theme_and_keymap_config_roundtrip(monkeypatch):
    monkeypatch.setenv("XIAOPU_THEME", "dracula")
    assert config.theme() == "dracula"
    config.set_theme("light")
    assert config.theme() == "light"
    monkeypatch.setenv("XIAOPU_KEYMAP", "minimal")
    assert config.keymap() == "minimal"
    config.set_keymap("default")
    assert config.keymap() == "default"
