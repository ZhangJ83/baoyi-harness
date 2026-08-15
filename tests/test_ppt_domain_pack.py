"""PPT Domain Pack tests: ontology, IR, verification policy, skills, transactions."""
from __future__ import annotations

import json
from pathlib import Path

import pytest

BENCH_DECK = Path(r"E:\project\agent\ppt-harness\benchmark_v0.1\ppt-eval\data\files\PowerPoint\3.pptx")


def test_task_ontology_is_exactly_eight_types():
    from domains.ppt.task_types import PPTTaskType

    assert {t.value for t in PPTTaskType} == {
        "atomic_edit", "atomic_style", "element_creation", "layout_reflow",
        "diagram_composition", "compose_from_slides", "source_grounded_build",
        "template_build",
    }


@pytest.mark.parametrize("instruction,expected", [
    ("Change the lecture number from Lecture 3 to Lecture 1", "atomic_edit"),
    ("Change the presentation title font size to 48pt", "atomic_style"),
    ("Create a text box on slide 2 with the phrase Important Note", "element_creation"),
    ("Convert the entire checklist on slide 8 into a two-column layout", "layout_reflow"),
    ("Create a flowchart using rounded rectangles and line arrows", "diagram_composition"),
    ("Combine slides 2 and 3 into a single slide with a table", "compose_from_slides"),
    ("基于 HTML/XLSX/JSON 生成一页四象限汇报", "source_grounded_build"),
    ("基于模板把思维导图内容改造成演示稿", "template_build"),
    # legacy frozen-benchmark labels map onto canonical types
    ("直接更新这份治理会 PPT，使用 XLSX 作为更新来源", "source_grounded_build"),
    ("Add a bullet point and resize the image to avoid overlap", "layout_reflow"),
])
def test_classify_ppt_task(instruction, expected):
    from domains.ppt.task_types import classify_ppt_task

    assert classify_ppt_task(instruction).value == expected


@pytest.mark.skipif(not BENCH_DECK.exists(), reason="benchmark deck not present")
def test_presentation_ir_from_real_deck():
    from domains.ppt.ir import from_pptx

    ir = from_pptx(BENCH_DECK)
    assert len(ir.slides) == 45
    joined = "\n".join(s.all_text() for s in ir.slides)
    assert "Dual Mode" in joined
    assert "shape" not in ir.summary().lower() or ir.summary()


def test_verification_policy_matches_design_examples():
    from domains.ppt.verification import verification_policy

    assert verification_policy("atomic_edit").required_kinds() == ("structural", "immutability")
    assert verification_policy("layout_reflow").required_kinds() == ("structural", "render", "visual")
    assert verification_policy("source_grounded_build").required_kinds() == (
        "content_grounding", "structural", "render", "visual")


def test_skills_registered_with_consistent_contracts():
    import domains.ppt  # noqa: F401
    from core.skill import get_skill
    from domains.ppt.skills import PPT_SKILLS
    from domains.ppt.verification import verification_policy

    assert set(PPT_SKILLS) == {
        "atomic_edit", "atomic_style", "element_creation", "layout_reflow",
        "diagram_composition", "compose_from_slides", "source_grounded_build",
        "template_build",
    }
    for name, spec in PPT_SKILLS.items():
        assert get_skill(name) is spec
        assert spec.capabilities
        assert spec.verification_contract == verification_policy(name).required_kinds()


def test_ppt_transaction_types():
    from domains.ppt.transaction import (
        PptAttributeChange,
        PptDelta,
        PptImmutabilityCertificate,
        PptMutationScope,
    )

    scope = PptMutationScope(slides={2, 3}, shapes={7})
    assert scope.fields == ("slides", "shapes", "properties")
    change = PptAttributeChange(slide=3, shape_id=7, property="fill",
                                before="FF66CC", after="00FF00")
    delta = PptDelta(added_slides=[4], changed_shapes=[7],
                     changed_properties=("fill",), attribute_changes=(change,))
    assert "added_slides=[4]" in delta.summarize()
    assert "slide=3 shape=7 fill" in delta.attribute_summary()
    cert = PptImmutabilityCertificate(artifact_ref="a", epoch=3, changed_slides=[1])
    assert cert.is_fresh(3)
    assert not cert.is_fresh(4)


def test_attribute_level_delta_detects_text_fill_position(tmp_path):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from domains.ppt.transaction import diff_decks

    def make(path, text, rgb, left):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(left, 914400, 914400 * 2, 914400)
        box.text_frame.text = text
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor.from_string(rgb)
        prs.save(str(path))
        return box.shape_id

    baseline = tmp_path / "base.pptx"
    modified = tmp_path / "mod.pptx"
    make(baseline, "hello", "FF66CC", 914400)
    shape_id = make(modified, "hello world", "00FF00", 914400 * 2)

    delta = diff_decks(baseline, modified)
    props = {c.property for c in delta.attribute_changes}
    assert {"text", "fill", "left"} <= props
    assert shape_id in delta.changed_shapes
    assert all(c.shape_id == shape_id for c in delta.attribute_changes)


def test_immutability_and_delta_subset_check(tmp_path):
    from pptx import Presentation

    from core.transaction import AllowedMutation, ImmutabilityPolicy
    from domains.ppt.transaction import (
        PptMutationScope,
        delta_within_mutation,
        diff_decks,
        verify_immutability,
    )

    def make(path, text):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(914400, 914400, 914400 * 2, 914400)
        box.text_frame.text = text
        prs.save(str(path))

    baseline = tmp_path / "base.pptx"
    modified = tmp_path / "mod.pptx"
    make(baseline, "keep me")
    make(modified, "changed me")

    # Non-target slides must stay untouched: slide 1 is not allowed -> violation.
    cert = verify_immutability(baseline, modified, allowed_slide_indexes=set())
    assert cert.passed is False
    assert cert.violations

    # Same change is fine when slide 1 is the declared mutation target.
    cert_ok = verify_immutability(baseline, modified, allowed_slide_indexes={1})
    assert cert_ok.passed is True

    delta = diff_decks(baseline, modified)
    mutation = AllowedMutation(
        scope=PptMutationScope(slides={1}, properties=("text",)),
        policy=ImmutabilityPolicy(allow=("text",), deny=("fill", "theme")),
    )
    ok, reasons = delta_within_mutation(delta, mutation)
    assert ok is True, reasons

    bad_policy = ImmutabilityPolicy(allow=("fill",), deny=("text",))
    bad_mutation = AllowedMutation(scope=PptMutationScope(slides={1}), policy=bad_policy)
    ok2, reasons2 = delta_within_mutation(delta, bad_mutation)
    assert ok2 is False
    assert any("denied" in r for r in reasons2)


def test_presentation_source_ir_from_workbuddy_task():
    task_dir = Path(r"E:\project\agent\ppt-harness\benchmark_v0.1\workbuddy\html-report-quadrant-ppt")
    if not task_dir.is_dir():
        pytest.skip("workbuddy task not present")
    from domains.ppt.intake import build_presentation_source_ir

    ir = build_presentation_source_ir(task_dir)
    assert ir.sources
    assert ir.page_structure or ir.slide_inventory
    assert ir.summary()


def test_tool_facade_is_vendor_neutral_and_capabilities_registered():
    import domains.ppt  # noqa: F401
    from core.capability import get_capability
    from core.tool import tools_for_capability
    from domains.ppt.tools import FACADE_TOOL_SPECS, PPT_TOOL_FACADE

    assert "presentation.read" in PPT_TOOL_FACADE
    get_capability("presentation.render")
    assert {t.name for t in FACADE_TOOL_SPECS} >= {"ppt_open", "ppt_render", "ppt_save"}
    assert tools_for_capability("presentation.render")
    joined = json.dumps(PPT_TOOL_FACADE).lower()
    for backend in ("python-pptx", "libreoffice", "soffice"):
        assert backend not in joined
