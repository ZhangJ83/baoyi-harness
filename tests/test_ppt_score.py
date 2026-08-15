from pathlib import Path

from pptx import Presentation

from benchmarks.ppt_score import score


def test_ppt_score_accepts_in_bounds_deck(tmp_path: Path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, 100000, 100000).text = "项目 限制"
    path = tmp_path / "ok.pptx"
    prs.save(path)
    result = score(path, min_slides=1, required_text=["项目", "限制"])
    assert result["score"] == 1.0
    assert result["checks"]["no_overflow"] is True


def test_ppt_score_rejects_geometry_overflow(tmp_path: Path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(prs.slide_width - 100, 0, 1000, 100000).text = "overflow"
    path = tmp_path / "overflow.pptx"
    prs.save(path)
    result = score(path)
    assert result["checks"]["no_overflow"] is False
    assert result["score"] < 1.0
