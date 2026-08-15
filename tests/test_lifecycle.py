import json
import os
from pathlib import Path
from tempfile import TemporaryDirectory
from unittest.mock import patch

import pytest
from pptx import Presentation

from agent.lifecycle import RecordMode, RunRecorder, discover_workspace
from agent.state import RunState
from agent.tools.registry import dispatch


class HarnessStub:
    def __init__(self, recorder=None):
        self.state = RunState()
        self.deck = None
        self.recorder = recorder
        self._done = None


def test_discovery_classifies_instructions_tasks_and_office_inputs():
    with TemporaryDirectory() as tmp:
        root = Path(tmp)
        (root / "XIAOPU.md").write_text("rules", encoding="utf-8")
        (root / "task_brief.txt").write_text("make slides", encoding="utf-8")
        (root / "data.xlsx").write_bytes(b"xlsx")
        result = discover_workspace(root)
    assert "XIAOPU.md" in result["instructions"]
    assert "task_brief.txt" in result["task_files"]
    assert any(item["path"] == "data.xlsx" for item in result["inputs"])


def test_existing_deck_uses_working_copy_preserves_original_and_records_provenance():
    with TemporaryDirectory() as tmp, patch.dict(os.environ, {"WORKSPACE": tmp}):
        root = Path(tmp)
        source = Presentation()
        source.slides.add_slide(source.slide_layouts[6])
        source.save(root / "source.pptx")
        original = (root / "source.pptx").read_bytes()
        recorder = RunRecorder("edit deck", "fake", "test", root)
        harness = HarnessStub(recorder)
        dispatch("open_deck", json.dumps({"path": "source.pptx"}), harness)
        dispatch("add_slide", json.dumps({"title": "Added", "bullets": ["A"]}), harness)
        dispatch("save_deck", json.dumps({"path": "final.pptx"}), harness)
        assert (root / "source.pptx").read_bytes() == original
        assert Path(harness.deck_working_path).is_file()
        assert (root / "final.pptx").is_file()
        assert "source-to-output" in recorder.provenance_path.read_text(encoding="utf-8")


def test_working_copy_is_idempotent_and_never_nests_working_suffixes():
    with TemporaryDirectory() as tmp, patch.dict(os.environ, {"WORKSPACE": tmp}):
        root = Path(tmp)
        source = Presentation()
        source.slides.add_slide(source.slide_layouts[6])
        source.save(root / "source.pptx")
        recorder = RunRecorder("edit deck", "fake", "test", root)
        first = recorder.working_copy(root / "source.pptx")
        second = recorder.working_copy(root / "source.pptx")
        third = recorder.working_copy(first)
        assert first == second == third
        assert first.name == "source.working.pptx"
        assert not list(recorder.work.glob("*.working.working.pptx"))
        rows = [json.loads(line) for line in recorder.steps_path.read_text(encoding="utf-8").splitlines()]
        assert sum(row["kind"] == "working_copy_created" for row in rows) == 1


def test_parallel_events_have_unique_ordered_sequences_and_valid_json():
    from concurrent.futures import ThreadPoolExecutor

    with TemporaryDirectory() as tmp:
        recorder = RunRecorder("parallel trace", "fake", "test", Path(tmp))
        with ThreadPoolExecutor(max_workers=8) as pool:
            list(pool.map(lambda i: recorder.event("parallel", index=i), range(100)))
        rows = [json.loads(line) for line in recorder.steps_path.read_text(encoding="utf-8").splitlines()]
        sequences = [row["sequence"] for row in rows]
        assert sequences == list(range(1, 101))


def test_record_modes_change_evidence_depth_without_changing_execution():
    with TemporaryDirectory() as tmp:
        root = Path(tmp)
        source = root / "source.txt"
        source.write_text("input", encoding="utf-8")
        minimal = RunRecorder("mode test", "fake", "test", root, mode="minimal")
        minimal.event("tool_started", arguments="private operational detail")
        minimal.check("structure", True, "passed")
        copied = minimal.working_copy(source)

        assert copied.read_text(encoding="utf-8") == "input"
        minimal_rows = [json.loads(line) for line in minimal.steps_path.read_text(encoding="utf-8").splitlines()]
        assert {row["kind"] for row in minimal_rows} == {"verification", "working_copy_created"}

        research = RunRecorder("mode test", "fake", "test", root, mode=RecordMode.RESEARCH)
        research.event("tool_started", arguments="x" * 5000)
        row = json.loads(research.steps_path.read_text(encoding="utf-8").splitlines()[0])
        assert row["kind"] == "tool_started"
        assert len(row["arguments"]) == 5000


def test_event_sink_failure_is_isolated_from_run_and_circuit_broken():
    class BrokenSink:
        def __init__(self):
            self.calls = 0

        def emit(self, row):
            self.calls += 1
            raise OSError("journal unavailable")

    with TemporaryDirectory() as tmp:
        sink = BrokenSink()
        recorder = RunRecorder("continue core work", "fake", "test", Path(tmp), event_sink=sink)
        recorder.event("tool_started", arguments="one")
        recorder.event("tool_finished", output="two")
        state = RunState()
        recorder.check("structural", True, "passed")
        recorder.finish("core task completed", state)

        assert sink.calls == 1
        assert recorder.completed
        assert recorder.recording_degraded
        assert recorder.record_failures == 1
        assert recorder.manifest["status"] == "completed"


def test_manifest_write_failure_does_not_block_completion():
    with TemporaryDirectory() as tmp:
        recorder = RunRecorder("best effort manifest", "fake", "test", Path(tmp))
        recorder.manifest_path = recorder.root  # writing text to a directory must fail
        recorder.check("structural", True, "passed")
        recorder.finish("completed despite observer failure", RunState())

        assert recorder.completed
        assert recorder.recording_degraded
        assert recorder.record_failures >= 1


def test_invalid_record_mode_falls_back_to_audit():
    with TemporaryDirectory() as tmp:
        recorder = RunRecorder("mode fallback", "fake", "test", Path(tmp), mode="invalid")
        assert recorder.mode is RecordMode.AUDIT


def test_recorder_directory_initialization_failure_is_observational(monkeypatch):
    original_mkdir = Path.mkdir
    failed = False

    def fail_first_recorder_mkdir(path, *args, **kwargs):
        nonlocal failed
        if not failed and ".xiaopu" in path.parts:
            failed = True
            raise OSError("journal volume unavailable")
        return original_mkdir(path, *args, **kwargs)

    with TemporaryDirectory() as tmp:
        monkeypatch.setattr(Path, "mkdir", fail_first_recorder_mkdir)
        recorder = RunRecorder("read-only work still runs", "fake", "test", Path(tmp))
        recorder.event("tool_started", arguments="continue")
        recorder.finish("completed", RunState())

        assert recorder.completed
        assert recorder.recording_degraded
        assert recorder.record_failures >= 1


def test_failed_verification_opens_bounded_repair_then_requires_reverification():
    harness = HarnessStub()
    dispatch("new_deck", json.dumps({"title": "Broken"}), harness)
    shape = next(shape for shape in harness.deck.slides[0].shapes if shape.has_text_frame and shape.text.strip() == "Broken")
    shape.left = harness.deck.slide_width - shape.width // 2
    assert "crosses slide boundary" in dispatch("ppt_verify", "{}", harness)
    assert harness.state.last_verification_failed
    dispatch("set_shape_geometry", json.dumps({"slide_number": 1, "shape_id": shape.shape_id, "x": 1.0, "y": 2.0, "w": 8.0, "height": 1.0}), harness)
    assert harness.state.repair_attempts == 1
    with pytest.raises(ValueError, match="verification evidence|ppt_structural"):
        dispatch("finish", json.dumps({"summary": "done"}), harness)
    assert "no structural issues" in dispatch("ppt_verify", "{}", harness)


def test_repair_budget_is_enforced():
    harness = HarnessStub()
    dispatch("new_deck", json.dumps({"title": "Repair"}), harness)
    shape = next(shape for shape in harness.deck.slides[0].shapes if shape.has_text_frame and shape.text.strip() == "Repair")
    harness.state.last_verification_failed = True
    harness.state.unresolved_checks.add("ppt_structural")
    harness.state.max_repairs = 1
    dispatch("set_shape_geometry", json.dumps({"slide_number": 1, "shape_id": shape.shape_id, "x": 1.0, "y": 2.0, "w": 8.0, "height": 1.0}), harness)
    # The same verifier-feedback cycle may make further mutations before reverify.
    dispatch("set_shape_geometry", json.dumps({"slide_number": 1, "shape_id": shape.shape_id, "x": 1.1, "y": 2.0, "w": 8.0, "height": 1.0}), harness)
    # A fresh verifier failure reopens a second cycle; the budget now blocks.
    harness.state.last_verification_failed = True
    with pytest.raises(RuntimeError, match="repair budget exhausted"):
        dispatch("set_shape_geometry", json.dumps({"slide_number": 1, "shape_id": shape.shape_id, "x": 1.2, "y": 2.0, "w": 8.0, "height": 1.0}), harness)


def test_trajectory_redacts_provider_secret():
    with TemporaryDirectory() as tmp, patch.dict(os.environ, {"OPENAI_API_KEY": "sk-secret-trajectory-value"}, clear=False):
        recorder = RunRecorder("safe trace", "fake", "test", Path(tmp))
        recorder.event("tool_started", arguments="token=sk-secret-trajectory-value")
        trace = recorder.steps_path.read_text(encoding="utf-8")
    assert "sk-secret-trajectory-value" not in trace
    assert "[REDACTED]" in trace


def test_official_task_evaluator_runs_with_local_paths_and_records_evidence():
    with TemporaryDirectory() as tmp, patch.dict(os.environ, {"WORKSPACE": tmp}):
        root = Path(tmp)
        task = root / "tasks" / "demo"
        evaluator = task / "tests" / "grading" / "test_verify.py"
        evaluator.parent.mkdir(parents=True)
        (task / "output").mkdir()
        (task / "tests" / "gold").mkdir()
        (task / "output" / "final.pptx").write_bytes(b"artifact")
        evaluator.write_text(
            "import os\nfrom pathlib import Path\n"
            "def test_output_is_bound():\n    assert Path(os.environ['WB_BENCH_OUTPUT_PATH']).is_file()\n",
            encoding="utf-8",
        )
        recorder = RunRecorder("tasks/demo", "fake", "test", root)
        harness = HarnessStub(recorder)
        harness.state.facts.update({
            "official_evaluator": str(evaluator.relative_to(root)),
            "official_evaluator_present": "true",
            "required_output_pptx": str(Path("tasks/demo/output/final.pptx")),
        })
        result = dispatch("run_task_evaluator", "{}", harness)
        assert "passed" in result
        assert "task_evaluator" in {row.kind for row in harness.state.fresh_evidence()}


def test_official_task_evaluator_does_not_inherit_provider_secrets(monkeypatch):
    captured = {}

    class Completed:
        returncode = 0
        stdout = "passed"
        stderr = ""

    def fake_run(*args, **kwargs):
        captured.update(kwargs["env"])
        return Completed()

    with TemporaryDirectory() as tmp, patch.dict(
        os.environ,
        {"WORKSPACE": tmp, "OPENAI_API_KEY": "secret-a", "ANTHROPIC_API_KEY": "secret-b", "DEEPSEEK_API_KEY": "secret-c"},
        clear=False,
    ):
        root = Path(tmp)
        task = root / "tasks" / "demo"
        evaluator = task / "tests" / "grading" / "test_verify.py"
        evaluator.parent.mkdir(parents=True)
        (task / "output").mkdir()
        (task / "tests" / "gold").mkdir()
        evaluator.write_text("def test_ok(): assert True\n", encoding="utf-8")
        (task / "output" / "final.pptx").write_bytes(b"artifact")
        recorder = RunRecorder("tasks/demo", "fake", "test", root)
        harness = HarnessStub(recorder)
        harness.state.facts.update({
            "official_evaluator": str(evaluator.relative_to(root)),
            "official_evaluator_present": "true",
            "required_output_pptx": "tasks/demo/output/final.pptx",
        })
        monkeypatch.setattr("agent.tools.lifecycle_tools.subprocess.run", fake_run)
        dispatch("run_task_evaluator", "{}", harness)

    assert "OPENAI_API_KEY" not in captured
    assert "ANTHROPIC_API_KEY" not in captured
    assert "DEEPSEEK_API_KEY" not in captured


def test_finish_exports_missing_task_trajectory_contract_files():
    with TemporaryDirectory() as tmp:
        root = Path(tmp)
        task = root / "tasks" / "demo-task"
        task.mkdir(parents=True)
        (task / "TRAJECTORY_CAPTURE_CONTRACT.md").write_text("contract", encoding="utf-8")
        recorder = RunRecorder("完成 tasks/demo-task 的 PPT 修改", "fake-model", "test", root)
        state = RunState()
        state.record_evidence("ppt_structural", "passed")
        recorder.check("ppt_structural", True, "passed")
        recorder.finish("done", state)
        expected = {"steps.jsonl", "run_metadata.json", "plan.md", "reads.md", "checks.md", "repairs.md", "artifacts.md"}
        assert expected == {path.name for path in (task / "trajectory").iterdir()}
        metadata = json.loads((task / "trajectory" / "run_metadata.json").read_text(encoding="utf-8"))
        assert metadata["task_id"] == "demo-task"
        assert metadata["model"] == "fake-model"


def test_finish_refreshes_task_trajectory_latest_view():
    with TemporaryDirectory() as tmp:
        root = Path(tmp)
        task = root / "tasks" / "demo-task"
        trajectory = task / "trajectory"
        trajectory.mkdir(parents=True)
        (task / "TRAJECTORY_CAPTURE_CONTRACT.md").write_text("contract", encoding="utf-8")
        (trajectory / "run_metadata.json").write_text('{"session_id":"stale"}', encoding="utf-8")
        (trajectory / "steps.jsonl").write_text("stale\n", encoding="utf-8")

        recorder = RunRecorder("complete tasks/demo-task", "fresh-model", "test", root)
        state = RunState()
        recorder.finish("done", state)

        metadata = json.loads((trajectory / "run_metadata.json").read_text(encoding="utf-8"))
        assert metadata["session_id"] == recorder.run_id
        assert metadata["model"] == "fresh-model"
        assert "stale" not in (trajectory / "steps.jsonl").read_text(encoding="utf-8")
