import json
import math
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt

from benchmarks.validate_model_generated_ppt_eval import (
    icc_2_1,
    performance_summary,
    recompute_deck_checks,
    recompute_pixel_checks,
    validate,
    weighted_kappa,
)


ROOT = Path(__file__).resolve().parents[1]


def test_agreement_metrics_are_one_for_identical_nonconstant_scores():
    values = [1, 2, 3, 4, 5, 2, 4]
    assert weighted_kappa(values, values) == 1.0
    assert math.isclose(icc_2_1([(float(value), float(value)) for value in values]), 1.0)


def test_result_validator_fails_closed_without_36_raw_artifacts(tmp_path):
    mapping = tmp_path / "mapping.json"
    mapping.write_text(json.dumps({"artifacts": []}), encoding="utf-8")
    score1 = tmp_path / "r1.json"
    score2 = tmp_path / "r2.json"
    for path in (score1, score2):
        path.write_text(json.dumps({"locked_before_adjudication": True, "scores": []}), encoding="utf-8")
    result = validate(
        ROOT / "benchmarks/pptbench_model_eval_v2.json",
        tmp_path / "raw",
        mapping,
        [score1, score2],
    )
    assert result["valid"] is False
    assert result["model_generated"] is False
    assert "anonymous_mapping_task_system_mismatch" in result["errors"]
    assert result["rendered_decks"] == 0


def test_blind_scores_produce_task_paired_statistics_and_holm():
    tasks = [{"id": f"t{i}"} for i in range(12)]
    protocol = {"systems":["xiaopu","claude_code","codex"], "tasks":tasks}
    mapping = []
    reviewers = [{}, {}]
    values = {"xiaopu":5, "claude_code":3, "codex":4}
    for system in protocol["systems"]:
        for task in tasks:
            artifact_id = f"{system}-{task['id']}"
            mapping.append({"anonymous_id":artifact_id, "system":system, "task_id":task["id"]})
            for reviewer in reviewers:
                reviewer[artifact_id] = {"dimensions":{"content_fidelity":values[system], "layout":values[system]}}
    result = performance_summary(protocol, mapping, reviewers, ["content_fidelity","layout"])
    assert result["system_mean_blind_score"] == {"xiaopu":5.0,"claude_code":3.0,"codex":4.0}
    assert result["contrasts"]["xiaopu_minus_claude_code"]["mean_delta"] == 2.0
    assert result["contrasts"]["xiaopu_minus_codex"]["paired_bootstrap_ci95"] == [1.0, 1.0]
    assert result["contrasts"]["xiaopu_minus_claude_code"]["holm_adjusted_p"] <= 0.05


def test_independent_pptx_and_pixel_recompute_enforces_contract(tmp_path):
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.8), Inches(0.9))
    title.text_frame.paragraphs[0].add_run().text = "Decision brief"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(50)
    body = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(10), Inches(1.0))
    body.text_frame.paragraphs[0].add_run().text = "Required fact"
    body.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    slide.notes_slide.notes_text_frame.text = "[Sources] Internal evaluation task packet"
    deck_path = tmp_path / "deck.pptx"
    deck.save(deck_path)
    task = {"min_slides": 1, "max_slides": 1, "required_text": ["Required fact"]}
    audit = recompute_deck_checks(deck_path, task)
    assert audit["valid"] is True
    assert audit["source_notes_missing"] == 0
    image = Image.new("RGB", (640, 360), "white")
    ImageDraw.Draw(image).rectangle((100, 80, 540, 280), fill=(30, 60, 100))
    slide_path = tmp_path / "slide-1.png"
    image.save(slide_path)
    assert recompute_pixel_checks([slide_path], 1) == {
        "rendered_png_count": 1,
        "blank_or_near_uniform_count": 0,
        "pass": True,
    }
