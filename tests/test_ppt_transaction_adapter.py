from __future__ import annotations

import json
from pathlib import Path
from tempfile import TemporaryDirectory

import pytest
from pptx import Presentation
from pptx.util import Inches

from agent.action_transaction import PostconditionFailed, ScopeViolation, TransactionCancelled
from agent.lifecycle import RunRecorder
from agent.ppt_transaction_adapter import run_ppt_transaction, slide_fingerprints
from agent.state import RunState


class HarnessStub:
    def __init__(self, deck: Presentation, recorder=None, cancelled: bool = False):
        self.deck = deck
        self.recorder = recorder
        self.state = RunState()
        self._cancelled = cancelled

    def cancel_requested(self) -> bool:
        return self._cancelled


def _deck() -> Presentation:
    deck = Presentation()
    for title in ("One", "Two", "Three"):
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        shape.text = title
    return deck


def test_scoped_edit_persists_checkpoint_certificate_and_trace() -> None:
    with TemporaryDirectory() as tmp:
        root = Path(tmp)
        source = root / "source.pptx"
        _deck().save(source)
        source_bytes = source.read_bytes()
        recorder = RunRecorder("edit slide 2", "fake", "test", root)
        harness = HarnessStub(Presentation(source), recorder)
        harness.deck_source_path = source
        harness.deck_working_path = recorder.working_copy(source)

        def edit(deck: Presentation) -> str:
            deck.slides[1].shapes[0].text = "Two edited"
            return "ok"

        outcome = run_ppt_transaction(
            harness,
            allowed_slides=[2],
            requested_slides=[2],
            action=edit,
            postcondition=lambda deck, result: result == "ok" and "edited" in deck.slides[1].shapes[0].text,
            transaction_id="scoped-success",
        )

        assert outcome.certificate.passed
        assert outcome.certificate.changed_slides == (2,)
        assert outcome.certificate.unexpected_changed_slides == ()
        assert outcome.checkpoint_path.is_file()
        assert outcome.committed_deck_path.is_file()
        assert outcome.certificate_path.is_file()
        assert source.read_bytes() == source_bytes
        assert harness.state.ppt_affected_slides == {2}
        persisted = json.loads(outcome.certificate_path.read_text(encoding="utf-8"))
        assert persisted["schema"] == "xiaopu-ppt-scope-certificate-v1"
        kinds = [json.loads(line)["kind"] for line in recorder.steps_path.read_text(encoding="utf-8").splitlines()]
        assert "action_transaction" in kinds
        assert "ppt_scope_verified" in kinds


def test_unauthorized_slide_is_rejected_before_checkpoint() -> None:
    harness = HarnessStub(_deck())
    original = slide_fingerprints(harness.deck)

    with pytest.raises(ScopeViolation):
        run_ppt_transaction(
            harness,
            allowed_slides=[1],
            requested_slides=[2],
            action=lambda deck: setattr(deck.slides[1].shapes[0], "text", "forbidden"),
        )

    assert slide_fingerprints(harness.deck) == original
    assert harness.state.mutation_epoch == 0


def test_unrequested_slide_change_fails_postcondition_and_rolls_back() -> None:
    harness = HarnessStub(_deck())
    original = slide_fingerprints(harness.deck)

    def edit_wrong_slide(deck: Presentation) -> None:
        deck.slides[0].shapes[0].text = "unexpected"
        harness.state.record_change("deck:slide:1:text")

    with pytest.raises(PostconditionFailed):
        run_ppt_transaction(
            harness,
            allowed_slides=[1, 2],
            requested_slides=[2],
            action=edit_wrong_slide,
        )

    assert slide_fingerprints(harness.deck) == original
    assert harness.state.mutation_epoch == 0
    assert harness.state.changed_files == set()


def test_action_failure_restores_working_copy_and_live_deck() -> None:
    with TemporaryDirectory() as tmp:
        root = Path(tmp)
        working = root / "working.pptx"
        deck = _deck()
        deck.save(working)
        before = working.read_bytes()
        harness = HarnessStub(deck)
        harness.deck_working_path = working

        def fail(deck: Presentation) -> None:
            deck.slides[1].shapes[0].text = "partial"
            deck.save(working)
            raise RuntimeError("mutation failed")

        with pytest.raises(RuntimeError, match="mutation failed"):
            run_ppt_transaction(
                harness,
                allowed_slides=[2],
                requested_slides=[2],
                action=fail,
            )

        assert working.read_bytes() == before
        assert harness.deck.slides[1].shapes[0].text == "Two"


def test_cancelled_harness_does_not_create_checkpoint_or_mutate() -> None:
    with TemporaryDirectory() as tmp:
        recorder = RunRecorder("cancel", "fake", "test", Path(tmp))
        harness = HarnessStub(_deck(), recorder, cancelled=True)

        with pytest.raises(TransactionCancelled):
            run_ppt_transaction(
                harness,
                allowed_slides=[2],
                requested_slides=[2],
                action=lambda deck: setattr(deck.slides[1].shapes[0], "text", "no"),
                transaction_id="cancel-before-checkpoint",
            )

        assert harness.deck.slides[1].shapes[0].text == "Two"
        assert not (recorder.work / "transactions" / "cancel-before-checkpoint").exists()

