"""Runtime bridge tests: the three legacy entrypoints now route through the portable stack."""
from pathlib import Path
from unittest.mock import patch

import json


def test_intake_bridge_builds_portable_source_ir(tmp_path):
    from pptx import Presentation

    from agent.intake import prepare_task_brief
    from agent.state import RunState

    root = Path(tmp_path)
    task = root / "tasks" / "demo"
    (task / "input").mkdir(parents=True)
    (task / "instruction.md").write_text("Create output/final.pptx", encoding="utf-8")
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(task / "input" / "deck.pptx"))

    state = RunState()
    with patch("agent.intake.config.sandbox_root", return_value=root):
        brief = prepare_task_brief("tasks/demo", state)

    payload = json.loads(brief)
    assert payload.get("input_pptx") == str(Path("tasks/demo/input/deck.pptx"))
    assert "ppt_source_ir" in state.facts
    assert "slides=1" in state.facts["ppt_source_ir"]
    assert payload.get("content_brief")


def test_task_compiler_bridge_returns_portable_contract():
    from agent.task_compiler import portable_contract_for

    contract = portable_contract_for(
        "完成 tasks/board-material-update-timeline-excel",
        {
            "task_instruction": "以现有 PPTX 为基础稿，使用 XLSX 作为更新来源。",
            "manifest_task_id": "board-material-update-timeline-excel",
        },
        "",
    )
    assert contract.task_type == "source_grounded_build"
    assert "presentation.read" in contract.capabilities
    assert "content_grounding" in contract.verification.required_kinds()
    # C_static now closes over mutation scope as well.
    assert contract.mutation is not None
    assert "theme" in contract.mutation.deny


def test_task_index_bridge_keeps_ownership_ranking(tmp_path):
    from agent.task_index import resolve_task

    for agent, task_id in [("alpha", "demo"), ("xiaopuharness", "demo"), ("xiaopuharness - 副本", "demo")]:
        task = tmp_path / "agent_workspaces" / "full13" / agent / "tasks" / task_id
        task.mkdir(parents=True)
        (task / "instruction.md").write_text("do it", encoding="utf-8")
    assert resolve_task(tmp_path, "demo").name == "demo"
    assert "xiaopuharness" in str(resolve_task(tmp_path, "demo"))
