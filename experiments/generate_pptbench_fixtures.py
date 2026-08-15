"""Create deterministic structural PPT fixtures for scorer smoke tests.

These decks are not agent outputs and must never be reported as PPTBench
performance. They only prove that the manifest, artifact naming, and offline
structural scorer agree end to end.
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation


def make_deck(task: dict, out: Path) -> None:
    prs = Presentation()
    while len(prs.slides) < task["min_slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(0, 0, 10 * 914400, 1 * 914400)
        box.text = task["prompt"]
    text = " ".join(task.get("required_text", []))
    if text:
        shape = prs.slides[0].shapes.add_textbox(0, 914400, 10 * 914400, 1 * 914400)
        shape.text = text
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)


def main() -> None:
    p = argparse.ArgumentParser()
    p.add_argument("--manifest", type=Path, default=Path("benchmarks/pptbench_tasks.json"))
    p.add_argument("--out", type=Path, required=True)
    args = p.parse_args()
    manifest = json.loads(args.manifest.read_text(encoding="utf-8-sig"))
    for task in manifest["tasks"]:
        make_deck(task, args.out / f"{task['id']}.pptx")
    print(f"created {len(manifest['tasks'])} structural fixtures under {args.out}")


if __name__ == "__main__":
    main()
