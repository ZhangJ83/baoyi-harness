"""Build a fail-closed model-generated PPT evaluation report from raw artifacts."""
from __future__ import annotations

import argparse
import hashlib
import itertools
import json
import math
import random
import re
import zipfile
from pathlib import Path

from PIL import Image, ImageStat
from pptx import Presentation


def load(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8-sig"))


def digest(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _shape_text(shape) -> str:
    return (shape.text_frame.text or "").strip() if getattr(shape, "has_text_frame", False) else ""


def _run_sizes(shape) -> list[float]:
    return [
        run.font.size.pt
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.font.size is not None and run.text.strip()
    ]


def recompute_deck_checks(deck_path: Path, task: dict) -> dict:
    """Recompute structural contract fields from the PPTX instead of trusting reports."""
    deck = Presentation(str(deck_path))
    slide_count = len(deck.slides)
    all_text: list[str] = []
    overlap_count = 0
    overflow_count = 0
    boundary_count = 0
    empty_text_boxes = 0
    body_font_violations = 0
    title_font_violations = 0
    title_font_unknown = 0
    source_notes_missing = 0
    for slide_index, slide in enumerate(deck.slides):
        text_shapes = []
        for shape in slide.shapes:
            text = _shape_text(shape)
            if text:
                all_text.append(text)
                text_shapes.append(shape)
                body_font_violations += sum(size < 16 for size in _run_sizes(shape))
                sizes = _run_sizes(shape)
                biggest = max(sizes) if sizes else 18.0
                width_inches = shape.width / 914400
                height_inches = shape.height / 914400
                chars_per_line = max(1, int(width_inches * 96 / (biggest * 0.75)))
                lines = sum(max(1, math.ceil(len(line) / chars_per_line)) for line in text.split("\n"))
                if lines * biggest * 0.62 / 72 > height_inches * 0.95:
                    overflow_count += 1
            elif getattr(shape, "has_text_frame", False):
                empty_text_boxes += 1
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > deck.slide_width or shape.top + shape.height > deck.slide_height:
                boundary_count += 1
        for left_index, first in enumerate(text_shapes):
            for second in text_shapes[left_index + 1:]:
                overlap_width = min(first.left + first.width, second.left + second.width) - max(first.left, second.left)
                overlap_height = min(first.top + first.height, second.top + second.height) - max(first.top, second.top)
                if overlap_width > 0 and overlap_height > 0:
                    ratio = (overlap_width * overlap_height) / min(first.width * first.height, second.width * second.height)
                    overlap_count += ratio > 0.15
        candidates = [shape for shape in text_shapes if _shape_text(shape)]
        title = slide.shapes.title if slide.shapes.title is not None and _shape_text(slide.shapes.title) else (min(candidates, key=lambda shape: shape.top) if candidates else None)
        title_sizes = _run_sizes(title) if title is not None else []
        if not title_sizes:
            title_font_unknown += 1
        elif max(title_sizes) < (50 if slide_index == 0 else 35):
            title_font_violations += 1
        try:
            notes = slide.notes_slide.notes_text_frame.text.casefold()
        except Exception:
            notes = ""
        source_notes_missing += "[sources]" not in notes
    combined = "\n".join(all_text).casefold()
    placeholders = ("lorem ipsum", "placeholder", "todo:", "[title]", "[content]")
    unresolved = empty_text_boxes + sum(marker in combined for marker in placeholders)
    required_text_missing = [value for value in task["required_text"] if str(value).casefold() not in combined]
    checks = {
        "opens": True,
        "slide_count": slide_count,
        "slide_count_in_bounds": task["min_slides"] <= slide_count <= task["max_slides"],
        "required_text_missing": required_text_missing,
        "overlap_count": overlap_count,
        "overflow_count": overflow_count,
        "boundary_count": boundary_count,
        "unresolved_placeholders": unresolved,
        "body_font_violations": body_font_violations,
        "title_font_violations": title_font_violations,
        "title_font_unknown": title_font_unknown,
        "source_notes_missing": source_notes_missing,
    }
    checks["valid"] = bool(
        checks["slide_count_in_bounds"]
        and not required_text_missing
        and overlap_count == 0
        and overflow_count == 0
        and boundary_count == 0
        and unresolved == 0
        and body_font_violations == 0
        and title_font_violations == 0
        and title_font_unknown == 0
        and source_notes_missing == 0
    )
    return checks


def recompute_pixel_checks(slide_pngs: list[Path], expected_count: int) -> dict:
    blank_count = 0
    for path in slide_pngs:
        with Image.open(path) as image:
            variance = ImageStat.Stat(image.convert("L")).var[0]
        blank_count += variance < 8
    return {
        "rendered_png_count": len(slide_pngs),
        "blank_or_near_uniform_count": blank_count,
        "pass": len(slide_pngs) == expected_count and expected_count > 0 and blank_count == 0,
    }


def weighted_kappa(a: list[int], b: list[int]) -> float:
    n = len(a)
    if n == 0 or n != len(b):
        return float("nan")
    observed = sum(((x - y) / 4) ** 2 for x, y in zip(a, b)) / n
    ca = [a.count(i) / n for i in range(1, 6)]
    cb = [b.count(i) / n for i in range(1, 6)]
    expected = sum(ca[i - 1] * cb[j - 1] * (((i - j) / 4) ** 2) for i in range(1, 6) for j in range(1, 6))
    return 1.0 if expected == 0 and observed == 0 else (float("nan") if expected == 0 else 1 - observed / expected)


def icc_2_1(rows: list[tuple[float, float]]) -> float:
    n, k = len(rows), 2
    if n < 2:
        return float("nan")
    row_means = [sum(row) / k for row in rows]
    col_means = [sum(row[j] for row in rows) / n for j in range(k)]
    grand = sum(row_means) / n
    msr = k * sum((value - grand) ** 2 for value in row_means) / (n - 1)
    msc = n * sum((value - grand) ** 2 for value in col_means) / (k - 1)
    mse = sum((rows[i][j] - row_means[i] - col_means[j] + grand) ** 2 for i in range(n) for j in range(k)) / ((n - 1) * (k - 1))
    denominator = msr + (k - 1) * mse + k * (msc - mse) / n
    return float("nan") if denominator == 0 else (msr - mse) / denominator


def paired_bootstrap_ci(deltas: list[float], repeats: int = 5000) -> list[float]:
    rng = random.Random(20260811)
    values = []
    for _ in range(repeats):
        sample = [deltas[rng.randrange(len(deltas))] for _ in deltas]
        values.append(sum(sample) / len(sample))
    values.sort()
    return [values[int(0.025 * repeats)], values[min(repeats - 1, int(0.975 * repeats))]]


def exact_paired_permutation(deltas: list[float]) -> float:
    observed = abs(sum(deltas) / len(deltas))
    extreme = 0
    total = 2 ** len(deltas)
    for signs in itertools.product((-1, 1), repeat=len(deltas)):
        value = abs(sum(sign * delta for sign, delta in zip(signs, deltas)) / len(deltas))
        extreme += value >= observed - 1e-12
    return extreme / total


def holm_adjust(values: dict[str, float]) -> dict[str, float]:
    ordered = sorted(values.items(), key=lambda item: item[1])
    adjusted = {}
    running = 0.0
    m = len(ordered)
    for rank, (name, value) in enumerate(ordered):
        running = max(running, min(1.0, value * (m - rank)))
        adjusted[name] = running
    return adjusted


def performance_summary(protocol: dict, mapping_rows: list[dict], reviewers: list[dict[str, dict]], dimensions: list[str]) -> dict:
    task_scores: dict[tuple[str, str], float] = {}
    for row in mapping_rows:
        artifact_id = row["anonymous_id"]
        per_reviewer = [sum(reviewer[artifact_id]["dimensions"].values()) / len(dimensions) for reviewer in reviewers]
        task_scores[(row["system"], row["task_id"])] = sum(per_reviewer) / len(per_reviewer)
    means = {system: sum(task_scores[(system, task["id"])] for task in protocol["tasks"]) / len(protocol["tasks"]) for system in protocol["systems"]}
    contrasts = {}
    raw_p = {}
    for comparator in ("claude_code", "codex"):
        deltas = [task_scores[("xiaopu", task["id"])] - task_scores[(comparator, task["id"])] for task in protocol["tasks"]]
        name = f"xiaopu_minus_{comparator}"
        p_value = exact_paired_permutation(deltas)
        raw_p[name] = p_value
        contrasts[name] = {"mean_delta":sum(deltas) / len(deltas), "paired_bootstrap_ci95":paired_bootstrap_ci(deltas), "paired_permutation_p":p_value, "n_tasks":len(deltas)}
    adjusted = holm_adjust(raw_p)
    for name, value in adjusted.items():
        contrasts[name]["holm_adjusted_p"] = value
    return {"system_mean_blind_score":means, "contrasts":contrasts, "unit":"task"}


def validate(protocol_path: Path, raw_root: Path, mapping_path: Path, score_paths: list[Path]) -> dict:
    protocol = load(protocol_path)
    mapping = load(mapping_path)
    errors: list[str] = []
    dimensions = protocol["blind_review"]["dimensions"]
    task_by_id = {task["id"]: task for task in protocol["tasks"]}
    expected = {(system, task["id"]) for system in protocol["systems"] for task in protocol["tasks"]}
    rows = mapping.get("artifacts", [])
    keyed = {(row.get("system"), row.get("task_id")): row for row in rows}
    if set(keyed) != expected or len(rows) != len(expected):
        errors.append("anonymous_mapping_task_system_mismatch")
    anonymous_ids = [row.get("anonymous_id") for row in rows]
    if len(anonymous_ids) != len(set(anonymous_ids)) or any(not value for value in anonymous_ids):
        errors.append("anonymous_ids_invalid_or_duplicate")
    protocol_hash = digest(protocol_path)
    rendered = 0
    independent_audits: dict[str, dict] = {}
    for system, task_id in sorted(expected):
        folder = raw_root / system / task_id
        deck, pdf, montage = folder / "deck.pptx", folder / "deck.pdf", folder / "montage.png"
        structural, pixels, trace = folder / "structural_report.json", folder / "pixel_audit.json", folder / "generation_trace.json"
        required = [deck, pdf, montage, structural, pixels, trace]
        if any(not path.is_file() for path in required):
            errors.append(f"{system}/{task_id}:missing_required_artifact")
            continue
        slide_pngs = sorted((folder / "slides").glob("*.png"))
        if not zipfile.is_zipfile(deck) or pdf.read_bytes()[:4] != b"%PDF" or montage.read_bytes()[:8] != b"\x89PNG\r\n\x1a\n" or not slide_pngs:
            errors.append(f"{system}/{task_id}:invalid_artifact_signature_or_missing_slides")
            continue
        if any(path.read_bytes()[:8] != b"\x89PNG\r\n\x1a\n" for path in slide_pngs):
            errors.append(f"{system}/{task_id}:invalid_slide_png")
            continue
        try:
            structure, pixel, generation = load(structural), load(pixels), load(trace)
            deck_audit = recompute_deck_checks(deck, task_by_id[task_id])
            pixel_audit = recompute_pixel_checks(slide_pngs, deck_audit["slide_count"])
        except Exception as exc:
            errors.append(f"{system}/{task_id}:independent_recompute_failed:{type(exc).__name__}")
            continue
        independent_audits[f"{system}/{task_id}"] = {"deck": deck_audit, "pixels": pixel_audit}
        if deck_audit.get("valid") is not True:
            errors.append(f"{system}/{task_id}:independently_recomputed_structural_gate_failed")
        structural_fields = (
            "valid", "opens", "slide_count", "overlap_count", "overflow_count",
            "boundary_count", "unresolved_placeholders", "body_font_violations",
            "title_font_violations", "title_font_unknown", "source_notes_missing",
        )
        if any(structure.get(field) != deck_audit.get(field) for field in structural_fields):
            errors.append(f"{system}/{task_id}:structural_report_not_reproducible")
        if pixel_audit.get("pass") is not True:
            errors.append(f"{system}/{task_id}:independently_recomputed_pixel_gate_failed")
        if any(pixel.get(field) != pixel_audit.get(field) for field in ("pass", "rendered_png_count", "blank_or_near_uniform_count")):
            errors.append(f"{system}/{task_id}:pixel_report_not_reproducible")
        generation_budget = generation.get("budget", {})
        budget_caps = protocol["budget"]
        budget_valid = bool(
            generation_budget.get("within_budget") is True
            and generation_budget.get("provider_usage_authoritative") is True
            and generation_budget.get("caps") == budget_caps
            and isinstance(generation_budget.get("output_tokens"), int)
            and generation_budget["output_tokens"] <= budget_caps["max_generated_output_tokens"]
            and isinstance(generation_budget.get("covered_local_tool_calls"), int)
            and generation_budget["covered_local_tool_calls"] <= budget_caps["max_covered_local_tool_calls"]
            and isinstance(generation_budget.get("wall_seconds"), (int, float))
            and generation_budget["wall_seconds"] <= budget_caps["max_agent_wall_seconds"] + 1
        )
        if not (generation.get("model_generated") is True and generation.get("cell_complete") is True
                and generation.get("model") == protocol.get("model") and generation.get("system") == system
                and generation.get("task_id") == task_id and generation.get("protocol_sha256") == protocol_hash
                and budget_valid):
            errors.append(f"{system}/{task_id}:generation_trace_mismatch")
        row = keyed.get((system, task_id), {})
        if row.get("deck_sha256") != digest(deck):
            errors.append(f"{system}/{task_id}:deck_hash_mapping_mismatch")
        rendered += 1
    if len(score_paths) != 2:
        errors.append("exactly_two_reviewer_score_files_required")
    scores = [load(path) for path in score_paths] if len(score_paths) == 2 and all(path.is_file() for path in score_paths) else []
    by_reviewer: list[dict[str, dict]] = []
    for index, payload in enumerate(scores):
        if payload.get("schema") != "pptbench-v2-blind-review-form":
            errors.append(f"reviewer_{index + 1}:schema_invalid")
        if payload.get("locked_before_adjudication") is not True:
            errors.append(f"reviewer_{index + 1}:form_not_locked")
        for field in ("independent_from_generation", "reviewer_non_author", "conflicts_declared"):
            if payload.get(field) is not True:
                errors.append(f"reviewer_{index + 1}:{field}_missing")
        if not isinstance(payload.get("reviewer_id"), str) or not payload.get("reviewer_id"):
            errors.append(f"reviewer_{index + 1}:reviewer_id_missing")
        if not isinstance(payload.get("review_completed_at"), str) or not payload.get("review_completed_at"):
            errors.append(f"reviewer_{index + 1}:review_completed_at_missing")
        if not isinstance(payload.get("locked_at"), str) or not payload.get("locked_at"):
            errors.append(f"reviewer_{index + 1}:locked_at_missing")
        if not isinstance(payload.get("source_draft_sha256"), str) or not re.fullmatch(r"[0-9a-f]{64}", payload.get("source_draft_sha256", "")):
            errors.append(f"reviewer_{index + 1}:source_draft_hash_invalid")
        entries = payload.get("scores", [])
        lookup = {row.get("anonymous_id"): row for row in entries}
        if set(lookup) != set(anonymous_ids) or len(entries) != len(anonymous_ids):
            errors.append(f"reviewer_{index + 1}:score_set_mismatch")
        for artifact_id, row in lookup.items():
            values = row.get("dimensions", {})
            if set(values) != set(dimensions) or any(not isinstance(value, int) or value < 1 or value > 5 for value in values.values()):
                errors.append(f"reviewer_{index + 1}:{artifact_id}:invalid_dimensions")
        by_reviewer.append(lookup)
    kappas: dict[str, float] = {}
    icc = float("nan")
    performance = {}
    if len(scores) == 2 and scores[0].get("reviewer_id") == scores[1].get("reviewer_id"):
        errors.append("reviewer_ids_must_be_distinct")
    if len(by_reviewer) == 2 and not any("reviewer_" in error for error in errors):
        for dimension in dimensions:
            a = [by_reviewer[0][artifact_id]["dimensions"][dimension] for artifact_id in anonymous_ids]
            b = [by_reviewer[1][artifact_id]["dimensions"][dimension] for artifact_id in anonymous_ids]
            kappas[dimension] = weighted_kappa(a, b)
        overall = []
        for artifact_id in anonymous_ids:
            overall.append(tuple(sum(by_reviewer[r][artifact_id]["dimensions"].values()) / len(dimensions) for r in range(2)))
        icc = icc_2_1(overall)
        performance = performance_summary(protocol, rows, by_reviewer, dimensions)
    finite_agreement = math.isfinite(icc) and all(math.isfinite(value) for value in kappas.values())
    if scores and not finite_agreement:
        errors.append("agreement_not_finite")
    valid = not errors and rendered == 36
    return {
        "schema": "model-generated-ppt-evaluation-v1",
        "valid": valid,
        "protocol_sha256": protocol_hash,
        "protocol_frozen_before_generation": True,
        "model_generated": valid,
        "paired_task_set": set(keyed) == expected,
        "n_tasks": 12,
        "n_systems": 3,
        "rendered_decks": rendered,
        "independently_recomputed_decks": len(independent_audits),
        "independent_artifact_audits": independent_audits,
        "blind_review": {
            "complete": valid,
            "n_reviewers": len(scores),
            "agreement_metric": icc if math.isfinite(icc) else None,
            "icc_2_1": icc if math.isfinite(icc) else None,
            "quadratic_weighted_kappa": kappas,
            "forms_locked_before_adjudication": len(scores) == 2 and all(payload.get("locked_before_adjudication") is True for payload in scores),
            "raw_score_files_sha256": [digest(path) for path in score_paths if path.is_file()],
        },
        "performance": performance,
        "errors": errors,
        "claim_boundary": "Evaluation completeness and agreement only; performance superiority requires separate paired statistics.",
    }


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--protocol", type=Path, default=Path("benchmarks/pptbench_model_eval_v2.json"))
    parser.add_argument("--raw-root", type=Path, required=True)
    parser.add_argument("--mapping", type=Path, required=True)
    parser.add_argument("--scores", type=Path, action="append", required=True)
    parser.add_argument("--out", type=Path, required=True)
    args = parser.parse_args()
    result = validate(args.protocol, args.raw_root, args.mapping, args.scores)
    args.out.parent.mkdir(parents=True, exist_ok=True)
    args.out.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(result, ensure_ascii=False))
    raise SystemExit(0 if result["valid"] else 2)


if __name__ == "__main__":
    main()
