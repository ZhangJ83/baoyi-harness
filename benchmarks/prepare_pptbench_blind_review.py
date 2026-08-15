"""Create two identity-blind PPT review bundles from frozen raw artifacts."""
from __future__ import annotations

import argparse
import hashlib
import json
import random
import secrets
import shutil
from pathlib import Path

from pptx import Presentation

IDENTITY_MARKERS = ("xiaopu", "claude", "codex", "anthropic", "openai")


def _visible_text(deck: Path) -> str:
    prs = Presentation(str(deck))
    return "\n".join(
        shape.text for slide in prs.slides for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    ).casefold()


def _write(path: Path, value: dict) -> None:
    path.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def prepare(protocol_path: Path, raw_root: Path, out_root: Path, rng=None) -> dict:
    if out_root.exists():
        raise FileExistsError(f"blind review output already exists: {out_root}")
    protocol = json.loads(protocol_path.read_text(encoding="utf-8"))
    rng = rng or random.SystemRandom()
    expected = [(system, task) for system in protocol["systems"] for task in protocol["tasks"]]
    anonymous_ids: list[str] = []
    while len(anonymous_ids) < len(expected):
        value = "D-" + secrets.token_hex(4).upper()
        if value not in anonymous_ids:
            anonymous_ids.append(value)
    rng.shuffle(anonymous_ids)
    mapping_rows = []
    for (system, task), anonymous_id in zip(expected, anonymous_ids):
        source = raw_root / system / task["id"]
        deck = source / "deck.pptx"
        montage = source / "montage.png"
        slides = sorted((source / "slides").glob("*.png"))
        if not deck.is_file() or not montage.is_file() or not slides:
            raise FileNotFoundError(f"incomplete review artifact: {system}/{task['id']}")
        visible = _visible_text(deck)
        leaked = [marker for marker in IDENTITY_MARKERS if marker in visible]
        if leaked:
            raise ValueError(f"visible system identity leak in {system}/{task['id']}: {','.join(leaked)}")
        mapping_rows.append({
            "anonymous_id": anonymous_id, "system": system, "task_id": task["id"],
            "deck_sha256": hashlib.sha256(deck.read_bytes()).hexdigest(),
        })
    out_root.mkdir(parents=True)
    _write(out_root / "private_mapping.json", {
        "schema":"pptbench-v2-private-anonymous-mapping", "artifacts":mapping_rows,
        "never_distribute_with_reviewer_bundle":True,
    })
    orders = []
    for reviewer_index in (1, 2):
        reviewer = out_root / f"reviewer_{reviewer_index}"
        artifacts_root = reviewer / "artifacts"
        artifacts_root.mkdir(parents=True)
        order = list(anonymous_ids)
        rng.shuffle(order)
        if reviewer_index == 2 and order == orders[0] and len(order) > 1:
            order = order[1:] + order[:1]
        orders.append(order)
        by_id = {row["anonymous_id"]: row for row in mapping_rows}
        task_by_id = {task["id"]: task for task in protocol["tasks"]}
        for anonymous_id in order:
            row = by_id[anonymous_id]
            task = task_by_id[row["task_id"]]
            source = raw_root / row["system"] / row["task_id"]
            target = artifacts_root / anonymous_id
            (target / "slides").mkdir(parents=True)
            shutil.copy2(source / "montage.png", target / "montage.png")
            for slide in sorted((source / "slides").glob("*.png")):
                shutil.copy2(slide, target / "slides" / slide.name)
            _write(target / "task.json", {
                "anonymous_id":anonymous_id, "task_id":task["id"], "audience":task["audience"],
                "deck_job":task["deck_job"], "prompt":task["prompt"], "facts":task["facts"],
                "required_text":task["required_text"], "min_slides":task["min_slides"], "max_slides":task["max_slides"],
            })
        _write(reviewer / "order.json", {"schema":"pptbench-v2-review-order", "anonymous_ids":order})
        _write(reviewer / "review_form.draft.json", {
            "schema":"pptbench-v2-blind-review-form-draft", "reviewer_id":"",
            "independent_from_generation":False, "reviewer_non_author":False,
            "conflicts_declared":False, "review_completed_at":"",
            "locked_before_adjudication":False,
            "scores":[{"anonymous_id":value, "dimensions":{dimension:None for dimension in protocol["blind_review"]["dimensions"]}, "blocking_comment":""} for value in order],
        })
    public_text = "\n".join(
        path.read_text(encoding="utf-8")
        for reviewer_dir in out_root.glob("reviewer_*")
        for path in reviewer_dir.rglob("*.json")
    )
    if any(marker in public_text.casefold() for marker in IDENTITY_MARKERS):
        raise RuntimeError("identity marker leaked into public reviewer metadata")
    return {"schema":"pptbench-v2-blind-bundle-readiness", "valid":True, "n_artifacts":len(expected), "reviewer_orders_distinct":orders[0] != orders[1], "private_mapping":str(out_root / "private_mapping.json")}


def main() -> int:
    p = argparse.ArgumentParser()
    p.add_argument("--protocol", type=Path, required=True)
    p.add_argument("--raw-root", type=Path, required=True)
    p.add_argument("--out-root", type=Path, required=True)
    p.add_argument("--report", type=Path, required=True)
    a = p.parse_args()
    result = prepare(a.protocol.resolve(), a.raw_root.resolve(), a.out_root.resolve())
    a.report.parent.mkdir(parents=True, exist_ok=True)
    _write(a.report, result)
    print(json.dumps(result))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
