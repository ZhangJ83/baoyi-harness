"""Deterministic, provider-independent PPT structural scorer."""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation


def score(path: Path, min_slides: int = 1, required_text: list[str] | None = None) -> dict:
    prs = Presentation(str(path))
    texts = [shape.text for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
    corpus = "\n".join(texts).casefold()
    required = required_text or []
    slide_w, slide_h = prs.slide_width, prs.slide_height
    positioned = [
        shape
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "left") and hasattr(shape, "top")
    ]
    # Geometry invariant only; rendered visual quality is a separate check.
    no_overflow = all(
        shape.left >= 0
        and shape.top >= 0
        and shape.left + shape.width <= slide_w
        and shape.top + shape.height <= slide_h
        for shape in positioned
    )
    checks = {
        "opens": True,
        "min_slides": len(prs.slides) >= min_slides,
        "required_text": all(term.casefold() in corpus for term in required),
        "no_overflow": no_overflow,
    }
    passed = sum(checks.values())
    return {
        "artifact": str(path),
        "slides": len(prs.slides),
        "checks": checks,
        "score": passed / len(checks),
        "deterministic": True,
    }


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("deck", type=Path)
    parser.add_argument("--min-slides", type=int, default=1)
    parser.add_argument("--required-text", action="append", default=[])
    args = parser.parse_args()
    result = score(args.deck, args.min_slides, args.required_text)
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0 if result["score"] == 1.0 else 1


if __name__ == "__main__":
    raise SystemExit(main())
