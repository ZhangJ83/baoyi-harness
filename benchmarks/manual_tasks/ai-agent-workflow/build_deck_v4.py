"""AI Agent workflow deck v4.

Page 1: native PPTX pipeline with status chip, five stage cards with progress
bars, and a delivery bar.
Page 2: HTML-dashboard style with nav actions, KPI strip, 2x2 cards including
progress bars and three action buttons, plus footer.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = Path(r"E:\project\agent\xiaopu\workspace\iteration_runs\tasks\ai-agent-workflow\output\ai_agent_workflow.pptx")
OUT.parent.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(0x0F, 0x1A, 0x2B)
NAVY2 = RGBColor(0x18, 0x2B, 0x46)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xA8, 0xB8, 0xCE)
AMBER = RGBColor(0xF6, 0xB5, 0x44)
INK = RGBColor(0x22, 0x2A, 0x35)
GRAY = RGBColor(0x66, 0x70, 0x7D)
PAPER = RGBColor(0xF3, 0xF5, 0xF8)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xDF, 0xE5, 0xEC)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF9, 0x7A, 0x2E)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
CYAN = RGBColor(0x06, 0xB6, 0xD4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill, line=None, round_=False, radius=0.12):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    if round_:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    return sh


def text(slide, x, y, w, h, value, size, color, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, spacing=1.0, font="Microsoft YaHei"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(value.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return box


def arrow(slide, x1, y1, x2, y2, color=AMBER):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(2)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement("{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd",
                             {"type": "arrow", "w": "med", "len": "med"}))


def progress(slide, x, y, w, frac, color):
    rect(slide, x, y, w, 0.10, RGBColor(0x2A, 0x3D, 0x5A), None, True, 0.5)
    rect(slide, x, y, max(0.08, w * frac), 0.10, color, None, True, 0.5)


# ---------------- page 1 ----------------
s1 = prs.slides.add_slide(BLANK)
rect(s1, 0, 0, 13.333, 7.5, NAVY)
rect(s1, 0, 0, 13.333, 0.10, AMBER)
rect(s1, 0, 7.14, 13.333, 0.36, NAVY2)
text(s1, 0.85, 0.50, 8.2, 0.4, "AI AGENT · WORKFLOW OVERVIEW", 12, AMBER, True)
rect(s1, 10.65, 0.50, 1.85, 0.44, GREEN, None, True, 0.5)
text(s1, 10.65, 0.56, 1.85, 0.30, "● 系统就绪", 10.5, WHITE, True, PP_ALIGN.CENTER)
text(s1, 0.85, 0.88, 11.6, 1.0, "AI Agent 工作流程", 40, WHITE, True)
text(s1, 0.85, 1.92, 11.6, 0.55, "从任务输入到可验证交付的六阶段闭环，每个阶段都产出可审计证据", 17, MUTED)
text(s1, 0.85, 6.88, 11.6, 0.30, "交付物 · 结构化证据 · 可恢复会话", 11, MUTED, False, PP_ALIGN.CENTER)

steps = [
    ("01", "输入 Input", "接收指令、附件与工作区上下文", "产出：结构化任务对象", 1.0, AMBER),
    ("02", "理解 Understand", "任务分类、能力匹配与计划生成", "产出：执行合同与工具面", 1.0, BLUE),
    ("03", "规划 Plan", "拆解阶段、预检风险与依赖关系", "产出：执行计划", 0.9, CYAN),
    ("04", "执行 Execute", "工具调用、原子事务与权限审计", "产出：候选产物", 0.75, ORANGE),
    ("05", "验证 Verify", "结构、渲染、视觉与官方评估器", "产出：验证证据", 0.6, GREEN),
]
x0, y0, w, h, gap = 0.72, 3.05, 2.16, 2.86, 0.22
for i, (num, title, desc, deliver, frac, accent) in enumerate(steps):
    x = x0 + i * (w + gap)
    rect(s1, x, y0, w, h, NAVY2, RGBColor(0x2E, 0x42, 0x62), True, 0.09)
    chip = rect(s1, x + 0.18, y0 + 0.18, 0.54, 0.54, accent, None, True, 0.5)
    tf = chip.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = NAVY
    text(s1, x + 0.16, y0 + 0.92, w - 0.32, 0.42, title, 14.5, WHITE, True)
    text(s1, x + 0.16, y0 + 1.38, w - 0.32, 0.75, desc, 10.5, MUTED, spacing=1.04)
    rect(s1, x + 0.16, y0 + 2.20, w - 0.32, 0.34, RGBColor(0x21, 0x35, 0x55), None, True, 0.5)
    text(s1, x + 0.24, y0 + 2.26, w - 0.48, 0.22, deliver, 8.5, accent)
    progress(s1, x + 0.16, y0 + 2.66, w - 0.32, frac, accent)
    if i < len(steps) - 1:
        arrow(s1, x + w + 0.03, y0 + h / 2, x + w + gap - 0.03, y0 + h / 2)

rect(s1, 0.72, 6.28, 11.90, 0.56, NAVY2, RGBColor(0x2E, 0x42, 0x62), True, 0.5)
text(s1, 0.92, 6.40, 11.5, 0.34, "INPUT → UNDERSTAND → PLAN → EXECUTE → VERIFY → DELIVER", 12.5, AMBER, True, PP_ALIGN.CENTER)

# ---------------- page 2 ----------------
s2 = prs.slides.add_slide(BLANK)
rect(s2, 0, 0, 13.333, 7.5, PAPER)
rect(s2, 0, 0, 13.333, 0.78, CARD, BORDER)
rect(s2, 0.55, 0.20, 0.38, 0.38, BLUE, None, True, 0.5)
text(s2, 1.08, 0.16, 3.2, 0.45, "AI Agent Console", 15, INK, True)
text(s2, 4.70, 0.25, 2.0, 0.32, "● 系统在线", 10.5, GREEN, True)
rect(s2, 9.05, 0.19, 1.35, 0.42, CARD, BORDER, True, 0.5)
text(s2, 9.05, 0.27, 1.35, 0.26, "帮助文档", 10.5, GRAY, False, PP_ALIGN.CENTER)
rect(s2, 10.55, 0.19, 2.25, 0.42, BLUE, None, True, 0.5)
text(s2, 10.55, 0.26, 2.25, 0.30, "＋ 新建任务", 11.5, WHITE, True, PP_ALIGN.CENTER)

text(s2, 0.75, 0.98, 11.8, 0.52, "工作台 · 四阶段运行视图", 25, INK, True)
text(s2, 0.75, 1.48, 11.8, 0.38, "复杂任务被拆解为可验证阶段，状态、指标与证据实时可见", 12.5, GRAY)

kpis = [
    ("1.2s", "平均响应", BLUE), ("98.2%", "任务完成率", GREEN),
    ("99.7%", "工具调用成功率", ORANGE), ("100%", "证据覆盖率", PURPLE),
]
kw = 2.86; kgap = 0.22; kx = 0.75; ky = 1.98
for i, (value, label, color) in enumerate(kpis):
    x = kx + i * (kw + kgap)
    rect(s2, x, ky, kw, 0.88, CARD, BORDER, True, 0.10)
    rect(s2, x, ky, 0.08, 0.88, color)
    text(s2, x + 0.22, ky + 0.10, kw - 0.4, 0.34, value, 19, INK, True)
    text(s2, x + 0.22, ky + 0.50, kw - 0.4, 0.26, label, 10, GRAY)

cards = [
    ("INPUT", "输入 Input", "01", BLUE, "接收指令、附件与工作区上下文，完成意图澄清与队列排序。",
     ["自然语言与多模态输入", "自动发现任务相关文件", "生成结构化任务对象"], "12+", "输入类型", 1.0),
    ("UNDERSTAND", "理解 Understand", "02", GREEN, "解析任务语义，匹配类型、能力与执行合同。",
     ["任务分类与 Skill 选择", "能力与工具面解析", "风险与依赖预检"], "8", "任务类型", 1.0),
    ("EXECUTE", "执行 Execute", "03", ORANGE, "在阶段工具面内执行修改，全部变更走原子事务与审计。",
     ["阶段化工具准入", "原子批量事务", "失败自动回滚"], "100%", "事务回滚", 0.8),
    ("VERIFY", "验证 Verify", "04", PURPLE, "以结构、渲染、视觉与官方评估器四层证据确认交付质量。",
     ["反例定位与局部修复", "证据新鲜度管理", "交付前强制 gate"], "4", "证据层", 0.7),
]
positions = [(0.75, 3.02), (6.90, 3.02), (0.75, 5.72), (6.90, 5.72)]
cw, ch = 5.70, 2.48
for (tag, title, num, accent, desc, bullets, metric, metric_label, frac), (cx, cy) in zip(cards, positions):
    rect(s2, cx, cy, cw, ch, CARD, BORDER, True, 0.06)
    rect(s2, cx, cy, 0.09, ch, accent)
    rect(s2, cx + 0.30, cy + 0.22, 0.62, 0.62, accent, None, True, 0.18)
    text(s2, cx + 0.30, cy + 0.30, 0.62, 0.45, num, 14, WHITE, True, PP_ALIGN.CENTER)
    text(s2, cx + 1.08, cy + 0.20, 2.6, 0.40, title, 16, INK, True)
    rect(s2, cx + cw - 1.50, cy + 0.28, 1.20, 0.30, PAPER, None, True, 0.5)
    text(s2, cx + cw - 1.50, cy + 0.32, 1.20, 0.22, tag, 8.5, accent, True, PP_ALIGN.CENTER)
    text(s2, cx + 0.32, cy + 0.94, cw - 0.64, 0.58, desc, 11, GRAY, spacing=1.04)
    bullets = "\n".join("•  " + b for b in bullets)
    text(s2, cx + 0.32, cy + 1.54, cw - 1.30, 0.70, bullets, 10, INK, spacing=1.02)
    rect(s2, cx + cw - 1.10, cy + 1.54, 0.78, 0.70, PAPER, BORDER, True, 0.10)
    text(s2, cx + cw - 1.10, cy + 1.62, 0.78, 0.28, metric, 14, accent, True, PP_ALIGN.CENTER)
    text(s2, cx + cw - 1.10, cy + 1.94, 0.78, 0.22, metric_label, 7.5, GRAY, False, PP_ALIGN.CENTER)
    progress(s2, cx + 0.32, cy + 2.34, cw - 1.30, frac, accent)
    text(s2, cx + 0.32, cy + 2.46, 1.0, 0.18, "查看日志", 8.5, accent, True)
    text(s2, cx + 1.30, cy + 2.46, 1.0, 0.18, "查看证据", 8.5, accent, True)
    text(s2, cx + 2.28, cy + 2.46, 1.0, 0.18, "导出", 8.5, accent, True)

rect(s2, 0, 7.06, 13.333, 0.44, CARD, BORDER)
text(s2, 0.75, 7.13, 11.8, 0.30, "AI Agent Console · 运行状态正常 · 所有阶段证据已归档", 10, GRAY, False, PP_ALIGN.CENTER)

prs.save(OUT)
print("saved", OUT, "slides", len(prs.slides))
