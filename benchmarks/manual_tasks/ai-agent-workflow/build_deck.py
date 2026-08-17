"""AI Agent workflow deck v6.

Page 1: six-stage native pipeline (input -> deliver) with stage cards and
delivery closure. Page 2: compact HTML-dashboard layout from v5.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = Path(r"E:\project\agent\xiaopu\workspace\iteration_runs\tasks\ai-agent-workflow\output\ai_agent_workflow.pptx")
OUT.parent.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(0x0F, 0x1A, 0x2B)
NAVY2 = RGBColor(0x18, 0x2B, 0x46)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xA8, 0xB8, 0xCE)
AMBER = RGBColor(0xF6, 0xB5, 0x44)
INK = RGBColor(0x22, 0x2A, 0x35)
GRAY = RGBColor(0x66, 0x70, 0x7D)
PAPER = RGBColor(0xF3, 0xF5, 0xF8)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xDF, 0xE5, 0xEC)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF9, 0x7A, 0x2E)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
CYAN = RGBColor(0x06, 0xB6, 0xD4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill, line=None, round_=False, radius=0.12):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    if round_:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    return sh


def text(slide, x, y, w, h, value, size, color, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, spacing=1.0, font="Microsoft YaHei"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(value.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return box


def arrow(slide, x1, y1, x2, y2, color=AMBER):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(2)
    ln = c.line._get_or_add_ln()
    ln.append(ln.makeelement("{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd",
                             {"type": "arrow", "w": "med", "len": "med"}))


def progress(slide, x, y, w, frac, color):
    rect(slide, x, y, w, 0.10, RGBColor(0x2A, 0x3D, 0x5A), None, True, 0.5)
    rect(slide, x, y, max(0.08, w * frac), 0.10, color, None, True, 0.5)


# ---------------- page 1 ----------------
s1 = prs.slides.add_slide(BLANK)
rect(s1, 0, 0, 13.333, 7.5, NAVY)
rect(s1, 0, 0, 13.333, 0.10, AMBER)
rect(s1, 0, 7.14, 13.333, 0.36, NAVY2)
text(s1, 0.85, 0.50, 8.2, 0.4, "AI AGENT · WORKFLOW OVERVIEW", 12, AMBER, True)
rect(s1, 10.65, 0.50, 1.85, 0.44, GREEN, None, True, 0.5)
text(s1, 10.65, 0.56, 1.85, 0.30, "● 系统就绪", 10.5, WHITE, True, PP_ALIGN.CENTER)
text(s1, 0.85, 0.88, 11.6, 1.0, "AI Agent 工作流程", 40, WHITE, True)
text(s1, 0.85, 1.92, 11.6, 0.55, "从任务输入到可验证交付的六阶段闭环，每个阶段都产出可审计证据", 17, MUTED)
text(s1, 0.85, 6.88, 11.6, 0.30, "交付物 · 结构化证据 · 可恢复会话", 11, MUTED, False, PP_ALIGN.CENTER)

steps = [
    ("01", "输入 Input", "指令、附件与上下文", "结构化任务", 1.0, AMBER),
    ("02", "理解 Understand", "分类、能力与计划", "执行合同", 1.0, BLUE),
    ("03", "规划 Plan", "阶段、风险与依赖", "执行计划", 0.9, CYAN),
    ("04", "执行 Execute", "调用、事务与审计", "候选产物", 0.75, ORANGE),
    ("05", "验证 Verify", "结构、渲染与评估", "验证证据", 0.6, GREEN),
    ("06", "交付 Deliver", "保存、摘要与归档", "最终交付物", 0.0, PURPLE),
]
x0, y0, gap = 0.72, 3.05, 0.18
w = (13.333 - 2 * x0 - 5 * gap) / 6
h = 2.86
for i, (num, title, desc, deliver, frac, accent) in enumerate(steps):
    x = x0 + i * (w + gap)
    rect(s1, x, y0, w, h, NAVY2, RGBColor(0x2E, 0x42, 0x62), True, 0.10)
    chip = rect(s1, x + 0.14, y0 + 0.18, 0.46, 0.46, accent, None, True, 0.5)
    tf = chip.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = NAVY
    text(s1, x + 0.12, y0 + 0.86, w - 0.24, 0.40, title, 12.5, WHITE, True)
    text(s1, x + 0.12, y0 + 1.32, w - 0.24, 0.72, desc, 9.5, MUTED, spacing=1.02)
    rect(s1, x + 0.12, y0 + 2.12, w - 0.24, 0.34, RGBColor(0x21, 0x35, 0x55), None, True, 0.5)
    text(s1, x + 0.18, y0 + 2.18, w - 0.36, 0.22, deliver, 8, accent, True, PP_ALIGN.CENTER)
    progress(s1, x + 0.12, y0 + 2.58, w - 0.24, frac, accent)
    if i < len(steps) - 1:
        arrow(s1, x + w + 0.025, y0 + h / 2, x + w + gap - 0.025, y0 + h / 2)

rect(s1, 0.72, 6.28, 11.90, 0.56, NAVY2, RGBColor(0x2E, 0x42, 0x62), True, 0.5)
text(s1, 0.92, 6.40, 11.5, 0.34, "INPUT → UNDERSTAND → PLAN → EXECUTE → VERIFY → DELIVER", 12.5, AMBER, True, PP_ALIGN.CENTER)

# ---------------- page 2 ----------------
s2 = prs.slides.add_slide(BLANK)
rect(s2, 0, 0, 13.333, 7.5, PAPER)
rect(s2, 0, 0, 13.333, 0.78, CARD, BORDER)
rect(s2, 0.55, 0.20, 0.38, 0.38, BLUE, None, True, 0.5)
text(s2, 1.08, 0.16, 3.2, 0.45, "AI Agent Console", 15, INK, True)
tabs = [("总览", True), ("执行", False), ("验证", False), ("设置", False)]
tx = 4.55
for label, active in tabs:
    wtab = 0.72
    if active:
        rect(s2, tx, 0.27, wtab, 0.28, BLUE, None, True, 0.5)
        text(s2, tx, 0.31, wtab, 0.22, label, 10.5, WHITE, True, PP_ALIGN.CENTER)
    else:
        text(s2, tx, 0.31, wtab, 0.22, label, 10.5, GRAY, False, PP_ALIGN.CENTER)
    tx += wtab + 0.18
text(s2, 8.95, 0.25, 1.6, 0.32, "● 系统在线", 10.5, GREEN, True)
rect(s2, 10.55, 0.19, 2.25, 0.42, BLUE, None, True, 0.5)
text(s2, 10.55, 0.26, 2.25, 0.30, "＋ 新建任务", 11.5, WHITE, True, PP_ALIGN.CENTER)

text(s2, 0.75, 0.94, 11.8, 0.48, "工作台 · 四阶段运行视图", 23, INK, True)
text(s2, 0.75, 1.40, 11.8, 0.34, "复杂任务被拆解为可验证阶段，状态、指标与证据实时可见", 12, GRAY)

kpis = [
    ("1.2s", "平均响应", BLUE), ("98.2%", "任务完成率", GREEN),
    ("99.7%", "工具调用成功率", ORANGE), ("100%", "证据覆盖率", PURPLE),
]
kw = 2.86; kgap = 0.22; kx = 0.75; ky = 1.84; kh = 0.78
for i, (value, label, color) in enumerate(kpis):
    x = kx + i * (kw + kgap)
    rect(s2, x, ky, kw, kh, CARD, BORDER, True, 0.10)
    rect(s2, x, ky, 0.08, kh, color)
    text(s2, x + 0.22, ky + 0.08, kw - 0.4, 0.30, value, 18, INK, True)
    text(s2, x + 0.22, ky + 0.46, kw - 0.4, 0.24, label, 9.5, GRAY)

cards = [
    ("INPUT", "输入 Input", "01", BLUE, "接收指令、附件与工作区上下文，完成意图澄清与队列排序。",
     ["自然语言与多模态输入", "自动发现任务相关文件", "生成结构化任务对象"], "12+", "输入类型", 1.0),
    ("UNDERSTAND", "理解 Understand", "02", GREEN, "解析任务语义，匹配类型、能力与执行合同。",
     ["任务分类与 Skill 选择", "能力与工具面解析", "风险与依赖预检"], "8", "任务类型", 1.0),
    ("EXECUTE", "执行 Execute", "03", ORANGE, "在阶段工具面内执行修改，全部变更走原子事务与审计。",
     ["阶段化工具准入", "原子批量事务", "失败自动回滚"], "100%", "事务回滚", 0.8),
    ("VERIFY", "验证 Verify", "04", PURPLE, "以结构、渲染、视觉与官方评估器四层证据确认交付质量。",
     ["反例定位与局部修复", "证据新鲜度管理", "交付前强制 gate"], "4", "证据层", 0.7),
]
positions = [(0.75, 2.76), (6.90, 2.76), (0.75, 4.86), (6.90, 4.86)]
cw, ch = 5.70, 1.98
for (tag, title, num, accent, desc, bullets, metric, metric_label, frac), (cx, cy) in zip(cards, positions):
    rect(s2, cx, cy, cw, ch, CARD, BORDER, True, 0.07)
    rect(s2, cx, cy, 0.09, ch, accent)
    rect(s2, cx + 0.26, cy + 0.18, 0.50, 0.50, accent, None, True, 0.2)
    text(s2, cx + 0.26, cy + 0.24, 0.50, 0.36, num, 13, WHITE, True, PP_ALIGN.CENTER)
    text(s2, cx + 0.90, cy + 0.16, 2.4, 0.34, title, 14.5, INK, True)
    rect(s2, cx + cw - 1.42, cy + 0.22, 1.14, 0.28, PAPER, None, True, 0.5)
    text(s2, cx + cw - 1.42, cy + 0.25, 1.14, 0.20, tag, 8, accent, True, PP_ALIGN.CENTER)
    text(s2, cx + 0.28, cy + 0.76, cw - 0.56, 0.46, desc, 10.5, GRAY, spacing=1.02)
    bullets = "\n".join("•  " + b for b in bullets)
    text(s2, cx + 0.28, cy + 1.24, cw - 1.26, 0.58, bullets, 9.5, INK, spacing=0.98)
    rect(s2, cx + cw - 1.06, cy + 1.24, 0.76, 0.58, PAPER, BORDER, True, 0.10)
    text(s2, cx + cw - 1.06, cy + 1.30, 0.76, 0.24, metric, 12.5, accent, True, PP_ALIGN.CENTER)
    text(s2, cx + cw - 1.06, cy + 1.58, 0.76, 0.20, metric_label, 7, GRAY, False, PP_ALIGN.CENTER)
    progress(s2, cx + 0.28, cy + 1.92, cw - 1.26, frac, accent)
    text(s2, cx + 0.28, cy + 1.98, 1.0, 0.16, "查看日志", 7.5, accent, True)
    text(s2, cx + 1.24, cy + 1.98, 1.0, 0.16, "查看证据", 7.5, accent, True)
    text(s2, cx + 2.20, cy + 1.98, 1.0, 0.16, "导出", 7.5, accent, True)

rect(s2, 0, 7.06, 13.333, 0.44, CARD, BORDER)
text(s2, 0.75, 7.13, 11.8, 0.30, "AI Agent Console · 运行状态正常 · 所有阶段证据已归档", 10, GRAY, False, PP_ALIGN.CENTER)

prs.save(OUT)
print("saved", OUT, "slides", len(prs.slides))
