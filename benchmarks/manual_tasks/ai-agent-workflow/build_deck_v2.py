"""Build a richer, visually varied 2-page AI Agent workflow deck.

Page 1: native PPTX shapes — dark hero, 5-step horizontal flow, delivery bar.
Page 2: HTML-dashboard style — top bar, 2x2 cards, badges, metric chips, footer.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

OUT = Path(r"E:\project\agent\xiaopu\workspace\iteration_runs\tasks\ai-agent-workflow\output\ai_agent_workflow.pptx")
OUT.parent.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(0x10, 0x1B, 0x2D)
NAVY2 = RGBColor(0x18, 0x2A, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xB8, 0xC4, 0xD6)
AMBER = RGBColor(0xF5, 0xB3, 0x41)
INK = RGBColor(0x22, 0x2A, 0x35)
GRAY = RGBColor(0x6B, 0x74, 0x80)
PAPER = RGBColor(0xF4, 0xF6, 0xF9)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE3, 0xE8, 0xEE)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF9, 0x7A, 0x2E)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
RED = RGBColor(0xEF, 0x44, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None, round_=False, radius=0.12):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    if round_:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    return sh


def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
             font="Microsoft YaHei", anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_arrow(slide, x1, y1, x2, y2, color):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(2.0)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd",
        {"type": "arrow", "w": "med", "len": "med"},
    ))
    return conn


# ----------------------------------------------------------------------------- page 1
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, 13.333, 7.5, NAVY)
add_rect(s1, 0, 0, 13.333, 0.10, AMBER)
add_text(s1, 0.9, 0.55, 11.5, 0.4, "AI AGENT · WORKFLOW OVERVIEW", 12, AMBER, True)
add_text(s1, 0.9, 0.95, 11.5, 1.1, "AI Agent 工作流程", 42, WHITE, True)
add_text(s1, 0.9, 2.05, 11.5, 0.6, "从任务输入到可验证交付的六阶段闭环，每一步都沉淀为可审计证据", 18, MUTED)

steps = [
    ("01", "输入 Input", "接收指令、文件与上下文"),
    ("02", "理解 Understand", "任务分类、能力匹配与计划"),
    ("03", "规划 Plan", "生成执行合同与工具面"),
    ("04", "执行 Execute", "工具调用、事务与沙箱"),
    ("05", "验证 Verify", "结构、渲染、视觉与评估器"),
]
x0, y0, w, h, gap = 0.75, 3.35, 2.12, 2.55, 0.24
for i, (num, title, desc) in enumerate(steps):
    x = x0 + i * (w + gap)
    add_rect(s1, x, y0, w, h, NAVY2, line=RGBColor(0x2C, 0x3E, 0x5E), round_=True, radius=0.10)
    chip = add_rect(s1, x + 0.18, y0 + 0.18, 0.52, 0.52, AMBER, round_=True, radius=0.5)
    tf = chip.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = num
    run.font.size = Pt(13); run.font.bold = True; run.font.color.rgb = NAVY
    add_text(s1, x + 0.16, y0 + 0.90, w - 0.32, 0.45, title, 15, WHITE, True)
    add_text(s1, x + 0.16, y0 + 1.45, w - 0.32, 0.95, desc, 11.5, MUTED, line_spacing=1.05)
    if i < len(steps) - 1:
        add_arrow(s1, x + w + 0.035, y0 + h / 2, x + w + gap - 0.035, y0 + h / 2, AMBER)

add_rect(s1, 0.75, 6.35, 11.83, 0.62, NAVY2, line=RGBColor(0x2C, 0x3E, 0x5E), round_=True, radius=0.5)
add_text(s1, 0.95, 6.47, 11.4, 0.4, "INPUT → UNDERSTAND → PLAN → EXECUTE → VERIFY → DELIVER", 13, AMBER, True, PP_ALIGN.CENTER)

# ----------------------------------------------------------------------------- page 2
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, 13.333, 7.5, PAPER)
add_rect(s2, 0, 0, 13.333, 0.78, CARD, line=BORDER)
add_rect(s2, 0.55, 0.20, 0.38, 0.38, BLUE, round_=True, radius=0.5)
add_text(s2, 1.10, 0.16, 3.0, 0.45, "AI Agent Console", 15, INK, True)
add_text(s2, 4.55, 0.24, 2.2, 0.35, "● 系统在线", 11, GREEN, True)
btn = add_rect(s2, 10.55, 0.19, 2.25, 0.42, BLUE, round_=True, radius=0.5)
add_text(s2, 10.55, 0.26, 2.25, 0.30, "＋ 新建任务", 11.5, WHITE, True, PP_ALIGN.CENTER)

add_text(s2, 0.75, 1.02, 11.8, 0.55, "工作台 · 四阶段运行视图", 26, INK, True)
add_text(s2, 0.75, 1.60, 11.8, 0.4, "将复杂任务拆解为可验证的阶段，实时跟踪每一步的状态与证据", 13, GRAY)

cards = [
    ("INPUT", "输入 Input", "01", BLUE,
     "接收用户指令、附件与工作区上下文，完成意图澄清与队列排序。",
     ["支持自然语言与多模态输入", "自动发现任务相关文件", "生成结构化任务对象"]),
    ("UNDERSTAND", "理解 Understand", "02", GREEN,
     "解析任务语义，匹配任务类型、所需能力与执行合同。",
     ["任务分类与 Skill 选择", "能力与工具面解析", "风险与依赖预检"]),
    ("EXECUTE", "执行 Execute", "03", ORANGE,
     "在阶段工具面内执行修改，所有变更走原子事务与权限审计。",
     ["阶段化工具准入", "原子批量事务", "失败自动回滚"]),
    ("VERIFY", "验证 Verify", "04", PURPLE,
     "用结构、渲染、视觉与官方评估器四层证据确认交付质量。",
     ["反例定位与局部修复", "证据新鲜度管理", "交付前强制 gate"]),
]
positions = [(0.75, 2.25), (6.90, 2.25), (0.75, 4.68), (6.90, 4.68)]
cw, ch = 5.70, 2.28
for (tag, title, num, accent, desc, bullets), (cx, cy) in zip(cards, positions):
    add_rect(s2, cx, cy, cw, ch, CARD, line=BORDER, round_=True, radius=0.06)
    add_rect(s2, cx, cy, 0.09, ch, accent)
    add_rect(s2, cx + 0.30, cy + 0.22, 0.62, 0.62, accent, round_=True, radius=0.18)
    add_text(s2, cx + 0.30, cy + 0.30, 0.62, 0.45, num, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text(s2, cx + 1.10, cy + 0.20, 3.2, 0.40, title, 17, INK, True)
    tag_pill = add_rect(s2, cx + cw - 1.25, cy + 0.28, 0.95, 0.30, PAPER, round_=True, radius=0.5)
    add_text(s2, cx + cw - 1.25, cy + 0.32, 0.95, 0.22, tag, 9, accent, True, PP_ALIGN.CENTER)
    add_text(s2, cx + 0.32, cy + 0.94, cw - 0.64, 0.62, desc, 11.5, GRAY, line_spacing=1.05)
    bullet_text = "\n".join("•  " + b for b in bullets)
    add_text(s2, cx + 0.32, cy + 1.52, cw - 0.64, 0.66, bullet_text, 10.5, INK, line_spacing=1.02)

add_rect(s2, 0, 7.06, 13.333, 0.44, CARD, line=BORDER)
add_text(s2, 0.75, 7.13, 11.8, 0.3, "AI Agent Console · 运行状态正常 · 所有阶段证据已归档", 10, GRAY, False, PP_ALIGN.CENTER)

prs.save(OUT)
print("saved", OUT, "slides", len(prs.slides))
