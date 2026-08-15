"""Common, policy-independent evaluation for one real-controller PPT cell."""
from __future__ import annotations

import argparse
import hashlib
import json
import os
from pathlib import Path

from pptx import Presentation

from agent.state import RunState
from agent.tools.ppt_tools import _inspect_rendered, _render, _verify


def digest(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


class _Harness:
    def __init__(self, deck):
        self.deck = deck
        self.state = RunState()


def evaluate(workspace: Path, task: dict, deck_name: str = "deck.pptx") -> dict:
    workspace = workspace.resolve()
    deck_path = workspace / deck_name
    reports = workspace / "common_evaluation"
    reports.mkdir(parents=True, exist_ok=True)
    errors: list[str] = []
    checks = {
        "pptx_opens": False, "slide_bounds": False, "required_text": False,
        "pdf_complete": False, "png_complete": False,
        "no_unintended_overlap": False, "no_overflow": False,
        "no_placeholder": False, "fresh_evidence": False,
    }
    structural_text = "not run"
    pixel_text = "not run"
    slide_count = 0
    harness = None
    try:
        deck = Presentation(str(deck_path))
        harness = _Harness(deck)
        checks["pptx_opens"] = True
        slide_count = len(deck.slides)
        checks["slide_bounds"] = task["min_slides"] <= slide_count <= task["max_slides"]
        texts = [shape.text for slide in deck.slides for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
        combined = "\n".join(texts).casefold()
        checks["required_text"] = all(str(value).casefold() in combined for value in task["required_text"])
        forbidden = ("lorem ipsum", "placeholder", "todo:", "insert ", "[title]", "[content]")
        checks["no_placeholder"] = not any(value in combined for value in forbidden)
        structural_text = _verify(harness)
        structural_pass = structural_text == "Verification: no structural issues found."
        checks["no_unintended_overlap"] = structural_pass
        checks["no_overflow"] = structural_pass
    except Exception as exc:
        errors.append(f"structure:{type(exc).__name__}:{exc}")
    (reports / "structural_report.txt").write_text(structural_text + "\n", encoding="utf-8")
    if harness is not None:
        prior = os.environ.get("WORKSPACE")
        os.environ["WORKSPACE"] = str(workspace)
        try:
            _render(harness, deck_name, "common_evaluation/rendered")
            pixel_text = _inspect_rendered(harness, "common_evaluation/rendered")
            render_dir = reports / "rendered"
            pdf = render_dir / f"{Path(deck_name).stem}.pdf"
            pngs = sorted(path for path in render_dir.glob("*.png") if path.name.lower() != "montage.png")
            checks["pdf_complete"] = pdf.is_file() and pdf.stat().st_size > 0
            checks["png_complete"] = len(pngs) == slide_count and slide_count > 0
            kinds = {record.kind for record in harness.state.fresh_evidence()}
            checks["fresh_evidence"] = {"ppt_structural", "ppt_render", "ppt_visual"}.issubset(kinds)
            (reports / "slide_pngs_manifest.json").write_text(json.dumps({
                "count": len(pngs), "slides": [{"path": str(path.relative_to(workspace)).replace("\\", "/"), "sha256": digest(path)} for path in pngs]
            }, indent=2) + "\n", encoding="utf-8")
        except Exception as exc:
            errors.append(f"render:{type(exc).__name__}:{exc}")
        finally:
            if prior is None:
                os.environ.pop("WORKSPACE", None)
            else:
                os.environ["WORKSPACE"] = prior
    (reports / "pixel_audit.txt").write_text(pixel_text + "\n", encoding="utf-8")
    result = {
        "schema": "real-controller-cell-evaluation-v1",
        "evaluator_sha256": digest(Path(__file__)),
        "renderer_sha256": digest(Path(__file__).resolve().parents[1] / "agent/tools/ppt_tools.py"),
        "task_id": task["id"], "deck": deck_name, "slide_count": slide_count,
        "evaluation": checks, "artifact_success": all(checks.values()),
        "infrastructure_valid": not any(error.startswith("render:") for error in errors),
        "errors": errors,
    }
    (reports / "cell_evaluation.json").write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return result


def main() -> int:
    p = argparse.ArgumentParser()
    p.add_argument("--workspace", type=Path, required=True)
    p.add_argument("--task-json", type=Path, required=True)
    p.add_argument("--deck", default="deck.pptx")
    a = p.parse_args()
    task = json.loads(a.task_json.read_text(encoding="utf-8"))
    result = evaluate(a.workspace, task, a.deck)
    print(json.dumps(result, ensure_ascii=False))
    return 0 if result["infrastructure_valid"] else 2


if __name__ == "__main__":
    raise SystemExit(main())
