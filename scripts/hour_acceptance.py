"""One-command, low-disk acceptance pass for the Xiaopu harness.

This script deliberately does not modify harness implementation files and does
not duplicate benchmark packages.  It exercises the domain-independent action
transaction, one real file-tool task, and two already-produced benchmark PPTX
smokes.  The only durable writes are a small JSON and Markdown report under the
requested report directory.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import platform
import shutil
import sys
import tempfile
import zipfile
from contextlib import contextmanager
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable

from pptx import Presentation

from agent.action_transaction import (
    ActionScope,
    ActionTransaction,
    CancellationToken,
    PostconditionFailed,
    ScopeViolation,
    TransactionCancelled,
)
from agent.state import RunState
from agent.tools.registry import dispatch


ROOT = Path(__file__).resolve().parents[1]
BENCHMARK = ROOT.parent / "ppt-harness" / "benchmark_v0.1" / "workbuddy"
SMOKE_MANIFEST = ROOT / ".smoke" / "complex_facades" / "evidence.json"


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for block in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def _assert(condition: bool, message: str) -> None:
    if not condition:
        raise AssertionError(message)


def _scope(*requested: str) -> ActionScope[str]:
    return ActionScope.from_iterables(
        allowed={"file:a", "file:b"}, requested=requested
    )


def transaction_success() -> dict[str, Any]:
    state = {"value": "before"}
    events = []
    tx = ActionTransaction(
        scope=_scope("file:a"),
        checkpoint=lambda: dict(state),
        execute=lambda: state.update(value="after") or "result",
        postcondition=lambda value: value == "result" and state["value"] == "after",
        commit=lambda _value: state.update(committed=True),
        rollback=lambda snapshot, _error: state.clear() or state.update(snapshot),
        event_sink=events.append,
        transaction_id="accept-success",
    )
    result = tx.run()
    phases = [event.phase for event in events]
    _assert(result.status == "committed", "transaction did not commit")
    _assert(state == {"value": "after", "committed": True}, "commit state mismatch")
    _assert(phases[-1] == "committed", "missing committed event")
    return {"transaction_id": result.transaction_id, "phases": phases, "state": state}


def transaction_scope_denied() -> dict[str, Any]:
    callbacks: list[str] = []
    events = []
    tx = ActionTransaction(
        scope=_scope("file:outside"),
        checkpoint=lambda: callbacks.append("checkpoint"),
        execute=lambda: callbacks.append("execute"),
        postcondition=lambda _value: True,
        rollback=lambda _snapshot, _error: callbacks.append("rollback"),
        event_sink=events.append,
        transaction_id="accept-scope-denied",
    )
    try:
        tx.run()
    except ScopeViolation as exc:
        _assert(callbacks == [], "scope denial ran a user callback")
        return {
            "expected_exception": type(exc).__name__,
            "callbacks": callbacks,
            "phases": [event.phase for event in events],
        }
    raise AssertionError("out-of-scope transaction unexpectedly committed")


def transaction_execute_rollback() -> dict[str, Any]:
    state = ["before"]
    events = []

    def execute() -> None:
        state.append("partial mutation")
        raise ValueError("synthetic execution failure")

    tx = ActionTransaction(
        scope=_scope("file:a"),
        checkpoint=lambda: list(state),
        execute=execute,
        postcondition=lambda _value: True,
        rollback=lambda snapshot, _error: state.__setitem__(slice(None), snapshot),
        event_sink=events.append,
        transaction_id="accept-execute-rollback",
    )
    try:
        tx.run()
    except ValueError as exc:
        phases = [event.phase for event in events]
        _assert(state == ["before"], "execution failure did not restore checkpoint")
        _assert(phases[-1] == "rolled_back", "rollback lifecycle is incomplete")
        return {"expected_exception": type(exc).__name__, "phases": phases, "state": state}
    raise AssertionError("failing execution unexpectedly committed")


def transaction_validation_rollback() -> dict[str, Any]:
    state = {"value": 1}
    events = []
    tx = ActionTransaction(
        scope=_scope("file:a"),
        checkpoint=lambda: dict(state),
        execute=lambda: state.update(value=2) or 2,
        postcondition=lambda _value: False,
        rollback=lambda snapshot, _error: state.clear() or state.update(snapshot),
        event_sink=events.append,
        transaction_id="accept-postcondition-rollback",
    )
    try:
        tx.run()
    except PostconditionFailed as exc:
        phases = [event.phase for event in events]
        _assert(state == {"value": 1}, "postcondition failure did not restore checkpoint")
        _assert("postcondition_started" in phases, "postcondition was not evaluated")
        _assert(phases[-1] == "rolled_back", "postcondition rollback is incomplete")
        return {"expected_exception": type(exc).__name__, "phases": phases, "state": state}
    raise AssertionError("false postcondition unexpectedly committed")


def transaction_cancel_rollback() -> dict[str, Any]:
    token = CancellationToken()
    state = ["before"]
    events = []

    def execute() -> str:
        state.append("partial mutation")
        token.cancel()
        return "done"

    tx = ActionTransaction(
        scope=_scope("file:a"),
        checkpoint=lambda: list(state),
        execute=execute,
        postcondition=lambda _value: True,
        rollback=lambda snapshot, _error: state.__setitem__(slice(None), snapshot),
        cancellation=token,
        event_sink=events.append,
        transaction_id="accept-cancel-rollback",
    )
    try:
        tx.run()
    except TransactionCancelled as exc:
        phases = [event.phase for event in events]
        _assert(state == ["before"], "cancellation did not restore checkpoint")
        _assert("cancelled" in phases, "cancellation event is missing")
        _assert(phases[-1] == "rolled_back", "cancellation rollback is incomplete")
        return {"expected_exception": type(exc).__name__, "phases": phases, "state": state}
    raise AssertionError("cancelled transaction unexpectedly committed")


class _Harness:
    def __init__(self) -> None:
        self.state = RunState()
        self.deck = None
        self.recorder = None


@contextmanager
def _workspace(path: Path):
    previous = os.environ.get("WORKSPACE")
    os.environ["WORKSPACE"] = str(path)
    try:
        yield
    finally:
        if previous is None:
            os.environ.pop("WORKSPACE", None)
        else:
            os.environ["WORKSPACE"] = previous


def ordinary_file_atomic_edit(report_parent: Path) -> dict[str, Any]:
    """Exercise the real file-tool facade without writing outside E:."""

    with tempfile.TemporaryDirectory(prefix="xiaopu_acceptance_", dir=report_parent) as tmp:
        workspace = Path(tmp)
        first = workspace / "first.md"
        second = workspace / "second.md"
        first.write_text("Status: draft\n", encoding="utf-8")
        second.write_text("Decision: pending\n", encoding="utf-8")
        harness = _Harness()
        with _workspace(workspace):
            result = dispatch(
                "apply_edits",
                json.dumps(
                    {
                        "edits": [
                            {"path": "first.md", "old": "draft", "new": "final"},
                            {"path": "second.md", "old": "pending", "new": "approved"},
                        ]
                    }
                ),
                harness,
            )
            try:
                dispatch("read_file", json.dumps({"path": "../outside.txt"}), harness)
            except PermissionError as exc:
                escape_error = str(exc)
            else:
                raise AssertionError("file tool allowed workspace escape")
        _assert(first.read_text(encoding="utf-8") == "Status: final\n", "first edit missing")
        _assert(second.read_text(encoding="utf-8") == "Decision: approved\n", "second edit missing")
        _assert(harness.state.mutation_epoch == 1, "two-file edit did not commit once")
        return {
            "tool": "apply_edits",
            "result": result,
            "mutation_epoch": harness.state.mutation_epoch,
            "changed_paths": sorted(harness.state.changed_files),
            "workspace_escape": {"blocked": True, "error": escape_error},
        }


def _slide_text(slide) -> str:
    chunks: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            chunks.append(shape.text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                chunks.extend(cell.text for cell in row.cells if cell.text.strip())
    return "\n".join(chunks)


def _ppt_package_check(path: Path) -> tuple[Presentation, dict[str, Any]]:
    _assert(path.is_file(), f"missing PPT smoke: {path}")
    with zipfile.ZipFile(path) as archive:
        corrupt_member = archive.testzip()
        members = len(archive.infolist())
    _assert(corrupt_member is None, f"corrupt OOXML member: {corrupt_member}")
    deck = Presentation(path)
    relation_count = 0
    for slide in deck.slides:
        for relation in slide.part.rels.values():
            relation_count += 1
            if not relation.is_external:
                _ = relation.target_part
    return deck, {
        "path": str(path),
        "bytes": path.stat().st_size,
        "sha256": _sha256(path),
        "zip_members": members,
        "slide_relationships_resolved": relation_count,
    }


def _historical_smoke(case: str) -> dict[str, Any]:
    manifest = json.loads(SMOKE_MANIFEST.read_text(encoding="utf-8"))
    matches = [item for item in manifest["results"] if item["case"] == case]
    _assert(len(matches) == 1, f"missing unique historical smoke record: {case}")
    row = matches[0]
    structural = str(row.get("check", {}).get("structural", ""))
    _assert("no blocking structural issues" in structural, "historical structural gate was not passing")
    _assert(row.get("source_unchanged") is True, "historical source-integrity assertion missing")
    return row


def ppt_smoke_xmind() -> dict[str, Any]:
    output = ROOT / ".smoke" / "complex_facades" / "from_outline" / "community_sustainability_workshop_deck.pptx"
    template = BENCHMARK / "xmind-screenshot-template-ppt" / "input" / "template" / "community_workshop_green_template.pptx"
    source = BENCHMARK / "xmind-screenshot-template-ppt" / "input" / "source" / "community_sustainability_workshop.xmind"
    reference = BENCHMARK / "xmind-screenshot-template-ppt" / "input" / "reference" / "mindmap_screenshot_style_reference.png"
    deck, package = _ppt_package_check(output)
    text = "\n".join(_slide_text(slide) for slide in deck.slides)
    _assert(len(deck.slides) == 8, "XMind smoke must contain 8 slides")
    for required in ("社区低碳共创坊", "30 天承诺卡", "不收手机号"):
        _assert(required in text, f"XMind smoke is missing required content: {required}")
    for page in range(1, 9):
        _assert(f"{page} / 8" in text, f"XMind smoke is missing page marker {page} / 8")
    sources = [template, source, reference]
    return {
        **package,
        "slides": len(deck.slides),
        "semantic_assertions": ["community theme", "30-day commitment", "privacy boundary", "1-8 page markers"],
        "source_hashes_observed_now": {str(path): _sha256(path) for path in sources},
        "historical_run_evidence": _historical_smoke("xmind-template/from_outline"),
        "visual_gate": "not asserted: no trustworthy PNG render exists in this session",
    }


def ppt_smoke_governance() -> dict[str, Any]:
    output = ROOT / ".smoke" / "complex_facades" / "batch_updates" / "quarterly_governance_board_deck_partial_smoke.pptx"
    source_deck = BENCHMARK / "board-material-update-timeline-excel" / "input" / "quarterly_governance_board_deck.pptx"
    register = BENCHMARK / "board-material-update-timeline-excel" / "input" / "program_milestone_update_register.xlsx"
    deck, package = _ppt_package_check(output)
    _assert(len(deck.slides) == 11, "governance smoke must retain 11 slides")
    slide_1 = _slide_text(deck.slides[0])
    slide_9 = _slide_text(deck.slides[8])
    slide_10 = _slide_text(deck.slides[9])
    slide_11 = _slide_text(deck.slides[10])
    for label, text in (("slide 1", slide_1), ("slide 9", slide_9)):
        _assert("v1.8 当前" in text and "2026-07-01" in text, f"{label} lacks scoped update")
    for label, text in (("slide 10", slide_10), ("slide 11", slide_11)):
        _assert("v1.7 草稿" in text and "2026-06-12" in text, f"{label} lost preserved historical/template text")
    return {
        **package,
        "slides": len(deck.slides),
        "semantic_assertions": ["slides 1/9 updated", "slides 10/11 preserved", "11-slide count retained"],
        "source_hashes_observed_now": {
            str(source_deck): _sha256(source_deck),
            str(register): _sha256(register),
        },
        "historical_run_evidence": _historical_smoke("governance/batch_updates-representative"),
        "scope_note": "representative facade smoke; not a full benchmark submission",
        "visual_gate": "not asserted: no trustworthy PNG render exists in this session",
    }


def _run_case(name: str, fn: Callable[[], dict[str, Any]]) -> dict[str, Any]:
    try:
        evidence = fn()
    except BaseException as exc:
        return {
            "name": name,
            "status": "failed",
            "error": f"{type(exc).__name__}: {exc}",
        }
    return {"name": name, "status": "passed", "evidence": evidence}


def _markdown(report: dict[str, Any]) -> str:
    rows = ["| 验收项 | 结果 |", "|---|---|"]
    for case in report["cases"]:
        detail = case.get("error", "证据已写入 JSON")
        rows.append(f"| `{case['name']}` | {case['status']} - {detail} |")
    limitations = "\n".join(f"- {item}" for item in report["limitations"])
    return (
        "# 小朴一小时验收报告\n\n"
        f"- 生成时间：`{report['generated_at']}`\n"
        f"- 总体结果：**{report['status']}**\n"
        f"- 通过：`{report['summary']['passed']}/{report['summary']['total']}`\n\n"
        + "\n".join(rows)
        + "\n\n## 证据边界\n\n"
        + limitations
        + "\n"
    )


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--report-dir",
        type=Path,
        default=ROOT / ".acceptance" / "hourly" / "current",
        help="Small report directory; keep this on a drive with free space.",
    )
    args = parser.parse_args()
    report_dir = args.report_dir.resolve()
    report_dir.mkdir(parents=True, exist_ok=True)

    cases = [
        _run_case("transaction_success", transaction_success),
        _run_case("transaction_scope_denied", transaction_scope_denied),
        _run_case("transaction_execute_rollback", transaction_execute_rollback),
        _run_case("transaction_validation_rollback", transaction_validation_rollback),
        _run_case("transaction_cancel_rollback", transaction_cancel_rollback),
        _run_case("ordinary_file_atomic_edit", lambda: ordinary_file_atomic_edit(report_dir.parent)),
        _run_case("ppt_xmind_from_outline", ppt_smoke_xmind),
        _run_case("ppt_governance_batch_update", ppt_smoke_governance),
    ]
    passed = sum(case["status"] == "passed" for case in cases)
    report = {
        "schema": "xiaopu.hour-acceptance.v1",
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "status": "passed" if passed == len(cases) else "failed",
        "summary": {"passed": passed, "total": len(cases)},
        "environment": {
            "python": sys.version,
            "platform": platform.platform(),
            "free_bytes": {
                "C": shutil.disk_usage("C:\\").free,
                "E": shutil.disk_usage("E:\\").free,
            },
        },
        "cases": cases,
        "limitations": [
            "The two PPT checks reopen existing real benchmark smokes; they do not rerun a model trajectory.",
            "These historical XMind/governance checks do not claim fresh PNG rendering. A separate two-case PPT transaction smoke has artifact-tool PNG evidence under .acceptance/hourly/transaction-smoke/.",
            "The governance PPT is explicitly a representative batch-transaction smoke, not the full task submission.",
            "Historical source_unchanged and structural-gate claims are imported from the frozen smoke manifest; current source hashes are recorded separately.",
        ],
    }
    json_path = report_dir / "acceptance.json"
    markdown_path = report_dir / "acceptance.md"
    json_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    markdown_path.write_text(_markdown(report), encoding="utf-8")
    print(json.dumps({
        "status": report["status"],
        "passed": passed,
        "total": len(cases),
        "json": str(json_path),
        "markdown": str(markdown_path),
    }, ensure_ascii=False))
    return 0 if report["status"] == "passed" else 1


if __name__ == "__main__":
    raise SystemExit(main())
