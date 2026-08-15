"""Run the four one-hour representative PPT facade smokes and emit evidence JSON."""
from __future__ import annotations

import hashlib
import json
import os
import shutil
import sys
from pathlib import Path

from pptx import Presentation

from agent.state import RunState
from agent.tools.registry import dispatch


ROOT = Path(__file__).resolve().parents[1]
BENCH = ROOT.parent / "ppt-harness" / "benchmark_v0.1"
OUT = ROOT / ".smoke" / "complex_facades"


class Harness:
    def __init__(self):
        self.state = RunState()
        self.deck = None


def digest(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def slide_text(slide) -> str:
    return "\n".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))


def save_deck(h: Harness, path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    h.deck.save(path)


def smoke_outline() -> dict:
    source = BENCH / "workbuddy" / "xmind-screenshot-template-ppt" / "input" / "template" / "community_workshop_green_template.pptx"
    content_ir = json.loads((ROOT / ".pytest-tmp" / "xmind-from-outline.json").read_text(encoding="utf-8"))
    source_hash = digest(source)
    h = Harness()
    h.deck = Presentation(source)
    report = dispatch("ppt_compose", json.dumps(content_ir, ensure_ascii=False), h)
    output = OUT / "from_outline" / "community_sustainability_workshop_deck.pptx"
    save_deck(h, output)
    reopened = Presentation(output)
    texts = "\n".join(slide_text(slide) for slide in reopened.slides)
    placeholders = ("布局族", "占位文本", "布局样张仅供选择")
    assert len(reopened.slides) == 8
    assert all(term not in texts for term in placeholders)
    notes_texts = [
        "\n".join(shape.text for shape in slide.notes_slide.shapes if getattr(shape, "has_text_frame", False))
        for slide in reopened.slides
    ]
    assert all("[Sources]" in text for text in notes_texts)
    assert digest(source) == source_hash
    h.deck = reopened
    check = json.loads(dispatch("ppt_check", "{}", h))
    return {
        "case": "xmind-template/from_outline",
        "output": str(output), "slides": 8, "mutation_epoch": h.state.mutation_epoch,
        "source_unchanged": True, "compose": report, "check": check,
    }


def smoke_batch() -> dict:
    task = BENCH / "agent_workspaces" / "xiaopuharness_full13" / "tasks" / "board-material-update-timeline-excel"
    source = task / "input" / "quarterly_governance_board_deck.pptx"
    source_hash = digest(source)
    h = Harness()
    h.deck = Presentation(source)
    # Match ppt_open's existing-deck baseline contract for a direct facade smoke.
    from agent.tools.ppt_tools import _collect_structural_findings
    h.state.ppt_existing_deck = True
    h.state.ppt_baseline_captured = True
    h.state.ppt_baseline_findings = {
        item["key"]: item["severity"] for item in _collect_structural_findings(h)
    }
    before = len(h.deck.slides)
    draft_label = "v1.7 " + chr(0x8349) + chr(0x7A3F)
    current_label = "v1.8 " + chr(0x5F53) + chr(0x524D)
    report = dispatch("ppt_edit_text", json.dumps({
        "operation": "batch_updates",
        "updates": [
            {"operation": "replace", "slide_number": 1, "old": draft_label, "new": current_label},
            {"operation": "replace", "slide_number": 1, "old": "2026-06-12", "new": "2026-07-01"},
            {"operation": "replace", "slide_number": 9, "old": draft_label, "new": current_label},
            {"operation": "replace", "slide_number": 9, "old": "2026-06-12", "new": "2026-07-01"},
        ],
    }, ensure_ascii=False), h)
    output = OUT / "batch_updates" / "quarterly_governance_board_deck_partial_smoke.pptx"
    save_deck(h, output)
    reopened = Presentation(output)
    assert len(reopened.slides) == before == 11
    assert current_label in slide_text(reopened.slides[0]) and "2026-07-01" in slide_text(reopened.slides[0])
    assert digest(source) == source_hash
    h.deck = reopened
    check = json.loads(dispatch("ppt_check", "{}", h))
    return {
        "case": "governance/batch_updates-representative",
        "scope_note": "representative facade smoke, not the full task submission",
        "output": str(output), "slides": 11, "mutation_epoch": h.state.mutation_epoch,
        "affected_slides": sorted(h.state.ppt_affected_slides), "source_unchanged": True,
        "edit": report, "check": check,
    }


def main() -> int:
    OUT.mkdir(parents=True, exist_ok=True)
    results = [smoke_outline(), smoke_batch()]
    manifest = OUT / "evidence.json"
    manifest.write_text(json.dumps({"schema": "xiaopu.complex-facade-smoke.v1", "results": results}, ensure_ascii=False, indent=2), encoding="utf-8")
    print(manifest)
    for item in results:
        print(item["case"], item["output"])
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
